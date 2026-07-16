"""
gerar_pptx_cliente (api/pptx_brauna.py) — gera sem exceção e a barra de
alocação do slide "Alocação de Carteira" (_slide_alocacao) tem largura
proporcional correta ao % informado.

_slide_alocacao desenha, por linha/classe, um retângulo arredondado "track"
de largura fixa (col_w[1] = 5.20in) e sobrepõe uma barra colorida cuja largura
é `track_w * min(atual, escala_max) / escala_max` (escala_max = teto comum da
régua, arredondado para múltiplo de 10 acima do maior valor, piso 30). Este
teste reabre o .pptx gerado com python-pptx e confere:
  1) a proporção geométrica entre duas barras bate com a proporção dos % reais
     (a barra de 40% deve ter o dobro da largura da barra de 20%, não só a
     cor/texto do desvio — essa é a "conferência manual" que motivou o teste);
  2) a largura de cada barra bate com a fórmula de largura usada pelo próprio
     slide (mesma escala_max), como guarda de regressão do cálculo em si.
"""
import io
import math

import pytest
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE

from api.pptx_brauna import gerar_pptx_cliente, CLS_LABEL

# Mesmas constantes usadas em _slide_alocacao (api/pptx_brauna.py) — a coluna 2
# ("Atual vs. Modelo") é a área do track/barra.
BAR_TRACK_W = Inches(5.20)
BAR_H = Inches(0.16)

DADOS_BASE = {
    "nome": "Cliente PPTX Teste",
    "assessor": "Denis Kinoshita",
    "perfil": "moderada",
    "data_ref": "15/07/2026",
    "patrimonio": 500_000,
    "rent": {"mes": 0.5, "ano": 5.0, "12m": 10.0, "24m": 20.0, "cdi_pct": 80.0},
    # Só 3 classes com valor > 0 (as demais ficam 0/0 e são puladas pelo slide),
    # com o maior valor definido para dar escala_max = 50 (piso do arredondamento
    # de math.ceil(45/10)*10) e nenhuma delas batendo exatamente na escala (o que
    # deixaria a barra do tamanho do track e ambíguo para o teste).
    "composicao": {
        "pos_fixado": 10.0,
        "inflacao": 25.0,
        "pre_fixado": 45.0,
        "acoes": 0, "fiis": 0, "multimercado": 0, "internacional": 0,
        "alternativos": 0, "criptomoedas": 0,
    },
    "modelo": {
        "pos_fixado": 10.0, "inflacao": 25.0, "pre_fixado": 45.0,
        "acoes": 0, "fiis": 0, "multimercado": 0, "internacional": 0,
        "alternativos": 0, "criptomoedas": 0,
    },
    "desvios": [],
    "acoes": [],
    "fiis": [],
    "rf_ativos": [],
    "checklist": {},
    "cross_ativos": [],
    "score_servir": 0,
    "status": "ok",
    "objetivo": "Teste",
    "alertas": [],
    "nota_lider": "",
}


@pytest.fixture(scope="module")
def slide_alocacao():
    pptx_bytes = gerar_pptx_cliente(DADOS_BASE)
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 7  # capa, resumo, alocação, rv, rf, servir, próximos passos
    return prs.slides[2]  # ordem fixa definida em gerar_pptx_cliente


def _barras_por_linha(slide):
    """Retorna, na ordem das linhas (pos_fixado, inflacao, pre_fixado), a shape
    da barra 'atual' de cada linha — a 2ª shape arredondada de cada par
    track+barra (a 1ª é sempre o track, de largura fixa BAR_TRACK_W)."""
    candidatas = [
        s for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and s.auto_shape_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        and s.height == BAR_H
    ]
    barras = []
    i = 0
    while i < len(candidatas):
        track = candidatas[i]
        assert track.width == BAR_TRACK_W, "1ª rounded-rect de cada linha deveria ser o track (largura fixa)"
        i += 1
        if i < len(candidatas) and candidatas[i].width != BAR_TRACK_W:
            barras.append(candidatas[i])
            i += 1
        else:
            barras.append(None)  # linha sem barra (atual == 0)
    return barras


def test_pptx_gera_sem_excecao_e_tem_7_slides(slide_alocacao):
    assert slide_alocacao is not None


def test_barras_tem_largura_proporcional_ao_percentual_informado(slide_alocacao):
    barras = _barras_por_linha(slide_alocacao)
    assert len(barras) == 3
    assert all(b is not None for b in barras)  # as 3 linhas têm atual > 0

    largura_10, largura_25, largura_45 = (b.width for b in barras)

    # Proporção geométrica: a razão entre larguras deve bater com a razão entre
    # os percentuais reais (dentro da tolerância de arredondamento int() do EMU).
    assert largura_25 / largura_10 == pytest.approx(25.0 / 10.0, rel=0.02)
    assert largura_45 / largura_10 == pytest.approx(45.0 / 10.0, rel=0.02)

    # Guarda de regressão da fórmula usada pelo próprio slide: escala_max é o
    # múltiplo de 10 arredondado para cima do maior valor (piso 30).
    escala_max = max(30, math.ceil(45.0 / 10.0) * 10)  # = 50
    for largura, atual in zip((largura_10, largura_25, largura_45), (10.0, 25.0, 45.0)):
        esperado = max(Inches(0.05), int(BAR_TRACK_W * min(atual, escala_max) / escala_max))
        assert largura == esperado


def test_marcador_de_alvo_nao_ultrapassa_o_track(slide_alocacao):
    """Marcador dourado do alvo (modelo) é uma rect fina (largura 0.03in) —
    confere que nenhuma sai da área do track (bug clássico: alvo=100% jogaria
    o marcador para fora da barra)."""
    marcadores = [
        s for s in slide_alocacao.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and s.auto_shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE
        and s.width == Inches(0.03)
    ]
    assert len(marcadores) == 3
    for m in marcadores:
        assert m.left >= Inches(2.75) - 1  # bar_track_x (col_x[1])
        assert m.left <= Inches(2.75) + BAR_TRACK_W


def test_todas_as_9_classes_tem_label_mapeado():
    # Sanidade: gerar_pptx_cliente não deveria quebrar com uma classe fora de
    # CLS_LABEL (a chave "cripto" errada, por ex., cairia como rótulo cru).
    assert set(DADOS_BASE["composicao"].keys()) <= set(CLS_LABEL.keys())
