"""
pptx_brauna.py — Gerador de PPTX com identidade visual Braúna Investimentos
Usar: from api.pptx_brauna import gerar_pptx_cliente
"""

import io
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paleta de cores ────────────────────────────────────────────────────────────
# DECISÃO 17/07/2026: tema claro unificado. O PDF técnico (gerar_pdf, em
# api/index.py) e o Braúna 360° (gerar_brauna360_html / _B360_CSS, em
# api/index.py) migraram para tema claro em 13/07/2026. Este arquivo tinha
# ficado, isoladamente, em fundo escuro teal (#023B40) por uma decisão de
# 16/07/2026 ("premium para cliente") — revisada e revertida por Denis: os 3
# documentos entregues ao mesmo cliente (PDF + 360° + PPTX) precisam ter a
# MESMA identidade visual, não só os mesmos números. Os hex abaixo foram
# copiados literalmente das custom properties do _B360_CSS (fonte de verdade
# da paleta) — não são re-tonalizados como antes.
#
# Exceção proposital: a capa (_slide_capa) mantém fundo --marinho (escuro) com
# texto creme, replicando o padrão já usado no Braúna 360° para as seções
# .cover/.close — dá o toque premium sem contradizer o tema claro do miolo.
# É o ÚNICO slide com fundo escuro; os outros 6 usam --paper/--card.
#
# Notas de contraste (WCAG AA — 4.5:1 texto normal, 3:1 texto grande ≥14pt
# bold ou ≥18pt regular), conferidas à mão para cada par texto/fundo usado
# neste arquivo:
#   - Sobre fundo claro (--paper/--card): C_INK, C_MUTED, C_PETROLEO, C_GOLD
#     e os status C_GOOD/C_CRIT passam confortavelmente em qualquer tamanho
#     usado aqui. C_WARN (~3,4:1 sobre branco) só atinge o critério de "texto
#     grande" — por isso os dois lugares onde aparece como texto (coluna de
#     desvio da tabela de alocação e a linha de "Pendências" do slide Modelo
#     de Servir) usam 14pt bold. C_FAINT (~4,3:1) fica reservado só para a
#     marca d'água (14pt bold rotacionado, onde a exigência cai para 3:1);
#     rodapé/legendas/textos pequenos usam C_MUTED, que folga o 4.5:1 em
#     qualquer tamanho.
#   - Sobre fundo escuro (--marinho na capa; --petroleo nas bandas de tabela
#     e na linha "Total Renda Fixa"): SEMPRE texto creme (C_CREAM /
#     C_CREAM_SOFT / C_CREAM_MUTE). C_GOLD bruto sobre --marinho ou
#     --petroleo dá ~2,4:1 — falha até o critério de texto grande — por isso
#     título/nome/rótulos da capa e das bandas escuras nunca usam dourado ou
#     tinta, só creme. Dourado nessas áreas fica restrito a elementos
#     decorativos sem texto (ex.: a linha divisória da capa).
# C_GOOD/C_WARN/C_CRIT formam a "status palette" (ok/atenção/crítico) —
# reservada, nunca reaproveitada como cor categórica solta; todo uso de
# status neste arquivo vem acompanhado de ícone/seta/sinal, nunca só cor (ver
# _status_icon, _desvio_cor e os "+"/"▲"/"▼" nos textos — regra pré-existente,
# mantida).
C_PAPER      = RGBColor(0xF6, 0xF7, 0xF6)   # fundo de página (6 dos 7 slides)
C_CARD       = RGBColor(0xFF, 0xFF, 0xFF)   # fundo de cards/tabelas
C_INK        = RGBColor(0x17, 0x27, 0x1E)   # texto principal
C_MUTED      = RGBColor(0x4F, 0x5E, 0x66)   # texto secundário (labels, rodapé, legendas)
C_FAINT      = RGBColor(0x6E, 0x7C, 0x82)   # reservado à marca d'água — ver nota de contraste acima
C_LINE       = RGBColor(0xE4, 0xE8, 0xE7)   # bordas / track neutro de barra
C_LINE_SOFT  = RGBColor(0xEE, 0xF1, 0xF0)   # linha alternada de tabela (zebra)
C_PETROLEO   = RGBColor(0x0E, 0x3A, 0x54)   # accent primário — "Atual" em barras/gráficos, bandas de tabela
C_GOLD       = RGBColor(0x8A, 0x6A, 0x28)   # accent dourado — "Alvo", destaque, kicker (só sobre fundo claro)
C_GOLD_SOFT  = RGBColor(0xC7, 0xAC, 0x80)   # dourado suave — base do tingimento de C_ROW_HILITE
C_MARINHO    = RGBColor(0x12, 0x30, 0x47)   # navy — SÓ capa (_slide_capa)
C_CREAM      = RGBColor(0xED, 0xE7, 0xDA)   # texto sobre marinho/petróleo
C_CREAM_SOFT = RGBColor(0xF3, 0xEE, 0xE4)   # texto sobre marinho/petróleo (ênfase/título)
C_CREAM_MUTE = RGBColor(0xB9, 0xC7, 0xCD)   # texto sobre marinho/petróleo (secundário)
C_GOOD       = RGBColor(0x3E, 0x7A, 0x5E)   # status ok
C_WARN       = RGBColor(0xB4, 0x83, 0x3A)   # status atenção — usar em texto ≥14pt bold (ver nota de contraste)
C_CRIT       = RGBColor(0xA5, 0x53, 0x3E)   # status crítico
C_ROW_ALT    = C_LINE_SOFT                  # faixa alternada de linha (tabelas)
C_ROW_HILITE = RGBColor(0xF1, 0xEA, 0xDF)   # linha em destaque (maior desvio) — tingimento ~25% de C_GOLD_SOFT sobre branco

# Mapeamento classe → label legível
CLS_LABEL = {
    "pos_fixado":    "Pós Fixado",
    "inflacao":      "Inflação",
    "pre_fixado":    "Pré Fixado",
    "acoes":         "Ações",
    "fiis":          "FIIs",
    "multimercado":  "Multimercado",
    "alternativos":  "Alternativos",
    "internacional": "Internacional",
    "criptomoedas":  "Criptomoedas",
}

# ── Helpers de formatação ──────────────────────────────────────────────────────

def _format_brl(valor: float) -> str:
    """Formata valor em reais: 850000 → 'R$ 850 mil' | 1500000 → 'R$ 1,5 M'"""
    if valor >= 1_000_000:
        v = valor / 1_000_000
        if v == int(v):
            return f"R$ {int(v)} M"
        return f"R$ {v:.1f} M".replace(".", ",")
    elif valor >= 1_000:
        v = valor / 1_000
        if v == int(v):
            return f"R$ {int(v)} mil"
        return f"R$ {v:.0f} mil"
    return f"R$ {valor:,.0f}".replace(",", ".")


def _status_cor(status: str) -> RGBColor:
    """Retorna cor RGB conforme status: ok, atencao, critico."""
    m = {"ok": C_GOOD, "atencao": C_WARN, "critico": C_CRIT}
    return m.get(status, C_WARN)


def _status_label(status: str) -> str:
    m = {"ok": "OK", "atencao": "Atenção", "critico": "Crítico"}
    return m.get(status, status.title())


def _status_icon(status: str) -> str:
    """Símbolo que SEMPRE acompanha a cor de status (nunca só cor)."""
    m = {"ok": "✓", "atencao": "⚠", "critico": "✕"}
    return m.get(status, "•")


def _desvio_cor(desvio: float) -> RGBColor:
    """Verde ±2pp, âmbar ±2-5pp, vermelho >5pp."""
    abs_d = abs(desvio)
    if abs_d <= 2:
        return C_GOOD
    elif abs_d <= 5:
        return C_WARN
    return C_CRIT


def _desvio_seta(desvio: float) -> str:
    """Seta de direção do desvio — sinal visual além da cor (nunca só cor)."""
    if desvio > 0:
        return "▲"
    elif desvio < 0:
        return "▼"
    return "—"


def _data_proxima_analise(data_ref: str) -> str:
    """Retorna data 3 meses à frente da data_ref (formato DD/MM/AAAA)."""
    try:
        d = datetime.datetime.strptime(data_ref, "%d/%m/%Y")
        mes = d.month + 3
        ano = d.year + (mes - 1) // 12
        mes = ((mes - 1) % 12) + 1
        return datetime.date(ano, mes, min(d.day, 28)).strftime("%d/%m/%Y")
    except Exception:
        return "—"


# ── Primitivas de slide ────────────────────────────────────────────────────────

def _nova_prs() -> Presentation:
    """Cria apresentação 16:9 widescreen (13.33" × 7.50")."""
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.50)
    return prs


def _blank_slide(prs: Presentation):
    """Adiciona slide em branco (layout 6 = blank)."""
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def _slide_bg(slide, prs: Presentation, color: RGBColor = None):
    """Preenche o fundo do slide. Default = C_PAPER (tema claro do miolo); a
    capa (_slide_capa) é a única que passa C_MARINHO explicitamente — ver
    racional da "exceção proposital" no cabeçalho do módulo."""
    if color is None:
        color = C_PAPER
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, fill_color: RGBColor, line_color=None, radius=False):
    """Adiciona retângulo (ou arredondado) e retorna a shape."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1 if not radius else 5,  # MSO_SHAPE_TYPE.RECTANGLE=1, ROUNDED_RECT=5
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _add_text_box(slide, text, x, y, w, h,
                  font_size=14, bold=False, color=C_INK,
                  align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Adiciona caixa de texto e retorna o frame."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def _add_watermark(slide, prs: Presentation, on_dark: bool = False):
    """Adiciona '#SomosBraúna' vertical na direita + rodapé URL.

    on_dark=True (só a capa, fundo --marinho) usa creme; on_dark=False (os
    outros 6 slides, fundo --paper) usa C_FAINT para a marca d'água (14pt
    bold rotacionado — exigência de contraste WCAG cai para 3:1, ver nota no
    topo do arquivo) e C_MUTED para o rodapé (10pt, precisa do 4.5:1 cheio).
    """
    wm_color     = C_CREAM_MUTE if on_dark else C_FAINT
    footer_color = C_CREAM_MUTE if on_dark else C_MUTED

    # Marca d'água vertical
    wm = slide.shapes.add_textbox(
        Inches(12.35), Inches(0.5),
        Inches(0.9), Inches(6.5)
    )
    tf = wm.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "#SomosBraúna"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = wm_color
    run.font.name = "Calibri"
    # Rotacionar 90° via XML
    sp = wm._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is None:
        xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    xfrm.set("rot", "5400000")  # 90° em EMU units (60000 per degree)

    # Rodapé
    _add_text_box(
        slide, "braunainvestimentos.com.br",
        Inches(3.5), Inches(7.22),
        Inches(6.33), Inches(0.25),
        font_size=10, bold=True, color=footer_color,
        align=PP_ALIGN.CENTER
    )


def _add_header_band(slide, prs: Presentation, titulo: str, subtitulo: str = ""):
    """Banda de título em tema claro: card branco com borda sutil (--paper e
    --card são quase idênticos em luminância, a borda C_LINE é o que
    diferencia a banda da página), título em tinta e subtítulo em dourado
    ("kicker" — passa 4.5:1 sobre branco)."""
    # Retângulo do card do cabeçalho
    _add_rect(slide,
              Inches(0.09), Inches(0.07),
              Inches(12.2), Inches(1.1),
              C_CARD, line_color=C_LINE)
    # Título
    _add_text_box(
        slide, titulo,
        Inches(0.25), Inches(0.10),
        Inches(11.8), Inches(0.9),
        font_size=36, bold=True, color=C_INK,
        align=PP_ALIGN.LEFT
    )
    # Subtítulo
    if subtitulo:
        _add_text_box(
            slide, subtitulo,
            Inches(0.25), Inches(1.15),
            Inches(11.8), Inches(0.35),
            font_size=14, bold=True, color=C_GOLD,
            align=PP_ALIGN.LEFT
        )


# ── Slides ─────────────────────────────────────────────────────────────────────

def _slide_capa(prs: Presentation, dados: dict):
    """Slide 1 — Capa. Único slide em fundo escuro (--marinho); todo o texto
    usa tons creme (nunca dourado/tinta — ver nota de contraste no topo do
    arquivo). O dourado fica restrito à linha divisória, puramente decorativa."""
    slide = _blank_slide(prs)
    _slide_bg(slide, prs, C_MARINHO)
    _add_watermark(slide, prs, on_dark=True)

    # Nome empresa
    _add_text_box(
        slide, "BRAÚNA INVESTIMENTOS",
        Inches(1.5), Inches(1.8),
        Inches(10.3), Inches(1.0),
        font_size=48, bold=True, color=C_CREAM_SOFT,
        align=PP_ALIGN.CENTER
    )

    # Subtítulo institucional
    _add_text_box(
        slide, "Análise de Carteira de Investimentos",
        Inches(1.5), Inches(2.75),
        Inches(10.3), Inches(0.6),
        font_size=22, bold=False, color=C_CREAM,
        align=PP_ALIGN.CENTER
    )

    # Linha divisória — retângulo fino dourado (decorativo, sem texto sobreposto)
    _add_rect(slide,
              Inches(2.5), Inches(3.45),
              Inches(8.33), Inches(0.04),
              C_GOLD)

    # Nome do cliente
    nome = dados.get("nome", "Cliente")
    _add_text_box(
        slide, nome,
        Inches(1.5), Inches(3.55),
        Inches(10.3), Inches(0.75),
        font_size=30, bold=True, color=C_CREAM_SOFT,
        align=PP_ALIGN.CENTER
    )

    # Perfil | Data | Assessor
    perfil   = dados.get("perfil", "").title()
    data_ref = dados.get("data_ref", "")
    assessor = dados.get("assessor", "")
    info_line = f"Perfil: {perfil}   |   Ref.: {data_ref}   |   Assessor: {assessor}"
    _add_text_box(
        slide, info_line,
        Inches(1.0), Inches(4.4),
        Inches(11.0), Inches(0.5),
        font_size=14, color=C_CREAM_MUTE,
        align=PP_ALIGN.CENTER
    )


def _slide_resumo(prs: Presentation, dados: dict):
    """Slide 2 — Resumo Executivo com 4 KPI cards."""
    slide = _blank_slide(prs)
    _slide_bg(slide, prs)
    _add_watermark(slide, prs)
    _add_header_band(slide, prs, "Resumo Executivo")

    rent    = dados.get("rent", {})
    status  = dados.get("status", "atencao")
    score   = dados.get("score_servir", 0)
    pat     = dados.get("patrimonio", 0)
    r12m    = rent.get("12m", 0)
    cdi_pct = rent.get("cdi_pct", 0)

    cards = [
        {
            "label": "Patrimônio Total",
            "valor": _format_brl(pat),
            "sub":   "",
            "cor":   C_GOLD,
        },
        {
            # Delta sempre com seta + sinal explícito — nunca só a cor (▲/▼ +/-).
            "label": "Rentabilidade 12m",
            "valor": (
                f"{'▲' if r12m > 0 else '▼' if r12m < 0 else '●'} "
                f"{'+' if r12m > 0 else ''}{r12m:.2f}%".replace(".", ",")
            ),
            "sub":   f"{cdi_pct:.1f}% do CDI".replace(".", ","),
            "cor":   C_GOOD if r12m >= 0 else C_CRIT,
        },
        {
            "label": "Status da Carteira",
            "valor": f"{_status_icon(status)} {_status_label(status)}",
            "sub":   "",
            "cor":   _status_cor(status),
        },
        {
            "label": "Modelo de Servir",
            "valor": f"{score:.1f} / 5,0".replace(".", ","),
            "sub":   "Score geral",
            "cor":   C_PETROLEO,
        },
    ]

    card_w = Inches(2.8)
    card_h = Inches(2.4)
    gap    = Inches(0.22)
    start_x = Inches(0.55)
    start_y = Inches(1.65)

    for i, c in enumerate(cards):
        x = start_x + i * (card_w + gap)

        # Fundo do card (branco sobre --paper; borda sutil para diferenciar
        # as duas superfícies, que têm luminância quase idêntica).
        rect = _add_rect(slide, x, start_y, card_w, card_h, C_CARD, line_color=C_LINE, radius=False)
        # Borda de cor na borda superior (retângulo fino)
        _add_rect(slide, x, start_y, card_w, Inches(0.06), c["cor"])

        # Label
        _add_text_box(
            slide, c["label"],
            x + Inches(0.15), start_y + Inches(0.15),
            card_w - Inches(0.3), Inches(0.4),
            font_size=13, bold=False, color=C_MUTED,
            align=PP_ALIGN.LEFT
        )
        # Valor principal
        _add_text_box(
            slide, c["valor"],
            x + Inches(0.1), start_y + Inches(0.65),
            card_w - Inches(0.2), Inches(0.95),
            font_size=28, bold=True, color=c["cor"],
            align=PP_ALIGN.LEFT
        )
        # Sub-label
        if c["sub"]:
            _add_text_box(
                slide, c["sub"],
                x + Inches(0.15), start_y + Inches(1.65),
                card_w - Inches(0.3), Inches(0.4),
                font_size=12, color=C_MUTED,
                align=PP_ALIGN.LEFT
            )

    # Objetivo
    objetivo = dados.get("objetivo", "")
    if objetivo:
        _add_text_box(
            slide, f"Objetivo: {objetivo}",
            Inches(0.55), Inches(4.25),
            Inches(11.5), Inches(0.4),
            font_size=13, color=C_MUTED,
            align=PP_ALIGN.LEFT
        )

    # Rentabilidades adicionais
    rmes = rent.get("mes", 0)
    rano = rent.get("ano", 0)
    r24m = rent.get("24m", 0)
    perf_txt = (
        f"Mês: {rmes:.2f}%    "
        f"Ano: {rano:.2f}%    "
        f"12m: {r12m:.2f}%    "
        f"24m: {r24m:.2f}%"
    ).replace(".", ",")
    _add_text_box(
        slide, perf_txt,
        Inches(0.55), Inches(4.7),
        Inches(11.5), Inches(0.4),
        font_size=13, bold=True, color=C_INK,
        align=PP_ALIGN.LEFT
    )

    # Alertas
    alertas = dados.get("alertas", [])
    if alertas:
        alertas_txt = [a.get("mensagem", "") if isinstance(a, dict) else str(a) for a in alertas]
        alerta_str = "  ⚠  " + "   |   ".join(t for t in alertas_txt if t)
        _add_text_box(
            slide, alerta_str,
            Inches(0.55), Inches(5.3),
            Inches(11.5), Inches(0.5),
            font_size=12, bold=True, color=C_CRIT,
            align=PP_ALIGN.LEFT
        )


def _slide_alocacao(prs: Presentation, dados: dict):
    """Slide 3 — Alocação Atual vs Modelo.

    A comparação atual-vs-alvo ganhou um encoding visual direto: uma barra fina
    (mark spec da skill dataviz — ≤0.16in de espessura, pontas arredondadas,
    cresce de uma baseline única em 0%) mostra o % atual sobre um track neutro,
    com um marcador vertical dourado sobreposto na posição do modelo (alvo).
    Isso comunica a distância atual↔alvo geometricamente, não só pela cor do
    texto de desvio — a cor (verde/âmbar/vermelho) permanece, mas sempre
    acompanhada da seta ▲/▼ e do sinal +/-, nunca sozinha (ver _desvio_seta).
    Optei por desenho manual (não native add_chart do python-pptx) porque o
    layout é uma barra única por linha embutida numa tabela de 5 colunas — um
    gráfico nativo exigiria um plot separado e perderia o alinhamento por linha
    com Atual%/Modelo%/Desvio; o desenho manual também é o que já dá controle
    fino sobre o marcador de alvo sobreposto, que add_chart não expõe.

    Paleta: coluna "Atual %" em C_PETROLEO e "Modelo %" em C_GOLD — mapeamento
    direto dos papéis definidos na paleta (petróleo = Atual, dourado = Alvo).
    A coluna de Desvio usa 14pt (em vez de 13pt) para que o texto em C_WARN
    (~3,4:1 sobre branco) se qualifique como "texto grande" no critério WCAG
    AA — ver nota de contraste no topo do arquivo.
    """
    import math

    slide = _blank_slide(prs)
    _slide_bg(slide, prs)
    _add_watermark(slide, prs)
    _add_header_band(slide, prs, "Alocação de Carteira", "Atual vs. Modelo de Referência")

    composicao = dados.get("composicao", {})
    modelo     = dados.get("modelo", {})
    classes    = list(CLS_LABEL.keys())

    # Cabeçalho da tabela — coluna 2 vira a área da barra comparativa.
    headers  = ["Classe de Ativo", "Atual vs. Modelo", "Atual %", "Modelo %", "Desvio"]
    col_x    = [Inches(0.35), Inches(2.75), Inches(8.05), Inches(9.30), Inches(10.55)]
    col_w    = [Inches(2.30), Inches(5.20), Inches(1.15), Inches(1.15), Inches(1.50)]
    row_h    = Inches(0.38)
    start_y  = Inches(1.55)

    # Linha de cabeçalho
    _add_rect(slide,
              Inches(0.25), start_y,
              Inches(11.8), row_h,
              C_PETROLEO)
    for j, h in enumerate(headers):
        _add_text_box(
            slide, h,
            col_x[j], start_y + Inches(0.05),
            col_w[j], row_h - Inches(0.05),
            font_size=13, bold=True, color=C_CREAM_SOFT,
            align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        )

    # Linhas de dados
    maior_desvio_idx = -1
    maior_desvio_abs = -1
    rows_data = []
    for cls in classes:
        atual  = composicao.get(cls, 0)
        alvo   = modelo.get(cls, 0)
        if atual == 0 and alvo == 0:
            continue
        desvio = round(atual - alvo, 1)
        rows_data.append((cls, atual, alvo, desvio))
        if abs(desvio) > maior_desvio_abs:
            maior_desvio_abs = abs(desvio)
            maior_desvio_idx = len(rows_data) - 1

    # Escala comum da barra (uma única "baseline" em 0% para todas as linhas),
    # arredondada para um número "limpo" acima do maior valor, com piso de 30%
    # para não esticar demais as barras quando a carteira é bem diversificada.
    maior_valor = max(
        [max(atual, alvo) for _, atual, alvo, _ in rows_data] + [0]
    )
    escala_max = max(30, math.ceil(maior_valor / 10.0) * 10)

    bar_track_x = col_x[1]
    bar_track_w = col_w[1]
    bar_h       = Inches(0.16)

    for i, (cls, atual, alvo, desvio) in enumerate(rows_data):
        y = start_y + (i + 1) * row_h
        # Fundo alternado, destaque na linha com maior desvio
        if i == maior_desvio_idx:
            bg = C_ROW_HILITE
        elif i % 2 == 0:
            bg = C_ROW_ALT
        else:
            bg = C_CARD
        _add_rect(slide, Inches(0.25), y, Inches(11.8), row_h, bg)

        desvio_cor = _desvio_cor(desvio)
        seta       = _desvio_seta(desvio)
        sinal      = "+" if desvio > 0 else ""
        desvio_str = f"{seta} {sinal}{desvio:.1f} pp".replace(".", ",") if desvio != 0 else "—"

        # ── Barra comparativa: track neutro + barra do atual + marcador do alvo ──
        bar_y = y + (row_h - bar_h) // 2
        _add_rect(slide, bar_track_x, bar_y, bar_track_w, bar_h, C_LINE, radius=True)
        if atual > 0:
            largura_atual = max(Inches(0.05), int(bar_track_w * min(atual, escala_max) / escala_max))
            _add_rect(slide, bar_track_x, bar_y, largura_atual, bar_h, desvio_cor, radius=True)
        # Marcador de alvo (dourado) sobreposto na posição do modelo — encoding
        # direto da comparação atual↔alvo, além da cor do texto de desvio.
        # Puramente decorativo (sem texto sobreposto), então o dourado bruto
        # não tem restrição de contraste de texto aqui.
        marker_x = bar_track_x + int(bar_track_w * min(alvo, escala_max) / escala_max)
        marker_x = min(marker_x, bar_track_x + bar_track_w - Inches(0.02))
        _add_rect(
            slide, marker_x, bar_y - Inches(0.05),
            Inches(0.03), bar_h + Inches(0.10),
            C_GOLD
        )

        vals = [
            CLS_LABEL.get(cls, cls),
            "",  # coluna 2 é a barra desenhada acima, sem texto
            f"{atual:.1f}%".replace(".", ","),
            f"{alvo:.1f}%".replace(".", ","),
            desvio_str,
        ]
        cores = [C_INK, C_INK, C_PETROLEO, C_GOLD, desvio_cor]
        bolds = [False, False, False, False, True]

        for j, (val, cor, bold) in enumerate(zip(vals, cores, bolds)):
            if j == 1:
                continue  # coluna da barra: nenhum texto sobreposto
            fsize = 14 if j == 4 else 13  # coluna de desvio em 14pt (ver docstring)
            _add_text_box(
                slide, val,
                col_x[j], y + Inches(0.04),
                col_w[j], row_h - Inches(0.06),
                font_size=fsize, bold=bold, color=cor,
                align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            )

    # Legenda
    _add_text_box(
        slide,
        "Barra = % atual  |  Marcador ouro = % modelo (alvo)  |  "
        "Verde ±2pp · Âmbar ±2–5pp · Vermelho >5pp  (▲ acima / ▼ abaixo)",
        Inches(0.35), Inches(6.72),
        Inches(11.5), Inches(0.4),
        font_size=10, color=C_MUTED,
        align=PP_ALIGN.LEFT
    )


def _slide_rv(prs: Presentation, dados: dict):
    """Slide 4 — Renda Variável & FIIs."""
    slide = _blank_slide(prs)
    _slide_bg(slide, prs)
    _add_watermark(slide, prs)
    _add_header_band(slide, prs, "Renda Variável & FIIs")

    acoes = sorted(dados.get("acoes", []), key=lambda x: x.get("pct", 0), reverse=True)
    fiis  = sorted(dados.get("fiis",  []), key=lambda x: x.get("pct", 0), reverse=True)

    if not acoes and not fiis:
        _add_text_box(
            slide, "Sem exposição em Renda Variável nesta carteira.",
            Inches(0.5), Inches(3.0),
            Inches(11.5), Inches(0.6),
            font_size=18, color=C_MUTED,
            align=PP_ALIGN.CENTER
        )
        return

    y = Inches(1.45)
    row_h = Inches(0.42)

    def _bloco(titulo, itens, y_start):
        # Subtítulo do bloco
        _add_text_box(
            slide, titulo,
            Inches(0.4), y_start,
            Inches(11.0), Inches(0.35),
            font_size=14, bold=True, color=C_GOLD,
            align=PP_ALIGN.LEFT
        )
        y_cur = y_start + Inches(0.38)
        for idx, item in enumerate(itens):
            ticker = item.get("ticker", "")
            nome   = item.get("nome", "")
            pct    = item.get("pct", 0)
            bg = C_CARD if idx % 2 == 0 else C_ROW_ALT
            _add_rect(slide, Inches(0.25), y_cur, Inches(11.8), row_h, bg)
            # Ticker
            _add_text_box(
                slide, ticker,
                Inches(0.4), y_cur + Inches(0.04),
                Inches(1.5), row_h,
                font_size=14, bold=True, color=C_GOLD
            )
            # Nome
            _add_text_box(
                slide, nome,
                Inches(2.0), y_cur + Inches(0.04),
                Inches(7.5), row_h,
                font_size=13, color=C_INK
            )
            # Percentual
            _add_text_box(
                slide, f"{pct:.1f}%".replace(".", ","),
                Inches(9.7), y_cur + Inches(0.04),
                Inches(1.8), row_h,
                font_size=13, bold=True, color=C_PETROLEO,
                align=PP_ALIGN.RIGHT
            )
            y_cur += row_h
        return y_cur + Inches(0.15)

    if acoes:
        y = _bloco("Ações", acoes, y)
    if fiis:
        y = _bloco("Fundos de Investimento Imobiliário (FIIs)", fiis, y)


def _slide_rf(prs: Presentation, dados: dict):
    """Slide 5 — Renda Fixa."""
    slide = _blank_slide(prs)
    _slide_bg(slide, prs)
    _add_watermark(slide, prs)
    _add_header_band(slide, prs, "Renda Fixa", "Principais posições")

    rf_ativos = dados.get("rf_ativos", [])
    MAX_ITEMS = 8
    exibir    = rf_ativos[:MAX_ITEMS]
    restantes = len(rf_ativos) - MAX_ITEMS

    row_h   = Inches(0.42)
    start_y = Inches(1.55)

    # Cabeçalho
    _add_rect(slide, Inches(0.25), start_y, Inches(11.8), row_h, C_PETROLEO)
    for txt, x, w, align in [
        ("Ativo", Inches(0.4),  Inches(7.0), PP_ALIGN.LEFT),
        ("Saldo", Inches(7.5),  Inches(2.3), PP_ALIGN.RIGHT),
        ("Classe", Inches(9.9), Inches(1.9), PP_ALIGN.CENTER),
    ]:
        _add_text_box(
            slide, txt,
            x, start_y + Inches(0.04),
            w, row_h,
            font_size=13, bold=True, color=C_CREAM_SOFT, align=align
        )

    total_rf = 0
    for i, ativo in enumerate(exibir):
        y   = start_y + (i + 1) * row_h
        bg  = C_CARD if i % 2 == 0 else C_ROW_ALT
        _add_rect(slide, Inches(0.25), y, Inches(11.8), row_h, bg)

        nome   = ativo.get("nome", "")
        saldo  = ativo.get("saldo", 0)
        classe = ativo.get("classe", "")
        total_rf += saldo

        _add_text_box(
            slide, nome,
            Inches(0.4), y + Inches(0.04),
            Inches(7.0), row_h,
            font_size=13, color=C_INK
        )
        _add_text_box(
            slide, _format_brl(saldo),
            Inches(7.5), y + Inches(0.04),
            Inches(2.3), row_h,
            font_size=13, bold=True, color=C_GOLD,
            align=PP_ALIGN.RIGHT
        )
        _add_text_box(
            slide, CLS_LABEL.get(classe, classe),
            Inches(9.9), y + Inches(0.04),
            Inches(1.9), row_h,
            font_size=12, color=C_MUTED,
            align=PP_ALIGN.CENTER
        )

    # "e outros X ativos"
    extra_y = start_y + (len(exibir) + 1) * row_h + Inches(0.1)
    if restantes > 0:
        _add_text_box(
            slide, f"...e outros {restantes} ativos não listados",
            Inches(0.4), extra_y,
            Inches(8.0), Inches(0.35),
            font_size=12, color=C_MUTED
        )
        extra_y += Inches(0.4)

    # Total RF — linha em C_PETROLEO: texto SEMPRE creme aqui (dourado/tinta
    # falham o contraste sobre esse fundo escuro, ver nota no topo do arquivo).
    if total_rf > 0:
        _add_rect(slide, Inches(0.25), extra_y, Inches(11.8), Inches(0.44), C_PETROLEO)
        _add_text_box(
            slide, "Total Renda Fixa listada:",
            Inches(0.4), extra_y + Inches(0.04),
            Inches(7.5), Inches(0.38),
            font_size=13, bold=True, color=C_CREAM_SOFT
        )
        _add_text_box(
            slide, _format_brl(total_rf),
            Inches(7.5), extra_y + Inches(0.04),
            Inches(4.4), Inches(0.38),
            font_size=13, bold=True, color=C_CREAM_SOFT,
            align=PP_ALIGN.RIGHT
        )


def _slide_servir(prs: Presentation, dados: dict):
    """Slide 6 — Modelo de Servir."""
    slide = _blank_slide(prs)
    _slide_bg(slide, prs)
    _add_watermark(slide, prs)
    score = dados.get("score_servir", 0)
    score_str = f"{score:.1f}".replace(".", ",")
    _add_header_band(slide, prs, f"Modelo de Servir  ·  Score {score_str}/5,0")

    checklist = dados.get("checklist", {})

    # Mapeamento nome → label legível
    ITEM_LABELS = {
        "contato_ativo":         "Contato ativo (últimos 30 dias)",
        "open_investments":      "Open Investments habilitado",
        "previdencia":           "Previdência estruturada",
        "seguro_vida":           "Seguro de vida contratado",
        "planejamento_sucesso":  "Planejamento sucessório iniciado",
        "declaracao_ir":         "Declaração IR assessorada",
        "consorcio":             "Consórcio avaliado",
        "cambio":                "Câmbio / remessa avaliado",
        "credito_estruturado":   "Crédito estruturado avaliado",
        "rebalanceamento":       "Rebalanceamento em dia",
    }

    items = list(checklist.items())
    col_count = 2
    row_h     = Inches(0.47)
    col_w     = Inches(5.8)
    start_y   = Inches(1.6)
    start_x   = [Inches(0.35), Inches(6.35)]

    pendentes_criticos = []

    for idx, (key, val) in enumerate(items):
        col = idx % col_count
        row = idx // col_count
        x   = start_x[col]
        y   = start_y + row * row_h

        label  = ITEM_LABELS.get(key, key.replace("_", " ").title())
        ok     = bool(val)
        icone  = "✓" if ok else "✕"
        cor    = C_GOOD if ok else C_CRIT

        bg = C_CARD if row % 2 == 0 else C_ROW_ALT
        _add_rect(slide, x, y, col_w, row_h - Inches(0.04), bg)

        _add_text_box(
            slide, f"{icone}  {label}",
            x + Inches(0.1), y + Inches(0.04),
            col_w - Inches(0.2), row_h - Inches(0.1),
            font_size=13, bold=not ok, color=cor
        )

        if not ok:
            pendentes_criticos.append(label)

    # Recomendação se houver pendências — 14pt para C_WARN se qualificar como
    # "texto grande" no critério de contraste WCAG AA (ver nota no topo do arquivo).
    if pendentes_criticos:
        recom_y = start_y + ((len(items) + 1) // col_count + 1) * row_h
        recom_y = max(recom_y, Inches(5.85))
        _add_text_box(
            slide,
            f"Pendências: ativar — {', '.join(pendentes_criticos[:3])}",
            Inches(0.35), recom_y,
            Inches(11.8), Inches(0.6),
            font_size=14, bold=True, color=C_WARN
        )

    # Cross-sell
    cross = dados.get("cross_ativos", [])
    if cross:
        cross_str = "Cross-sell ativo: " + ", ".join(
            c.replace("_", " ").title() for c in cross
        )
        _add_text_box(
            slide, cross_str,
            Inches(0.35), Inches(6.55),
            Inches(11.5), Inches(0.38),
            font_size=11, color=C_MUTED
        )


def _slide_proximos_passos(prs: Presentation, dados: dict):
    """Slide 7 — Próximos Passos."""
    slide = _blank_slide(prs)
    _slide_bg(slide, prs)
    _add_watermark(slide, prs)
    _add_header_band(slide, prs, "Próximos Passos")

    acoes = []

    # 1. Desvios de alocação
    desvios = dados.get("desvios", [])
    for d in sorted(desvios, key=lambda x: abs(x.get("desvio", 0)), reverse=True):
        dev = d.get("desvio", 0)
        if abs(dev) > 2:
            label  = d.get("label", "")
            verbo  = "Aumentar" if dev < 0 else "Reduzir"
            pp_str = f"{abs(dev):.1f} pp".replace(".", ",")
            acoes.append(f"{verbo} {label}: {pp_str} fora do modelo")
        if len(acoes) >= 3:
            break

    # 2. Checklist pendente
    checklist = dados.get("checklist", {})
    ITEM_LABELS = {
        "contato_ativo":         "Contato ativo",
        "open_investments":      "Open Investments",
        "previdencia":           "Previdência",
        "seguro_vida":           "Seguro de vida",
        "planejamento_sucesso":  "Planejamento sucessório",
        "declaracao_ir":         "Declaração IR",
        "consorcio":             "Consórcio",
        "cambio":                "Câmbio",
        "credito_estruturado":   "Crédito estruturado",
        "rebalanceamento":       "Rebalanceamento",
    }
    for key, val in checklist.items():
        if not val and len(acoes) < 5:
            acoes.append(f"Ativar: {ITEM_LABELS.get(key, key)}")

    # 3. Cross-sell oportunidades
    cross_ativos = dados.get("cross_ativos", [])
    cross_todos  = ["previdencia", "consorcio", "seguro_vida", "cambio"]
    for item in cross_todos:
        if item not in cross_ativos and len(acoes) < 6:
            acoes.append(f"Oportunidade cross-sell: {item.replace('_', ' ').title()}")

    # 4. Nota do líder
    nota = dados.get("nota_lider", "")
    if nota and len(acoes) < 6:
        acoes.append(f"Nota: {nota}")

    # Renderizar lista numerada
    row_h   = Inches(0.55)
    start_y = Inches(1.55)
    for i, acao in enumerate(acoes[:6]):
        y  = start_y + i * row_h
        bg = C_CARD if i % 2 == 0 else C_ROW_ALT
        _add_rect(slide, Inches(0.25), y, Inches(11.8), row_h - Inches(0.04), bg)
        # Número
        _add_text_box(
            slide, str(i + 1),
            Inches(0.4), y + Inches(0.06),
            Inches(0.5), row_h,
            font_size=18, bold=True, color=C_GOLD,
            align=PP_ALIGN.CENTER
        )
        # Texto
        _add_text_box(
            slide, acao,
            Inches(1.0), y + Inches(0.06),
            Inches(10.8), row_h,
            font_size=14, color=C_INK
        )

    # Data próxima análise
    data_ref     = dados.get("data_ref", "")
    proxima_data = _data_proxima_analise(data_ref)
    _add_text_box(
        slide,
        f"Próxima análise sugerida: {proxima_data}",
        Inches(0.35), Inches(6.85),
        Inches(11.5), Inches(0.4),
        font_size=13, bold=True, color=C_GOLD,
        align=PP_ALIGN.RIGHT
    )


# ── Função principal ───────────────────────────────────────────────────────────

def gerar_pptx_cliente(dados: dict) -> bytes:
    """
    Gera um PPTX com visual Braúna Investimentos (7 slides).

    Parâmetros: dict com chaves nome, assessor, perfil, data_ref, patrimonio,
    rent, composicao, modelo, desvios, acoes, fiis, rf_ativos, checklist,
    cross_ativos, score_servir, status, objetivo, alertas, nota_lider.

    Retorna: bytes do arquivo .pptx prontos para send_file().
    """
    prs = _nova_prs()

    _slide_capa(prs, dados)
    _slide_resumo(prs, dados)
    _slide_alocacao(prs, dados)
    _slide_rv(prs, dados)
    _slide_rf(prs, dados)
    _slide_servir(prs, dados)
    _slide_proximos_passos(prs, dados)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── Teste local ────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    dados_exemplo = {
        "nome":      "João Paulo Machado",
        "assessor":  "Denis Kinoshita",
        "perfil":    "moderada",
        "data_ref":  "29/05/2025",
        "patrimonio": 850_000,
        "rent": {
            "mes":     0.72,
            "ano":     4.10,
            "12m":     11.20,
            "24m":     22.80,
            "cdi_pct": 75.9,
        },
        "composicao": {
            "pos_fixado":    42,
            "inflacao":      20,
            "pre_fixado":    9,
            "acoes":         8,
            "fiis":          3,
            "multimercado":  5,
            "internacional": 11,
            "alternativos":  0,
            "criptomoedas":  0,
        },
        "modelo": {
            "pos_fixado":    44,
            "inflacao":      23,
            "pre_fixado":    10,
            "acoes":         5,
            "fiis":          1.5,
            "multimercado":  6,
            "internacional": 9,
            "alternativos":  0,
            "criptomoedas":  0,
        },
        "desvios": [
            {"label": "Inflação",      "real": 20, "alvo": 23, "desvio": -3},
            {"label": "Ações",         "real": 8,  "alvo": 5,  "desvio":  3},
            {"label": "Internacional", "real": 11, "alvo": 9,  "desvio":  2},
            {"label": "Pós Fixado",    "real": 42, "alvo": 44, "desvio": -2},
        ],
        "acoes": [
            {"ticker": "BBAS3", "nome": "Banco do Brasil",     "pct": 4.2},
            {"ticker": "ITUB4", "nome": "Itaú Unibanco",       "pct": 2.1},
            {"ticker": "VALE3", "nome": "Vale S.A.",            "pct": 1.7},
        ],
        "fiis": [
            {"ticker": "XPML11", "nome": "XP Malls",           "pct": 1.8},
            {"ticker": "HGLG11", "nome": "CSHG Logística",     "pct": 1.2},
        ],
        "rf_ativos": [
            {"nome": "Tesouro IPCA+ 2035",   "saldo": 95_000,  "classe": "inflacao"},
            {"nome": "LCI Banco Inter 97% CDI", "saldo": 75_000, "classe": "pos_fixado"},
            {"nome": "CDB Nubank 110% CDI",  "saldo": 60_000,  "classe": "pos_fixado"},
            {"nome": "Tesouro Selic 2027",   "saldo": 50_000,  "classe": "pos_fixado"},
            {"nome": "Debênture ENGIE 2029", "saldo": 45_000,  "classe": "pre_fixado"},
            {"nome": "CRA Raízen 2028",      "saldo": 40_000,  "classe": "inflacao"},
        ],
        "checklist": {
            "contato_ativo":        True,
            "open_investments":     False,
            "previdencia":          False,
            "seguro_vida":          True,
            "planejamento_sucesso": False,
            "declaracao_ir":        True,
            "rebalanceamento":      True,
        },
        "cross_ativos":  ["previdencia", "consorcio"],
        "score_servir":  3.5,
        "status":        "atencao",
        "objetivo":      "Crescimento patrimonial com proteção",
        "alertas":       ["Inflação abaixo do modelo por 3 meses"],
        "nota_lider":    "Agendar reunião semestral de revisão",
    }

    resultado = gerar_pptx_cliente(dados_exemplo)

    destino = "./teste_brauna.pptx"
    with open(destino, "wb") as f:
        f.write(resultado)

    print(f"Gerado com sucesso: {destino}  ({len(resultado):,} bytes)")
