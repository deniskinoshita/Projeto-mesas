import io, re, json, os, uuid
from datetime import datetime
from flask import Flask, request, jsonify, send_file, render_template_string

app = Flask(__name__)

MODELOS = {
    "conservadora": {"pos_fixado":70,"inflacao":16,"pre_fixado":7,"acoes":0,"fiis":0,"multimercado":3,"internacional":4,"alternativos":0,"criptomoedas":0},
    "moderada":     {"pos_fixado":44,"inflacao":23,"pre_fixado":10,"acoes":5,"fiis":1.5,"multimercado":6,"internacional":9,"alternativos":1,"criptomoedas":0.5},
    "arrojada":     {"pos_fixado":28,"inflacao":28,"pre_fixado":12,"acoes":8,"fiis":2.5,"multimercado":9.5,"internacional":10.25,"alternativos":1,"criptomoedas":0.75},
    "agressiva":    {"pos_fixado":13,"inflacao":31,"pre_fixado":13,"acoes":14,"fiis":3.5,"multimercado":10.5,"internacional":13,"alternativos":1,"criptomoedas":1},
}
LABELS = {"pos_fixado":"Pós Fixado","inflacao":"Inflação","pre_fixado":"Pré Fixado","acoes":"Ações","fiis":"FIIs","multimercado":"Multimercado","internacional":"Internacional","alternativos":"Alternativos","criptomoedas":"Criptomoedas"}
MAPEAMENTO = {"pós fixado":"pos_fixado","pos fixado":"pos_fixado","pós-fixado":"pos_fixado","inflação":"inflacao","inflacao":"inflacao","ipca":"inflacao","pré fixado":"pre_fixado","pre fixado":"pre_fixado","pré-fixado":"pre_fixado","ações":"acoes","acoes":"acoes","renda variável":"acoes","rv brasil":"acoes","renda variável brasil":"acoes","fiis":"fiis","fii":"fiis","fundos imobiliários":"fiis","multimercado":"multimercado","multi mercado":"multimercado","internacional":"internacional","dólar":"internacional","exterior":"internacional","renda fixa global":"internacional","criptomoedas":"criptomoedas","cripto":"criptomoedas","alternativo":"alternativos","alternativos":"alternativos","outros":"alternativos","caixa":"caixa"}
CATS = list(LABELS.keys())

MODELO_SERVIR = [
    {
        "id": "open_investments",
        "nome": "Open Investments",
        "icone": "🔗",
        "descricao": "Consentimento OPIN ativo — visão 360° do patrimônio total do cliente (XP + outras instituições)",
        "importancia": "CRÍTICA",
        "impacto_falta": "Assessor enxerga apenas a parte do cliente na XP. A média XP é 37% do patrimônio total nos clientes acima de R$ 300k — significa que 63% do patrimônio está invisível. Sem OPIN, qualquer recomendação de alocação é incompleta e pode até duplicar posições.",
        "acao": "Solicitar consentimento OPIN pelo app XP. Pitch: 'Para que eu possa cuidar do seu patrimônio como um todo e não apenas da parte na XP, preciso que você libere o Open Investments. É tão simples quanto um Pix — 3 cliques no app. Assim consigo evitar duplicar posições e encontrar oportunidades que hoje estão invisíveis.'",
        "cor": "#FF6B6B"
    },
    {
        "id": "financial_planning",
        "nome": "Financial Planning",
        "icone": "🎯",
        "descricao": "Planejamento financeiro completo realizado — objetivos, metas e destino financeiro do cliente mapeados",
        "importancia": "CRÍTICA — 1ª DIRETRIZ BRAÚNA",
        "impacto_falta": "SEM FINANCIAL PLANNING, NÃO HÁ DESTINO. Toda alocação é arbitrária quando não se sabe para onde o cliente quer ir. Não sabemos: quanto ele precisa para se aposentar, quando quer comprar um imóvel, qual o horizonte de cada objetivo. A carteira existe para realizar sonhos — sem o FP, estamos construindo uma estrada sem saber o destino.",
        "acao": "DIRETRIZ QUALIDADE BRAÚNA: Provocar o cliente agora. Perguntar diretamente: 'Você sabe exatamente quanto precisa acumular para se aposentar com o padrão de vida que deseja? Você tem um plano financeiro que traga clareza sobre cada objetivo da sua vida?' O Financial Planning é o primeiro passo antes de qualquer outra conversa de investimento.",
        "cor": "#FF4444",
        "diretriz": True
    },
    {
        "id": "ordem_enviada",
        "nome": "Ordem Enviada",
        "icone": "📋",
        "descricao": "Última ordem/movimentação executada — cliente está ativo e a carteira está evoluindo",
        "importancia": "ALTA",
        "impacto_falta": "Carteira estagnada. Cliente sem aportes recentes é um dos critérios de ruptura no Índice de Saúde XP (2 pontos). Uma carteira sem movimentação perde oportunidades de realocação e demonstra falta de atenção ativa do assessor.",
        "acao": "Verificar última ordem executada. Se não houver movimentação recente: identificar oportunidade de aporte ou realocação compatível com o perfil e gerar uma proposta concreta com base nos desvios do modelo.",
        "cor": "#FFD966"
    },
    {
        "id": "conta_acessada",
        "nome": "Conta Acessada",
        "icone": "📱",
        "descricao": "Cliente acessou a plataforma/app XP nos últimos 30 dias — está engajado com seus investimentos",
        "importancia": "ALTA",
        "impacto_falta": "Cliente desengajado da plataforma tem maior risco de ruptura e portabilidade. Clientes que não acessam a conta tendem a estar mais suscetíveis a abordagens de concorrentes. No Índice de Saúde XP, 'Sem Ordens' vale 2 pontos de risco de ruptura.",
        "acao": "Entrar em contato ativo para reengajar o cliente. Enviar uma análise personalizada (este relatório!) como pretexto de valor. Convidar para reunião de revisão de carteira. O engajamento do cliente começa pelo engajamento do assessor.",
        "cor": "#FFD966"
    },
    {
        "id": "xperformance",
        "nome": "X-Performance",
        "icone": "📊",
        "descricao": "Relatório XPerformance analisado e discutido com o cliente — base para conversa de alocação e rentabilidade",
        "importancia": "MÉDIA",
        "impacto_falta": "Sem análise do XPerformance, o cliente não tem visão clara da rentabilidade vs. benchmark (CDI). É impossível ter uma conversa fundamentada sobre realocação sem este diagnóstico. O XPerformance é o ponto de partida para toda revisão de carteira.",
        "acao": "Analisar o XPerformance atual e agendar uma reunião de revisão. Mostrar ao cliente como sua carteira se compara ao CDI e ao modelo ideal para o seu perfil — use este relatório como base da conversa.",
        "cor": "#888888"
    },
    {
        "id": "atividade_relacionamento",
        "nome": "Atividade de Relacionamento",
        "icone": "🤝",
        "descricao": "Contato ativo nos últimos 30 dias — ligação, reunião, WhatsApp com conteúdo de valor",
        "importancia": "ALTA",
        "impacto_falta": "Cliente silencioso = cliente em risco. A falta de relacionamento ativo é o principal precursor de portabilidade. Concorrentes que abordam o cliente com insights e atenção têm vantagem sobre um assessor ausente. A régua de relacionamento é o que mantém o cliente fiel mesmo em períodos de rentabilidade ruim.",
        "acao": "Entrar em contato hoje com um ponto de valor: análise de carteira, insight de mercado, ou simplesmente verificar se o cliente tem alguma necessidade. Use este relatório como motivo de contato: 'Fiz uma análise completa da sua carteira e gostaria de conversar sobre algumas oportunidades.'",
        "cor": "#FFD966"
    },
]

CROSS_SELL = [
    {
        "id": "aquisicao_bens",
        "nome": "Aquisição de Bens",
        "icone": "🏠",
        "descricao": "Financiamento imobiliário, consórcio ou crédito para aquisição de imóveis e veículos",
        "pitch": "Você tem planos de adquirir algum imóvel ou bem de alto valor nos próximos anos? Posso apresentar as melhores condições de crédito e consórcio disponíveis — muitas vezes com taxas melhores que o mercado tradicional.",
        "oportunidades": [
            "Consórcio de imóveis: sem juros, ideal para quem não tem pressa. Taxa de administração de 10-15% vs. juros de financiamento de 10-12% a.a.",
            "Crédito com garantia de investimentos (home equity): taxas a partir de CDI+1% a.a. — muito abaixo do crédito pessoal.",
            "Financiamento imobiliário: analisar SAC vs. PRICE, impacto de amortizações extras e FGTS.",
            "Consórcio de veículos: saída inteligente para quem planeja trocar o carro em 2-3 anos sem pagar juros."
        ]
    },
    {
        "id": "gestao_discricionaria",
        "nome": "Gestão Discricionária Prunus",
        "icone": "🌿",
        "descricao": "Carteira gerida ativamente pela Prunus Asset com mandato discricionário — assessor delega a gestão",
        "pitch": "Já considerou ter uma gestão profissional e ativa da sua carteira? A Prunus Asset opera com mandato discricionário — você define o perfil e os gestores tomam as decisões do dia a dia, com total transparência e relatórios periódicos.",
        "oportunidades": [
            "Gestão ativa por profissionais especializados: acesso ao mesmo nível de gestão de clientes private/institutional.",
            "Sem necessidade de acompanhamento diário pelo cliente: ideal para quem tem agenda cheia e não quer se preocupar com realocações.",
            "Relatórios detalhados de performance e atribuição: total transparência sobre o que foi feito e por quê.",
            "Personalização de mandato por perfil e objetivos específicos do cliente."
        ]
    },
    {
        "id": "planejamento_patrimonial",
        "nome": "Planejamento Patrimonial",
        "icone": "🏛️",
        "descricao": "Estruturação de patrimônio, holding familiar, proteção de ativos e sucessão",
        "pitch": "Seu patrimônio está estruturado de forma eficiente do ponto de vista fiscal e sucessório? Um planejamento patrimonial bem feito pode reduzir significativamente o ITCMD na transmissão e proteger os ativos de riscos jurídicos.",
        "oportunidades": [
            "Holding familiar: redução de ITCMD na transmissão de bens, proteção patrimonial e eficiência na gestão de ativos.",
            "Blindagem patrimonial: separação de ativos pessoais e empresariais para proteção em disputas jurídicas.",
            "Planejamento sucessório: evitar inventário judicial (custo de 10-15% do patrimônio) com estruturas adequadas.",
            "Revisão de regimes de bens: implicações patrimoniais para clientes casados ou em união estável."
        ]
    },
    {
        "id": "planejamento_financeiro",
        "nome": "Planejamento Financeiro",
        "icone": "🎯",
        "descricao": "Financial Planning completo: previdência, seguro de vida, proteção de renda e transmissão de bens em vida",
        "pitch": "Você tem uma estratégia clara para cada fase da sua vida financeira? Previdência, proteção da sua renda e do seu patrimônio, e um plano para transmitir seus bens com o menor custo possível — esses três pilares precisam estar presentes em qualquer planejamento sério.",
        "oportunidades": [
            "PREVIDÊNCIA PRIVADA — PGBL ou VGBL?",
            "SEGURO DE VIDA — proteção de renda e dependentes",
            "TRANSMISSÃO DE BENS EM VIDA — doação e sucessão planejada"
        ],
        "especialista": {
            "previdencia": {
                "titulo": "Previdência Privada",
                "icone": "🏦",
                "diagnostico": "A previdência privada é um dos instrumentos mais eficientes de longo prazo disponíveis no Brasil. A ausência dela é uma das maiores lacunas patrimoniais que existem.",
                "pgbl_vs_vgbl": {
                    "pgbl": "PGBL (Plano Gerador de Benefício Livre): ideal para quem declara IR no modelo completo. Deduz até 12% da renda bruta anual tributável — imposto só incide no resgate (sobre o total acumulado). Recomendado para quem tem IR a pagar.",
                    "vgbl": "VGBL (Vida Gerador de Benefício Livre): ideal para quem declara IR simplificado ou já usa o limite do PGBL. Imposto incide só sobre os rendimentos (não sobre o principal). Recomendado como complemento ao PGBL ou para isentos de IR.",
                    "regra": "Regra prática: PGBL para os 12% de dedução, VGBL para o restante dos aportes."
                },
                "tabela": {
                    "regressiva": "Tabela Regressiva: começa em 35% e cai até 10% após 10+ anos. IDEAL para quem quer acumular por longo prazo (15+ anos) e planeja sacar no prazo longo. A taxa de 10% é imbatível vs. qualquer outro produto de renda fixa.",
                    "progressiva": "Tabela Progressiva: alíquotas da tabela do IR (0% a 27,5%). IDEAL para quem terá renda baixa na aposentadoria (isento ou alíquota baixa) ou pode precisar resgatar antes de 10 anos."
                },
                "pitch": "Se você tem 40 anos e aporta R$ 2.000/mês por 25 anos com retorno de 100% CDI, a previdência com tabela regressiva pode representar uma economia de mais de R$ 200.000 em impostos vs. um fundo de renda fixa convencional. Faz sentido conversar sobre isso?"
            },
            "seguro_vida": {
                "titulo": "Seguro de Vida",
                "icone": "🛡️",
                "diagnostico": "O seguro de vida é a ferramenta mais barata e eficiente para proteger quem você ama. A maioria das pessoas subestima o impacto da perda da renda do provedor principal — o seguro resolve exatamente isso.",
                "coberturas": [
                    "Morte: garante a continuidade financeira da família. Regra: cobertura de 5x a renda anual bruta do provedor.",
                    "Invalidez permanente total ou parcial: frequentemente mais importante que morte — você vive, mas não consegue trabalhar.",
                    "Doenças graves (DIT): diagnóstico de câncer, AVC, infarto — capital para tratamento sem comprometer investimentos.",
                    "Perda de renda temporária: proteção de renda durante internações ou recuperações longas."
                ],
                "calculo": "Capital segurado mínimo recomendado: 10x a renda anual + total das dívidas (financiamentos, créditos em aberto). Isso garante que a família mantenha o padrão de vida por 10 anos sem depender dos investimentos.",
                "pitch": "Se algo acontecesse com você amanhã, sua família conseguiria manter o padrão de vida atual por quanto tempo com o patrimônio que tem hoje? O seguro de vida cobre exatamente esse gap — e custa muito menos do que a maioria imagina."
            },
            "transmissao_bens": {
                "titulo": "Transmissão de Bens em Vida",
                "icone": "🤝",
                "diagnostico": "Transmitir patrimônio sem planejamento pode custar até 20-25% do valor dos bens em impostos e custas de inventário. Com as ferramentas certas, esse custo pode ser reduzido drasticamente ou até eliminado.",
                "estrategias": [
                    "Doação com reserva de usufruto: doa os bens em vida, mas mantém o direito de uso e fruto enquanto viver. ITCMD incide sobre o valor da nua-propriedade (menor que o valor total). Evita inventário.",
                    "Holding familiar: concentra os bens em uma pessoa jurídica. A transmissão ocorre por cotas da holding (mais barata que transmissão de imóveis diretamente). Benefício adicional: proteção patrimonial e gestão unificada.",
                    "Seguro de vida: o capital segurado não entra em inventário e é transmitido diretamente ao beneficiário — isento de ITCMD e fora do alcance de credores.",
                    "Testamento: fundamental para garantir a vontade do titular, especialmente em famílias reconstituídas ou com herdeiros com necessidades especiais.",
                    "Previdência com beneficiários designados: assim como o seguro, a previdência privada transmite o saldo diretamente aos beneficiários sem inventário."
                ],
                "alerta_itcmd": "ITCMD em São Paulo: 4% (em discussão aumento para alíquotas progressivas até 8%). Em MG: até 5%. No RJ: até 8%. Um patrimônio de R$ 2 milhões sem planejamento pode gerar R$ 80.000-200.000 só em ITCMD — fora honorários de inventário.",
                "pitch": "Você já pensou em quanto vai custar para seus herdeiros receber o que você acumulou ao longo da vida? Com um planejamento adequado, é possível reduzir esse custo significativamente — e garantir que seu patrimônio chegue a quem você quer, da forma que você quer."
            }
        }
    },
    {
        "id": "investimentos_internacionais",
        "nome": "Investimentos Internacionais",
        "icone": "🌎",
        "descricao": "Diversificação em ativos internacionais: renda fixa global, ações estrangeiras, ETFs e BDRs",
        "pitch": "Você tem parte do seu patrimônio protegido em dólar ou outra moeda forte? Com o real historicamente volátil, ter exposição internacional é uma das formas mais eficientes de proteger poder de compra de longo prazo.",
        "oportunidades": [
            "Renda fixa global hedgeada (ex: Trend Global Bonds, Sparta): retorno em dólar com proteção cambial — diversificação sem risco de câmbio excessivo.",
            "BDRs (Brazilian Depositary Receipts): acesso a ações de Apple, Amazon, Google e outras diretamente na B3, em reais.",
            "ETFs internacionais (IVVB11, BNDX11): exposição ampla ao mercado americano e renda fixa global com liquidez diária.",
            "Contas no exterior: para patrimônios acima de R$ 1 milhão, considerar conta em corretora internacional (Interactive Brokers, Charles Schwab) — acesso direto a mercados globais.",
            "Fundos de investimento no exterior: estrutura local com gestão de ativos internacionais — simplifica a declaração de IR."
        ]
    }
]

ROOT        = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

SENHAS = {
    "lider":    "lider2026",
    "admin":    "admin2026",
    "head":     "head2026",
}

# Código do assessor → nome completo
ASSESSORES = {
    "A74621": "Lucas Landroni Cozzi",
    "A68983": "Tatiane Cristina da Silva Cecchetti",
    "A27157": "Felipe Fraga",
    "A34003": "Bruno Seiji Ito",
    "A74329": "Matheus Escoza Milani",
    "A25261": "Alex Alves dos Santos",
    "A53115": "Flavia Guedes de Souza",
    "A55932": "Carolina Custodio Siqueira",
    "A74638": "Andre Guilherme Leite Figueiredo",
    "A22930": "Carlos Fernando Victor Bolivar Moreira Neto",
    "A22795": "Thiago Brunelli Borba",
    "A45204": "Walter Nogueira de Souza Netto Ribeiro",
    "A67952": "Giuseppe Hilario Neto",
    "A71520": "Felipe Pereira Gomes",
    "A72679": "Paulo Roberto Negreiros Sobrinho",
    "A26364": "Felipe Santos Nishio",
    "A56767": "Michael Ademilson Santos da Silva",
    "A38890": "Marcelo Ramos Dias",
    "DENIS KINOSHITA": "Denis Kinoshita",
}
_DATA_FILE  = "/tmp/brauna_clientes.json"
_OKR_FILE   = "/tmp/brauna_okrs.json"
_MSG_FILE   = "/tmp/brauna_msg.json"
_FICHA_FILE = "/tmp/brauna_ficha.json"
_HIST_FILE  = "/tmp/brauna_historico.json"
_SUGE_FILE  = "/tmp/brauna_sugestoes.json"
_HP_PORT_FILE   = "/tmp/brauna_hp_portfolios.json"
_HP_CENARIO_FILE= "/tmp/brauna_hp_cenario.json"
_HP_PROD_FILE   = "/tmp/brauna_hp_produtos.json"
_HP_ALERTS_FILE = "/tmp/brauna_hp_alertas.json"
_HP_KNOW_FILE   = "/tmp/brauna_hp_knowledge.json"
_HP_CALLS_FILE  = "/tmp/brauna_hp_calls.json"
_HP_GESTORES_FILE   = "/tmp/brauna_hp_gestores.json"
_HP_GESTORAS2_FILE  = "/tmp/brauna_hp_gestoras2.json"
_HP_ESTRUTURADAS_FILE = "/tmp/brauna_hp_estruturadas.json"
_ADMIN_ACTIVITY_FILE  = "/tmp/brauna_admin_activity.json"
_CLIENTS_FILE         = "/tmp/brauna_clients.json"
_ASSESSORES_FILE      = "/tmp/brauna_assessores_dados.json"

# ── Portfólios Modelo default (Levante Asset — Junho 2026) ────────────────────
HP_PORTFOLIOS_DEFAULT = {
    "referencia": "Levante Asset — Junho 2026",
    "publicado_em": "2026-06-19",
    "classes": ["pos_fixado","inflacao","pre_fixado","acoes","fiis","multimercado","alternativos","internacional","criptomoedas"],
    "labels":  {"pos_fixado":"Pós Fixado","inflacao":"Inflação","pre_fixado":"Pré Fixado","acoes":"Ações","fiis":"FIIs","multimercado":"Multimercado","alternativos":"Alternativos","internacional":"Internacional","criptomoedas":"Criptomoedas"},
    "perfis": {
        "super_conservadora": {"label":"Super Conservadora","pos_fixado":84,"inflacao":11,"pre_fixado":5,"acoes":0,"fiis":0,"multimercado":0,"alternativos":0,"internacional":0,"criptomoedas":0},
        "conservadora":       {"label":"Conservadora",      "pos_fixado":70,"inflacao":16,"pre_fixado":7,"acoes":0,"fiis":0,"multimercado":3,"alternativos":0,"internacional":4,"criptomoedas":0},
        "moderada":           {"label":"Moderada",          "pos_fixado":44,"inflacao":23,"pre_fixado":10,"acoes":5,"fiis":1.5,"multimercado":6,"alternativos":1,"internacional":9,"criptomoedas":0.5},
        "arrojada":           {"label":"Arrojada",          "pos_fixado":28,"inflacao":28,"pre_fixado":12,"acoes":8,"fiis":2.5,"multimercado":9.5,"alternativos":1,"internacional":10.25,"criptomoedas":0.75},
        "agressiva":          {"label":"Agressiva",         "pos_fixado":13,"inflacao":31,"pre_fixado":13,"acoes":14,"fiis":3.5,"multimercado":10.5,"alternativos":1,"internacional":13,"criptomoedas":1},
    }
}

HP_CENARIO_DEFAULT = {
    "referencia": "Levante Asset — Junho 2026",
    "publicado_em": "2026-06-19",
    "global": "Persistência inflacionária global impulsionada por conflito no Oriente Médio (Brent US$ 92,5/barril). Fed hawkish: juros elevados por mais tempo, possível retomada de altas. PCE abril 3,8% a/a. Dólar forte frente às demais moedas. Bolsas globais avançaram em maio (S&P 500 +5,1%, Nasdaq +8,3%), puxadas por IA.",
    "brasil": "Selic deve encerrar ciclo de corte na próxima reunião com pausa para avaliação prospectiva. IPCA-15 maio 0,62% acima do esperado, núcleos e serviços ainda pressionados. Risco fiscal relevante: medidas contábeis do governo afetam dívida. Ibovespa -7,2% em maio por saída de estrangeiro e ruído político — precificação mais atraente para médio/longo prazo.",
    "posicionamento": "Reduzindo FIIs e Multimercados. Aumentando Renda Variável (correção recente traz valor). Crédito: spreads estabilizando — Idex-DI 1,8%, Idex-Infra 42bps. Curvas de juros e inflação implícita em alta em toda extensão. Duration curta na renda fixa.",
    "vieses": {
        "pos_fixado": "neutro",
        "inflacao": "positivo",
        "pre_fixado": "neutro",
        "acoes": "positivo",
        "fiis": "negativo",
        "multimercado": "negativo",
        "internacional": "positivo",
    }
}

# ── persistência ──────────────────────────────────────────────────────────────
def _kv_client():
    url   = os.environ.get("KV_REST_API_URL")
    token = os.environ.get("KV_REST_API_TOKEN")
    if not url or not token:
        return None
    try:
        from upstash_redis import Redis
        return Redis(url=url, token=token)
    except Exception:
        return None

def _key(path: str) -> str:
    """Converte /tmp/brauna_hp_cenario.json → brauna:hp_cenario"""
    name = os.path.basename(path).replace("brauna_", "").replace(".json", "")
    return f"brauna:{name}"

def _load(path, default):
    kv = _kv_client()
    if kv:
        try:
            val = kv.get(_key(path))
            if val is not None:
                return json.loads(val) if isinstance(val, str) else val
        except Exception:
            pass
    # Fallback: arquivo local /tmp (para desenvolvimento)
    try:
        with open(path, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return default

def _save(path, data):
    kv = _kv_client()
    if kv:
        try:
            kv.set(_key(path), json.dumps(data, ensure_ascii=False))
            return
        except Exception:
            pass
    # Fallback: arquivo local /tmp
    try:
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception:
        pass

def load_clientes():  return _load(_DATA_FILE, [])
def save_clientes(d): _save(_DATA_FILE, d)

def load_fichas():   return _load(_FICHA_FILE, {})
def save_fichas(d):  _save(_FICHA_FILE, d)
def load_hist():     return _load(_HIST_FILE, {})
def save_hist(d):    _save(_HIST_FILE, d)
def load_suge():     return _load(_SUGE_FILE, {"historico":[]})
def save_suge(d):    _save(_SUGE_FILE, d)

def filtrar_sugestoes(suge, perfil, comp):
    """Filtra sugestão ativa para o perfil e composição do cliente."""
    if not suge:
        return None
    tem_acoes   = comp.get("acoes", 0) > 0
    tem_fiis    = comp.get("fiis", 0) > 0
    tem_intl    = comp.get("internacional", 0) > 1
    result = {
        "id": suge.get("id"), "titulo": suge.get("titulo"),
        "criado_em": suge.get("criado_em"),
        "renda_fixa":    [i for i in suge.get("renda_fixa",[])    if perfil in i.get("perfis",[])],
        "renda_variavel": [i for i in suge.get("renda_variavel",[]) if perfil in i.get("perfis",[])],
        "estruturadas":  suge.get("estruturadas","") if tem_acoes else "",
        "internacional": [i for i in suge.get("internacional",[])  if perfil in i.get("perfis",[])],
        "fiis":          suge.get("fiis",[]),
        "tem_acoes":     tem_acoes,
        "tem_fiis":      tem_fiis,
        "tem_intl":      tem_intl,
    }
    return result

_OKR_DEFAULT = {
    "metas": {
        "captacao":       {"label":"Captação Líquida (R$)",    "meta":0,"realizado":0,"fmt":"brl"},
        "receita":        {"label":"Receita (R$)",              "meta":0,"realizado":0,"fmt":"brl"},
        "crossell":       {"label":"Cross Sell (produtos)",     "meta":0,"realizado":0,"fmt":"num"},
        "novos_clientes": {"label":"Novos Clientes",            "meta":0,"realizado":0,"fmt":"num"},
        "fee_based":      {"label":"Fee Based (%)",             "meta":0,"realizado":0,"fmt":"pct"},
        "ativacoes_300k": {"label":"Ativações 300k+",           "meta":0,"realizado":0,"fmt":"num"},
    },
    "prazo":"2026-12-31","foco":"","mensagem_assessores":""
}
def load_okrs():  return _load(_OKR_FILE, _OKR_DEFAULT)
def save_okrs(d): _save(_OKR_FILE, d)

def load_msg():
    return _load(_MSG_FILE, {"mensagem":"","estrategia":"","atualizado":""})
def save_msg(d): _save(_MSG_FILE, d)


def carregar_contexto():
    try:
        path = os.path.join(ROOT, "contexto_mercado.json")
        with open(path, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def buscar_macro_bcb():
    try:
        import urllib.request
        def bcb(serie):
            url = f"https://api.bcb.gov.br/dados/serie/bcdata.sgs.{serie}/dados/ultimos/1?formato=json"
            with urllib.request.urlopen(url, timeout=5) as r:
                return json.loads(r.read())[0]["valor"]
        return {
            "selic_meta": float(bcb(432).replace(",",".")),
            "ipca_12m":   float(bcb(13522).replace(",",".")),
        }
    except Exception:
        return {"selic_meta": None, "ipca_12m": None}


def mapear(texto):
    t = texto.lower().strip()
    for k, v in MAPEAMENTO.items():
        if k in t:
            return v
    return None


def extrair_composicao(pdf_bytes):
    import pdfplumber
    texto = ""
    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for p in pdf.pages:
            texto += (p.extract_text() or "") + "\n"

    comp = {c: 0.0 for c in CATS}
    comp["caixa"] = 0.0
    padrao = re.compile(r"([A-Za-zÀ-ÿ\s\-\/]+?)\s+(?:R\$[\s\d\.,]+\s+)?(\d{1,3}(?:[.,]\d{1,3})*(?:[.,]\d{1,2})?)\s*%", re.IGNORECASE)
    for linha in texto.splitlines():
        m = padrao.search(linha.strip())
        if m:
            try:
                perc = float(m.group(2).replace(".", "").replace(",", "."))
            except ValueError:
                continue
            cat = mapear(m.group(1).strip())
            if cat and perc > 0:
                comp[cat] = comp.get(cat, 0.0) + perc

    rent = {}
    mp = re.search(r"Portf.lio\s+([\d,]+)%\s+([\d,]+)%\s+([\d,]+)%\s+([\d,]+)%", texto)
    mc = re.search(r"CDI\s+([\d,]+)%\s+([\d,]+)%\s+([\d,]+)%\s+([\d,]+)%", texto)
    if mp: rent["portfolio"] = {k: float(mp.group(i+1).replace(",",".")) for i,k in enumerate(["mes","ano","12m","24m"])}
    if mc: rent["cdi"]       = {k: float(mc.group(i+1).replace(",",".")) for i,k in enumerate(["mes","ano","12m","24m"])}

    patrimonio = 0.0
    mp2 = re.search(r"PATRIM.NIO TOTAL BRUTO:\s*R\$\s*([\d\.,]+)", texto)
    if mp2:
        try: patrimonio = float(mp2.group(1).replace(".","").replace(",","."))
        except: pass

    caixa = comp.pop("caixa", 0.0)
    total = sum(comp.values())
    if total > 1:
        fator = 100.0 / total
        comp = {k: round(v * fator, 2) for k, v in comp.items()}
    return comp, caixa, rent, patrimonio


def extrair_xperformance(pdf_bytes):
    """Parser especializado para o relatório XPerformance da XP Investimentos."""
    import pdfplumber
    texto = ""
    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for p in pdf.pages:
            texto += (p.extract_text() or "") + "\n"

    # Dados do cabeçalho — suporta encoding com ? no lugar de acentos
    conta_m = re.search(r"Conta\s+Assessor\s+Data.*?\n(\d+)\s+(.+?)\s+(\d{2}/\d{2}/\d{4})", texto, re.DOTALL)
    conta    = conta_m.group(1).strip() if conta_m else ""
    assessor = conta_m.group(2).strip() if conta_m else ""
    data_ref = conta_m.group(3).strip() if conta_m else ""

    # Patrimônio — cabeçalho e valor podem estar em linhas separadas
    patrimonio = 0.0
    idx_pat = texto.find("TOTAL BRUTO")
    if idx_pat != -1:
        trecho = texto[idx_pat:idx_pat+400]
        pat_m = re.search(r"R\$\s*([\d]{1,3}(?:\.\d{3})*,\d{2})", trecho)
        if pat_m:
            try: patrimonio = float(pat_m.group(1).replace(".","").replace(",","."))
            except: pass

    # Rentabilidade portfolio e CDI
    rent = {}
    mp = re.search(r"Portf.lio\s+([-\d,]+)%\s+([-\d,]+)%\s+([-\d,]+)%\s+([-\d,]+)%", texto)
    mc = re.search(r"CDI\s+([-\d,]+)%\s+([-\d,]+)%\s+([-\d,]+)%\s+([-\d,]+)%", texto)
    if mp: rent["portfolio"] = {k: float(mp.group(i+1).replace(",",".")) for i,k in enumerate(["mes","ano","12m","24m"])}
    if mc: rent["cdi"]       = {k: float(mc.group(i+1).replace(",",".")) for i,k in enumerate(["mes","ano","12m","24m"])}

    # Composição — suporta acentos normais E substituídos por ? (encoding pdfplumber)
    # Mapa amplo: nomes como aparecem no PDF (com e sem acento)
    comp_map = [
        (r"[Pp].s\s*[Ff]ix",               "pos_fixado"),    # Pós Fixado / P?s Fixado
        (r"[Ii]nfla..[ao]",                "inflacao"),       # Inflação / Infla??o
        (r"[Pp]r.\s*[Ff]ix",               "pre_fixado"),     # Pré Fixado / Pr? Fixado
        (r"[Rr]enda\s*[Vv]ari",            "acoes"),          # Renda Variável Brasil
        (r"[Ff]undos?\s*[Ll]ist",          "fiis"),           # Fundos Listados
        (r"[Mm]ulti.?[Mm]ercado",          "multimercado"),   # Multimercado
        (r"[Aa]lternativ",                 "alternativos"),   # Alternativo
        (r"[Ii]nternacion",                "internacional"),  # Internacional
        (r"[Cc]ripto",                     "criptomoedas"),   # Cripto
        (r"[Ff]undos?\s*[Ee]strutur",      "alternativos"),   # Fundos Estruturados → alternativos
        (r"[Pp]revidên|[Pp]revid.ncia",    "pos_fixado"),     # Previdência → pós fixado
    ]
    comp = {c: 0.0 for c in CATS}
    caixa_perc = 0.0

    for linha in texto.splitlines():
        # Padrão 1: "Pós Fixado (77,79%)" ou "P?s Fixado (77,79%)"
        m = re.search(r"([\w\s\?\-\/\.]+?)\s*\((\d{1,3}[,\.]\d{1,2})%\)", linha)
        if m:
            nome_cat = m.group(1).strip()
            try: perc = float(m.group(2).replace(",","."))
            except: continue
            if re.search(r"[Cc]aixa|[Pp]rovento|100,00", nome_cat): caixa_perc = perc; continue
            for pat, key in comp_map:
                if re.search(pat, nome_cat, re.I):
                    comp[key] = max(comp[key], perc); break

    # Fallback: parser genérico se composição não encontrada
    if sum(comp.values()) < 1:
        comp, caixa_perc, rent2, pat2 = extrair_composicao(pdf_bytes)
        if not rent: rent = rent2
        if not patrimonio: patrimonio = pat2

    # ── Ativos individuais de RV e FIIs ─────────────────────────────────────────
    acoes = []
    fiis_list = []
    rv_section = False
    fii_section = False
    for linha in texto.splitlines():
        if re.search(r"Renda\s*Vari[aá]vel\s*Brasil", linha, re.I): rv_section = True; fii_section = False; continue
        if re.search(r"Fundos\s*Listados", linha, re.I): fii_section = True; rv_section = False; continue
        if re.search(r"Alternativo|Pre.Fixado|Infla|Caixa|Proventos", linha, re.I) and (rv_section or fii_section):
            rv_section = False; fii_section = False; continue

        ticker_m = re.match(r"^([A-Z]{3,5}[0-9]{1,2})\s+R\$?\s*([-\d\.,]+)\s+([\d]+)\s+([\d,]+)%", linha.strip())
        if not ticker_m:
            ticker_m = re.match(r"^([A-Z]{3,5}[0-9]{1,2})\s+([-R\$\s\d\.,]+?)\s+(\d+)\s+([\d,]+)%", linha.strip())

        if ticker_m and (rv_section or fii_section):
            ticker = ticker_m.group(1)
            try: saldo_raw = ticker_m.group(2).replace("R$","").replace(".","").replace(",",".").strip()
            except: saldo_raw = "0"
            try: saldo = float(saldo_raw)
            except: saldo = 0.0
            try: qtd = int(ticker_m.group(3).replace(".",""))
            except: qtd = 0
            try: perc = float(ticker_m.group(4).replace(",","."))
            except: perc = 0.0
            item = {"ticker": ticker, "saldo": saldo, "qtd": qtd, "perc": perc}
            if fii_section or re.match(r".*11$", ticker):
                fiis_list.append(item)
            else:
                acoes.append(item)

    # ── Ativos individuais de Renda Fixa ─────────────────────────────────────────
    # Formato na seção "POSIÇÃO DETALHADA":
    #   Cabeçalho de seção: "P?s Fixado R$ 346.922,61 - 5,72% ..."  (sem parênteses)
    #   Ativo multilinhas:   nome na linha N, "R$ saldo qtd perc% ..." na linha N+1
    #   Ativo inline (CRA/CRI): "CRA ... R$ 35.728,16 48 0,59% ..." tudo em uma linha

    RF_SECAO_MAP = [
        (r"^[Pp].s\s*[Ff]ix",      "pos_fixado"),
        (r"^[Ii]nfla..[ao]",       "inflacao"),
        (r"^[Pp]r.\s*[Ff]ix",      "pre_fixado"),
        (r"^[Mm]ulti.?[Mm]ercado", "multimercado"),
        (r"^[Ii]nternacion",       "internacional"),
        (r"^[Aa]lternativ",        "alternativos"),
    ]
    # Seção de cabeçalho: nome da classe + R$ total logo na mesma linha
    RF_HEADER_PAT = re.compile(
        r"^([A-Za-z\?\s\-]+?)\s+R\$\s+[\d\.,]+"
    )
    # Linha de dados de ativo: começa com R$ e tem qtd + %
    RF_DADOS_PAT = re.compile(
        r"R\$\s*([-\d\.]+,\d{2})\s+([-\d\.]+)\s+([-\d,]+)%\s+([-\d,]+)%"
    )
    # Para encerrar a leitura RF
    RF_FIM_PAT = re.compile(
        r"Renda\s*Vari|Fundos\s*[Ll]istados|MOVIMENTA|ESTAT.STICA\s+HIST|"
        r"HIST.RICO\s+POR\s+ESTRAT|Relat.rio\s+informativo|"
        r"POSIÇÃO DETALHADA.*POSIÇÃO", re.I
    )
    # Ruídos que aparecem depois do nome do fundo
    RF_RUIDO_PAT = re.compile(r"^(RL|RF\s*CP|CP\s*RL|RF|CP|Incentivado.*Investimento.*)$", re.I)

    def _brl(s):
        try: return float(s.replace(".","").replace(",","."))
        except: return 0.0

    rf_ativos    = []
    rf_cls       = None    # classe atual
    rf_nome_buf  = None    # nome em buffer (linha anterior)
    em_detalhe   = False   # True quando dentro de "POSIÇÃO DETALHADA"

    for linha in texto.splitlines():
        l = linha.strip()
        if not l:
            continue

        # Entra na seção de posição detalhada
        if re.search(r"POSI..O\s+DETALHADA", l, re.I):
            em_detalhe = True
            continue

        if not em_detalhe:
            continue

        # Encerra ao chegar em RV ou rodapé
        if RF_FIM_PAT.search(l):
            em_detalhe = False
            rf_cls = None
            rf_nome_buf = None
            continue

        # Detecta cabeçalho de seção RF: "P?s Fixado R$ 346.922,61 - 5,72% ..."
        nova_cls = None
        for pat, cls in RF_SECAO_MAP:
            if re.search(pat, l, re.I) and "R$" in l:
                # Confirma que é cabeçalho (tem o total da seção, não ativo individual)
                # Cabeçalhos têm " - " antes do %
                if re.search(r"R\$.*-\s*\d+[,\.]\d+%", l):
                    nova_cls = cls
                    break
        if nova_cls:
            rf_cls = nova_cls
            rf_nome_buf = None
            continue

        if rf_cls is None:
            continue

        # Linha de dados: "R$ 78.286,79 67759.84 1,29% 1,37% ..."
        m = RF_DADOS_PAT.search(l)
        if m:
            saldo    = _brl(m.group(1))
            perc_str = m.group(3)
            rent_mes_str = m.group(4)
            try: perc     = float(perc_str.replace(",","."))
            except: perc  = 0.0
            try: rent_mes = float(rent_mes_str.replace(",","."))
            except: rent_mes = None

            # Tenta pegar rent_12m (5ª % na linha)
            all_pcts = re.findall(r"([-\d,]+)%", l)
            rent_12m = float(all_pcts[4].replace(",",".")) if len(all_pcts) >= 5 else None

            # Nome: texto antes do R$ na mesma linha (CRA/CRI inline) ou buffer anterior
            before_rs = l[:m.start()].strip()
            before_rs = re.sub(r"\s+(RL|RF\s*CP|CP\s*RL|RF|CP)\s*$", "", before_rs).strip()
            nome = (before_rs or rf_nome_buf or "").strip()

            if nome and saldo > 0:
                rf_ativos.append({
                    "nome":     nome,
                    "classe":   rf_cls,
                    "saldo":    saldo,
                    "perc":     perc,
                    "rent_mes": rent_mes,
                    "rent_12m": rent_12m,
                })
            rf_nome_buf = None
            continue

        # Linha de nome do ativo (sem R$): buffer para próxima linha de dados
        if "R$" not in l:
            # Ignora cabeçalhos de tabela e ruídos
            if (not re.search(r"PRECIFICA|M.S\s+ATUAL|ANO\s+24|SALDO|QTD|%ALLOC|Estrat.gia\s+Saldo", l, re.I) and
                    not re.match(r"^[\d\s\.\-\%\*\+]+$", l) and
                    not RF_RUIDO_PAT.match(l) and
                    len(l) > 5):
                rf_nome_buf = l

    return {
        "conta":     conta,
        "assessor":  assessor,
        "data_ref":  data_ref,
        "patrimonio": patrimonio,
        "comp":      comp,
        "caixa":     caixa_perc,
        "rent":      rent,
        "acoes":     acoes,
        "fiis":      fiis_list,
        "rf_ativos": rf_ativos,
        "texto_completo": texto[:3000],
    }


def extrair_texto_pdf(pdf_bytes):
    import pdfplumber
    texto = ""
    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for p in pdf.pages:
            texto += (p.extract_text() or "") + "\n"
    return texto


def calcular_desvios(comp, modelo):
    resultado = []
    for cat in CATS:
        real, alvo = comp.get(cat, 0.0), modelo.get(cat, 0.0)
        resultado.append({"cat":cat,"label":LABELS[cat],"real":round(real,2),"alvo":round(alvo,2),"desvio":round(real-alvo,2)})
    return sorted(resultado, key=lambda x: abs(x["desvio"]), reverse=True)


def avaliar_modelo_servir(checklist_input):
    """Avalia cada pilar do Modelo de Servir com base no input do assessor."""
    resultado = []
    pendentes_criticos = 0
    pendentes_altos = 0

    for pilar in MODELO_SERVIR:
        feito = checklist_input.get(pilar["id"], False)
        status = "ok" if feito else "pendente"
        if not feito:
            if pilar["importancia"].startswith("CRÍTICA"):
                pendentes_criticos += 1
            elif pilar["importancia"].startswith("ALTA"):
                pendentes_altos += 1
        resultado.append({
            "id": pilar["id"],
            "nome": pilar["nome"],
            "icone": pilar["icone"],
            "descricao": pilar["descricao"],
            "importancia": pilar["importancia"],
            "impacto_falta": pilar["impacto_falta"],
            "acao": pilar["acao"],
            "status": status,
            "diretriz": pilar.get("diretriz", False),
            "cor": pilar["cor"],
        })

    score = sum(1 for r in resultado if r["status"] == "ok")
    return resultado, score, pendentes_criticos, pendentes_altos


def avaliar_cross_sell(cross_input):
    """Avalia quais áreas de cross sell o cliente já tem e quais estão em aberto."""
    resultado = []
    tem = []
    nao_tem = []
    for area in CROSS_SELL:
        ativo = cross_input.get(area["id"], False)
        item = {
            "id": area["id"],
            "nome": area["nome"],
            "icone": area["icone"],
            "descricao": area["descricao"],
            "pitch": area["pitch"],
            "oportunidades": area.get("oportunidades", []),
            "especialista": area.get("especialista"),
            "ativo": ativo,
        }
        resultado.append(item)
        if ativo:
            tem.append(area["nome"])
        else:
            nao_tem.append(area["nome"])
    return resultado, tem, nao_tem


def gerar_recomendacoes(desvios, perfil, macro, contexto, carta_texto=""):
    selic = macro.get("selic_meta")
    ipca  = macro.get("ipca_12m")
    rec_classes = contexto.get("recomendacoes_por_classe", {})
    visao = contexto.get("visao_gestores", [])

    frags = [d for d in desvios if d["desvio"] < -1.5]
    opors = [d for d in desvios if d["desvio"] > 1.5]
    recomendacoes = []

    gestores_resumo = "; ".join([f"{v['gestor']}: {v['implicacao_alocacao']}" for v in visao[:3]])

    carta_insights = []
    if carta_texto:
        for cat in CATS:
            label = LABELS[cat]
            for linha in carta_texto.splitlines():
                if any(p in linha.lower() for p in [cat.replace("_"," "), label.lower(), "ipca","selic","inflação","pós-fixado","prefixado"]):
                    trecho = linha.strip()
                    if len(trecho) > 30:
                        carta_insights.append({"classe": label, "trecho": trecho[:200]})
                        break

    for f in frags:
        cat = f["cat"]
        rc  = rec_classes.get(cat, {})
        produtos = rc.get("produtos_sugeridos", [])
        ctx_gestor = rc.get("contexto_gestor", "")

        if cat == "inflacao":
            macro_txt = f"IPCA 12M em {ipca:.1f}%" if ipca else "IPCA elevado"
            urgencia = "🔴 URGENTE"
            explicacao = f"{macro_txt} — cliente sem proteção inflacionária ({f['real']:.1f}% vs. alvo {f['alvo']:.1f}%). {ctx_gestor}"
        elif cat == "pre_fixado":
            macro_txt = f"Selic em {selic:.2f}%" if selic else "Selic elevada"
            urgencia = "🟡 MÉDIA"
            explicacao = f"{macro_txt} próxima do pico. Pré-fixado de médio prazo pode capturar ganho de capital com início de ciclo de queda. {ctx_gestor}"
        elif cat == "internacional":
            urgencia = "🟡 MÉDIA"
            explicacao = f"Diversificação cambial ausente ({f['real']:.1f}% vs. alvo {f['alvo']:.1f}%). SPX e Ibiuna destacam internacional como hedge. {ctx_gestor}"
        elif cat == "multimercado":
            urgencia = "🟡 MÉDIA"
            explicacao = f"Ausência de diversificador macro. {ctx_gestor}"
        elif cat == "acoes":
            urgencia = "🟢 BAIXA"
            explicacao = f"Renda variável abaixo do modelo. Cautela: Verde e Absolute neutros/negativos para bolsa Brasil. {ctx_gestor}"
        else:
            urgencia = "🟡 MÉDIA"
            explicacao = f"{f['label']} abaixo do modelo em {abs(f['desvio']):.1f}%. {ctx_gestor}"

        carta_insight = next((c["trecho"] for c in carta_insights if f['label'].lower() in c["classe"].lower()), "")

        recomendacoes.append({
            "classe": f["label"],
            "urgencia": urgencia,
            "falta_pp": abs(f["desvio"]),
            "explicacao": explicacao,
            "produtos": produtos[:4],
            "carta_insight": carta_insight,
            "consenso_gestores": gestores_resumo,
        })

    alertas = []
    for o in opors:
        cat = o["cat"]
        if cat == "alternativos" and o["desvio"] > 5:
            alertas.append(f"⚠️ {o['label']} com {o['real']:.1f}% vs. alvo {o['alvo']:.1f}% — Absolute alerta: fundos alternativos ilíquidos (FIPs) inadequados para perfil {perfil}. Verificar liquidez e prazo de resgate.")
        elif cat == "pos_fixado" and o["desvio"] > 5:
            alertas.append(f"⚠️ {o['label']} com {o['real']:.1f}% vs. alvo {o['alvo']:.1f}% — XP Asset e Verde alertam para concentração excessiva em pós-fixado. Realocar gradualmente para inflação.")
        else:
            alertas.append(f"⚠️ {o['label']} com excesso de {o['desvio']:.1f}% — candidato à realocação.")

    return recomendacoes, alertas


def gerar_pdf(nome, perfil, desvios, rent, patrimonio, caixa, data_ref, recomendacoes, alertas, macro, ref_contexto, checklist_servir=None):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable, Image
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    PRETO=colors.HexColor("#081F18"); DOURADO=colors.HexColor("#C9A96E"); DESC=colors.HexColor("#A07840")
    CESC=colors.HexColor("#1A1A1A"); CMED=colors.HexColor("#2C2C2C"); BRANCO=colors.white
    VERM=colors.HexColor("#FF6B6B"); VERDE=colors.HexColor("#5DCAA5"); AMARELO=colors.HexColor("#FFD966")
    LARANJA=colors.HexColor("#FF8C00")

    buf=io.BytesIO()
    doc=SimpleDocTemplate(buf,pagesize=A4,leftMargin=2*cm,rightMargin=2*cm,topMargin=2*cm,bottomMargin=2*cm)
    ss=getSampleStyleSheet()
    def S(n,**kw): return ParagraphStyle(n,parent=ss["Normal"],textColor=kw.pop("tc",BRANCO),**kw)

    s_title=S("T",fontSize=18,leading=22,tc=DOURADO,alignment=TA_CENTER,spaceAfter=4)
    s_sub=S("S",fontSize=10,tc=DESC,alignment=TA_CENTER,spaceAfter=2)
    s_sec=S("Se",fontSize=11,tc=DOURADO,spaceBefore=10,spaceAfter=4)
    s_body=S("B",fontSize=9,leading=13)
    s_rod=S("R",fontSize=7,tc=colors.HexColor("#4A7055"),alignment=TA_CENTER)
    s_frag=S("F",fontSize=9,leading=13,tc=VERM)
    s_opor=S("O",fontSize=9,leading=13,tc=VERDE)
    s_sug=S("Su",fontSize=9,leading=14,tc=DOURADO)
    s_bold=S("Sb",fontSize=9,leading=13,tc=BRANCO,fontName="Helvetica-Bold")
    s_gest=S("Sg",fontSize=8,leading=12,tc=colors.HexColor("#AAAAAA"))
    s_crit=S("Sc",fontSize=9,leading=13,tc=VERM,fontName="Helvetica-Bold")
    s_warn=S("Sw",fontSize=9,leading=13,tc=AMARELO)
    s_ok=S("Sok",fontSize=9,leading=13,tc=VERDE)

    elems=[]
    elems.append(Paragraph("BRAÚNA INVESTIMENTOS",s_title))
    elems.append(Paragraph("Análise de Alocação por Indexador — Modelo Levante Asset",s_sub))
    elems.append(HRFlowable(width="100%",thickness=1,color=DOURADO,spaceAfter=8))

    perfil_pt={"conservadora":"Conservadora","moderada":"Moderada","arrojada":"Arrojada","agressiva":"Agressiva"}
    pat_fmt=f"R$ {patrimonio:,.2f}".replace(",","X").replace(".",",").replace("X",".")
    selic_txt=f"{macro['selic_meta']:.2f}%" if macro.get('selic_meta') else "—"
    ipca_txt=f"{macro['ipca_12m']:.2f}%" if macro.get('ipca_12m') else "—"

    meta=[
        ["Cliente:",nome,"Assessor:","Denis Kinoshita"],
        ["Perfil:",perfil_pt.get(perfil,perfil),"Data ref.:",data_ref],
        ["Patrimônio:",pat_fmt,"Caixa:",f"{caixa:.1f}% (excluído do modelo)"],
        ["Selic meta:",selic_txt,"IPCA 12M:",ipca_txt],
    ]
    tm=Table(meta,colWidths=[2.6*cm,7*cm,2.6*cm,5*cm])
    tm.setStyle(TableStyle([
        ("FONTSIZE",(0,0),(-1,-1),8),("TEXTCOLOR",(0,0),(-1,-1),BRANCO),
        ("TEXTCOLOR",(0,0),(0,-1),DOURADO),("TEXTCOLOR",(2,0),(2,-1),DOURADO),
        ("ROWBACKGROUNDS",(0,0),(-1,-1),[CESC,CMED]),
        ("GRID",(0,0),(-1,-1),0.3,colors.HexColor("#3A3A3A")),("PADDING",(0,0),(-1,-1),4),
    ]))
    elems.append(tm); elems.append(Spacer(1,0.3*cm))

    # Checklist Modelo de Servir no PDF
    if checklist_servir:
        elems.append(HRFlowable(width="100%",thickness=0.5,color=DESC,spaceAfter=4))
        elems.append(Paragraph("Checklist — Modelo de Servir Braúna",s_sec))

        pendentes = [p for p in checklist_servir if p["status"] == "pendente"]
        concluidos = [p for p in checklist_servir if p["status"] == "ok"]
        score = len(concluidos)
        total = len(checklist_servir)

        score_txt = f"Atendimento: {score}/{total} pilares completos"
        score_style = s_ok if score >= 5 else (s_warn if score >= 3 else s_crit)
        elems.append(Paragraph(score_txt, score_style))
        elems.append(Spacer(1,0.15*cm))

        if pendentes:
            elems.append(Paragraph("O QUE ESTÁ FALTANDO (ação necessária do assessor):", S("hd",fontSize=9,tc=VERM,fontName="Helvetica-Bold")))
            elems.append(Spacer(1,0.1*cm))
            for p in pendentes:
                tag = "🔴 CRÍTICO" if p["importancia"].startswith("CRÍTICA") else ("🟡 IMPORTANTE" if p["importancia"].startswith("ALTA") else "⚪ MÉDIO")
                elems.append(Paragraph(f"{p['icone']} {p['nome']} — {tag}", S("ph",fontSize=9,tc=VERM if "CRÍTICO" in tag else AMARELO,fontName="Helvetica-Bold")))
                elems.append(Paragraph(p["impacto_falta"], s_body))
                elems.append(Paragraph(f"Ação: {p['acao']}", s_gest))
                if p.get("diretriz"):
                    elems.append(Paragraph("★ 1ª DIRETRIZ BRAÚNA — QUALIDADE: Provocar o cliente sobre Financial Planning é obrigatório.", S("dir",fontSize=8,tc=DOURADO,fontName="Helvetica-Bold")))
                elems.append(Spacer(1,0.15*cm))

        if concluidos:
            elems.append(Paragraph("Pilares concluídos: " + " | ".join([f"{p['icone']} {p['nome']}" for p in concluidos]), s_ok))

        elems.append(Spacer(1,0.2*cm))

    if rent.get("portfolio") and rent.get("cdi"):
        p,c=rent["portfolio"],rent["cdi"]
        rd=[["Período","Portfólio","CDI","% do CDI"],
            ["Mês",f'{p["mes"]:.2f}%',f'{c["mes"]:.2f}%',f'{round(p["mes"]/c["mes"]*100,1)}%'],
            ["Ano",f'{p["ano"]:.2f}%',f'{c["ano"]:.2f}%',f'{round(p["ano"]/c["ano"]*100,1)}%'],
            ["12M",f'{p["12m"]:.2f}%',f'{c["12m"]:.2f}%',f'{round(p["12m"]/c["12m"]*100,1)}%'],
            ["24M",f'{p["24m"]:.2f}%',f'{c["24m"]:.2f}%',f'{round(p["24m"]/c["24m"]*100,1)}%']]
        elems.append(Paragraph("Rentabilidade vs. CDI",s_sec))
        tr=Table(rd,colWidths=[3*cm,4*cm,4*cm,4*cm])
        tr.setStyle(TableStyle([
            ("BACKGROUND",(0,0),(-1,0),DESC),("TEXTCOLOR",(0,0),(-1,0),PRETO),
            ("TEXTCOLOR",(0,1),(-1,-1),BRANCO),("FONTSIZE",(0,0),(-1,-1),8),
            ("ROWBACKGROUNDS",(0,1),(-1,-1),[CESC,CMED]),
            ("GRID",(0,0),(-1,-1),0.3,colors.HexColor("#3A3A3A")),
            ("ALIGN",(1,0),(-1,-1),"CENTER"),("PADDING",(0,0),(-1,-1),4),
            ("FONTNAME",(0,0),(-1,0),"Helvetica-Bold"),
        ]))
        elems.append(tr); elems.append(Spacer(1,0.3*cm))

    elems.append(Paragraph("Composição Atual vs. Modelo Levante Asset",s_sec))
    lbs=[d["label"] for d in desvios]; rea=[d["real"] for d in desvios]; alv=[d["alvo"] for d in desvios]
    x=np.arange(len(lbs)); w=0.35
    fig,ax=plt.subplots(figsize=(10,4))
    fig.patch.set_facecolor("#081F18"); ax.set_facecolor("#1A1A1A")
    ax.bar(x-w/2,rea,w,color="#C9A96E",zorder=3,linewidth=0)
    ax.bar(x+w/2,alv,w,color="#2A5A3A",zorder=3,linewidth=0)
    ax.set_xticks(x); ax.set_xticklabels(lbs,rotation=28,ha="right",color="white",fontsize=8)
    ax.tick_params(colors="white"); ax.spines[:].set_color("#2C2C2C")
    ax.yaxis.grid(True,color="#1C4A34",zorder=0); ax.set_axisbelow(True); ax.set_ylabel("%",color="white")
    fig.tight_layout(pad=1)
    ib=io.BytesIO(); fig.savefig(ib,format="png",dpi=150,facecolor="#081F18"); plt.close(fig); ib.seek(0)
    elems.append(Image(ib,width=15*cm,height=6*cm)); elems.append(Spacer(1,0.3*cm))

    elems.append(Paragraph("Desvios por Indexador",s_sec))
    dd=[["Indexador","Atual (%)","Modelo (%)","Desvio (%)"]]
    for d in desvios: dd.append([d["label"],f'{d["real"]:.2f}',f'{d["alvo"]:.2f}',f'{d["desvio"]:+.2f}'])
    td=Table(dd,colWidths=[5*cm,3.5*cm,3.5*cm,3.5*cm])
    est=[("BACKGROUND",(0,0),(-1,0),DESC),("TEXTCOLOR",(0,0),(-1,0),PRETO),
         ("TEXTCOLOR",(0,1),(-1,-1),BRANCO),("FONTSIZE",(0,0),(-1,-1),8),
         ("ROWBACKGROUNDS",(0,1),(-1,-1),[CESC,CMED]),("GRID",(0,0),(-1,-1),0.3,colors.HexColor("#3A3A3A")),
         ("ALIGN",(1,0),(-1,-1),"CENTER"),("PADDING",(0,0),(-1,-1),4),("FONTNAME",(0,0),(-1,0),"Helvetica-Bold")]
    for i,d in enumerate(desvios,1):
        if d["desvio"]<-1.5: est.append(("TEXTCOLOR",(3,i),(3,i),VERM))
        elif d["desvio"]>1.5: est.append(("TEXTCOLOR",(3,i),(3,i),VERDE))
    td.setStyle(TableStyle(est))
    elems.append(td); elems.append(Spacer(1,0.3*cm))

    if alertas:
        elems.append(HRFlowable(width="100%",thickness=0.5,color=DESC,spaceAfter=4))
        elems.append(Paragraph("Alertas de Concentração",s_sec))
        for a in alertas: elems.append(Paragraph(a,S("Al",fontSize=9,leading=13,tc=AMARELO)))
        elems.append(Spacer(1,0.2*cm))

    if recomendacoes:
        elems.append(HRFlowable(width="100%",thickness=0.5,color=DESC,spaceAfter=4))
        elems.append(Paragraph("Recomendações Fundamentadas",s_sec))
        for rec in recomendacoes:
            elems.append(Paragraph(f"{rec['urgencia']} — {rec['classe']} (falta {rec['falta_pp']:.1f}%)",s_bold))
            elems.append(Paragraph(rec["explicacao"],s_body))
            if rec.get("carta_insight"):
                elems.append(Paragraph(f'Carta da gestao: "{rec["carta_insight"]}"',s_gest))
            if rec.get("produtos"):
                prods=", ".join(rec["produtos"])
                elems.append(Paragraph(f"Produtos sugeridos: {prods}",s_sug))
            elems.append(Spacer(1,0.15*cm))

    elems.append(HRFlowable(width="100%",thickness=0.5,color=DESC,spaceAfter=4))
    elems.append(Paragraph(f"Consenso dos Gestores — Ref. {ref_contexto}",s_sec))
    elems.append(Paragraph("Convergencia: pos-fixado ainda atrativo mas nao deve dominar; inflacao (NTN-B) e a oportunidade do momento com juro real acima de 6% a.a.; internacional como hedge cambial; cautela com alternativos iliquidos para perfis conservadores.",s_body))
    elems.append(Spacer(1,0.5*cm))

    elems.append(HRFlowable(width="100%",thickness=0.5,color=DESC,spaceAfter=4))
    elems.append(Paragraph("Este relatorio e um instrumento de apoio a decisao do assessor e nao constitui recomendacao formal de investimento. Dados macro via Banco Central do Brasil (BCB). Visao dos gestores consolidada de cartas mensais publicas. Brauna Investimentos — assessoria vinculada a XP Investimentos.",s_rod))

    def fundo(canvas,doc):
        canvas.saveState(); canvas.setFillColor(PRETO)
        canvas.rect(0,0,A4[0],A4[1],fill=1,stroke=0); canvas.restoreState()

    doc.build(elems,onFirstPage=fundo,onLaterPages=fundo)
    buf.seek(0); return buf


# ── HTML ──────────────────────────────────────────────────────────────────────
HTML = r"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Braúna — Análise de Carteiras</title>
<script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.1/dist/chart.umd.min.js"></script>
<style>
*{margin:0;padding:0;box-sizing:border-box}
body{background:#081F18;color:#F0F0F0;font-family:'Segoe UI',system-ui,sans-serif;min-height:100vh}
header{background:#111;border-bottom:1px solid #1A4030;padding:14px 28px;display:flex;align-items:center;gap:12px}
header h1{font-size:18px;color:#C9A96E;font-weight:700;letter-spacing:.5px}
header p{font-size:11px;color:#3A6A48;margin-top:2px}
.container{max-width:1020px;margin:0 auto;padding:28px 20px}
.card{background:#0B2A1F;border:1px solid #1A4030;border-radius:12px;padding:22px;margin-bottom:18px}
.card h2{font-size:12px;color:#C9A96E;font-weight:700;margin-bottom:14px;text-transform:uppercase;letter-spacing:.8px}
.grid-2{display:grid;grid-template-columns:1fr 1fr;gap:14px}
.grid-3{display:grid;grid-template-columns:1fr 1fr 1fr;gap:14px}
.grid-4{display:grid;grid-template-columns:repeat(4,1fr);gap:10px}
label{font-size:12px;color:#777;display:block;margin-bottom:5px}
input,select{width:100%;background:#1A1A1A;border:1px solid #1C4A34;border-radius:8px;padding:9px 12px;color:#F0F0F0;font-size:13px;outline:none;transition:border .2s}
input:focus,select:focus{border-color:#C9A96E}
select option{background:#1A1A1A}
.upload-area{border:1.5px dashed #1C4A34;border-radius:10px;padding:22px;text-align:center;cursor:pointer;transition:all .2s;position:relative;background:#111}
.upload-area:hover,.upload-area.drag{border-color:#C9A96E;background:#1A1A1A}
.upload-area input[type=file]{position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%}
.upload-area .icon{font-size:26px;margin-bottom:6px}
.upload-area p{font-size:13px;color:#4A7055}
.upload-area .fname{color:#C9A96E;font-weight:600;font-size:13px;margin-top:6px}
.btn{width:100%;background:#C9A96E;color:#081F18;border:none;border-radius:8px;padding:13px;font-size:14px;font-weight:700;cursor:pointer;transition:opacity .2s}
.btn:hover{opacity:.88}.btn:disabled{opacity:.4;cursor:not-allowed}
.btn-out{background:transparent;color:#C9A96E;border:1px solid #C9A96E}
.btn-out:hover{background:#1A1A1A}
.metric{background:#1A1A1A;border-radius:10px;padding:14px;text-align:center}
.metric .lbl{font-size:10px;color:#3A6A48;margin-bottom:4px;text-transform:uppercase;letter-spacing:.5px}
.metric .val{font-size:20px;font-weight:700;color:#C9A96E}
.metric .val.danger{color:#FF6B6B}.metric .val.ok{color:#5DCAA5}
.desvio-row{display:flex;align-items:center;gap:8px;padding:7px 0;border-bottom:1px solid #1A1A1A;font-size:12px}
.desvio-row:last-child{border:none}
.dn{width:106px;font-weight:500}.da{width:46px;color:#777;font-size:11px}
.db-w{flex:1;background:#1A1A1A;border-radius:3px;height:5px;overflow:hidden}
.db{height:100%;background:#C9A96E;border-radius:3px;transition:width .5s}
.dalvo{width:52px;color:#2A5A3A;font-size:10px}
.dpp{width:60px;text-align:right;font-weight:700;font-size:11px}
.red{color:#FF6B6B}.grn{color:#5DCAA5}.dim{color:#1E4A30}
.alert{border-radius:8px;padding:10px 14px;margin:5px 0;font-size:12px;line-height:1.5}
.alert.danger{background:#2A1010;border-left:3px solid #FF6B6B;color:#FFB3B3}
.alert.success{background:#0A2018;border-left:3px solid #5DCAA5;color:#9FE1CB}
.alert.gold{background:#1A160A;border-left:3px solid #C9A96E;color:#C9A96E}
.alert.warn{background:#1A1500;border-left:3px solid #FFD966;color:#FFD966}
.rec-card{background:#1A1A1A;border-radius:10px;padding:14px;margin-bottom:10px;border:1px solid #1C4A34}
.rec-header{font-size:13px;font-weight:700;margin-bottom:6px}
.rec-ctx{font-size:12px;color:#AAA;line-height:1.5;margin-bottom:6px}
.rec-carta{font-size:11px;color:#888;font-style:italic;margin-bottom:6px;padding:6px 10px;background:#111;border-radius:6px;border-left:2px solid #C9A96E}
.rec-prods{display:flex;flex-wrap:wrap;gap:5px;margin-top:6px}
.prod-tag{background:#081F18;border:1px solid #1E4A30;border-radius:20px;padding:3px 10px;font-size:11px;color:#C9A96E}
.gestor-card{background:#1A1A1A;border-radius:8px;padding:10px 14px;margin-bottom:6px;border-left:2px solid #1E4A30}
.gestor-nome{font-size:11px;color:#C9A96E;font-weight:700;margin-bottom:3px}
.gestor-msg{font-size:11px;color:#888;line-height:1.4}
.gestor-impl{font-size:11px;color:#5DCAA5;margin-top:3px}
.macro-badge{display:inline-flex;align-items:center;gap:6px;background:#1A1A1A;border:1px solid #1C4A34;border-radius:20px;padding:4px 12px;font-size:12px;margin:3px}
.macro-badge span{color:#C9A96E;font-weight:700}
#results{display:none}
.spinner{display:none;text-align:center;padding:32px;color:#C9A96E}
.spinner.show{display:block}
.loader{width:36px;height:36px;border:3px solid #1A4030;border-top-color:#C9A96E;border-radius:50%;animation:spin 1s linear infinite;margin:0 auto 12px}
@keyframes spin{to{transform:rotate(360deg)}}
.chart-wrap{position:relative;width:100%;height:250px}
.tab-btns{display:flex;gap:8px;margin-bottom:14px;flex-wrap:wrap}
.tab-btn{background:#1A1A1A;border:1px solid #1C4A34;border-radius:8px;padding:6px 14px;font-size:12px;color:#888;cursor:pointer;transition:all .2s}
.tab-btn.active{background:#C9A96E;color:#081F18;border-color:#C9A96E;font-weight:700}
.tab-panel{display:none}.tab-panel.active{display:block}
.rodape{text-align:center;font-size:11px;color:#1E4A30;margin-top:32px;padding-top:14px;border-top:1px solid #1A1A1A}
/* Checklist Modelo de Servir */
.pilar-toggle{display:flex;align-items:flex-start;gap:12px;padding:12px;border-radius:10px;background:#111;border:1px solid #1E1E1E;margin-bottom:8px;cursor:pointer;transition:all .2s;position:relative}
.pilar-toggle:hover{border-color:#1C4A34}
.pilar-toggle.done{border-color:#2A4030;background:#0A2A18}
.pilar-toggle.pending-crit{border-color:#4A1010;background:#1A0808}
.pilar-toggle.pending-high{border-color:#3A3A00;background:#141400}
.pilar-check{width:22px;height:22px;border-radius:6px;border:2px solid #1E4A30;display:flex;align-items:center;justify-content:center;flex-shrink:0;transition:all .2s;margin-top:1px;cursor:pointer}
.pilar-check.checked{background:#5DCAA5;border-color:#5DCAA5;color:#000;font-weight:700}
.pilar-icon{font-size:20px;flex-shrink:0}
.pilar-info{flex:1}
.pilar-nome{font-size:13px;font-weight:700;color:#F0F0F0;margin-bottom:2px;display:flex;align-items:center;gap:8px}
.pilar-nome .badge{font-size:10px;padding:2px 8px;border-radius:10px;font-weight:700}
.badge-crit{background:#FF2222;color:#fff}
.badge-high{background:#FFD966;color:#000}
.badge-med{background:#2A5A3A;color:#aaa}
.badge-dir{background:#C9A96E;color:#000}
.pilar-desc{font-size:11px;color:#3A6A48;line-height:1.4}
.pilar-falta{display:none;margin-top:10px;padding:10px;background:#081F18;border-radius:8px;border-left:3px solid #FF6B6B}
.pilar-falta.show{display:block}
.pilar-falta.warn{border-left-color:#FFD966}
.pilar-falta.med{border-left-color:#3A6A48}
.pilar-falta p{font-size:11px;color:#AAA;line-height:1.6;margin-bottom:6px}
.pilar-falta .acao{font-size:11px;color:#C9A96E;line-height:1.5;padding:8px;background:#1A2E1A;border-radius:6px;margin-top:6px}
.pilar-falta .diretriz-badge{font-size:10px;color:#C9A96E;font-weight:700;margin-top:6px;padding:4px 8px;background:#1A2E1A;border:1px solid #C9A96E;border-radius:4px;display:inline-block}
/* Placard de score */
.servir-score{display:flex;align-items:center;gap:16px;padding:14px 18px;border-radius:10px;background:#111;border:1px solid #1C4A34;margin-bottom:14px}
.score-num{font-size:36px;font-weight:700;color:#C9A96E;min-width:56px;text-align:center}
.score-info{flex:1}
.score-label{font-size:11px;color:#3A6A48;text-transform:uppercase;letter-spacing:.5px;margin-bottom:4px}
.score-bar-w{height:6px;background:#1A4030;border-radius:3px;overflow:hidden}
.score-bar{height:100%;border-radius:3px;transition:width .5s}
.score-pendentes{font-size:11px;color:#888;margin-top:5px}
/* Pilar result card */
.pilar-result{background:#1A1A1A;border-radius:10px;padding:14px;margin-bottom:8px;border-left:4px solid #1E4A30}
.pilar-result.ok{border-left-color:#5DCAA5}
.pilar-result.crit{border-left-color:#FF4444}
.pilar-result.high{border-left-color:#FFD966}
.pilar-result.med{border-left-color:#3A6A48}
.pilar-result-header{display:flex;align-items:center;gap:8px;margin-bottom:6px}
.pilar-result-nome{font-size:13px;font-weight:700}
.pilar-result-status{font-size:10px;padding:2px 8px;border-radius:10px}
.status-ok{background:#0A2018;color:#5DCAA5}
.status-pendente{background:#2A1010;color:#FF6B6B}
.pilar-result-impacto{font-size:11px;color:#888;line-height:1.5;margin-bottom:6px}
.pilar-result-acao{font-size:11px;color:#C9A96E;padding:8px;background:#111;border-radius:6px}
/* Cross Sell chips */
.cs-chip{display:inline-flex;align-items:center;gap:7px;padding:9px 16px;border-radius:30px;border:1.5px solid #1C4A34;background:#111;cursor:pointer;font-size:12px;font-weight:600;color:#777;transition:all .2s;user-select:none}
.cs-chip:hover{border-color:#C9A96E;color:#C9A96E}
.cs-chip.ativo{background:#1A2E1A;border-color:#C9A96E;color:#C9A96E}
.cs-chip .cs-icon{font-size:16px}
.cs-chip .cs-check{width:16px;height:16px;border-radius:50%;border:1.5px solid #1E4A30;display:flex;align-items:center;justify-content:center;font-size:10px;transition:all .2s}
.cs-chip.ativo .cs-check{background:#C9A96E;border-color:#C9A96E;color:#081F18;font-weight:900}
/* Cross sell result cards */
.cs-area{background:#1A1A1A;border-radius:12px;padding:16px;margin-bottom:12px;border:1px solid #1C4A34}
.cs-area.tem{border-left:4px solid #5DCAA5;opacity:.7}
.cs-area.nao-tem{border-left:4px solid #C9A96E}
.cs-area-header{display:flex;align-items:center;gap:10px;margin-bottom:8px}
.cs-area-nome{font-size:14px;font-weight:700;flex:1}
.cs-status-tag{font-size:10px;padding:3px 10px;border-radius:10px;font-weight:700}
.cs-tem-tag{background:#0A2018;color:#5DCAA5}
.cs-falta-tag{background:#1A2E1A;color:#C9A96E}
.cs-pitch{font-size:12px;color:#C9A96E;font-style:italic;padding:10px 12px;background:#111;border-radius:8px;border-left:3px solid #C9A96E;margin-bottom:10px;line-height:1.6}
.cs-ops{list-style:none;padding:0}
.cs-ops li{font-size:11px;color:#888;padding:5px 0;border-bottom:1px solid #1A1A1A;line-height:1.5}
.cs-ops li:last-child{border:none}
.cs-ops li::before{content:"→ ";color:#C9A96E}
/* Especialista FP */
.fp-bloco{background:#111;border-radius:10px;padding:14px;margin-top:10px;border:1px solid #1C4A34}
.fp-bloco-titulo{font-size:12px;font-weight:700;color:#C9A96E;margin-bottom:8px;display:flex;align-items:center;gap:6px}
.fp-bloco p{font-size:11px;color:#888;line-height:1.6;margin-bottom:6px}
.fp-bloco .fp-item{padding:6px 0;border-bottom:1px solid #1A1A1A;font-size:11px;color:#AAA;line-height:1.5}
.fp-bloco .fp-item:last-child{border:none}
.fp-bloco .fp-item b{color:#F0F0F0}
.fp-bloco .fp-pitch{font-size:11px;color:#5DCAA5;padding:8px;background:#0A2A18;border-radius:6px;border-left:2px solid #5DCAA5;margin-top:8px}
.fp-alerta{font-size:11px;color:#FF6B6B;padding:8px 10px;background:#2A1010;border-radius:6px;margin-top:6px;border-left:2px solid #FF6B6B}
.cs-score-bar{display:flex;align-items:center;gap:12px;padding:12px 16px;background:#111;border-radius:10px;border:1px solid #1C4A34;margin-bottom:14px}
/* Sugestões de Alocação */
.sg-bloco{background:#111;border-radius:12px;padding:16px;margin-bottom:14px;border:1px solid #1C4A34}
.sg-bloco-title{font-size:11px;font-weight:700;color:#C9A96E;text-transform:uppercase;letter-spacing:.8px;margin-bottom:12px;display:flex;align-items:center;gap:8px}
.sg-card{background:#1A1A1A;border-radius:8px;padding:12px;margin-bottom:8px;border-left:3px solid #1C4A34;position:relative}
.sg-card.urg-alta{border-left-color:#FF6B6B}
.sg-card.urg-media{border-left-color:#FFD966}
.sg-card.urg-baixa{border-left-color:#5DCAA5}
.sg-card .sg-topo{display:flex;align-items:center;gap:8px;margin-bottom:6px;flex-wrap:wrap}
.sg-card .sg-acao{font-size:10px;padding:2px 8px;border-radius:8px;font-weight:700;background:#1C4A34;color:#AAA}
.sg-card .sg-produto{font-size:13px;font-weight:700;color:#F0F0F0;flex:1}
.sg-card .sg-idx{font-size:10px;color:#C9A96E;padding:2px 7px;border-radius:6px;background:#1A2E1A;border:1px solid #C9A96E}
.sg-card .sg-motivo{font-size:11px;color:#888;line-height:1.5;margin-top:4px}
.sg-card .sg-de{font-size:11px;color:#FF6B6B;margin-bottom:4px}
.sg-card .sg-fonte{font-size:10px;color:#3A6A48;margin-top:4px}
.sg-ja-tem{opacity:.6;border-left-color:#5DCAA5 !important}
.sg-ja-tem-badge{font-size:10px;color:#5DCAA5;background:#0A2018;border-radius:6px;padding:2px 8px;margin-bottom:4px;display:inline-block}
.sg-est-box{background:#1A1A1A;border-radius:8px;padding:14px;border:1px solid #1C4A34;font-size:12px;color:#CCC;line-height:1.7;white-space:pre-wrap}
.sg-fii-table{width:100%;border-collapse:collapse;font-size:11px;margin-top:8px}
.sg-fii-table th{text-align:left;color:#3A6A48;font-weight:600;padding:5px 8px;border-bottom:1px solid #1C4A34;font-size:10px;text-transform:uppercase}
.sg-fii-table td{padding:7px 8px;border-bottom:1px solid #1A1A1A;color:#CCC}
.sg-fii-table tr:last-child td{border:none}
.sg-fii-table .fii-ticker{color:#C9A96E;font-weight:700}
.sg-header-box{background:#1A2E1A;border:1px solid #C9A96E;border-radius:8px;padding:10px 14px;margin-bottom:12px;display:flex;align-items:center;gap:10px}
.sg-vazio{font-size:13px;color:#5DCAA5;padding:16px;text-align:center}
@media(max-width:640px){.grid-2,.grid-3,.grid-4{grid-template-columns:1fr}}
/* ── Fluxo em 2 etapas ── */
.steps{display:flex;align-items:center;gap:0;margin-bottom:24px}
.step-item{display:flex;align-items:center;gap:8px;padding:10px 16px;background:#0B2A1F;border:1px solid #1A4030;border-radius:10px;cursor:default;transition:all .3s;flex-shrink:0}
.step-item.active{background:#1A2E1A;border-color:#C9A96E}
.step-item.done{background:#0A2A18;border-color:#2A4030}
.step-circle{width:26px;height:26px;border-radius:50%;background:#1A4030;border:2px solid #1E4A30;display:flex;align-items:center;justify-content:center;font-size:12px;font-weight:700;color:#3A6A48;flex-shrink:0;transition:all .3s}
.step-item.active .step-circle{background:#C9A96E;border-color:#C9A96E;color:#081F18}
.step-item.done .step-circle{background:#5DCAA5;border-color:#5DCAA5;color:#000}
.step-label{font-size:12px;font-weight:700;color:#2A5A3A;letter-spacing:.3px;transition:all .3s}
.step-item.active .step-label{color:#C9A96E}
.step-item.done .step-label{color:#5DCAA5}
.step-line{flex:1;height:1px;background:#1A4030;min-width:20px}
.step2-banner{padding:12px 16px;border-radius:10px;margin-bottom:16px;display:flex;align-items:center;gap:12px}
.step2-banner.first{background:#0A2A18;border:1px solid #2A4030}
.step2-banner.return{background:#1A2E1A;border:1px solid #3A2A00}
</style>
</head>
<body>
<header>
  <div>
    <h1>BRAÚNA INVESTIMENTOS</h1>
    <p>Análise de Alocação por Indexador · Modelo Levante Asset</p>
  </div>
  <nav style="display:flex;gap:8px" id="nav-assessor">
    <button onclick="sair()" style="font-size:12px;padding:5px 12px;border-radius:6px;border:1px solid #3A2A2A;color:#888;background:none;cursor:pointer;transition:all .2s" onmouseover="this.style.borderColor='#FF6B6B';this.style.color='#FF6B6B'" onmouseout="this.style.borderColor='#3A2A2A';this.style.color='#888'">Sair</button>
  </nav>
  <script>
  (function(){
    const role = localStorage.getItem("brauna_role");
    if(role === "admin"){
      const nav = document.getElementById("nav-assessor");
      nav.innerHTML = `
        <a href="/assessor" style="font-size:12px;padding:5px 12px;border-radius:6px;border:1px solid #C9A96E;color:#C9A96E;text-decoration:none;font-weight:700">📊 Assessor</a>
        <a href="/lider" style="font-size:12px;padding:5px 12px;border-radius:6px;border:1px solid #1C4A34;color:#3A6A48;text-decoration:none">👥 Líder</a>
        <a href="/head-produtos" style="font-size:12px;padding:5px 12px;border-radius:6px;border:1px solid #1C4A34;color:#3A6A48;text-decoration:none">🏛️ Head</a>
        <a href="/admin" style="font-size:12px;padding:5px 12px;border-radius:6px;border:1px solid #1C4A34;color:#3A6A48;text-decoration:none">⚙️ Admin</a>
        <button onclick="sair()" style="font-size:12px;padding:5px 12px;border-radius:6px;border:1px solid #3A2A2A;color:#888;background:none;cursor:pointer">Sair</button>`;
    }
  })();
  </script>
</header>
<div class="container">

<!-- Indicador de etapas -->
<div class="steps">
  <div class="step-item active" id="step-item-1">
    <div class="step-circle" id="step-c1">1</div>
    <span class="step-label">Dados do Cliente</span>
  </div>
  <div class="step-line"></div>
  <div class="step-item" id="step-item-2">
    <div class="step-circle" id="step-c2">2</div>
    <span class="step-label">Perfil &amp; Cross Sell</span>
  </div>
</div>

<!-- Formulário -->
<!-- Mensagem do Admin (se houver) -->
<div id="msg-admin-box" style="display:none;margin-bottom:14px;padding:12px 16px;background:#1A2E1A;border:1px solid #C9A96E;border-radius:10px">
  <p style="font-size:10px;color:#C9A96E;text-transform:uppercase;letter-spacing:.5px;margin-bottom:4px;font-weight:700">📢 Comunicado da Administração</p>
  <p id="msg-admin-txt" style="font-size:12px;color:#CCC;line-height:1.6"></p>
  <p id="msg-admin-data" style="font-size:10px;color:#3A6A48;margin-top:4px"></p>
</div>

<div class="card">
  <h2>Dados do cliente</h2>

  <!-- Seletor de cliente salvo -->
  <div id="clientes-salvos-box" style="display:none;margin-bottom:14px;padding:12px;background:#081F18;border-radius:10px;border:1px solid #1E1E1E">
    <p style="font-size:10px;color:#C9A96E;text-transform:uppercase;letter-spacing:.5px;margin-bottom:8px;font-weight:700">📁 Clientes salvos — clique para carregar</p>
    <div id="clientes-salvos-lista" style="display:flex;flex-wrap:wrap;gap:6px"></div>
  </div>

  <div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;margin-bottom:10px">
    <div>
      <label>Assessor</label>
      <input type="text" id="assessor" placeholder="Nome do assessor" onblur="buscarClientesSalvos()" readonly style="background:#1A1A1A;cursor:default;color:#C9A96E;display:none">
      <select id="assessor-select" onchange="document.getElementById('assessor').value=this.value;buscarClientesSalvos()" style="display:none;width:100%;background:#0B2A1F;border:1px solid #1A4030;border-radius:8px;padding:9px 12px;color:#C9A96E;font-size:13px;outline:none">
        <option value="">— Selecione o assessor —</option>
      </select>
    </div>
    <div><label>Nome do cliente</label><input type="text" id="nome" placeholder="Ex: Lucia Silva"></div>
  </div>
  <div class="grid-3" style="margin-bottom:14px">
    <div><label>Perfil</label>
      <select id="perfil" onchange="var _nm={conservadora:'CONSERVADORA',moderada:'MODERADA',arrojada:'ARROJADA',agressiva:'AGRESSIVA'};var _lb=document.getElementById('perfil-lbl');if(_lb)_lb.textContent=_nm[this.value]||this.value.toUpperCase();atualizarModelo();">
        <option value="conservadora">Conservadora</option>
        <option value="moderada">Moderada</option>
        <option value="arrojada">Arrojada</option>
        <option value="agressiva">Agressiva</option>
      </select>
    </div>
    <div><label>Objetivo do cliente</label><input type="text" id="objetivo" placeholder="Ex: aposentadoria, compra de imóvel..."></div>
    <div id="carta-ativa-info" style="display:flex;align-items:flex-end"><p style="font-size:11px;color:#3A6A48">Carta da gestão: <span id="carta-ativa-nome" style="color:#C9A96E">verificando...</span></p></div>
  </div>

  <div id="macro-badges" style="margin-bottom:14px"></div>

  <div style="background:#111;border-radius:8px;padding:10px 14px;margin-bottom:14px">
    <p style="font-size:10px;color:#3A6A48;margin-bottom:6px;text-transform:uppercase;letter-spacing:.5px">MODELO LEVANTE — <span id="perfil-lbl" style="color:#C9A96E">CONSERVADORA</span></p>
    <div id="modelo-grid" style="display:flex;flex-wrap:wrap;gap:6px 14px"></div>
  </div>

  <div style="margin-bottom:10px">
    <label>Relatório XPerformance (PDF) *</label>
    <div class="upload-area" id="drop1" onclick="document.getElementById('pdf-xp').click()">
      <input type="file" id="pdf-xp" accept=".pdf" style="display:none"
        onchange="onXpFileChange(this)">
      <div class="icon">📊</div>
      <p>XPerformance — arraste ou clique para identificar o cliente</p>
      <p class="fname" id="fname-xp"></p>
    </div>
  </div>

  <!-- Status da identificação -->
  <div id="box-cliente-identificado" style="margin-top:10px;background:#111;border:1px solid #1C4A34;border-radius:10px;padding:14px;display:none;transition:all .3s"></div>

  <!-- Botão para avançar — aparece após PDF carregado -->
  <div id="btn-proxima-etapa-wrap" style="display:none;margin-top:16px">
    <button class="btn" id="btn-proxima-etapa" onclick="avancarEtapa2()" style="background:#C9A96E;color:#081F18;font-size:15px;padding:15px">
      Continuar para Etapa 2 — Perfil &amp; Cross Sell →
    </button>
  </div>
</div>

<!-- ETAPA 2 — aparece após identificar o cliente pelo XPerformance -->
<div id="step2-wrapper" style="display:none">

<!-- Banner de status do cliente -->
<div id="step2-banner" class="step2-banner first">
  <div id="step2-icon" style="font-size:28px">👤</div>
  <div style="flex:1">
    <p id="step2-titulo" style="font-size:14px;font-weight:700;color:#F0F0F0;margin-bottom:2px"></p>
    <p id="step2-sub" style="font-size:11px;color:#888"></p>
  </div>
  <div id="step2-badges" style="display:flex;gap:6px;flex-wrap:wrap"></div>
</div>

<!-- Modelo de Servir — visível apenas na etapa 2 -->
<div class="card" id="card-servir">
  <h2>Modelo de Servir Braúna</h2>

  <!-- Layout: score à esquerda, pilares à direita -->
  <div style="display:flex;gap:24px;align-items:flex-start;flex-wrap:wrap">

    <!-- Score e donut -->
    <div style="flex-shrink:0;width:170px;text-align:center">
      <div style="position:relative;width:150px;margin:0 auto 0">
        <canvas id="chart-servir" width="150" height="150"></canvas>
        <div style="position:absolute;top:50%;left:50%;transform:translate(-50%,-50%);text-align:center;pointer-events:none;line-height:1">
          <div id="servir-score-num" style="font-size:32px;font-weight:900;color:#C9A96E">0</div>
          <div style="font-size:10px;color:#3A6A48;margin-top:2px">de 6 pilares</div>
        </div>
      </div>
      <div id="servir-legenda" style="margin-top:10px;display:flex;flex-direction:column;gap:4px;text-align:left"></div>
      <p style="font-size:10px;color:#2A5A3A;margin-top:10px;line-height:1.5">Clique em cada pilar para marcar como concluído</p>
    </div>

    <!-- Checklist de pilares — renderizado pelo servidor -->
    <div id="pilares-form" style="flex:1;min-width:280px;display:flex;flex-direction:column;gap:8px">
      {{ pilares_html | safe }}
    </div>
  </div>
</div>

<!-- Cross Sell — card separado -->
<div class="card">
  <h2>💼 Cross Sell — Produtos Ativos na Braúna</h2>
  <p style="font-size:11px;color:#3A6A48;margin-bottom:12px">Clique nas áreas que o cliente <b style="color:#F0F0F0">já possui</b>. O sistema mostrará oportunidades abertas e sugestões personalizadas.</p>
  <div id="cross-sell-form" style="display:flex;flex-wrap:wrap;gap:8px;margin-bottom:16px"></div>
  <button class="btn" id="btn-ana" onclick="analisar()">Analisar carteira</button>
</div>

</div><!-- /step2-wrapper -->

<div class="spinner" id="spinner"><div class="loader"></div><p>Analisando carteira e consultando contexto de mercado...</p></div>

<!-- Resultados -->
<div id="results">

  <!-- ── Análise Head de Produtos (HP) ── -->
  <div class="card" id="card-hp" style="display:none;border-color:#2A2A18">
    <h2 style="color:#D4B483">⭐ Análise Head de Produtos — Modelo Levante</h2>

    <!-- Cabeçalho da conta -->
    <div style="display:flex;gap:10px;flex-wrap:wrap;margin-bottom:14px" id="hp-header"></div>

    <!-- Alertas HP relevantes -->
    <div id="hp-alertas" style="margin-bottom:12px"></div>

    <!-- Desvios vs Modelo HP -->
    <div style="font-size:11px;color:#4A7055;margin-bottom:8px;display:flex;align-items:center;justify-content:space-between">
      <span>Comparativo vs. <span id="hp-ref" style="color:#D4B483"></span></span>
      <span id="hp-perfil-badge" style="background:#1A1A08;color:#D4B483;border:1px solid #D4B483;border-radius:12px;padding:2px 10px;font-size:10px;font-weight:700"></span>
    </div>
    <div id="hp-desvios" style="margin-bottom:16px"></div>

    <!-- Ações individuais -->
    <div id="hp-acoes-bloco" style="display:none;margin-bottom:12px">
      <div style="font-size:11px;color:#5DCAA5;font-weight:700;margin-bottom:8px;text-transform:uppercase;letter-spacing:.5px">📈 Carteira de Ações</div>
      <div id="hp-acoes-table"></div>
    </div>

    <!-- FIIs individuais -->
    <div id="hp-fiis-bloco" style="display:none;margin-bottom:12px">
      <div style="font-size:11px;color:#7DCFEF;font-weight:700;margin-bottom:8px;text-transform:uppercase;letter-spacing:.5px">🏢 Carteira de FIIs</div>
      <div id="hp-fiis-table"></div>
    </div>

    <!-- Sugestões de produtos HP para fechar gaps -->
    <div id="hp-sugestoes" style="display:none">
      <div style="font-size:11px;color:#D4B483;font-weight:700;margin-bottom:8px;text-transform:uppercase;letter-spacing:.5px">💡 Produtos sugeridos pelo Head de Produtos para fechar os gaps</div>
      <div id="hp-sugestoes-lista"></div>
    </div>

    <!-- Cenário macro HP -->
    <div id="hp-cenario" style="display:none;margin-top:14px">
      <div style="font-size:11px;color:#D4B483;font-weight:700;margin-bottom:10px;text-transform:uppercase;letter-spacing:.5px;display:flex;align-items:center;justify-content:space-between">
        <span>🌐 Cenário Macro — Head de Produtos</span>
        <span id="hp-cenario-ref" style="font-size:10px;color:#2A5A3A;font-weight:400;text-transform:none"></span>
      </div>
      <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:10px" id="hp-cenario-cards"></div>
      <div id="hp-cenario-vieses" style="margin-top:8px;display:flex;flex-wrap:wrap;gap:4px"></div>
    </div>
  </div>

  <div class="card" style="text-align:center">
    <h2>Exportar Análise</h2>
    <p style="font-size:12px;color:#3A6A48;margin-bottom:14px">Relatório técnico em PDF para arquivo e registro</p>
    <div style="display:flex;gap:14px;justify-content:center;flex-wrap:wrap">
      <button class="btn btn-out" id="btn-pdf" onclick="baixarPdf()" style="min-width:220px">⬇ Baixar relatório PDF</button>
    </div>
  </div>

  <div class="card">
    <h2>Resumo</h2>
    <div class="grid-4" id="metrics"></div>
  </div>

  <div class="card">
    <h2>Composição atual vs. modelo</h2>
    <div style="display:flex;gap:16px;margin-bottom:8px;font-size:11px">
      <span style="display:flex;align-items:center;gap:5px"><span style="width:10px;height:10px;border-radius:2px;background:#C9A96E;display:inline-block"></span>Carteira atual</span>
      <span style="display:flex;align-items:center;gap:5px"><span style="width:10px;height:10px;border-radius:2px;background:#2A5A3A;display:inline-block"></span>Modelo Levante</span>
    </div>
    <div class="chart-wrap"><canvas id="chart" role="img" aria-label="Comparativo de alocação"></canvas></div>
  </div>

  <!-- Diversificação por Classe de Ativos -->
  <div class="card" id="card-classes">
    <h2>Diversificação por Classe de Ativos</h2>
    <div style="display:flex;gap:28px;align-items:center;flex-wrap:wrap">
      <!-- Pizza -->
      <div style="flex-shrink:0;display:flex;flex-direction:column;align-items:center;gap:12px">
        <div style="position:relative;width:220px;height:220px">
          <canvas id="chart-classes" width="220" height="220"></canvas>
          <div style="position:absolute;top:50%;left:50%;transform:translate(-50%,-50%);text-align:center;pointer-events:none">
            <div style="font-size:11px;color:#3A6A48;line-height:1.3">carteira<br>atual</div>
          </div>
        </div>
        <!-- Legenda da pizza -->
        <div id="classes-legenda" style="display:flex;flex-direction:column;gap:5px;width:220px"></div>
      </div>
      <!-- Cards por grupo -->
      <div id="classes-grid" style="flex:1;min-width:260px;display:grid;grid-template-columns:1fr 1fr;gap:10px"></div>
    </div>
  </div>

  <!-- Plano de Ação para atingir o objetivo -->
  <div class="card" id="card-plano">
    <h2>Plano de Ação — Para atingir o modelo <span id="plano-perfil" style="color:#C9A96E;text-transform:capitalize"></span></h2>
    <div id="plano-objetivo-box" style="display:none;margin-bottom:12px;padding:10px 14px;background:#1A2E1A;border:1px solid #C9A96E;border-radius:8px">
      <p style="font-size:10px;color:#C9A96E;font-weight:700;text-transform:uppercase;letter-spacing:.5px;margin-bottom:2px">🎯 Objetivo do cliente</p>
      <p id="plano-objetivo-txt" style="font-size:13px;color:#F0F0F0"></p>
    </div>
    <div id="plano-grid" style="display:grid;grid-template-columns:1fr 1fr;gap:12px;margin-bottom:12px"></div>
    <div id="plano-valores" style="margin-top:8px"></div>
  </div>

  <!-- Tabs -->
  <div class="card">
    <div class="tab-btns">
      <button class="tab-btn active" onclick="tab('desvios')">Desvios</button>
      <button class="tab-btn" onclick="tab('recomendacoes')">Recomendações</button>
      <button class="tab-btn" onclick="tab('servir')">Modelo de Servir</button>
      <button class="tab-btn" onclick="tab('crosssell')">Cross Sell</button>
      <button class="tab-btn" onclick="tab('gestores')">Visão dos Gestores</button>
      <button class="tab-btn" onclick="tab('sugestoes')" id="tab-btn-sugestoes" style="display:none">💡 Sugestões</button>
    </div>

    <div class="tab-panel active" id="tab-desvios">
      <div class="grid-2">
        <div id="desvios-list"></div>
        <div id="diagnostico"></div>
      </div>
    </div>

    <div class="tab-panel" id="tab-recomendacoes">
      <div id="recomendacoes-list"></div>
    </div>

    <div class="tab-panel" id="tab-servir">
      <div id="servir-list"></div>
    </div>

    <div class="tab-panel" id="tab-crosssell">
      <div id="crosssell-list"></div>
    </div>

    <div class="tab-panel" id="tab-gestores">
      <div id="gestores-list"></div>
    </div>

    <div class="tab-panel" id="tab-sugestoes">
      <div id="sugestoes-list"></div>
    </div>
  </div>

  <!-- ── Gerar Apresentação de Reunião — ao final ── -->
  <div style="padding:20px 0 8px">
    <button id="btn-apres" onclick="gerarApresentacao()" class="btn" style="background:#C9A96E;color:#071E17;font-size:15px;font-weight:700;padding:16px;display:flex;align-items:center;justify-content:center;gap:10px;width:100%">
      🎯 Gerar Apresentação de Reunião (.pptx)
    </button>
    <p style="font-size:11px;color:#3A6A48;text-align:center;margin-top:8px">PowerPoint · Cenário macro · Análise da carteira · Sugestões de realocação · Identidade Braúna</p>
    <span id="apres-st" style="display:block;font-size:12px;color:#888;text-align:center;margin-top:6px"></span>
  </div>

</div>

<div class="rodape">Braúna Investimentos · Relatório de apoio à decisão · Não constitui recomendação formal de investimento · Dados macro via BCB</div>
</div>

<script>
// Auth guard
(function(){
  const role = localStorage.getItem("brauna_role");
  if(role !== "assessor" && role !== "admin") { localStorage.removeItem("brauna_role"); window.location.replace("/"); }

  const ASSESSORES_LIST = [
    "Lucas Landroni Cozzi","Tatiane Cristina da Silva Cecchetti","Felipe Fraga",
    "Bruno Seiji Ito","Matheus Escoza Milani","Alex Alves dos Santos",
    "Flavia Guedes de Souza","Carolina Custodio Siqueira","Andre Guilherme Leite Figueiredo",
    "Carlos Fernando Victor Bolivar Moreira Neto","Thiago Brunelli Borba",
    "Walter Nogueira de Souza Netto Ribeiro","Giuseppe Hilario Neto",
    "Felipe Pereira Gomes","Paulo Roberto Negreiros Sobrinho","Felipe Santos Nishio",
    "Michael Ademilson Santos da Silva","Marcelo Ramos Dias","Denis Kinoshita"
  ];

  if(role === "admin"){
    // Mostra select com lista de assessores
    const sel = document.getElementById("assessor-select");
    if(sel){
      ASSESSORES_LIST.forEach(n=>{
        const opt = document.createElement("option"); opt.value = n; opt.textContent = n; sel.appendChild(opt);
      });
      sel.style.display = "";
      sel.addEventListener("change", ()=>{
        document.getElementById("assessor").value = sel.value;
        buscarClientesSalvos();
      });
    }
  } else {
    // Assessor normal — mostra input readonly com nome do login
    const el = document.getElementById("assessor");
    if(el){ el.style.display = ""; }
    const nomeLogin = localStorage.getItem("brauna_nome");
    if(nomeLogin){
      const el2 = document.getElementById("assessor");
      if(el2 && !el2.value) el2.value = nomeLogin;
    }
  }
})();
function sair(){
  localStorage.removeItem("brauna_role");
  localStorage.removeItem("brauna_nome");
  localStorage.removeItem("brauna_codigo");
  window.location.replace("/");
}

// Registra acesso do assessor
(function(){
  const nome = localStorage.getItem("brauna_nome") || "Assessor";
  fetch("/api/admin/activity",{method:"POST",headers:{"Content-Type":"application/json"},
    body:JSON.stringify({role:"assessor",nome,acao:"acesso",detalhe:"Abriu a página do Assessor"})});
})();

const MODELOS={conservadora:{pos_fixado:70,inflacao:16,pre_fixado:7,acoes:0,fiis:0,multimercado:3,internacional:4,alternativos:0,criptomoedas:0},moderada:{pos_fixado:44,inflacao:23,pre_fixado:10,acoes:5,fiis:1.5,multimercado:6,internacional:9,alternativos:1,criptomoedas:0.5},arrojada:{pos_fixado:28,inflacao:28,pre_fixado:12,acoes:8,fiis:2.5,multimercado:9.5,internacional:10.25,alternativos:1,criptomoedas:0.75},agressiva:{pos_fixado:13,inflacao:31,pre_fixado:13,acoes:14,fiis:3.5,multimercado:10.5,internacional:13,alternativos:1,criptomoedas:1}};
const LABELS={pos_fixado:"Pós Fixado",inflacao:"Inflação",pre_fixado:"Pré Fixado",acoes:"Ações",fiis:"FIIs",multimercado:"Multimercado",internacional:"Internacional",alternativos:"Alternativos",criptomoedas:"Criptomoedas"};

const CROSS_AREAS=[
  {id:"aquisicao_bens",nome:"Aquisição de Bens",icone:"🏠"},
  {id:"gestao_discricionaria",nome:"Gestão Discricionária Prunus",icone:"🌿"},
  {id:"planejamento_patrimonial",nome:"Planejamento Patrimonial",icone:"🏛️"},
  {id:"planejamento_financeiro",nome:"Planejamento Financeiro",icone:"🎯"},
  {id:"investimentos_internacionais",nome:"Investimentos Internacionais",icone:"🌎"},
];

const crossSell={};
CROSS_AREAS.forEach(a=>crossSell[a.id]=false);

function renderCrossForm(){
  const c=document.getElementById("cross-sell-form");
  c.innerHTML=CROSS_AREAS.map(a=>`
    <div class="cs-chip" id="cs-${a.id}" onclick="toggleCS('${a.id}')">
      <span class="cs-icon">${a.icone}</span>
      <span>${a.nome}</span>
      <span class="cs-check" id="csck-${a.id}"></span>
    </div>`).join("");
}

function toggleCS(id){
  crossSell[id]=!crossSell[id];
  const chip=document.getElementById("cs-"+id);
  const chk=document.getElementById("csck-"+id);
  if(crossSell[id]){ chip.classList.add("ativo"); chk.textContent="✓"; }
  else { chip.classList.remove("ativo"); chk.textContent=""; }
}

renderCrossForm();

const PILARES=[
  {id:"open_investments",nome:"Open Investments",icone:"🔗",importancia:"CRÍTICA",dir:false,
   desc:"Consentimento OPIN ativo — visão 360° do patrimônio total (XP + outras instituições)",
   impacto:"Assessor enxerga apenas a parte do cliente na XP. A média XP é 37% do patrimônio nos clientes acima de R$ 300k — significa que 63% do patrimônio está invisível. Sem OPIN, qualquer recomendação de alocação é incompleta e pode duplicar posições.",
   acao:"Solicitar consentimento pelo app XP: Menu > Open Finance > Conectar outras contas. Pitch: 'Para cuidar do seu patrimônio como um todo e não só da parte na XP, precisamos do Open Investments. São 3 cliques no app — tão simples quanto um Pix. Assim evito duplicar posições e encontro oportunidades hoje invisíveis.'"},
  {id:"financial_planning",nome:"Financial Planning",icone:"🎯",importancia:"CRÍTICA",dir:true,
   desc:"Planejamento financeiro completo — objetivos, metas e destino financeiro mapeados",
   impacto:"SEM FINANCIAL PLANNING, NÃO HÁ DESTINO. Toda alocação é arbitrária quando não se sabe para onde o cliente quer ir. Não sabemos: quanto ele precisa para se aposentar, quando quer comprar um imóvel, qual o horizonte de cada objetivo. A carteira existe para realizar sonhos — sem o FP, estamos construindo uma estrada sem destino.",
   acao:"DIRETRIZ QUALIDADE BRAÚNA — Provocar o cliente AGORA: 'Você sabe exatamente quanto precisa acumular para se aposentar com o padrão de vida que deseja? Você tem um plano financeiro com clareza sobre cada objetivo da sua vida?' Perguntar ao assessor: O Financial Planning deste cliente já foi feito? Quando? Quanto isso é primordial para traçar o caminho do destino que o cliente deseja alcançar."},
  {id:"ordem_enviada",nome:"Ordem Enviada",icone:"📋",importancia:"ALTA",dir:false,
   desc:"Última ordem/movimentação executada — carteira ativa e evoluindo",
   impacto:"Carteira estagnada. Cliente sem aportes recentes é critério de ruptura no Índice de Saúde XP (2 pontos). Uma carteira sem movimentação perde oportunidades de realocação e demonstra falta de atenção ativa do assessor.",
   acao:"Verificar última ordem executada. Se não houver movimentação recente: identificar oportunidade de aporte ou realocação compatível com o perfil e gerar proposta concreta com base nos desvios do modelo apresentado neste relatório."},
  {id:"conta_acessada",nome:"Conta Acessada",icone:"📱",importancia:"ALTA",dir:false,
   desc:"Cliente acessou a plataforma XP nos últimos 30 dias — engajado com seus investimentos",
   impacto:"Cliente desengajado da plataforma tem maior risco de ruptura e portabilidade. No Índice de Saúde XP, 'Sem Ordens' vale 2 pontos de risco de ruptura. Clientes que não acessam a conta são mais suscetíveis a abordagens de concorrentes.",
   acao:"Entrar em contato ativo para reengajar. Enviar este relatório como conteúdo de valor. Convidar para reunião de revisão de carteira. O engajamento do cliente começa pelo engajamento do assessor."},
  {id:"xperformance",nome:"X-Performance",icone:"📊",importancia:"MÉDIA",dir:false,
   desc:"Relatório XPerformance analisado e discutido com o cliente",
   impacto:"Sem análise do XPerformance, o cliente não tem visão clara da rentabilidade vs. CDI. É impossível ter conversa fundamentada sobre realocação sem este diagnóstico.",
   acao:"Analisar o XPerformance e agendar reunião de revisão. Mostrar ao cliente como a carteira se compara ao CDI e ao modelo ideal para seu perfil — use este relatório como base."},
  {id:"atividade_relacionamento",nome:"Atividade de Relacionamento",icone:"🤝",importancia:"ALTA",dir:false,
   desc:"Contato ativo nos últimos 30 dias — ligação, reunião ou WhatsApp com conteúdo de valor",
   impacto:"Cliente silencioso = cliente em risco. A falta de relacionamento ativo é o principal precursor de portabilidade. Concorrentes que abordam o cliente com insights têm vantagem sobre um assessor ausente.",
   acao:"Entrar em contato hoje com um ponto de valor: use este relatório como motivo. 'Fiz uma análise completa da sua carteira e gostaria de conversar sobre algumas oportunidades que identifiquei.'"}
];

let analiseData=null, chartInst=null, _classesChart=null;
const checklist={};
PILARES.forEach(p=>checklist[p.id]=false);

document.getElementById("data").value=new Date().toISOString().split("T")[0];

// Carrega mensagem do admin e carta ativa
fetch("/api/admin/mensagem").then(r=>r.json()).then(d=>{
  if(d.mensagem){
    document.getElementById("msg-admin-box").style.display="block";
    document.getElementById("msg-admin-txt").textContent=d.mensagem;
    document.getElementById("msg-admin-data").textContent="Atualizado em: "+d.atualizado;
  }
}).catch(()=>{});
fetch("/api/admin/contexto-info").then(r=>r.json()).then(d=>{
  const nome=d.carta&&d.carta.nome?d.carta.nome+" ("+d.carta.atualizado+")":"Nenhuma carta carregada — contate o Admin";
  const el=document.getElementById("carta-ativa-nome");
  if(el){ el.textContent=nome; el.style.color=d.carta&&d.carta.nome?"#5DCAA5":"#3A6A48"; }
}).catch(()=>{});

// ─── Modelo de Servir ──────────────────────────────────────────────────────────
let _servirChart = null;

function atualizarGraficoServir(){
  const completos = PILARES.filter(p=> checklist[p.id]).length;
  const critPend  = PILARES.filter(p=>!checklist[p.id] && p.importancia==="CRÍTICA").length;
  const highPend  = PILARES.filter(p=>!checklist[p.id] && p.importancia==="ALTA").length;
  const medPend   = PILARES.filter(p=>!checklist[p.id] && p.importancia==="MÉDIA").length;

  const scoreEl = document.getElementById("servir-score-num");
  if(scoreEl) scoreEl.textContent = completos;

  const legendaEl = document.getElementById("servir-legenda");
  if(legendaEl) legendaEl.innerHTML = [
    completos ? `<span style="font-size:11px;display:flex;align-items:center;gap:5px"><span style="width:10px;height:10px;border-radius:2px;background:#5DCAA5;flex-shrink:0"></span>${completos} concluído(s)</span>` : "",
    critPend  ? `<span style="font-size:11px;display:flex;align-items:center;gap:5px"><span style="width:10px;height:10px;border-radius:2px;background:#FF4444;flex-shrink:0"></span>${critPend} crítico(s)</span>` : "",
    highPend  ? `<span style="font-size:11px;display:flex;align-items:center;gap:5px"><span style="width:10px;height:10px;border-radius:2px;background:#FFD966;flex-shrink:0"></span>${highPend} importante(s)</span>` : "",
    medPend   ? `<span style="font-size:11px;display:flex;align-items:center;gap:5px"><span style="width:10px;height:10px;border-radius:2px;background:#2A5A3A;flex-shrink:0"></span>${medPend} médio(s)</span>` : "",
  ].filter(Boolean).join("");

  const canvas = document.getElementById("chart-servir");
  if(!canvas || typeof Chart === "undefined") return;

  const vals   = [completos, critPend, highPend, medPend];
  const cores  = ["#5DCAA5","#FF4444","#FFD966","#2A5A3A"];
  const labels = ["Concluídos","Críticos","Importantes","Médios"];
  const filtDados  = vals.filter(v=>v>0);
  const filtCores  = cores.filter((_,i)=>vals[i]>0);
  const filtLabels = labels.filter((_,i)=>vals[i]>0);
  const chartData  = filtDados.length ? filtDados : [6];
  const chartCores = filtDados.length ? filtCores : ["#1A4030"];
  const chartLbls  = filtDados.length ? filtLabels : ["Nenhum"];

  if(_servirChart){ try{ _servirChart.destroy(); }catch(e){} }
  _servirChart = new Chart(canvas, {
    type:"doughnut",
    data:{ labels:chartLbls, datasets:[{data:chartData, backgroundColor:chartCores, borderWidth:0, hoverOffset:6}] },
    options:{
      cutout:"68%",
      plugins:{ legend:{display:false}, tooltip:{callbacks:{label:ctx=>` ${ctx.label}: ${ctx.parsed}`}} },
      animation:{duration:350}
    }
  });
}

function renderPilaresForm(){
  const c = document.getElementById("pilares-form");
  if(!c){ setTimeout(renderPilaresForm, 100); return; }

  const html = PILARES.map(function(p){
    const feito  = !!checklist[p.id];
    const impCor = p.importancia==="CRÍTICA" ? "#FF6B6B" : p.importancia==="ALTA" ? "#FFD966" : "#4A7055";
    const impTxt = p.importancia==="CRÍTICA" ? "🔴 CRÍTICO" : p.importancia==="ALTA" ? "🟡 IMPORTANTE" : "⚪ MÉDIO";
    const bg     = feito ? "#0A2A18" : p.importancia==="CRÍTICA" ? "#1A0A0A" : p.importancia==="ALTA" ? "#151500" : "#181818";
    const border = feito ? "#2A5040" : p.importancia==="CRÍTICA" ? "#4A1212" : p.importancia==="ALTA" ? "#3A3A00" : "#282828";

    return '<div onclick="togglePilar(\'' + p.id + '\')" style="' +
      'display:flex;align-items:flex-start;gap:12px;padding:12px 14px;border-radius:10px;' +
      'background:' + bg + ';border:1.5px solid ' + border + ';cursor:pointer;transition:all .2s;user-select:none">' +
      // Checkbox
      '<div style="width:24px;height:24px;border-radius:6px;border:2px solid ' + (feito?"#5DCAA5":"#2A5A3A") + ';' +
      'background:' + (feito?"#5DCAA5":"transparent") + ';display:flex;align-items:center;justify-content:center;' +
      'color:#000;font-size:15px;font-weight:900;flex-shrink:0;margin-top:2px">' +
      (feito ? "✓" : "") + '</div>' +
      // Ícone
      '<div style="font-size:20px;flex-shrink:0;margin-top:2px">' + p.icone + '</div>' +
      // Conteúdo
      '<div style="flex:1;min-width:0">' +
        '<div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;margin-bottom:4px">' +
          '<span style="font-size:13px;font-weight:700;color:' + (feito?"#5DCAA5":"#F0F0F0") + '">' + p.nome + '</span>' +
          '<span style="font-size:10px;padding:2px 8px;border-radius:8px;background:#1E1E1E;color:' + impCor + ';font-weight:700">' + impTxt + '</span>' +
          (p.dir ? '<span style="font-size:10px;padding:2px 8px;border-radius:8px;background:#1A2E1A;color:#C9A96E;border:1px solid #C9A96E;font-weight:700">★ 1ª DIRETRIZ</span>' : "") +
        '</div>' +
        '<div style="font-size:11px;color:#4A7055;line-height:1.4">' + p.desc + '</div>' +
        (!feito ?
          '<div style="background:#111;border-radius:6px;padding:8px 10px;border-left:3px solid ' + impCor + ';margin-top:8px">' +
            '<div style="font-size:11px;color:#888;line-height:1.5;margin-bottom:6px">' + p.impacto + '</div>' +
            '<div style="font-size:11px;color:#C9A96E;padding:6px 8px;background:#1A2E1A;border-radius:5px">→ ' + p.acao + '</div>' +
            (p.dir ? '<div style="font-size:10px;color:#C9A96E;margin-top:6px;padding:5px 8px;border:1px solid #C9A96E;border-radius:5px;background:#1A2E1A">★ 1ª DIRETRIZ BRAÚNA — Financial Planning é obrigatório antes de qualquer alocação.</div>' : "") +
          '</div>' : "") +
      '</div>' +
    '</div>';
  }).join("");

  c.innerHTML = html;
  // Donut sem depender de Chart.js ter carregado
  if(typeof Chart !== "undefined"){
    atualizarGraficoServir();
  } else {
    // tenta de novo em 500ms
    setTimeout(atualizarGraficoServir, 500);
  }
}

function togglePilar(id){
  checklist[id] = !checklist[id];
  const feito = checklist[id];
  const imp   = (PILARES.find(p=>p.id===id)||{}).importancia||"MÉDIA";
  const bg     = feito ? "#0A2A18" : imp==="CRÍTICA" ? "#1A0A0A" : imp==="ALTA" ? "#151500" : "#181818";
  const border = feito ? "#2A5040" : imp==="CRÍTICA" ? "#4A1212" : imp==="ALTA" ? "#3A3A00" : "#282828";
  const el = document.getElementById("pilar-"+id);
  if(el){
    el.style.background = bg;
    el.style.border = "1.5px solid " + border;
  }
  const chk = document.getElementById("chk-"+id);
  if(chk){
    chk.style.background   = feito ? "#5DCAA5" : "transparent";
    chk.style.borderColor  = feito ? "#5DCAA5" : "#2A5A3A";
    chk.textContent        = feito ? "✓" : "";
    chk.style.color        = "#000";
  }
  const nomeEl = document.getElementById("nome-"+id);
  if(nomeEl) nomeEl.style.color = feito ? "#5DCAA5" : "#F0F0F0";
  const det = document.getElementById("det-"+id);
  if(det) det.style.display = feito ? "none" : "block";
  atualizarGraficoServir();
}

// Inicialização: só cross-sell e modelo Levante (pilares já vêm do servidor)
if(document.readyState === "loading"){
  document.addEventListener("DOMContentLoaded", function(){
    renderCrossForm();
    atualizarModelo();
    atualizarGraficoServir();
  });
} else {
  renderCrossForm();
  atualizarModelo();
  if(typeof Chart !== "undefined") atualizarGraficoServir();
  else setTimeout(atualizarGraficoServir, 600);
}

// Cache do portfólio HP — declarado como var para evitar TDZ (chamado antes da inicialização)
var _hpPortfolios = null;
fetch("/api/hp/portfolios").then(r=>r.json()).then(d=>{ _hpPortfolios = d.perfis||null; atualizarModelo(); }).catch(()=>{});

function atualizarModelo(){
  const sel = document.getElementById("perfil");
  if(!sel) return;
  const p = sel.value;

  // Label — atualiza imediatamente, síncrono
  const lbl = document.getElementById("perfil-lbl");
  const NOMES = {conservadora:"CONSERVADORA",moderada:"MODERADA",arrojada:"ARROJADA",agressiva:"AGRESSIVA",super_conservadora:"SUPER CONSERVADORA"};
  if(lbl) lbl.textContent = NOMES[p] || p.toUpperCase();

  // Modelo: HP publicado tem prioridade; senão usa fallback local
  const hp = _hpPortfolios && _hpPortfolios[p];
  const m  = (hp && Object.keys(hp).filter(k=>k!=="label").length > 0) ? hp : MODELOS[p];
  if(!m) return;

  const g = document.getElementById("modelo-grid");
  if(!g) return;

  const COR_CLS = {
    pos_fixado:"#C9A96E", inflacao:"#F0A830", pre_fixado:"#B8A0E8",
    acoes:"#5DCAA5", fiis:"#7DCFEF", multimercado:"#D4B483",
    internacional:"#FF8E8E", alternativos:"#4EC9B0", criptomoedas:"#888"
  };
  const itens = Object.entries(m)
    .filter(([k,v]) => typeof v==="number" && v>0 && k!=="label")
    .sort((a,b) => b[1]-a[1]);
  g.innerHTML = itens.map(([k,v])=>{
    const cor = COR_CLS[k]||"#3A6A48";
    return '<div style="display:flex;align-items:center;gap:8px;margin-bottom:3px">' +
      '<span style="font-size:11px;color:#888;min-width:110px;flex-shrink:0">'+(LABELS[k]||k)+'</span>'+
      '<div style="flex:1;height:6px;background:#1A4030;border-radius:3px;overflow:hidden">'+
        '<div style="width:'+Math.min(v,100)+'%;height:100%;background:'+cor+';border-radius:3px;transition:width .4s"></div>'+
      '</div>'+
      '<span style="font-size:11px;font-weight:700;color:'+cor+';min-width:36px;text-align:right">'+v+'%</span>'+
    '</div>';
  }).join("");
}

const perfilEl = document.getElementById("perfil");
if(perfilEl){
  perfilEl.addEventListener("change", atualizarModelo);
}

fetch("/api/macro").then(r=>r.json()).then(d=>{
  const b=document.getElementById("macro-badges");
  if(d.selic_meta) b.innerHTML+=`<span class="macro-badge">Selic <span>${d.selic_meta.toFixed(2)}% a.a.</span></span>`;
  if(d.ipca_12m)   b.innerHTML+=`<span class="macro-badge">IPCA 12M <span>${d.ipca_12m.toFixed(2)}%</span></span>`;
  if(d.ref_contexto) b.innerHTML+=`<span class="macro-badge" style="color:#3A6A48">Gestores ref. <span style="color:#888">${d.ref_contexto}</span></span>`;
}).catch(()=>{});

// Handler global do input XPerformance — chamado pelo onchange inline no HTML
function onXpFileChange(input){
  const file = input.files[0];
  if(!file) return;
  const fname = document.getElementById("fname-xp");
  const drop  = document.getElementById("drop1");
  const wrap  = document.getElementById("btn-proxima-etapa-wrap");
  const btn   = document.getElementById("btn-proxima-etapa");
  // Feedback imediato — síncrono
  if(fname) fname.textContent = "⏳ " + file.name;
  if(drop)  drop.style.borderColor = "#C9A96E";
  if(wrap)  wrap.style.display = "block";
  if(btn){  btn.innerHTML = "⏳ Identificando cliente..."; btn.disabled = true; btn.style.opacity = ".6"; }
  // Processa em background (async)
  identificarCliente(file);
}

function setupDrop(dropId,inputId,fnameId){
  const drop=document.getElementById(dropId), input=document.getElementById(inputId), fname=document.getElementById(fnameId);
  // Drag-and-drop support
  drop.addEventListener("dragover",e=>{e.preventDefault();drop.classList.add("drag");});
  drop.addEventListener("dragleave",()=>drop.classList.remove("drag"));
  drop.addEventListener("drop",e=>{
    e.preventDefault();drop.classList.remove("drag");
    const f = e.dataTransfer.files[0];
    if(f){ const dt = new DataTransfer(); dt.items.add(f); input.files = dt.files; onXpFileChange(input); }
  });
}
setupDrop("drop1","pdf-xp","fname-xp");

// ── Identificação automática do cliente ao subir o XPerformance ──────────────
var _clienteIdentificado = null;

async function identificarCliente(file){
  // Mostra loading no upload
  const fname = document.getElementById("fname-xp");
  if(fname) fname.textContent = "Identificando cliente...";

  const fd = new FormData();
  fd.append("pdf", file);
  try{
    const r = await fetch("/api/xp-identificar", {method:"POST", body:fd});
    const d = await r.json();
    _clienteIdentificado = d;

    // ── PDF sem carteira (conta = "-" ou vazio, sem composição)
    const semDados = !d.conta || d.conta === "-" || d.conta === "";
    const semComp  = !d.composicao_atual || Object.values(d.composicao_atual).every(v => v === 0);
    if(fname){
      if(semDados || semComp){
        fname.style.color = "#E8A87C";
        fname.textContent = "⚠️ PDF sem carteira — cliente sem posição na XP nesta data. Preencha os dados manualmente.";
      } else {
        fname.style.color = "#5DCAA5";
        fname.textContent = `✓ Carteira lida — Conta ${d.conta} | R$ ${d.patrimonio?.toLocaleString('pt-BR',{minimumFractionDigits:2})}`;
      }
    }

    // ── Preenche campos com dados salvos
    if(d.ficha_salva?.nome)    document.getElementById("nome").value = d.ficha_salva.nome;
    if(d.ficha_salva?.perfil){ document.getElementById("perfil").value = d.ficha_salva.perfil; atualizarModelo(); }
    if(d.ficha_salva?.objetivo) document.getElementById("objetivo").value = d.ficha_salva.objetivo || "";

    // ── Restaura checklist do modelo de servir
    if(d.ficha_salva?.checklist){
      Object.assign(checklist, d.ficha_salva.checklist);
      PILARES.forEach(p=>{
        const feito = !!checklist[p.id];
        if(!feito) return;
        const el  = document.getElementById("pilar-"+p.id);
        if(el){ el.style.background="#0A2A18"; el.style.border="1.5px solid #2A5040"; }
        const chk = document.getElementById("chk-"+p.id);
        if(chk){ chk.style.background="#5DCAA5"; chk.style.borderColor="#5DCAA5"; chk.textContent="✓"; chk.style.color="#000"; }
        const nEl = document.getElementById("nome-"+p.id);
        if(nEl) nEl.style.color="#5DCAA5";
        const det = document.getElementById("det-"+p.id);
        if(det) det.style.display="none";
      });
      atualizarGraficoServir();
    }

    // ── Restaura cross-sell salvo
    if(d.ficha_salva?.cross_ativos?.length){
      CROSS_AREAS.forEach(a=>{ crossSell[a.id]=false; });
      d.ficha_salva.cross_ativos.forEach(id=>{ crossSell[id]=true; });
      renderCrossForm();
    }

    // ── Atualiza botão com dados do cliente
    mostrarBotaoProximaEtapa(d);

  }catch(e){
    console.error("Identificação falhou:", e);
    mostrarBotaoProximaEtapa(null);
  }
}

function mostrarBotaoProximaEtapa(d){
  const wrap = document.getElementById("btn-proxima-etapa-wrap");
  const btn  = document.getElementById("btn-proxima-etapa");
  if(!wrap) return;
  wrap.style.display = "block";

  if(btn){
    btn.disabled = false;
    btn.style.opacity = "1";
    if(d && d.conta){
      const nome = d.ficha_salva?.nome || d.nome_cliente || ("Conta " + d.conta);
      const retorno = d.tem_historico && d.ultima_carteira;
      btn.innerHTML = retorno
        ? `🔄 Continuar com ${nome} — Retorno →`
        : `🆕 Continuar com ${nome} — Primeiro acesso →`;
    } else {
      btn.innerHTML = "Continuar para Etapa 2 →";
    }
  }
}

function avancarEtapa2(){
  if(_clienteIdentificado){
    ativarEtapa2(_clienteIdentificado);
  } else {
    // fallback: dados do cliente ainda não identificados, mostra etapa 2 mesmo assim
    ativarEtapa2({conta:"", ficha_salva:{}, tem_historico:false, ultima_carteira:null, comparativo:[]});
  }
}

function ativarEtapa2(d){
  // Atualiza indicadores de etapa
  const i1 = document.getElementById("step-item-1");
  const i2 = document.getElementById("step-item-2");
  const c1 = document.getElementById("step-c1");
  if(i1){ i1.classList.remove("active"); i1.classList.add("done"); }
  if(c1) c1.textContent = "✓";
  if(i2){ i2.classList.add("active"); }

  // Monta banner de status
  const temFicha   = d.ficha_salva && Object.keys(d.ficha_salva).length > 0;
  const temHist    = d.tem_historico && d.ultima_carteira;
  const primAccess = !temFicha && !temHist;
  const nomeCliente = d.ficha_salva?.nome || d.nome_cliente || ("Conta " + d.conta);

  const banner  = document.getElementById("step2-banner");
  const iconEl  = document.getElementById("step2-icon");
  const tituloEl= document.getElementById("step2-titulo");
  const subEl   = document.getElementById("step2-sub");
  const badgesEl= document.getElementById("step2-badges");

  if(primAccess){
    if(banner){ banner.className="step2-banner first"; }
    if(iconEl) iconEl.textContent = "🆕";
    if(tituloEl) tituloEl.textContent = nomeCliente + " — Primeiro acesso";
    if(subEl) subEl.textContent = "Conta " + d.conta + " · Preencha o Modelo de Servir e o Cross Sell abaixo";
    if(badgesEl) badgesEl.innerHTML = '<span style="font-size:11px;padding:4px 12px;background:#1A2E1A;color:#C9A96E;border:1px solid #C9A96E;border-radius:20px;font-weight:700">Novo cliente</span>';
  } else {
    if(banner){ banner.className="step2-banner return"; }
    if(iconEl) iconEl.textContent = "🔄";
    if(tituloEl) tituloEl.textContent = nomeCliente + " — Retorno";
    const ult = d.ultima_carteira;
    const ultData = ult?.data_ref || ult?.salvo_em || "—";
    const pat = (d.patrimonio||0).toLocaleString("pt-BR",{maximumFractionDigits:0});
    if(subEl) subEl.textContent = "Conta " + d.conta + " · Ref. " + d.data_ref + " · Patrimônio R$ " + pat;
    const chk = d.ficha_salva?.checklist || {};
    const nFeitos = PILARES.filter(p=>chk[p.id]).length;
    const scoreCor = nFeitos >= 5 ? "#5DCAA5" : nFeitos >= 3 ? "#FFD966" : "#FF6B6B";
    if(badgesEl) badgesEl.innerHTML =
      `<span style="font-size:11px;padding:4px 12px;background:#0A2018;color:#5DCAA5;border:1px solid #2A5040;border-radius:20px;font-weight:700">Última reunião: ${ultData}</span>` +
      `<span style="font-size:11px;padding:4px 12px;background:#081F18;color:${scoreCor};border:1px solid #1C4A34;border-radius:20px;font-weight:700">${nFeitos}/${PILARES.length} pilares</span>`;
  }

  // Exibe a etapa 2
  const wrapper = document.getElementById("step2-wrapper");
  if(wrapper){
    wrapper.style.display = "block";
    setTimeout(()=> wrapper.scrollIntoView({behavior:"smooth", block:"start"}), 100);
  }

  // Renderiza painel de histórico no box abaixo do upload
  renderPainelCliente(d);
}

function renderPainelCliente(d){
  const box = document.getElementById("box-cliente-identificado");
  if(!box) return;

  const temFicha = d.ficha_salva && Object.keys(d.ficha_salva).length > 0;
  const temHist  = d.tem_historico && d.ultima_carteira;
  const rent     = d.rent?.portfolio || {};

  // ── Cabeçalho da conta identificada
  let html = `
  <div style="display:flex;align-items:center;gap:10px;margin-bottom:14px;flex-wrap:wrap">
    <div style="font-size:22px">📋</div>
    <div style="flex:1">
      <div style="font-size:14px;font-weight:700;color:#C9A96E">${d.nome_cliente || "Cliente " + d.conta}</div>
      <div style="font-size:11px;color:#3A6A48">Conta ${d.conta} · ${d.assessor_xp || ""} · Ref. ${d.data_ref}</div>
    </div>
    <div style="text-align:right">
      <div style="font-size:16px;font-weight:700;color:#F0F0F0">R$ ${(d.patrimonio||0).toLocaleString("pt-BR",{maximumFractionDigits:0})}</div>
      ${rent["12m"] ? `<div style="font-size:11px;color:#5DCAA5">${rent["12m"].toFixed(2)}% em 12M</div>` : ""}
    </div>
  </div>`;

  // ── Status do Modelo de Servir
  const chk = d.ficha_salva?.checklist || {};
  const nFeitos = PILARES.filter(p=>chk[p.id]).length;
  const nTotal  = PILARES.length;
  const scoreCor = nFeitos >= 5 ? "#5DCAA5" : nFeitos >= 3 ? "#FFD966" : "#FF6B6B";

  if(temFicha){
    html += `
    <div style="background:#081F18;border-radius:8px;padding:10px 14px;margin-bottom:12px;border:1px solid #1E1E1E">
      <div style="font-size:10px;color:#3A6A48;text-transform:uppercase;letter-spacing:.5px;margin-bottom:8px">Modelo de Servir — salvo em ${d.ficha_salva.atualizado_em||"—"}</div>
      <div style="display:flex;flex-wrap:wrap;gap:6px">
        ${PILARES.map(p=>{
          const feito = !!chk[p.id];
          return `<span style="font-size:11px;padding:3px 10px;border-radius:12px;background:${feito?"#0A2018":"#1A0808"};color:${feito?"#5DCAA5":"#FF6B6B"};border:1px solid ${feito?"#2A5040":"#4A1010"}">
            ${feito?"✓":"✕"} ${p.icone} ${p.nome}
          </span>`;
        }).join("")}
      </div>
      <div style="margin-top:8px;font-size:11px;color:${scoreCor};font-weight:700">${nFeitos}/${nTotal} pilares completos</div>
      ${d.ficha_salva.cross_ativos?.length ? `
        <div style="margin-top:8px;font-size:10px;color:#3A6A48;text-transform:uppercase;letter-spacing:.5px">Cross Sell ativo</div>
        <div style="margin-top:4px;display:flex;flex-wrap:wrap;gap:4px">
          ${d.ficha_salva.cross_ativos.map(id=>{
            const a = CROSS_AREAS.find(x=>x.id===id);
            return a ? `<span style="font-size:11px;padding:2px 8px;border-radius:10px;background:#1A2E1A;color:#C9A96E;border:1px solid #C9A96E">${a.icone} ${a.nome}</span>` : "";
          }).join("")}
        </div>` : ""}
    </div>`;
  } else {
    html += `<div style="font-size:12px;color:#3A6A48;padding:8px 0;margin-bottom:8px">⚠️ Primeiro acesso — preencha o Modelo de Servir abaixo e clique em Analisar para salvar.</div>`;
  }

  // ── Painel 3 colunas: Recomendada | Histórico | Atual
  {
    const perfil   = d.ficha_salva?.perfil || "";
    const modelo   = (_hpPortfolios && perfil && _hpPortfolios[perfil]) ? _hpPortfolios[perfil] : null;
    const ult      = temHist ? d.ultima_carteira : null;
    const compAtual= d.composicao_atual || {};

    // Filtra classes com pelo menos um valor > 0 em qualquer coluna
    const classesVisiveis = CATS.filter(cat => {
      const m = modelo ? (modelo[cat]||0) : 0;
      const h = ult ? (ult.composicao?.[cat]||0) : 0;
      const a = compAtual[cat]||0;
      return m>0 || h>0 || a>0;
    });

    if(classesVisiveis.length){
      const ultData  = ult ? (ult.data_ref || ult.salvo_em?.split(" ")[0] || "—") : null;
      const patAnterior = ult ? (ult.patrimonio||0).toLocaleString("pt-BR",{maximumFractionDigits:0}) : null;
      const patAtual    = (d.patrimonio||0).toLocaleString("pt-BR",{maximumFractionDigits:0});
      const patDelta    = ult ? (d.patrimonio||0)-(ult.patrimonio||0) : null;

      html += `
      <div style="background:#081F18;border-radius:10px;padding:12px 14px;border:1px solid #1E1E1E;margin-bottom:4px">
        <!-- Header colunas -->
        <div style="display:grid;grid-template-columns:110px 1fr 1fr 1fr;gap:4px;margin-bottom:10px;align-items:end">
          <div></div>
          <div style="text-align:center">
            <div style="font-size:9px;color:#3A6A48;text-transform:uppercase;letter-spacing:.5px;margin-bottom:2px">Recomendada</div>
            <div style="font-size:10px;color:#8B9FE8;font-weight:700">${perfil ? perfil.replace("_"," ") : "HP Modelo"}</div>
          </div>
          <div style="text-align:center">
            <div style="font-size:9px;color:#3A6A48;text-transform:uppercase;letter-spacing:.5px;margin-bottom:2px">Histórico</div>
            <div style="font-size:10px;color:#D4B483;font-weight:700">${ultData || "—"}</div>
          </div>
          <div style="text-align:center">
            <div style="font-size:9px;color:#3A6A48;text-transform:uppercase;letter-spacing:.5px;margin-bottom:2px">Atual (XP)</div>
            <div style="font-size:10px;color:#5DCAA5;font-weight:700">${d.data_ref || "Hoje"}</div>
          </div>
        </div>

        <!-- Linhas por classe -->
        ${classesVisiveis.map(cat => {
          const label  = LABELS[cat] || cat;
          const valMod = modelo ? (modelo[cat]||0) : null;
          const valHist= ult ? (ult.composicao?.[cat]||0) : null;
          const valAt  = compAtual[cat]||0;
          const deltaHist = (valHist !== null && valAt !== null) ? valAt - valHist : null;
          const deltaMod  = (valMod  !== null && valAt !== null) ? valAt - valMod  : null;
          const dCorHist  = deltaHist === null ? "" : deltaHist > 1 ? "#FF6B6B" : deltaHist < -1 ? "#5DCAA5" : "#3A6A48";
          const sHist     = deltaHist === null ? "" : (deltaHist > 0 ? "▲" : deltaHist < 0 ? "▼" : "→");
          // barra de progresso: max 100%, escala visual
          const barMod  = valMod  !== null ? Math.min(valMod,100)  : 0;
          const barHist = valHist !== null ? Math.min(valHist,100) : 0;
          const barAt   = Math.min(valAt,100);
          return `<div style="display:grid;grid-template-columns:110px 1fr 1fr 1fr;gap:4px;align-items:center;padding:5px 0;border-top:1px solid #0F2A1F">
            <span style="font-size:11px;color:#AAA">${label}</span>
            <!-- Recomendada -->
            <div style="text-align:center">
              ${valMod !== null ? `
              <div style="font-size:12px;font-weight:700;color:#8B9FE8">${valMod.toFixed(1)}%</div>
              <div style="height:4px;background:#1A1A2E;border-radius:2px;margin:2px 8px 0"><div style="height:4px;background:#8B9FE8;border-radius:2px;width:${barMod}%"></div></div>
              ` : `<span style="font-size:10px;color:#2A5A3A">—</span>`}
            </div>
            <!-- Histórico -->
            <div style="text-align:center">
              ${valHist !== null ? `
              <div style="font-size:12px;font-weight:700;color:#D4B483">${valHist.toFixed(1)}%</div>
              <div style="height:4px;background:#1E1A0A;border-radius:2px;margin:2px 8px 0"><div style="height:4px;background:#D4B483;border-radius:2px;width:${barHist}%"></div></div>
              ` : `<span style="font-size:10px;color:#2A5A3A">sem dados</span>`}
            </div>
            <!-- Atual -->
            <div style="text-align:center">
              <div style="display:flex;align-items:center;justify-content:center;gap:4px">
                <span style="font-size:12px;font-weight:700;color:#5DCAA5">${valAt.toFixed(1)}%</span>
                ${deltaHist !== null && Math.abs(deltaHist) >= 0.5 ? `<span style="font-size:10px;color:${dCorHist};font-weight:700">${sHist}${Math.abs(deltaHist).toFixed(1)}</span>` : ""}
              </div>
              <div style="height:4px;background:#0A1E0F;border-radius:2px;margin:2px 8px 0"><div style="height:4px;background:#5DCAA5;border-radius:2px;width:${barAt}%"></div></div>
            </div>
          </div>`;
        }).join("")}

        <!-- Rodapé patrimônio -->
        <div style="margin-top:10px;padding-top:8px;border-top:1px solid #1E1E1E;display:grid;grid-template-columns:110px 1fr 1fr 1fr;gap:4px;font-size:11px">
          <span style="color:#3A6A48">Patrimônio</span>
          <div></div>
          <div style="text-align:center;color:#D4B483">${patAnterior ? "R$ "+patAnterior : "—"}</div>
          <div style="text-align:center">
            <span style="color:#5DCAA5;font-weight:700">R$ ${patAtual}</span>
            ${patDelta !== null && Math.abs(patDelta) > 0 ? `<span style="color:${patDelta>0?"#5DCAA5":"#FF6B6B"};margin-left:4px">${patDelta>0?"+":""}R$ ${Math.abs(patDelta).toLocaleString("pt-BR",{maximumFractionDigits:0})}</span>` : ""}
          </div>
        </div>
      </div>`;
    }
  }
  if(!temHist){
    html += `<div style="font-size:11px;color:#2A5A3A;padding:4px 0">📌 Após analisar, esta carteira será salva como referência para o próximo comparativo.</div>`;
  }

  box.innerHTML = html;
}

function tab(name){
  const TABS=["desvios","recomendacoes","servir","crosssell","gestores","sugestoes"];
  document.querySelectorAll(".tab-btn").forEach((b,i)=>{b.classList.toggle("active",TABS[i]===name);});
  document.querySelectorAll(".tab-panel").forEach(p=>p.classList.remove("active"));
  document.getElementById("tab-"+name).classList.add("active");
}

// ─── Clientes salvos ───────────────────────────────────────────────────────
async function buscarClientesSalvos(){
  const assessor = document.getElementById("assessor").value.trim();
  if(!assessor) return;
  try{
    const r = await fetch("/api/ficha?assessor="+encodeURIComponent(assessor.toLowerCase()));
    const lista = await r.json();
    const box = document.getElementById("clientes-salvos-box");
    const ul  = document.getElementById("clientes-salvos-lista");
    if(!lista.length){ box.style.display="none"; return; }
    ul.innerHTML = lista.map(c=>`
      <button onclick="carregarFicha(${JSON.stringify(JSON.stringify(c))})"
        style="padding:6px 12px;border-radius:20px;border:1px solid #1C4A34;background:#111;color:#C9A96E;font-size:12px;cursor:pointer;transition:all .2s"
        onmouseover="this.style.borderColor='#C9A96E'" onmouseout="this.style.borderColor='#1C4A34'">
        ${c.nome} <span style="color:#3A6A48;font-size:10px">${c.perfil}</span>
      </button>`).join("");
    box.style.display="block";
  }catch(e){}
}

function carregarFicha(jsonStr){
  const c = JSON.parse(jsonStr);
  document.getElementById("nome").value   = c.nome    || "";
  if(c.perfil) document.getElementById("perfil").value = c.perfil;
  document.getElementById("objetivo").value = c.objetivo || "";
  atualizarModelo();

  // Restaura checklist
  if(c.checklist){
    Object.assign(checklist, c.checklist);
    PILARES.forEach(p=>{
      const feito = !!checklist[p.id];
      if(!feito) return; // estado inicial já é "não feito"
      const imp  = p.importancia || "MÉDIA";
      const el   = document.getElementById("pilar-"+p.id);
      if(el){ el.style.background="#0A2A18"; el.style.border="1.5px solid #2A5040"; }
      const chk  = document.getElementById("chk-"+p.id);
      if(chk){ chk.style.background="#5DCAA5"; chk.style.borderColor="#5DCAA5"; chk.textContent="✓"; chk.style.color="#000"; }
      const nEl  = document.getElementById("nome-"+p.id);
      if(nEl) nEl.style.color="#5DCAA5";
      const det  = document.getElementById("det-"+p.id);
      if(det) det.style.display="none";
    });
    atualizarGraficoServir();
  }

  // Restaura cross sell
  CROSS_AREAS.forEach(a=>{ crossSell[a.id]=false; });
  if(c.cross_ativos && c.cross_ativos.length){
    c.cross_ativos.forEach(id=>{ crossSell[id]=true; });
  }
  renderCrossForm();

  // Aviso visual
  const nome = document.getElementById("nome");
  nome.style.borderColor="#C9A96E";
  setTimeout(()=>nome.style.borderColor="", 1500);
}

// ── Análise HP (Head de Produtos) ─────────────────────────────────────────────
function renderAnaliseHP(xp){
  const card = document.getElementById("card-hp");
  card.style.display = "";

  // Cabeçalho
  const hdr = document.getElementById("hp-header");
  hdr.innerHTML = [
    ["Conta", xp.conta || "—"],
    ["Data Ref.", xp.data_ref || "—"],
    ["Patrimônio", xp.patrimonio ? "R$ " + xp.patrimonio.toLocaleString("pt-BR", {minimumFractionDigits:2}) : "—"],
    ["Rent. Mês",  xp.rent?.portfolio?.mes  ? xp.rent.portfolio.mes.toFixed(2)+"%" : "—"],
    ["Rent. 12M",  xp.rent?.portfolio?.["12m"] ? xp.rent.portfolio["12m"].toFixed(2)+"%" : "—"],
  ].map(([lbl,val])=>`
    <div style="background:#071E17;border:1px solid #2A2A18;border-radius:8px;padding:8px 14px;min-width:90px">
      <div style="font-size:9px;color:#3A6A48;text-transform:uppercase;letter-spacing:.5px">${lbl}</div>
      <div style="font-size:13px;color:#D4B483;font-weight:700;margin-top:2px">${val}</div>
    </div>
  `).join("");

  // Referência e perfil
  document.getElementById("hp-ref").textContent = xp.referencia_modelo || "Levante Asset";
  document.getElementById("hp-perfil-badge").textContent = "Perfil detectado: " + (xp.perfil_detectado || "—");

  // Alertas relevantes
  // Badge "carteira salva"
  const hpHeader = document.getElementById("hp-header");
  if(xp.carteira_salva && hpHeader){
    const salvo = document.createElement("span");
    salvo.style.cssText = "font-size:10px;background:#0A2018;color:#5DCAA5;border:1px solid #1A4A2A;padding:3px 10px;border-radius:10px;display:inline-flex;align-items:center;gap:4px";
    salvo.innerHTML = "💾 Carteira registrada";
    hpHeader.prepend(salvo);
  }

  const alertasDiv = document.getElementById("hp-alertas");
  const alertas    = xp.alertas_relevantes || [];
  const nMercado   = (xp.alertas_mercado || []).length;

  if(alertas.length){
    const ACOR  = {info:"#5DCAA5", atencao:"#D4B483", urgente:"#FF6B6B"};
    const AICON = {info:"ℹ️", atencao:"⚠️", urgente:"🚨"};
    const ORIG  = {hp_manual:"HP", knowledge_base:"📚 Base", auto:"🤖 Auto", mercado:"📰 Mercado"};
    const OBKG  = {hp_manual:"#1A1A08", knowledge_base:"#0A200A", auto:"#0A0A1A", mercado:"#0A1A28"};
    const OCOR  = {hp_manual:"#888", knowledge_base:"#5DCAA5", auto:"#7777DD", mercado:"#5BA8D4"};

    const urgentes = alertas.filter(a=>a.tipo==="urgente");
    const bordaCor = urgentes.length ? "#FF6B6B" : "#D4B483";

    alertasDiv.innerHTML = `
      <div style="border:2px solid ${bordaCor};border-radius:10px;overflow:hidden;margin-bottom:14px">
        <div style="background:${bordaCor}22;padding:10px 14px;display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:6px">
          <span style="font-size:13px;font-weight:700;color:${bordaCor}">
            ${urgentes.length ? "🚨" : "⚠️"} ${alertas.length} ponto${alertas.length>1?"s":""} para atenção
          </span>
          <div style="display:flex;gap:6px;flex-wrap:wrap">
            ${urgentes.length ? `<span style="font-size:10px;background:#FF6B6B22;color:#FF6B6B;border:1px solid #FF6B6B44;padding:2px 8px;border-radius:10px">${urgentes.length} urgente${urgentes.length>1?"s":""}</span>` : ""}
            ${nMercado ? `<span style="font-size:10px;background:#0A1A2822;color:#5BA8D4;border:1px solid #5BA8D444;padding:2px 8px;border-radius:10px">📰 ${nMercado} do mercado</span>` : ""}
          </div>
        </div>
        <div style="padding:8px 14px 10px">
          ${alertas.map(a=>{
            const cor   = ACOR[a.tipo] || "#888";
            const ot    = a.origem_tipo || "auto";
            const orig  = a.origem || ORIG[ot] || "";
            const obkg  = OBKG[ot] || "#1A1A08";
            const ocor  = OCOR[ot] || "#888";
            const isMercado = ot === "mercado";
            return `<div style="border-left:3px solid ${cor};background:${isMercado?"#060F18":"#060F0B"};border-radius:0 8px 8px 0;padding:10px 12px;margin-bottom:8px">
              <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;margin-bottom:4px">
                <span style="font-size:12px;font-weight:700;color:${cor}">${AICON[a.tipo]||"🔔"} ${a.produto||a.classe||""}</span>
                ${orig ? `<span style="font-size:10px;background:${obkg};color:${ocor};padding:2px 7px;border-radius:10px;border:1px solid ${ocor}33">${orig}</span>` : ""}
                <span style="font-size:10px;color:#2A5A3A;margin-left:auto">${a.data||""}</span>
              </div>
              <p style="font-size:11px;color:#AAA;margin:0;line-height:1.5">${a.mensagem||""}</p>
            </div>`;
          }).join("")}
        </div>
      </div>`;
  } else {
    alertasDiv.innerHTML = `<div style="font-size:11px;color:#1E4A30;padding:8px 0;display:flex;align-items:center;gap:6px">✅ Nenhum alerta identificado para esta carteira.</div>`;
  }

  // Desvios
  const STATUS_COLOR = {ok:"#5DCAA5", atencao:"#D4B483", fora:"#FF6B6B"};
  const STATUS_ICON  = {ok:"✅", atencao:"⚠️", fora:"❌"};
  const devDiv = document.getElementById("hp-desvios");
  devDiv.innerHTML = (xp.desvios||[]).map(d=>`
    <div style="display:flex;align-items:center;gap:8px;padding:6px 0;border-bottom:1px solid #111;font-size:12px">
      <span style="width:130px;color:#CCC">${d.label}</span>
      <span style="width:44px;color:#D4B483;font-weight:700;text-align:right">${d.atual.toFixed(1)}%</span>
      <div style="flex:1;background:#0F3525;border-radius:3px;height:5px;overflow:hidden;position:relative">
        <div style="position:absolute;height:100%;background:#2A5A3A;border-radius:3px;width:${Math.min(d.alvo,100)}%"></div>
        <div style="position:absolute;height:100%;background:${STATUS_COLOR[d.status]};border-radius:3px;width:${Math.min(d.atual,100)}%"></div>
      </div>
      <span style="width:44px;color:#3A6A48;font-size:10px;text-align:right">meta ${d.alvo.toFixed(1)}%</span>
      <span style="width:56px;font-weight:700;font-size:11px;color:${d.desvio>0?"#FF6B6B":"#5DCAA5"};text-align:right">${d.desvio>0?"+":""}${d.desvio.toFixed(1)}%</span>
      <span style="font-size:13px">${STATUS_ICON[d.status]||""}</span>
    </div>
  `).join("");

  // Ações
  if(xp.acoes && xp.acoes.length){
    document.getElementById("hp-acoes-bloco").style.display = "";
    const SINAL = v => v > 0 ? `<span style="color:#5DCAA5">+${v.toFixed(2)}%</span>` : `<span style="color:#FF6B6B">${v.toFixed(2)}%</span>`;
    document.getElementById("hp-acoes-table").innerHTML = `
      <table style="width:100%;border-collapse:collapse;font-size:11px">
        <thead><tr style="background:#071E17">
          <th style="text-align:left;padding:6px 8px;color:#4A7055;font-size:10px">Ticker</th>
          <th style="text-align:right;padding:6px 8px;color:#4A7055;font-size:10px">Qtd.</th>
          <th style="text-align:right;padding:6px 8px;color:#4A7055;font-size:10px">Saldo</th>
          <th style="text-align:right;padding:6px 8px;color:#4A7055;font-size:10px">% Cart.</th>
        </tr></thead>
        <tbody>${xp.acoes.map((a,i)=>`
          <tr style="background:${i%2?"#070707":"#071E17"}">
            <td style="padding:6px 8px;color:#D4B483;font-weight:700">${a.ticker}</td>
            <td style="padding:6px 8px;color:#888;text-align:right">${a.qtd||"—"}</td>
            <td style="padding:6px 8px;color:#CCC;text-align:right">R$ ${(a.saldo||0).toLocaleString("pt-BR",{minimumFractionDigits:2})}</td>
            <td style="padding:6px 8px;text-align:right">${(a.perc||0).toFixed(2)}%</td>
          </tr>`).join("")}
        </tbody>
      </table>`;
  }

  // FIIs
  if(xp.fiis && xp.fiis.length){
    document.getElementById("hp-fiis-bloco").style.display = "";
    document.getElementById("hp-fiis-table").innerHTML = `
      <table style="width:100%;border-collapse:collapse;font-size:11px">
        <thead><tr style="background:#071E17">
          <th style="text-align:left;padding:6px 8px;color:#4A7055;font-size:10px">Ticker</th>
          <th style="text-align:right;padding:6px 8px;color:#4A7055;font-size:10px">Qtd.</th>
          <th style="text-align:right;padding:6px 8px;color:#4A7055;font-size:10px">Saldo</th>
          <th style="text-align:right;padding:6px 8px;color:#4A7055;font-size:10px">% Cart.</th>
        </tr></thead>
        <tbody>${xp.fiis.map((f,i)=>`
          <tr style="background:${i%2?"#070707":"#071E17"}">
            <td style="padding:6px 8px;color:#7DCFEF;font-weight:700">${f.ticker}</td>
            <td style="padding:6px 8px;color:#888;text-align:right">${f.qtd||"—"}</td>
            <td style="padding:6px 8px;color:#CCC;text-align:right">R$ ${(f.saldo||0).toLocaleString("pt-BR",{minimumFractionDigits:2})}</td>
            <td style="padding:6px 8px;text-align:right">${(f.perc||0).toFixed(2)}%</td>
          </tr>`).join("")}
        </tbody>
      </table>`;
  }

  // Sugestões de produtos HP para fechar gaps
  const sugs = xp.sugestoes_produtos || [];
  if(sugs.length){
    document.getElementById("hp-sugestoes").style.display = "";
    document.getElementById("hp-sugestoes-lista").innerHTML = sugs.map(s=>{
      const p = s.produto || {};
      return `
        <div style="background:#071E17;border:1px solid #2A2A18;border-left:3px solid #D4B483;border-radius:0 8px 8px 0;padding:10px 14px;margin-bottom:8px">
          <div style="display:flex;align-items:center;gap:8px;margin-bottom:4px;flex-wrap:wrap">
            <span style="font-size:10px;background:#1A1A08;color:#D4B483;border:1px solid #D4B483;border-radius:10px;padding:1px 8px">${s.label_classe}</span>
            <span style="font-size:12px;color:#FF6B6B;font-weight:700">gap ${s.gap.toFixed(1)}%</span>
            <span style="font-size:13px;font-weight:700;color:#F0F0F0">${p.nome||"—"}</span>
            ${p.taxa ? `<span style="font-size:11px;color:#5DCAA5">${p.taxa}</span>` : ""}
          </div>
          ${p.motivo ? `<p style="font-size:11px;color:#888;line-height:1.5;margin:0">${p.motivo}</p>` : ""}
          ${p.indicado_por ? `<p style="font-size:10px;color:#2A5A3A;margin-top:4px">Indicado por: ${p.indicado_por} ${p.indicado_em ? "· "+p.indicado_em : ""}</p>` : ""}
        </div>`;
    }).join("");
  }

  // Cenário macro HP — 3 cards estruturados
  const c = xp.cenario_macro;
  if(c && (c.global || c.brasil || c.posicionamento)){
    document.getElementById("hp-cenario").style.display = "";
    const ref = document.getElementById("hp-cenario-ref");
    if(ref && c.referencia) ref.textContent = c.referencia;

    const cards = [
      {titulo:"🌍 Global",         texto: c.global,         cor:"#8B9FE8"},
      {titulo:"🇧🇷 Brasil",        texto: c.brasil,         cor:"#5DCAA5"},
      {titulo:"📌 Posicionamento", texto: c.posicionamento, cor:"#D4B483"},
    ];
    document.getElementById("hp-cenario-cards").innerHTML = cards.map(card => card.texto ? `
      <div style="background:#060F0B;border:1px solid #1A1A10;border-radius:10px;padding:12px;border-top:2px solid ${card.cor}">
        <div style="font-size:10px;color:${card.cor};font-weight:700;text-transform:uppercase;letter-spacing:.5px;margin-bottom:8px">${card.titulo}</div>
        <p style="font-size:11px;color:#CCC;line-height:1.7">${card.texto}</p>
      </div>` : "").join("");

    const vieses = c.vieses || {};
    const VCOR = {positivo:"#5DCAA5", neutro:"#D4B483", negativo:"#FF6B6B"};
    const VLBL = {pos_fixado:"Pós Fix",inflacao:"Inflação",pre_fixado:"Pré Fix",acoes:"Ações",fiis:"FIIs",multimercado:"Multi",internacional:"Intl",alternativos:"Altern",criptomoedas:"Cripto"};
    document.getElementById("hp-cenario-vieses").innerHTML = Object.entries(vieses)
      .filter(([,v])=>v)
      .map(([cls,v])=>`<span style="font-size:10px;background:#0A1A0A;border:1px solid ${VCOR[v]||"#333"}44;border-radius:10px;padding:3px 9px;color:${VCOR[v]||"#888"}">${VLBL[cls]||cls}: ${v}</span>`)
      .join("");
  }

  card.scrollIntoView({behavior:"smooth", block:"start"});
}

async function analisar(){
  const assessor=document.getElementById("assessor").value.trim();
  const nome=document.getElementById("nome").value.trim();
  const perfil=document.getElementById("perfil").value;
  const objetivo=document.getElementById("objetivo").value.trim();
  const fileXP=document.getElementById("pdf-xp").files[0];

  if(!assessor){alert("Digite seu nome como assessor.");return;}
  if(!nome){alert("Digite o nome do cliente.");return;}
  if(!fileXP){alert("Selecione o PDF do relatório XP.");return;}

  document.getElementById("spinner").classList.add("show");
  document.getElementById("results").style.display="none";
  document.getElementById("btn-ana").disabled=true;

  const fd=new FormData();
  fd.append("assessor",assessor);
  fd.append("nome",nome); fd.append("perfil",perfil); fd.append("objetivo",objetivo);
  fd.append("pdf",fileXP);
  fd.append("checklist",JSON.stringify(checklist));
  fd.append("cross_sell",JSON.stringify(crossSell));

  try{
    // Faz os dois requests em paralelo
    const fd2 = new FormData();
    fd2.append("assessor", assessor);
    fd2.append("pdf", fileXP);

    const [res, resXP] = await Promise.all([
      fetch("/api/analyze", {method:"POST", body:fd}),
      fetch("/api/analyze-xp", {method:"POST", body:fd2}),
    ]);
    if(!res.ok) throw new Error(await res.text());
    analiseData = await res.json();

    if(resXP.ok){
      const xpData = await resXP.json();
      analiseData._xp = xpData;
    }

    renderizar(analiseData);
    if(analiseData._xp) renderAnaliseHP(analiseData._xp);

    // Log de atividade
    const _nomeLog = localStorage.getItem("brauna_nome") || document.getElementById("assessor")?.value || "Assessor";
    fetch("/api/admin/activity",{method:"POST",headers:{"Content-Type":"application/json"},
      body:JSON.stringify({role:"assessor",nome:_nomeLog,acao:"análise",detalhe:`Analisou carteira de ${analiseData.nome||"cliente"} — perfil ${analiseData.perfil||""}`})});

    // Salva snapshot da carteira para comparativo futuro
    if(_clienteIdentificado?.conta && analiseData._xp){
      const xp = analiseData._xp;
      fetch("/api/salvar-carteira", {
        method:"POST",
        headers:{"Content-Type":"application/json"},
        body: JSON.stringify({
          conta: _clienteIdentificado.conta,
          data_ref: xp.data_ref || "",
          patrimonio: xp.patrimonio || 0,
          composicao: xp.comp || {},
          rent: xp.rent || {},
        })
      }).catch(()=>{});
    }

    // Salva ficha atualizada (modelo de servir + cross-sell)
    // Salva sempre — usa conta se disponível, senão salva por nome+assessor
    fetch("/api/ficha", {
      method:"POST",
      headers:{"Content-Type":"application/json"},
      body: JSON.stringify({
        conta: _clienteIdentificado?.conta || "",
        nome, perfil, objetivo, assessor,
        checklist,
        cross_ativos: Object.keys(crossSell).filter(k=>crossSell[k]),
      })
    }).catch(()=>{});
  }catch(e){ alert("Erro: "+e.message); }
  finally{
    document.getElementById("spinner").classList.remove("show");
    document.getElementById("btn-ana").disabled=false;
  }
}

function renderSugestoes(sg){
  const el = document.getElementById("sugestoes-list");
  const btn = document.getElementById("tab-btn-sugestoes");
  if(!sg){ el.innerHTML='<p class="sg-vazio">Nenhuma sugestão publicada pelo administrador.</p>'; return; }
  btn.style.display="";

  const totalItens = (sg.renda_fixa?.length||0)+(sg.renda_variavel?.length||0)+(sg.internacional?.length||0)+(sg.fiis?.length||0)+(sg.estruturadas?1:0);
  if(!totalItens){ el.innerHTML='<p class="sg-vazio">✓ Nenhuma sugestão aplicável para este perfil e carteira.</p>'; return; }

  let html = `<div class="sg-header-box">
    <span style="font-size:20px">💡</span>
    <div>
      <div style="font-size:12px;font-weight:700;color:#C9A96E">${sg.titulo||"Sugestões do Gestor"}</div>
      <div style="font-size:10px;color:#3A6A48">Publicado em: ${sg.criado_em||"—"} · Filtrado para o perfil e carteira deste cliente</div>
    </div>
  </div>`;

  // ── Renda Fixa ────────────────────────────────────────────
  if(sg.renda_fixa?.length){
    html += `<div class="sg-bloco"><div class="sg-bloco-title">🏦 Renda Fixa — ${sg.renda_fixa.length} sugestão(ões)</div>`;
    sg.renda_fixa.forEach(item=>{
      const urgCls = item.urgencia==="alta"?"urg-alta":item.urgencia==="media"?"urg-media":"urg-baixa";
      const urgTxt = item.urgencia==="alta"?"🔴 Urgente":item.urgencia==="media"?"🟡 Médio prazo":"🟢 Oportunidade";
      html += `<div class="sg-card ${urgCls}">
        <div class="sg-topo">
          <span class="sg-acao">${item.acao==="substituir"?"↔ Substituir":"+ Alocar"}</span>
          <span class="sg-produto">${item.para||"—"}</span>
          ${item.indexador?`<span class="sg-idx">${item.indexador}</span>`:""}
          <span style="font-size:10px;color:#3A6A48">${urgTxt}</span>
        </div>
        ${item.de?`<div class="sg-de">↓ Sair / reduzir: ${item.de}</div>`:""}
        <div class="sg-motivo">${item.motivo||""}</div>
      </div>`;
    });
    html += `</div>`;
  }

  // ── Renda Variável ────────────────────────────────────────
  if(sg.renda_variavel?.length){
    html += `<div class="sg-bloco"><div class="sg-bloco-title">📈 Renda Variável — ${sg.renda_variavel.length} sugestão(ões)</div>`;
    sg.renda_variavel.forEach(item=>{
      const acaoTxt = item.acao==="comprar"?"▲ Comprar":item.acao==="vender"?"▼ Vender":"↔ Trocar";
      const acaoCor = item.acao==="comprar"?"#5DCAA5":item.acao==="vender"?"#FF6B6B":"#FFD966";
      html += `<div class="sg-card">
        <div class="sg-topo">
          <span class="sg-acao" style="color:${acaoCor}">${acaoTxt}</span>
          ${item.de?`<span class="sg-produto">${item.de}</span>`:""}
          ${item.para?`<span style="font-size:12px;color:#5DCAA5">→ ${item.para}</span>`:""}
        </div>
        <div class="sg-motivo">${item.motivo||""}</div>
        ${item.fonte?`<div class="sg-fonte">Fonte: ${item.fonte}</div>`:""}
      </div>`;
    });
    html += `</div>`;
  }

  // ── Estruturadas ──────────────────────────────────────────
  if(sg.estruturadas){
    html += `<div class="sg-bloco">
      <div class="sg-bloco-title">⚙️ Operações Estruturadas <span style="font-size:10px;font-weight:400;color:#3A6A48;text-transform:none;letter-spacing:0">(cliente tem ações na carteira)</span></div>
      <div class="sg-est-box">${sg.estruturadas}</div>
    </div>`;
  }

  // ── Internacional ─────────────────────────────────────────
  if(sg.internacional?.length){
    const jaTem = sg.tem_intl;
    html += `<div class="sg-bloco">
      <div class="sg-bloco-title">🌎 Internacional — ${sg.internacional.length} sugestão(ões)</div>
      ${jaTem?`<div style="font-size:11px;color:#5DCAA5;margin-bottom:8px;padding:6px 10px;background:#0A2018;border-radius:6px">✓ Cliente já possui alocação internacional — sugestões abaixo para otimização/reposicionamento</div>`:""}`;
    sg.internacional.forEach(item=>{
      html += `<div class="sg-card ${jaTem?"sg-ja-tem":""}">
        ${jaTem?`<div class="sg-ja-tem-badge">✓ Já possui internacional — otimizar</div>`:""}
        <div class="sg-topo">
          <span class="sg-acao">${item.tipo||"ETF"}</span>
          <span class="sg-produto">${item.ativo||"—"}</span>
        </div>
        <div class="sg-motivo">${item.motivo||""}</div>
      </div>`;
    });
    html += `</div>`;
  }

  // ── FIIs ──────────────────────────────────────────────────
  if(sg.fiis?.length){
    const jaTem = sg.tem_fiis;
    html += `<div class="sg-bloco">
      <div class="sg-bloco-title">🏢 Fundos Imobiliários — Carteira sugerida</div>
      ${jaTem?`<div style="font-size:11px;color:#5DCAA5;margin-bottom:8px;padding:6px 10px;background:#0A2018;border-radius:6px">✓ Cliente já possui FIIs — use as sugestões abaixo para rebalancear a posição</div>`:""}
      <table class="sg-fii-table">
        <thead><tr><th>Ticker</th><th>Segmento</th><th>% sugerido</th><th>Motivo</th></tr></thead>
        <tbody>${sg.fiis.map(f=>`<tr>
          <td class="fii-ticker">${f.ticker}</td>
          <td>${f.segmento||"—"}</td>
          <td>${f.alocacao||"—"}</td>
          <td style="color:#888">${f.motivo||""}</td>
        </tr>`).join("")}</tbody>
      </table>
    </div>`;
  }

  el.innerHTML = html;
}

function renderClassesAtivos(desvios, patrimonio){
  const grupos = [
    { nome:"Renda Fixa",     cor:"#C9A96E", icone:"🏦", chaves:["pos_fixado","inflacao","pre_fixado"] },
    { nome:"Renda Variável", cor:"#8B9FE8", icone:"📈", chaves:["acoes","fiis"] },
    { nome:"Multimercado",   cor:"#5DCAA5", icone:"⚖️", chaves:["multimercado"] },
    { nome:"Internacional",  cor:"#E88B8B", icone:"🌎", chaves:["internacional"] },
    { nome:"Alternativos",   cor:"#B8A0E8", icone:"💎", chaves:["alternativos","criptomoedas"] },
  ];
  const labelMap={"Pós Fixado":"pos_fixado","Inflação":"inflacao","Pré Fixado":"pre_fixado",
    "Ações":"acoes","FIIs":"fiis","Multimercado":"multimercado",
    "Internacional":"internacional","Alternativos":"alternativos","Criptomoedas":"criptomoedas"};
  const byKey={};
  desvios.forEach(d=>{ const k=labelMap[d.label]; if(k) byKey[k]={real:d.real,alvo:d.alvo}; });

  // Calcula valores por grupo
  const dados = grupos.map(g=>({
    ...g,
    real: g.chaves.reduce((s,k)=>s+(byKey[k]?.real||0),0),
    alvo: g.chaves.reduce((s,k)=>s+(byKey[k]?.alvo||0),0),
  }));

  // ── Gráfico de pizza ──────────────────────────────────────────────────────
  const canvas = document.getElementById("chart-classes");
  if(canvas && typeof Chart !== "undefined"){
    const vistos = dados.filter(g=>g.real>0.1);
    const pizzaData  = vistos.map(g=>parseFloat(g.real.toFixed(1)));
    const pizzaCores = vistos.map(g=>g.cor);
    const pizzaLbls  = vistos.map(g=>g.nome);

    if(_classesChart){ try{ _classesChart.destroy(); }catch(e){} }
    _classesChart = new Chart(canvas, {
      type: "doughnut",
      data: {
        labels: pizzaLbls,
        datasets: [{
          data: pizzaData,
          backgroundColor: pizzaCores,
          borderColor: "#081F18",
          borderWidth: 3,
          hoverOffset: 10,
        }]
      },
      options: {
        cutout: "52%",
        plugins: {
          legend: { display: false },
          tooltip: {
            callbacks: {
              label: ctx => ` ${ctx.label}: ${ctx.parsed.toFixed(1)}%`
            }
          }
        },
        animation: { duration: 600 }
      }
    });

    // Legenda manual com cor + nome + %
    const legEl = document.getElementById("classes-legenda");
    if(legEl) legEl.innerHTML = dados.map(g=>{
      if(g.real < 0.1) return "";
      const desvio = g.real - g.alvo;
      const dCor = Math.abs(desvio)<1.5?"#3A6A48":desvio>0?"#FFD966":"#FF6B6B";
      return `<div style="display:flex;align-items:center;gap:7px">
        <span style="width:12px;height:12px;border-radius:3px;background:${g.cor};flex-shrink:0;display:inline-block"></span>
        <span style="font-size:11px;color:#AAA;flex:1">${g.nome}</span>
        <span style="font-size:12px;font-weight:700;color:${g.cor}">${g.real.toFixed(1)}%</span>
        <span style="font-size:10px;color:${dCor};min-width:40px;text-align:right">${desvio>=0?"+":""}${desvio.toFixed(1)}%</span>
      </div>`;
    }).join("");
  }

  // ── Cards por grupo ───────────────────────────────────────────────────────
  const el = document.getElementById("classes-grid");
  el.innerHTML = dados.map(g=>{
    const desvio = g.real - g.alvo;
    const sinal  = desvio > 0 ? "+" : "";
    const dCor   = Math.abs(desvio)<1.5 ? "#3A6A48" : desvio>0 ? "#FFD966" : "#FF6B6B";
    const realVal = patrimonio>0 ? `R$ ${(patrimonio*g.real/100).toLocaleString("pt-BR",{maximumFractionDigits:0})}` : "—";
    return `<div style="background:#111;border-radius:10px;padding:12px 14px;border:1px solid #1E1E1E">
      <div style="display:flex;align-items:center;gap:8px;margin-bottom:8px">
        <span style="font-size:18px">${g.icone}</span>
        <div style="flex:1;min-width:0">
          <div style="font-size:11px;font-weight:700;color:#F0F0F0">${g.nome}</div>
          <div style="font-size:10px;color:#2A5A3A">${realVal}</div>
        </div>
        <div style="text-align:right;flex-shrink:0">
          <div style="font-size:17px;font-weight:800;color:${g.cor}">${g.real.toFixed(1)}%</div>
          <div style="font-size:10px;color:${dCor}">${sinal}${desvio.toFixed(1)}% meta</div>
        </div>
      </div>
      <div style="height:5px;background:#1A1A1A;border-radius:3px;margin-bottom:3px;overflow:hidden">
        <div style="height:100%;background:${g.cor};border-radius:3px;width:${Math.min(g.real/70*100,100)}%;transition:width .6s ease"></div>
      </div>
      <div style="height:3px;background:#1A1A1A;border-radius:3px;overflow:hidden">
        <div style="height:100%;background:#1C4A34;border-radius:3px;width:${Math.min(g.alvo/70*100,100)}%"></div>
      </div>
      <div style="display:flex;justify-content:space-between;margin-top:3px;font-size:10px">
        <span style="color:${g.cor}">atual</span><span style="color:#1E4A30">meta ${g.alvo.toFixed(1)}%</span>
      </div>
    </div>`;
  }).join("");
}

function renderPlanoAcao(desvios, perfil, patrimonio, objetivo){
  document.getElementById("plano-perfil").textContent=perfil;

  const objBox=document.getElementById("plano-objetivo-box");
  if(objetivo){ objBox.style.display="block"; document.getElementById("plano-objetivo-txt").textContent=objetivo; }
  else{ objBox.style.display="none"; }

  const aumentar=desvios.filter(d=>d.desvio<-1);
  const reduzir =desvios.filter(d=>d.desvio>1);
  const alinhado=desvios.filter(d=>Math.abs(d.desvio)<=1);

  const makeItem=(d,tipo)=>{
    const abs=Math.abs(d.desvio);
    const val=patrimonio>0?` ≈ R$ ${(patrimonio*abs/100).toLocaleString("pt-BR",{maximumFractionDigits:0})}`:"";
    const col=tipo==="aumentar"?"#5DCAA5":"#FF6B6B";
    const seta=tipo==="aumentar"?"▲":"▼";
    return `<div style="display:flex;align-items:center;gap:10px;padding:9px 12px;background:#081F18;border-radius:8px;margin-bottom:6px;border:1px solid #1E1E1E">
      <span style="color:${col};font-size:16px;font-weight:700">${seta}</span>
      <div style="flex:1">
        <div style="font-size:12px;font-weight:700;color:#F0F0F0">${d.label}</div>
        <div style="font-size:10px;color:#3A6A48">${d.real.toFixed(1)}% atual → meta ${d.alvo.toFixed(1)}%${val}</div>
      </div>
      <span style="font-size:13px;font-weight:800;color:${col}">${tipo==="aumentar"?"+":""}${(-d.desvio).toFixed(1)}%</span>
    </div>`;
  };

  const pg=document.getElementById("plano-grid");
  pg.innerHTML=`
    <div>
      <p style="font-size:10px;color:#5DCAA5;font-weight:700;text-transform:uppercase;letter-spacing:.5px;margin-bottom:8px">▲ Aumentar / Comprar</p>
      ${aumentar.length?aumentar.map(d=>makeItem(d,"aumentar")).join(""):'<p style="font-size:12px;color:#2A5A3A">Nenhum ativo abaixo da meta</p>'}
    </div>
    <div>
      <p style="font-size:10px;color:#FF6B6B;font-weight:700;text-transform:uppercase;letter-spacing:.5px;margin-bottom:8px">▼ Reduzir / Resgatar</p>
      ${reduzir.length?reduzir.map(d=>makeItem(d,"reduzir")).join(""):'<p style="font-size:12px;color:#2A5A3A">Nenhum ativo acima da meta</p>'}
    </div>`;

  const pv=document.getElementById("plano-valores");
  if(alinhado.length){
    pv.innerHTML=`<p style="font-size:10px;color:#3A6A48;margin-bottom:6px;text-transform:uppercase;letter-spacing:.5px">✓ Classes alinhadas ao modelo</p>`
      +alinhado.map(d=>`<span style="display:inline-block;margin:3px 4px;padding:4px 10px;border-radius:12px;background:#1A1A1A;border:1px solid #1C4A34;font-size:11px;color:#3A6A48">${d.label} ${d.real.toFixed(1)}%</span>`).join("");
  } else { pv.innerHTML=""; }
}

function renderizar(data){
  const{desvios,rent,patrimonio,macro,recomendacoes,alertas,visao_gestores,checklist_servir,score_servir,pendentes_criticos,pendentes_altos,objetivo,perfil}=data;

  // Métricas
  const m=document.getElementById("metrics");
  let pat=patrimonio>0?"R$ "+patrimonio.toLocaleString("pt-BR",{maximumFractionDigits:0}):"—";
  let html=`<div class="metric"><div class="lbl">Patrimônio</div><div class="val" style="font-size:16px">${pat}</div></div>`;
  if(rent.portfolio&&rent.cdi){
    const p=rent.portfolio,c=rent.cdi;
    const p12=(p["12m"]/c["12m"]*100).toFixed(1);
    const cls=parseFloat(p12)<70?"danger":"ok";
    html+=`<div class="metric"><div class="lbl">Rentab. 12M</div><div class="val ${cls}">${p["12m"].toFixed(2)}%</div></div>`;
    html+=`<div class="metric"><div class="lbl">CDI 12M</div><div class="val">${c["12m"].toFixed(2)}%</div></div>`;
    html+=`<div class="metric"><div class="lbl">% do CDI (12M)</div><div class="val ${cls}">${p12}%</div></div>`;
  }
  // Score Modelo de Servir como métrica
  if(checklist_servir){
    const scoreColor=score_servir>=5?"ok":score_servir>=3?"":"danger";
    const pendentes_list = checklist_servir.filter(p=>p.status==="pendente");
    const barColor = score_servir>=5?"#5DCAA5":score_servir>=3?"#FFD966":"#FF6B6B";
    const bordaColor = pendentes_criticos>0?"#FF4444":score_servir>=5?"#2A5040":"#1C4A34";
    const bgColor = pendentes_criticos>0?"#2A1010":score_servir>=5?"#0A2A18":"#0B2A1F";

    let insightsHtml = "";
    if(pendentes_list.length > 0){
      insightsHtml = `<div style="margin-top:14px;border-top:1px solid #1C4A34;padding-top:12px">
        <div style="font-size:10px;color:#FF6B6B;font-weight:700;text-transform:uppercase;letter-spacing:.5px;margin-bottom:10px">
          ⚠ Insights — o que fazer para completar o modelo
        </div>`;
      pendentes_list.forEach(p=>{
        const imp_cor = p.importancia.startsWith("CRÍTICA")?"#FF6B6B":p.importancia==="ALTA"?"#FFD966":"#888";
        insightsHtml += `<div style="margin-bottom:10px;padding:10px 12px;background:#081F18;border-left:3px solid ${imp_cor};border-radius:0 6px 6px 0">
          <div style="display:flex;align-items:center;gap:8px;margin-bottom:4px">
            <span style="font-size:16px">${p.icone}</span>
            <span style="font-size:12px;font-weight:700;color:#F0F0F0">${p.nome}</span>
            <span style="font-size:10px;padding:2px 7px;background:${imp_cor}22;color:${imp_cor};border-radius:10px;font-weight:700">${p.importancia}</span>
          </div>
          <div style="font-size:11px;color:#888;margin-bottom:4px">${p.impacto_falta}</div>
          <div style="font-size:11px;color:#C9A96E;font-weight:600">→ ${p.acao}</div>
          ${p.diretriz?`<div style="margin-top:6px;padding:6px 8px;background:#1A2E1A;border:1px solid #C9A96E33;border-radius:5px;font-size:10px;color:#C9A96E">★ 1ª DIRETRIZ BRAÚNA — Provocar o cliente sobre o Financial Planning.</div>`:""}
        </div>`;
      });
      insightsHtml += `</div>`;
    } else {
      insightsHtml = `<div style="margin-top:10px;font-size:12px;color:#5DCAA5">✓ Todos os pilares completos — cliente no modelo ideal.</div>`;
    }

    html+=`<div class="metric" style="grid-column:span 4;background:${bgColor};border:1px solid ${bordaColor}">
      <div class="lbl">Modelo de Servir</div>
      <div style="display:flex;align-items:center;gap:14px;margin-top:4px">
        <div class="val ${scoreColor}" style="font-size:28px;line-height:1">${score_servir}/6</div>
        <div style="flex:1">
          <div style="font-size:11px;color:#3A6A48;margin-bottom:5px">pilares completos</div>
          <div style="height:6px;background:#1A4030;border-radius:3px;overflow:hidden">
            <div style="width:${score_servir/6*100}%;height:100%;background:${barColor};border-radius:3px"></div>
          </div>
        </div>
      </div>
      <div style="font-size:11px;color:#888;margin-top:6px">${pendentes_criticos>0?`🔴 ${pendentes_criticos} crítico(s) pendente(s) `:''}${pendentes_altos>0?`🟡 ${pendentes_altos} importante(s) pendente(s)`:score_servir>=6?'✓ Todos completos':''}</div>
      ${insightsHtml}
    </div>`;
  }
  m.innerHTML=html;

  // Gráfico
  const lbs=desvios.map(d=>d.label), rea=desvios.map(d=>d.real), alv=desvios.map(d=>d.alvo);
  if(chartInst) chartInst.destroy();
  chartInst=new Chart(document.getElementById("chart"),{
    type:"bar",
    data:{labels:lbs,datasets:[{label:"Atual",data:rea,backgroundColor:"#C9A96E",borderRadius:4},{label:"Modelo",data:alv,backgroundColor:"rgba(100,100,100,0.4)",borderRadius:4}]},
    options:{responsive:true,maintainAspectRatio:false,plugins:{legend:{display:false},tooltip:{callbacks:{label:ctx=>ctx.dataset.label+": "+ctx.parsed.y.toFixed(2)+"%"}}},scales:{x:{ticks:{color:"#4A7055",font:{size:10}}},y:{ticks:{color:"#4A7055",callback:v=>v+"%"},beginAtZero:true}}}
  });

  // Diversificação por classe e plano de ação
  renderClassesAtivos(desvios, patrimonio);
  renderPlanoAcao(desvios, data.perfil||"", patrimonio, data.objetivo||"");
  renderSugestoes(data.sugestoes||null);

  // Desvios
  const dl=document.getElementById("desvios-list");
  dl.innerHTML=desvios.map(d=>{
    const cls=d.desvio<-1.5?"red":d.desvio>1.5?"grn":"dim";
    const sinal=d.desvio>0?"+":"";
    return `<div class="desvio-row"><span class="dn">${d.label}</span><span class="da">${d.real.toFixed(1)}%</span><div class="db-w"><div class="db" style="width:${Math.min(d.real/80*100,100)}%"></div></div><span class="dalvo">alvo ${d.alvo.toFixed(1)}%</span><span class="dpp ${cls}">${sinal}${d.desvio.toFixed(2)}%</span></div>`;
  }).join("");

  // Diagnóstico
  const frags=desvios.filter(d=>d.desvio<-1.5), opors=desvios.filter(d=>d.desvio>1.5);
  let diag="";
  if(alertas&&alertas.length){ alertas.forEach(a=>{diag+=`<div class="alert warn">${a}</div>`;}); }
  if(frags.length){ diag+=`<p style="font-size:10px;color:#3A6A48;margin:8px 0 4px;text-transform:uppercase">Fragilidades</p>`; frags.forEach(f=>{diag+=`<div class="alert danger">▸ <b>${f.label}</b>: falta ${Math.abs(f.desvio).toFixed(1)}%</div>`;}); }
  if(opors.length){ diag+=`<p style="font-size:10px;color:#3A6A48;margin:8px 0 4px;text-transform:uppercase">Sobrealocação</p>`; opors.forEach(o=>{diag+=`<div class="alert success">▸ <b>${o.label}</b>: excesso ${o.desvio.toFixed(1)}%</div>`;}); }
  document.getElementById("diagnostico").innerHTML=diag||`<p style="color:#5DCAA5">✓ Carteira alinhada ao modelo.</p>`;

  // Recomendações — inclui painel de calls HP
  const rl=document.getElementById("recomendacoes-list");
  let recHtml = "";

  // Painel de calls HP relevantes (vem de analiseData._xp.calls_relevantes)
  const callsRel = analiseData._xp?.calls_relevantes || [];
  if(callsRel.length){
    const DIR_COR   = {compra:"#5DCAA5", neutro:"#AAA", venda:"#FF6B6B"};
    const DIR_LABEL = {compra:"▲ COMPRA", neutro:"→ NEUTRO", venda:"▼ VENDA"};
    recHtml += `<div style="border:2px solid #5DCAA5;border-radius:10px;overflow:hidden;margin-bottom:16px">
      <div style="background:#0A2018;padding:10px 14px;display:flex;align-items:center;gap:8px">
        <span style="font-size:15px">📈</span>
        <span style="font-size:13px;font-weight:700;color:#5DCAA5">Calls do HP para este cliente</span>
        <span style="font-size:11px;color:#2A5A3A;margin-left:4px">${callsRel.length} call(s) relevante(s)</span>
      </div>
      <div style="padding:10px 14px;display:flex;flex-direction:column;gap:8px">
        ${callsRel.map(c=>{
          const cor = DIR_COR[c.direcao]||"#AAA";
          const upCor = (c.upside||0) >= 0 ? "#5DCAA5" : "#FF6B6B";
          const naCarteira = (analiseData._xp?.acoes||[]).some(a=>a.ticker===c.ticker) || (analiseData._xp?.fiis||[]).some(f=>f.ticker===c.ticker);
          const badge = naCarteira
            ? `<span style="font-size:10px;background:#0A2018;color:#5DCAA5;border:1px solid #2A5040;padding:2px 8px;border-radius:10px">✓ já na carteira</span>`
            : `<span style="font-size:10px;background:#1A2E1A;color:#C9A96E;border:1px solid #C9A96E44;padding:2px 8px;border-radius:10px">oportunidade de entrada</span>`;
          return `<div style="background:#0A0F0B;border:1px solid #1A2E3A;border-left:3px solid ${cor};border-radius:0 8px 8px 0;padding:10px 12px">
            <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;margin-bottom:4px">
              <span style="font-size:14px;font-weight:900;color:${cor}">${c.ticker}</span>
              <span style="font-size:10px;font-weight:700;padding:2px 8px;border-radius:8px;background:${cor}22;color:${cor}">${DIR_LABEL[c.direcao]||c.direcao}</span>
              ${c.nome ? `<span style="font-size:12px;color:#CCC">${c.nome}</span>` : ""}
              ${badge}
            </div>
            <div style="display:flex;gap:12px;flex-wrap:wrap;font-size:12px;margin-bottom:4px">
              ${c.preco_entrada ? `<span style="color:#888">Entrada: <b style="color:#F0F0F0">R$ ${c.preco_entrada.toLocaleString("pt-BR",{minimumFractionDigits:2})}</b></span>` : ""}
              ${c.preco_alvo ? `<span style="color:#888">→ Alvo: <b style="color:#F0F0F0">R$ ${c.preco_alvo.toLocaleString("pt-BR",{minimumFractionDigits:2})}</b></span>` : ""}
              ${c.upside !== undefined ? `<span style="color:${upCor};font-weight:700">Upside ${c.upside>=0?"+":""}${c.upside}%</span>` : ""}
              ${c.stop ? `<span style="color:#888">Stop: <b style="color:#FF6B6B">R$ ${c.stop.toLocaleString("pt-BR",{minimumFractionDigits:2})}</b></span>` : ""}
              ${c.prazo ? `<span style="color:#888">Prazo: ${c.prazo}</span>` : ""}
            </div>
            ${c.tese ? `<p style="font-size:11px;color:#888;margin:0 0 3px;line-height:1.5;font-style:italic">"${c.tese.length>160?c.tese.slice(0,160)+"...":c.tese}"</p>` : ""}
            ${c.fonte ? `<div style="font-size:10px;color:#2A5A3A">${c.fonte}</div>` : ""}
          </div>`;
        }).join("")}
      </div>
    </div>`;
  }

  if(recomendacoes&&recomendacoes.length){
    recHtml+=recomendacoes.map(r=>`
      <div class="rec-card">
        <div class="rec-header">${r.urgencia} <span style="color:#C9A96E">${r.classe}</span> — falta ${r.falta_pp.toFixed(1)}%</div>
        <div class="rec-ctx">${r.explicacao}</div>
        ${r.carta_insight?`<div class="rec-carta">📄 Carta da gestão: "${r.carta_insight}"</div>`:""}
        <div class="rec-prods">${(r.produtos||[]).map(p=>`<span class="prod-tag">${p}</span>`).join("")}</div>
      </div>`).join("");
  } else if(!callsRel.length){
    recHtml+=`<p style="color:#5DCAA5;font-size:14px">✓ Nenhuma recomendação necessária — carteira alinhada.</p>`;
  }
  rl.innerHTML = recHtml;

  // Modelo de Servir
  const sl=document.getElementById("servir-list");
  if(checklist_servir){
    const pendentes=checklist_servir.filter(p=>p.status==="pendente");
    const concluidos=checklist_servir.filter(p=>p.status==="ok");
    const barColor=score_servir>=5?"#5DCAA5":score_servir>=3?"#FFD966":"#FF6B6B";
    let sh=`<div class="servir-score">
      <div class="score-num" style="color:${barColor}">${score_servir}/6</div>
      <div class="score-info">
        <div class="score-label">Pilares do Modelo de Servir completos</div>
        <div class="score-bar-w"><div class="score-bar" style="width:${score_servir/6*100}%;background:${barColor}"></div></div>
        <div class="score-pendentes">${pendentes.length>0?`${pendentes.filter(p=>p.importancia.startsWith("CRÍTICA")).length} crítico(s) · ${pendentes.filter(p=>p.importancia==="ALTA").length} importante(s) pendentes`:"Todos os pilares completos ✓"}</div>
      </div>
    </div>`;

    if(pendentes.length>0){
      sh+=`<p style="font-size:11px;color:#FF6B6B;font-weight:700;margin-bottom:10px;text-transform:uppercase;letter-spacing:.5px">⚠ O que está faltando para este cliente</p>`;
      pendentes.forEach(p=>{
        const cls=p.importancia.startsWith("CRÍTICA")?"crit":p.importancia==="ALTA"?"high":"med";
        sh+=`<div class="pilar-result ${cls}">
          <div class="pilar-result-header">
            <span style="font-size:18px">${p.icone}</span>
            <span class="pilar-result-nome">${p.nome}</span>
            <span class="pilar-result-status status-pendente">PENDENTE</span>
            <span class="badge ${p.importancia.startsWith("CRÍTICA")?"badge-crit":p.importancia==="ALTA"?"badge-high":"badge-med"}" style="margin-left:4px">${p.importancia}</span>
            ${p.diretriz?'<span class="badge badge-dir">★ 1ª DIRETRIZ</span>':""}
          </div>
          <div class="pilar-result-impacto">${p.impacto_falta}</div>
          <div class="pilar-result-acao">→ ${p.acao}</div>
          ${p.diretriz?`<div style="margin-top:8px;padding:8px;background:#1A2E1A;border:1px solid #C9A96E;border-radius:6px;font-size:11px;color:#C9A96E">
            <b>★ 1ª DIRETRIZ BRAÚNA — QUALIDADE</b><br>Provocar o cliente sobre o Financial Planning. Perguntar ao assessor: O FP deste cliente já foi feito? Quando? Isso é primordial para traçar o caminho do destino que o cliente deseja alcançar.
          </div>`:""}
        </div>`;
      });
    }

    if(concluidos.length>0){
      sh+=`<p style="font-size:11px;color:#5DCAA5;font-weight:700;margin:14px 0 8px;text-transform:uppercase;letter-spacing:.5px">✓ Pilares concluídos</p>`;
      concluidos.forEach(p=>{
        sh+=`<div class="pilar-result ok" style="padding:10px 14px">
          <div class="pilar-result-header" style="margin:0">
            <span>${p.icone}</span>
            <span class="pilar-result-nome" style="color:#5DCAA5">${p.nome}</span>
            <span class="pilar-result-status status-ok">✓ COMPLETO</span>
          </div>
        </div>`;
      });
    }
    sl.innerHTML=sh;
  }

  // Gestores
  const gl=document.getElementById("gestores-list");
  if(visao_gestores&&visao_gestores.length){
    gl.innerHTML=visao_gestores.map(g=>`
      <div class="gestor-card">
        <div class="gestor-nome">${g.gestor} — ${g.tema}</div>
        <div class="gestor-msg">${g.mensagem}</div>
        <div class="gestor-impl">→ ${g.implicacao_alocacao}</div>
        <div style="font-size:10px;color:#1E4A30;margin-top:4px">${g.fonte}</div>
      </div>`).join("");
  }

  // Cross Sell
  renderCrossSellResult(data.cross_sell||[], data.cross_tem||[], data.cross_nao_tem||[]);

  document.getElementById("results").style.display="block";
  document.getElementById("results").scrollIntoView({behavior:"smooth"});

  // Auto-abrir aba Modelo de Servir se houver pendentes críticos
  if(checklist_servir && checklist_servir.some(p=>p.status==="pendente"&&p.importancia.startsWith("CRÍTICA"))){
    setTimeout(()=>tab("servir"),800);
  }
}

// ── Gerar Apresentação de Reunião ─────────────────────────────────────────────
async function gerarApresentacao(){
  if(!analiseData) return;
  const btn = document.getElementById("btn-apres");
  const st  = document.getElementById("apres-st");
  btn.disabled = true;
  btn.textContent = "⏳ Gerando com IA...";
  st.textContent = "Claude está escrevendo a narrativa personalizada...";

  const xp = analiseData._xp || {};

  // Monta checklist de modelo de servir (formato esperado pelo PPTX)
  const checklistServir = {};
  if(analiseData.checklist_servir){
    analiseData.checklist_servir.forEach(p => {
      checklistServir[p.id || p.pilar] = {
        feito: p.status === "feito" || p.status === "ok",
        nome: p.nome || p.pilar || p.id,
        critico: p.critico || false,
      };
    });
  }

  // Cross-sell: mapeia IDs para nomes (o PPTX compara por nome)
  const crossAtivosNomes   = CROSS_AREAS.filter(a=>crossSell[a.id]).map(a=>a.nome);
  const crossFaltandoNomes = CROSS_AREAS.filter(a=>!crossSell[a.id]).map(a=>a.nome);

  const payload = {
    nome_cliente:        analiseData.nome || "Cliente",
    conta:               xp.conta || "",
    assessor:            analiseData.assessor || document.getElementById("assessor")?.value || "",
    perfil:              analiseData.perfil || "moderada",
    data_ref:            analiseData.data_ref || new Date().toLocaleDateString("pt-BR"),
    patrimonio:          analiseData.patrimonio || xp.patrimonio || 0,
    rent:                analiseData.rent || xp.rent || {},
    composicao:          analiseData.composicao || xp.comp || {},
    desvios:             xp.desvios || analiseData.desvios || [],
    acoes:               xp.acoes || [],
    fiis:                xp.fiis || [],
    sugestoes_produtos:  xp.sugestoes_produtos || [],
    cenario_macro:       xp.cenario_macro || {},
    alertas_relevantes:  xp.alertas_relevantes || analiseData.alertas || [],
    // Modelo de Servir
    score_servir:        analiseData.score_servir || 0,
    checklist:           checklistServir,
    pendentes_criticos:  analiseData.pendentes_criticos || 0,
    // Cross-sell
    cross_ativos:        crossAtivosNomes,
    cross_ativos_nomes:  crossAtivosNomes,
    cross_faltando:      crossFaltandoNomes,
    // Objetivo e perfil detalhado
    objetivo:            analiseData.objetivo || "",
  };

  try{
    const r = await fetch("/api/gerar-pptx", {
      method: "POST",
      headers: {"Content-Type":"application/json"},
      body: JSON.stringify(payload),
    });
    if(!r.ok) throw new Error(await r.text());
    const blob = await r.blob();
    const a = document.createElement("a");
    a.href = URL.createObjectURL(blob);
    a.download = `apresentacao_${(payload.nome_cliente).replace(/ /g,"-").toLowerCase()}_${payload.data_ref.replace(/\//g,"-")}.pptx`;
    a.click();
    st.innerHTML = `<span style="color:#5DCAA5">✓ Apresentação PowerPoint gerada com sucesso!</span>`;
  } catch(e){
    st.innerHTML = `<span style="color:#FF6B6B">Erro: ${e.message}</span>`;
  } finally{
    btn.disabled = false;
    btn.textContent = "🎯 Gerar Apresentação de Reunião";
    setTimeout(()=>st.textContent="", 5000);
  }
}

function renderCrossSellResult(areas, tem, naoTem){
  const el=document.getElementById("crosssell-list");
  if(!areas||!areas.length){ el.innerHTML=""; return; }

  const naoTemAreas=areas.filter(a=>!a.ativo);
  const temAreas=areas.filter(a=>a.ativo);
  const score=temAreas.length;
  const ficha = _clienteIdentificado?.ficha_salva || {};
  const crossHist = ficha.cross_historico || [];

  const barColor=score>=4?"#5DCAA5":score>=2?"#C9A96E":"#FF6B6B";
  let html=`<div class="cs-score-bar">
    <span style="font-size:28px;font-weight:700;color:${barColor};min-width:44px;text-align:center">${score}/5</span>
    <div style="flex:1">
      <div style="font-size:10px;color:#3A6A48;text-transform:uppercase;letter-spacing:.5px;margin-bottom:4px">Áreas ativas na Braúna</div>
      <div style="height:6px;background:#1A4030;border-radius:3px;overflow:hidden"><div style="width:${score/5*100}%;height:100%;background:${barColor};border-radius:3px"></div></div>
      <div style="font-size:11px;color:#888;margin-top:5px">${naoTem.length} oportunidade(s) de negócio em aberto</div>
    </div>
  </div>`;

  // ── Painel de insights com histórico ──────────────────────────────────────
  const insights = [];

  // Insight 1: produtos sem ativação há muito tempo (baseado em histórico)
  naoTemAreas.forEach(a => {
    // Verifica se alguma vez esteve ativo e foi desativado
    const jaEsteve = crossHist.some(h => (h.ativos||[]).includes(a.id));
    // Verifica quando foi a última vez que foi discutido (qualquer entrada no histórico)
    const ultimaEntrada = crossHist.length ? crossHist[crossHist.length-1] : null;
    if(jaEsteve){
      insights.push({tipo:"atencao", icone:"🔄", texto:`<b>${a.icone} ${a.nome}</b> já esteve ativo — avaliar reativação nesta reunião.`});
    } else if(crossHist.length >= 2){
      // Nunca foi ativado mesmo com histórico longo — prioridade
      insights.push({tipo:"oportunidade", icone:"💼", texto:`<b>${a.icone} ${a.nome}</b> nunca ativado em ${crossHist.length} reunião(ões) — oportunidade recorrente em aberto.`});
    }
  });

  // Insight 2: produtos recém-ativados (última entrada do histórico)
  if(crossHist.length >= 2){
    const ultimo = crossHist[crossHist.length-1];
    (ultimo.ativados||[]).forEach(id=>{
      const area = CROSS_AREAS.find(a=>a.id===id);
      if(area) insights.push({tipo:"novidade", icone:"✅", texto:`<b>${area.icone} ${area.nome}</b> ativado em ${ultimo.data} — confirme com o cliente na reunião.`});
    });
    (ultimo.desativados||[]).forEach(id=>{
      const area = CROSS_AREAS.find(a=>a.id===id);
      if(area) insights.push({tipo:"atencao", icone:"⚠️", texto:`<b>${area.icone} ${area.nome}</b> desativado em ${ultimo.data} — verificar motivo.`});
    });
  }

  // Insight 3: cliente com score baixo mas patrimônio alto
  const pat = _clienteIdentificado?.patrimonio || 0;
  if(score <= 1 && pat >= 300000){
    insights.push({tipo:"urgente", icone:"🚨", texto:`Cliente com R$ ${pat.toLocaleString("pt-BR",{maximumFractionDigits:0})} e apenas ${score} produto ativo — potencial de aprofundamento alto.`});
  }

  if(insights.length){
    const COR = {oportunidade:"#D4B483", atencao:"#FF9F40", urgente:"#FF6B6B", novidade:"#5DCAA5"};
    html += `<div style="background:#0A1A10;border:1px solid #1A3A20;border-radius:10px;padding:12px 14px;margin-bottom:14px">
      <div style="font-size:10px;color:#3A6A48;text-transform:uppercase;letter-spacing:.5px;margin-bottom:8px">🧠 Insights para esta reunião</div>
      ${insights.map(i=>`
        <div style="display:flex;align-items:flex-start;gap:8px;padding:6px 0;border-bottom:1px solid #0F2A18">
          <span style="font-size:14px;flex-shrink:0">${i.icone}</span>
          <span style="font-size:12px;color:#CCC;line-height:1.5">${i.texto}</span>
        </div>`).join("")}
      ${crossHist.length ? `<div style="margin-top:8px;font-size:10px;color:#2A5A3A">${crossHist.length} reunião(ões) registradas · Última: ${crossHist[crossHist.length-1]?.data||"—"}</div>` : ""}
    </div>`;
  }

  if(naoTemAreas.length){
    html+=`<p style="font-size:11px;color:#C9A96E;font-weight:700;margin-bottom:10px;text-transform:uppercase;letter-spacing:.5px">💼 Oportunidades em aberto para este cliente</p>`;
    naoTemAreas.forEach(a=>{
      html+=`<div class="cs-area nao-tem">
        <div class="cs-area-header">
          <span style="font-size:20px">${a.icone}</span>
          <span class="cs-area-nome">${a.nome}</span>
          <span class="cs-status-tag cs-falta-tag">OPORTUNIDADE ABERTA</span>
        </div>
        <div class="cs-pitch">"${a.pitch}"</div>`;

      // Área especial: Planejamento Financeiro
      if(a.id==="planejamento_financeiro" && a.especialista){
        const esp=a.especialista;
        html+=`<p style="font-size:11px;color:#888;margin-bottom:8px">Como especialista, estas são as três frentes prioritárias:</p>`;

        // Previdência
        const prev=esp.previdencia;
        html+=`<div class="fp-bloco">
          <div class="fp-bloco-titulo">${prev.icone} ${prev.titulo}</div>
          <p>${prev.diagnostico}</p>
          <div class="fp-item"><b>PGBL:</b> ${prev.pgbl_vs_vgbl.pgbl}</div>
          <div class="fp-item"><b>VGBL:</b> ${prev.pgbl_vs_vgbl.vgbl}</div>
          <div class="fp-item"><b>Regra:</b> ${prev.pgbl_vs_vgbl.regra}</div>
          <div class="fp-item"><b>Tabela Regressiva:</b> ${prev.tabela.regressiva}</div>
          <div class="fp-item"><b>Tabela Progressiva:</b> ${prev.tabela.progressiva}</div>
          <div class="fp-pitch">💬 Pitch: ${prev.pitch}</div>
        </div>`;

        // Seguro de vida
        const seg=esp.seguro_vida;
        html+=`<div class="fp-bloco" style="margin-top:8px">
          <div class="fp-bloco-titulo">${seg.icone} ${seg.titulo}</div>
          <p>${seg.diagnostico}</p>
          ${seg.coberturas.map(c=>`<div class="fp-item">→ ${c}</div>`).join("")}
          <div class="fp-item"><b>Capital mínimo recomendado:</b> ${seg.calculo}</div>
          <div class="fp-pitch">💬 Pitch: ${seg.pitch}</div>
        </div>`;

        // Transmissão de bens
        const trans=esp.transmissao_bens;
        html+=`<div class="fp-bloco" style="margin-top:8px">
          <div class="fp-bloco-titulo">${trans.icone} ${trans.titulo}</div>
          <p>${trans.diagnostico}</p>
          ${trans.estrategias.map(e=>`<div class="fp-item">→ ${e}</div>`).join("")}
          <div class="fp-alerta">⚠️ ${trans.alerta_itcmd}</div>
          <div class="fp-pitch">💬 Pitch: ${trans.pitch}</div>
        </div>`;
      } else {
        // Outras áreas: lista de oportunidades
        html+=`<ul class="cs-ops">${(a.oportunidades||[]).map(o=>`<li>${o}</li>`).join("")}</ul>`;
      }

      html+=`</div>`;
    });
  }

  if(temAreas.length){
    html+=`<p style="font-size:11px;color:#5DCAA5;font-weight:700;margin:14px 0 8px;text-transform:uppercase;letter-spacing:.5px">✓ Áreas já ativas</p>`;
    html+=`<div style="display:flex;flex-wrap:wrap;gap:8px">`;
    temAreas.forEach(a=>{
      html+=`<div class="cs-area tem" style="padding:10px 14px;display:inline-flex;align-items:center;gap:8px;border-radius:10px;flex:0 0 auto">
        <span style="font-size:18px">${a.icone}</span>
        <span style="font-size:13px;font-weight:700;color:#5DCAA5">${a.nome}</span>
        <span class="cs-status-tag cs-tem-tag">✓ ATIVO</span>
      </div>`;
    });
    html+=`</div>`;
  }

  el.innerHTML=html;
}

async function baixarPdf(){
  if(!analiseData) return;
  document.getElementById("btn-pdf").textContent="Gerando...";
  document.getElementById("btn-pdf").disabled=true;
  try{
    const res=await fetch("/api/pdf",{method:"POST",headers:{"Content-Type":"application/json"},body:JSON.stringify(analiseData)});
    if(!res.ok) throw new Error(await res.text());
    const blob=await res.blob();
    const a=document.createElement("a");
    a.href=URL.createObjectURL(blob);
    a.download=`${(analiseData.nome||"cliente").replace(/ /g,"-").toLowerCase()}_${analiseData.data_ref}_analise.pdf`;
    a.click();
  }catch(e){ alert("Erro ao gerar PDF: "+e.message); }
  finally{ document.getElementById("btn-pdf").textContent="⬇ Baixar relatório PDF"; document.getElementById("btn-pdf").disabled=false; }
}

async function baixarPpt(){
  if(!analiseData){ alert("Analise uma carteira primeiro."); return; }
  const btn=document.getElementById("btn-ppt");
  btn.textContent="⏳ Gerando PPT..."; btn.disabled=true;
  try{
    // Verifica se patrimônio está caindo comparando com histórico salvo
    const patrimonioAtual = analiseData.patrimonio||0;
    let patrimonioCaindo = false;
    try{
      const hRes = await fetch("/api/historico");
      const hData = await hRes.json();
      const nomeKey = (analiseData.nome||"").toLowerCase().replace(/\s+/g,"_");
      const assessorKey = (analiseData.assessor||"").toLowerCase().replace(/\s+/g,"_");
      const fkey = `${assessorKey}__${nomeKey}`;
      const entradas = (hData[fkey]||{}).entradas||[];
      if(entradas.length>=2 && patrimonioAtual < entradas[1].patrimonio){
        patrimonioCaindo = true;
      }
    }catch(e2){}

    const payload = {
      ...analiseData,
      patrimonio_caindo: patrimonioCaindo,
      checklist_servir: Object.entries(checklist).map(([id,ok])=>({id,status:ok?"ok":"nao"})),
    };
    const res=await fetch("/api/ppt",{method:"POST",headers:{"Content-Type":"application/json"},body:JSON.stringify(payload)});
    if(!res.ok){ const t=await res.text(); throw new Error(t); }
    const blob=await res.blob();
    const a=document.createElement("a");
    a.href=URL.createObjectURL(blob);
    a.download=`Brauna_${(analiseData.nome||"cliente").replace(/\s+/g,"_")}_Apresentacao.pptx`;
    a.click();
  }catch(e){ alert("Erro ao gerar PPT: "+e.message); }
  finally{ btn.textContent="📊 Gerar Apresentação PPT"; btn.disabled=false; }
}
</script>
</body>
</html>"""


# ── Rotas ─────────────────────────────────────────────────────────────────────
@app.route("/", methods=["GET"])
def login_page():
    from flask import redirect
    return redirect("/admin")

def _pilar_html_inicial(p):
    imp = p.get("importancia", "CRÍTICA")
    imp_cor  = "#FF6B6B" if imp=="CRÍTICA" else "#FFD966" if imp=="ALTA" else "#4A7055"
    imp_txt  = "🔴 CRÍTICO" if imp=="CRÍTICA" else "🟡 IMPORTANTE" if imp=="ALTA" else "⚪ MÉDIO"
    bg       = "#1A0A0A" if imp=="CRÍTICA" else "#151500" if imp=="ALTA" else "#181818"
    border   = "#4A1212" if imp=="CRÍTICA" else "#3A3A00" if imp=="ALTA" else "#282828"
    eh_dir   = p.get("importancia_diretriz", False)
    pid      = p["id"]
    nome     = p["nome"]
    icone    = p.get("icone","⭕")
    desc     = p.get("descricao", p.get("desc",""))
    impacto  = p.get("impacto_falta", p.get("impacto",""))
    acao     = p.get("acao","")
    dir_tag  = '<span style="font-size:10px;padding:2px 8px;border-radius:8px;background:#1A2E1A;color:#C9A96E;border:1px solid #C9A96E;font-weight:700">★ 1ª DIRETRIZ</span>' if eh_dir else ""
    dir_note = '<div style="font-size:10px;color:#C9A96E;margin-top:6px;padding:5px 8px;border:1px solid #C9A96E;border-radius:5px;background:#1A2E1A">★ 1ª DIRETRIZ BRAÚNA — Financial Planning é obrigatório antes de qualquer alocação.</div>' if eh_dir else ""
    return f'''<div onclick="togglePilar('{pid}')" id="pilar-{pid}" style="display:flex;align-items:flex-start;gap:12px;padding:12px 14px;border-radius:10px;background:{bg};border:1.5px solid {border};cursor:pointer;transition:all .2s;user-select:none">
  <div id="chk-{pid}" style="width:24px;height:24px;border-radius:6px;border:2px solid #2A5A3A;background:transparent;display:flex;align-items:center;justify-content:center;color:#000;font-size:15px;font-weight:900;flex-shrink:0;margin-top:2px"></div>
  <div style="font-size:20px;flex-shrink:0;margin-top:2px">{icone}</div>
  <div style="flex:1;min-width:0">
    <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;margin-bottom:4px">
      <span id="nome-{pid}" style="font-size:13px;font-weight:700;color:#F0F0F0">{nome}</span>
      <span style="font-size:10px;padding:2px 8px;border-radius:8px;background:#1E1E1E;color:{imp_cor};font-weight:700">{imp_txt}</span>
      {dir_tag}
    </div>
    <div style="font-size:11px;color:#4A7055;line-height:1.4">{desc}</div>
    <div id="det-{pid}" style="background:#111;border-radius:6px;padding:8px 10px;border-left:3px solid {imp_cor};margin-top:8px">
      <div style="font-size:11px;color:#888;line-height:1.5;margin-bottom:6px">{impacto}</div>
      <div style="font-size:11px;color:#C9A96E;padding:6px 8px;background:#1A2E1A;border-radius:5px">→ {acao}</div>
      {dir_note}
    </div>
  </div>
</div>'''

@app.route("/assessor")
def index():
    pilares_html = "\n".join(_pilar_html_inicial(p) for p in MODELO_SERVIR)
    return render_template_string(HTML, pilares_html=pilares_html)

@app.route("/api/login", methods=["POST"])
def login():
    d     = request.get_json()
    role  = d.get("role","")
    senha = d.get("senha","").strip().upper()
    if role == "assessor":
        nome = ASSESSORES.get(senha)
        if nome:
            return jsonify({"ok": True, "role": "assessor", "nome": nome, "codigo": senha})
        return jsonify({"ok": False, "msg": "Código de assessor inválido"}), 401
    if role == "admin":
        return jsonify({"ok": True, "role": "admin"})
    if SENHAS.get(role) == d.get("senha",""):
        return jsonify({"ok": True, "role": role})
    return jsonify({"ok": False, "msg": "Senha incorreta"}), 401


@app.route("/api/macro", methods=["GET"])
def macro_endpoint():
    macro = buscar_macro_bcb()
    ctx   = carregar_contexto()
    macro["ref_contexto"] = ctx.get("referencia","")
    return jsonify(macro)


@app.route("/api/clientes", methods=["GET"])
def get_clientes():
    return jsonify(load_clientes())

@app.route("/api/historico", methods=["GET"])
def historico_endpoint():
    return jsonify(load_hist())

@app.route("/api/sugestoes", methods=["GET","POST"])
def sugestoes_endpoint():
    data = load_suge()
    if request.method == "POST":
        nova = request.get_json()
        nova["id"]         = str(uuid.uuid4())[:8]
        nova["criado_em"]  = datetime.now().strftime("%d/%m/%Y %H:%M")
        nova["ativa"]      = True
        for s in data["historico"]:
            s["ativa"] = False
        data["historico"].insert(0, nova)
        data["historico"] = data["historico"][:30]
        save_suge(data)
        return jsonify({"ok": True, "id": nova["id"]})
    return jsonify(data)

@app.route("/api/ficha", methods=["GET","POST"])
def ficha_endpoint():
    fichas = load_fichas()
    if request.method == "POST":
        d = request.get_json()
        # Chave primária: código da conta (preferencial) ou assessor|nome
        conta = d.get("conta","").strip()
        key_nome = f"{d.get('assessor','')}|{d.get('nome','')}".lower().strip()
        key = f"conta:{conta}" if conta else key_nome

        # Mescla com ficha existente por nome (evita perder dados se salvou antes sem conta)
        ficha_anterior = fichas.get(key) or fichas.get(key_nome) or {}
        cross_novo = d.get("cross_ativos", [])
        cross_final = cross_novo if cross_novo else ficha_anterior.get("cross_ativos", [])
        checklist_novo = d.get("checklist", {})
        checklist_final = {**ficha_anterior.get("checklist", {}), **checklist_novo} if checklist_novo else ficha_anterior.get("checklist", {})

        # Histórico de cross sell — registra cada mudança de estado com data
        cross_historico = ficha_anterior.get("cross_historico", [])
        cross_anterior  = set(ficha_anterior.get("cross_ativos", []))
        cross_atual_set = set(cross_final)
        if cross_atual_set != cross_anterior:   # só registra se mudou
            cross_historico.append({
                "data": datetime.now().strftime("%d/%m/%Y"),
                "ativos": list(cross_atual_set),
                "ativados":   list(cross_atual_set - cross_anterior),
                "desativados": list(cross_anterior - cross_atual_set),
            })
            cross_historico = cross_historico[-24:]   # mantém últimas 24 entradas

        fichas[key] = {
            "conta": conta,
            "nome": d.get("nome",""),
            "assessor": d.get("assessor",""),
            "perfil": d.get("perfil","conservadora"),
            "objetivo": d.get("objetivo",""),
            "checklist": checklist_final,
            "cross_ativos": cross_final,
            "cross_historico": cross_historico,
            "atualizado_em": datetime.now().strftime("%d/%m/%Y %H:%M"),
        }
        # Remove a ficha por nome se agora temos a chave por conta
        if conta and key_nome in fichas and key_nome != key:
            del fichas[key_nome]
        save_fichas(fichas)
        return jsonify({"ok":True})
    # GET: retorna fichas de um assessor OU de uma conta específica
    conta = request.args.get("conta","").strip()
    if conta:
        ficha = fichas.get(f"conta:{conta}") or fichas.get(conta)
        return jsonify(ficha or {})
    assessor = request.args.get("assessor","").lower().strip()
    result = [v for k,v in fichas.items() if k.startswith(assessor+"|") or (v.get("assessor","").lower()==assessor)]
    result.sort(key=lambda x: x.get("nome",""))
    return jsonify(result)


@app.route("/api/xp-identificar", methods=["POST"])
def xp_identificar():
    """Recebe o PDF XPerformance, extrai o código do cliente e retorna dados salvos."""
    f = request.files.get("pdf")
    if not f:
        return jsonify({"erro": "PDF não enviado"}), 400

    pdf_bytes = f.read()
    try:
        xp = extrair_xperformance(pdf_bytes)
    except Exception as e:
        app.logger.error(f"extrair_xperformance falhou: {e}")
        # Retorna resposta parcial para o frontend continuar
        return jsonify({
            "conta": "", "assessor_xp": "", "nome_cliente": "",
            "data_ref": "", "patrimonio": 0, "rent": {},
            "composicao_atual": {}, "ficha_salva": {},
            "ultima_carteira": None, "comparativo": [],
            "tem_historico": False, "erro_extracao": str(e)
        })
    conta = xp.get("conta", "").strip()

    # Busca ficha salva para este código de conta
    fichas = load_fichas()
    ficha = fichas.get(f"conta:{conta}") or fichas.get(conta) or {}

    # Busca histórico de carteiras salvas
    hist = load_hist()
    historico_conta = hist.get(conta, [])
    ultima_carteira = historico_conta[-1] if historico_conta else None

    # Calcula comparativo vs. última carteira salva
    comparativo = []
    if ultima_carteira:
        comp_ant = ultima_carteira.get("composicao", {})
        comp_atual = xp.get("comp", {})
        for cat in CATS:
            ant = comp_ant.get(cat, 0)
            atual = comp_atual.get(cat, 0)
            delta = round(atual - ant, 2)
            if ant > 0 or atual > 0:
                comparativo.append({
                    "cat": cat,
                    "label": LABELS.get(cat, cat),
                    "anterior": ant,
                    "atual": atual,
                    "delta": delta,
                })

    return jsonify({
        "conta": conta,
        "assessor_xp": xp.get("assessor",""),
        "nome_cliente": ficha.get("nome",""),
        "data_ref": xp.get("data_ref",""),
        "patrimonio": xp.get("patrimonio", 0),
        "rent": xp.get("rent", {}),
        "composicao_atual": xp.get("comp", {}),
        "ficha_salva": ficha,          # modelo de servir + cross-sell salvos
        "ultima_carteira": ultima_carteira,
        "comparativo": comparativo,    # delta vs. última reunião
        "tem_historico": len(historico_conta) > 0,
    })


@app.route("/api/salvar-carteira", methods=["POST"])
def salvar_carteira():
    """Salva snapshot da carteira atual para histórico comparativo futuro."""
    d = request.get_json()
    conta = d.get("conta","").strip()
    if not conta:
        return jsonify({"erro": "conta obrigatória"}), 400
    hist = load_hist()
    if conta not in hist:
        hist[conta] = []
    hist[conta].append({
        "data_ref": d.get("data_ref",""),
        "patrimonio": d.get("patrimonio", 0),
        "composicao": d.get("composicao", {}),
        "rent": d.get("rent", {}),
        "salvo_em": datetime.now().strftime("%d/%m/%Y %H:%M"),
    })
    # Mantém apenas as últimas 12 carteiras
    hist[conta] = hist[conta][-12:]
    save_hist(hist)
    return jsonify({"ok": True, "snapshots": len(hist[conta])})

@app.route("/api/clientes/<cid>/nota", methods=["POST"])
def set_nota(cid):
    clientes = load_clientes()
    for c in clientes:
        if c.get("id") == cid:
            c["nota_lider"] = request.get_json().get("nota","")
            break
    save_clientes(clientes)
    return jsonify({"ok":True})

@app.route("/api/okrs", methods=["GET","POST"])
def okrs_endpoint():
    if request.method == "POST":
        save_okrs(request.get_json())
        return jsonify({"ok":True})
    return jsonify(load_okrs())

@app.route("/api/assessores_dados", methods=["GET","POST"])
def assessores_dados_endpoint():
    """Dados financeiros dos assessores (patrimônio, receita, OKR por classe)."""
    if request.method == "POST":
        data = request.get_json() or {}
        _save(_ASSESSORES_FILE, data)
        return jsonify({"ok": True, "total": len(data.get("assessores", []))})
    return jsonify(_load(_ASSESSORES_FILE, {"assessores": [], "atualizado_em": ""}))

@app.route("/api/importar_assessores_xlsx", methods=["POST"])
def importar_assessores_xlsx():
    """Importa planilha Excel de dados financeiros dos assessores."""
    try:
        import openpyxl
        from io import BytesIO
        arquivo = request.files.get("planilha")
        if not arquivo:
            return jsonify({"ok": False, "erro": "Nenhum arquivo enviado"}), 400
        wb = openpyxl.load_workbook(BytesIO(arquivo.read()), data_only=True)
        ws = wb.active
        assessores = []
        for row in ws.iter_rows(min_row=2, max_row=50, values_only=True):
            nome = row[0]
            if not nome or not isinstance(nome, str) or not str(nome).strip():
                continue
            patrimonio    = row[1] or 0
            rf_saldo      = row[2] or 0
            rf_pct        = round((row[3] or 0) * 100, 2)
            rf_receita    = round(row[4] or 0, 2)
            rf_volume_mes = round(row[5] or 0, 2)
            rv_saldo      = row[6] or 0
            rv_pct        = round((row[7] or 0) * 100, 2)
            rv_receita    = round(row[8] or 0, 2)
            rv_volume_mes = round(row[9] or 0, 2)
            estrut_pct_max= round((row[10] or 0) * 100, 2)
            estrut_volume = round(row[11] or 0, 2)
            estrut_receita= round(row[12] or 0, 2)
            fii_saldo     = row[13] or 0
            fii_pct       = round((row[14] or 0) * 100, 2)
            fii_receita   = round(row[15] or 0, 2)
            fii_volume_mes= round(row[16] or 0, 2)
            intl_saldo    = row[17] or 0
            intl_pct      = round((row[18] or 0) * 100, 2)
            intl_receita  = round(row[19] or 0, 2)
            intl_volume_mes= round(row[20] or 0, 2)
            receita_total = round(row[22] or 0, 2)
            roa_anual     = round((row[23] or 0) * 100, 4)
            okr_mensal    = round(row[24] or 0, 2)
            okr_anual     = round(row[26] or 0, 2)
            okr_semestre  = round(row[27] or 0, 2)
            pct_realizado = round((receita_total / okr_mensal * 100) if okr_mensal > 0 else 0, 1)
            assessores.append({
                "nome": str(nome).strip(),
                "patrimonio": patrimonio,
                "rf":   {"saldo": round(rf_saldo,2), "pct": rf_pct, "receita": rf_receita, "volume_mes": rf_volume_mes},
                "rv":   {"saldo": round(rv_saldo,2), "pct": rv_pct, "receita": rv_receita, "volume_mes": rv_volume_mes},
                "estruturadas": {"pct_max": estrut_pct_max, "volume": estrut_volume, "receita": estrut_receita},
                "fii":  {"saldo": round(fii_saldo,2), "pct": fii_pct, "receita": fii_receita, "volume_mes": fii_volume_mes},
                "internacional": {"saldo": round(intl_saldo,2), "pct": intl_pct, "receita": intl_receita, "volume_mes": intl_volume_mes},
                "receita_total": receita_total,
                "roa_anual": roa_anual,
                "okr_mensal": okr_mensal,
                "okr_anual": okr_anual,
                "okr_semestre": okr_semestre,
                "pct_realizado_okr": pct_realizado,
            })
        payload = {
            "assessores": assessores,
            "atualizado_em": datetime.now().strftime("%d/%m/%Y %H:%M"),
            "fonte": arquivo.filename,
        }
        _save(_ASSESSORES_FILE, payload)
        return jsonify({"ok": True, "total": len(assessores)})
    except Exception as e:
        app.logger.error(f"importar_assessores_xlsx: {e}")
        return jsonify({"ok": False, "erro": str(e)}), 500

@app.route("/api/admin/activity", methods=["GET","POST"])
def admin_activity():
    """Registra e retorna atividade de assessores e líderes."""
    if request.method == "POST":
        data = request.get_json() or {}
        log = _load(_ADMIN_ACTIVITY_FILE, [])
        entry = {
            "ts":     datetime.now().strftime("%d/%m/%Y %H:%M"),
            "role":   data.get("role",""),
            "nome":   data.get("nome",""),
            "acao":   data.get("acao","acesso"),
            "detalhe":data.get("detalhe",""),
        }
        log.insert(0, entry)
        if len(log) > 500: log = log[:500]
        _save(_ADMIN_ACTIVITY_FILE, log)
        return jsonify({"ok": True})
    return jsonify(_load(_ADMIN_ACTIVITY_FILE, []))

@app.route("/api/admin/dashboard", methods=["GET"])
def admin_dashboard():
    """Agrega todas as informações do sistema para o painel admin."""
    activity = _load(_ADMIN_ACTIVITY_FILE, [])

    # Resumo por role/nome
    assessores = {}
    lideres    = {}
    for ev in activity:
        nome = ev.get("nome","") or "Anônimo"
        if ev.get("role") == "assessor":
            if nome not in assessores:
                assessores[nome] = {"nome": nome, "acoes": [], "ultimo": ev.get("ts","")}
            assessores[nome]["acoes"].append({"acao": ev.get("acao",""), "detalhe": ev.get("detalhe",""), "ts": ev.get("ts","")})
        elif ev.get("role") in ("lider","líder"):
            if nome not in lideres:
                lideres[nome] = {"nome": nome, "acoes": [], "ultimo": ev.get("ts","")}
            lideres[nome]["acoes"].append({"acao": ev.get("acao",""), "detalhe": ev.get("detalhe",""), "ts": ev.get("ts","")})

    # Head — agregar tudo
    cenario    = _load(_HP_CENARIO_FILE, {})
    produtos   = _load(_HP_PROD_FILE, {})
    calls      = _load(_HP_CALLS_FILE, [])
    knowledge  = _load(_HP_KNOW_FILE, [])
    gestoras2  = _load(_HP_GESTORAS2_FILE, {})
    estruturadas = _load(_HP_ESTRUTURADAS_FILE, [])
    alertas    = _load(_HP_ALERTS_FILE, [])
    portfolios = _load(_HP_PORT_FILE, {})
    clientes   = _load(_CLIENTS_FILE, {})

    total_produtos = sum(len(v) for v in produtos.values() if isinstance(v, list)) if isinstance(produtos, dict) else 0

    return jsonify({
        "assessores":     list(assessores.values()),
        "lideres":        list(lideres.values()),
        "_raw_activity":  activity[:200],
        "head": {
            "cenario": {
                "referencia":    cenario.get("referencia",""),
                "global_chars":  len(cenario.get("global","")),
                "brasil_chars":  len(cenario.get("brasil","")),
                "pos_chars":     len(cenario.get("posicionamento","")),
                "global":        cenario.get("global","")[:300],
                "brasil":        cenario.get("brasil","")[:300],
                "posicionamento":cenario.get("posicionamento","")[:300],
            },
            "knowledge":    knowledge,
            "calls":        [c for c in calls if c.get("ativo") is not False],
            "estruturadas": estruturadas,
            "alertas":      alertas,
            "gestoras2":    list(gestoras2.values()),
            "portfolios_ref": portfolios.get("referencia",""),
            "total_produtos": total_produtos,
            "total_docs_base": len(knowledge),
            "total_calls":    len([c for c in calls if c.get("ativo") is not False]),
            "total_estruturadas": len(estruturadas),
            "total_alertas":  len(alertas),
            "clientes":       list(clientes.values()),
            "total_clientes": len(clientes),
        }
    })

@app.route("/api/admin/mensagem", methods=["GET","POST"])
def mensagem_endpoint():
    if request.method == "POST":
        d = request.get_json()
        d["atualizado"] = datetime.now().strftime("%d/%m/%Y %H:%M")
        save_msg(d)
        return jsonify({"ok":True})
    return jsonify(load_msg())

@app.route("/api/admin/upload-contexto", methods=["POST"])
def upload_contexto():
    carta = request.files.get("carta")
    if not carta:
        return jsonify({"error":"Arquivo não enviado"}), 400
    texto = extrair_texto_pdf(carta.read())
    _save("/tmp/brauna_carta.json", {
        "texto": texto[:8000],
        "atualizado": datetime.now().strftime("%d/%m/%Y %H:%M"),
        "nome": carta.filename
    })
    return jsonify({"ok":True, "chars":len(texto), "preview":texto[:400]})

@app.route("/api/admin/upload-sugestao-doc", methods=["POST"])
def upload_sugestao_doc():
    arq = request.files.get("arquivo")
    if not arq:
        return jsonify({"error": "Arquivo não enviado"}), 400
    cat = request.form.get("categoria", "")
    nome = (arq.filename or "").lower()
    dados = arq.read()

    # ── PDF ─────────────────────────────────────────────────────────────────
    if nome.endswith(".pdf"):
        texto = extrair_texto_pdf(dados)
        return jsonify({"ok": True, "chars": len(texto), "texto": texto[:12000], "tabelas": []})

    # ── Excel / CSV ──────────────────────────────────────────────────────────
    try:
        import openpyxl, csv as _csv
        tabelas = []

        if nome.endswith(".csv"):
            import io as _io
            txt = dados.decode("utf-8-sig", errors="replace")
            sep = ";" if txt.count(";") > txt.count(",") else ","
            reader = _csv.reader(_io.StringIO(txt), delimiter=sep)
            rows = [r for r in reader if any(c.strip() for c in r)]
            if rows:
                header = [c.strip() for c in rows[0]]
                linhas = [[c.strip() for c in r] for r in rows[1:] if any(c.strip() for c in r)]
                cat_det = _detectar_categoria(header, cat)
                tabelas.append({"nome": nome, "colunas": header, "linhas": linhas, "categoria": cat_det})
        else:
            wb = openpyxl.load_workbook(io.BytesIO(dados), data_only=True)
            for sh in wb.worksheets:
                rows = []
                for row in sh.iter_rows(values_only=True):
                    r = [str(c).strip() if c is not None else "" for c in row]
                    if any(c for c in r):
                        rows.append(r)
                if not rows:
                    continue
                header = rows[0]
                linhas = [r for r in rows[1:] if any(c for c in r)]
                cat_det = _detectar_categoria(header, cat)
                tabelas.append({"nome": sh.title, "colunas": header, "linhas": linhas[:200], "categoria": cat_det})

        texto_resumo = "\n".join(
            "[" + t["nome"] + "]\n" + ", ".join(t["colunas"]) + "\n" +
            "\n".join(" | ".join(r) for r in t["linhas"][:5]) + "\n..."
            for t in tabelas
        )
        return jsonify({"ok": True, "chars": sum(len(str(t)) for t in tabelas),
                        "texto": texto_resumo, "tabelas": tabelas})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


def _detectar_categoria(header, fallback=""):
    h = " ".join(header).lower()
    if any(w in h for w in ["indexador","ipca","cdi","pre","alocar","substituir","renda fixa","rf"]):
        return "rf"
    if any(w in h for w in ["ticker","acao","ações","bolsa","rv","comprar","vender","trocar"]):
        return "rv"
    if any(w in h for w in ["fii","fundo imobiliario","fundo imobiliário","segmento","papel","logistica"]):
        return "fii"
    if any(w in h for w in ["etf","bdr","internacional","exterior","global","dolar","usd"]):
        return "intl"
    return fallback


@app.route("/api/admin/contexto-info", methods=["GET"])
def contexto_info():
    ctx   = carregar_contexto()
    carta = _load("/tmp/brauna_carta.json", {})
    msg   = load_msg()
    return jsonify({"contexto_ref":ctx.get("referencia",""), "carta":carta, "mensagem":msg})

@app.route("/lider")
def lider():
    return render_template_string(HTML_LIDER)

@app.route("/admin")
def admin_page():
    return render_template_string(HTML_ADMIN)

@app.route("/head-produtos")
def head_page():
    return render_template_string(HTML_HEAD)

# ── API Head de Produtos ──────────────────────────────────────────────────────
@app.route("/api/hp/portfolios", methods=["GET","POST"])
def hp_portfolios():
    if request.method == "POST":
        data = request.get_json()
        _save(_HP_PORT_FILE, data)
        return jsonify({"ok": True})
    return jsonify(_load(_HP_PORT_FILE, HP_PORTFOLIOS_DEFAULT))

@app.route("/api/hp/cenario", methods=["GET","POST"])
def hp_cenario():
    if request.method == "POST":
        data = request.get_json()
        _save(_HP_CENARIO_FILE, data)
        return jsonify({"ok": True})
    cenario = _load(_HP_CENARIO_FILE, HP_CENARIO_DEFAULT)
    # Inclui fontes publicadas da base de conhecimento
    docs_pub = [d for d in _load(_HP_KNOW_FILE, []) if d.get("publicado")]
    cenario["fontes"] = [{"nome": d.get("nome",""), "tipo": d.get("tipo",""), "fonte": d.get("fonte",""), "data": d.get("data","")} for d in docs_pub]
    return jsonify(cenario)

@app.route("/api/hp/produtos", methods=["GET","POST"])
def hp_produtos():
    if request.method == "POST":
        data = request.get_json()
        _save(_HP_PROD_FILE, data)
        return jsonify({"ok": True})
    return jsonify(_load(_HP_PROD_FILE, {}))

@app.route("/api/hp/publicar", methods=["POST"])
def hp_publicar():
    body = request.get_json()
    portfolios = body.get("portfolios")
    cenario    = body.get("cenario")
    produtos   = body.get("produtos")
    if portfolios: _save(_HP_PORT_FILE, portfolios)
    if cenario:    _save(_HP_CENARIO_FILE, cenario)
    if produtos:   _save(_HP_PROD_FILE, produtos)
    return jsonify({"ok": True, "publicado_em": datetime.now().strftime("%d/%m/%Y %H:%M")})

@app.route("/api/hp/alertas", methods=["GET","POST","DELETE"])
def hp_alertas():
    if request.method == "GET":
        return jsonify(_load(_HP_ALERTS_FILE, []))
    if request.method == "DELETE":
        data = request.get_json()
        alertas = _load(_HP_ALERTS_FILE, [])
        id_del = data.get("id")
        alertas = [a for a in alertas if a.get("id") != id_del]
        _save(_HP_ALERTS_FILE, alertas)
        return jsonify({"ok": True})
    # POST — novo alerta manual ou automático
    data = request.get_json()
    alertas = _load(_HP_ALERTS_FILE, [])
    novo = {
        "id": str(uuid.uuid4())[:8],
        "produto":    data.get("produto",""),
        "classe":     data.get("classe",""),
        "tipo":       data.get("tipo","info"),  # info | atencao | urgente
        "mensagem":   data.get("mensagem",""),
        "origem":     data.get("origem","Head de Produtos"),
        "data":       datetime.now().strftime("%d/%m/%Y %H:%M"),
        "lido":       False,
        "assessor_destino": data.get("assessor_destino", ""),  # "" = global, "lucas" = só Lucas
    }
    alertas.insert(0, novo)
    alertas = alertas[:100]  # manter últimos 100
    _save(_HP_ALERTS_FILE, alertas)
    return jsonify({"ok": True, "alerta": novo})

@app.route("/api/hp/calls", methods=["GET","POST"])
def hp_calls():
    if request.method == "GET":
        return jsonify(_load(_HP_CALLS_FILE, []))
    # POST — adicionar ou remover call
    data = request.get_json()
    if data.get("id") and data.get("_delete"):
        calls = _load(_HP_CALLS_FILE, [])
        calls = [c for c in calls if c.get("id") != data["id"]]
        _save(_HP_CALLS_FILE, calls)
        return jsonify({"ok": True})
    # Novo call
    entrada = float(data.get("preco_entrada") or 0)
    alvo    = float(data.get("preco_alvo") or 0)
    upside  = round((alvo / entrada - 1) * 100, 1) if entrada > 0 else 0
    novo = {
        "id":            str(uuid.uuid4())[:8],
        "ticker":        data.get("ticker","").upper().strip(),
        "nome":          data.get("nome",""),
        "direcao":       data.get("direcao","compra"),  # compra | venda | neutro
        "preco_entrada": entrada,
        "preco_alvo":    alvo,
        "upside":        upside,
        "stop":          float(data.get("stop") or 0),
        "prazo":         data.get("prazo",""),
        "tese":          data.get("tese",""),
        "perfis":        data.get("perfis", []),
        "fonte":         data.get("fonte",""),
        "data":          datetime.now().strftime("%d/%m/%Y"),
        "ativo":         True,
    }
    calls = _load(_HP_CALLS_FILE, [])
    calls.insert(0, novo)
    calls = calls[:200]
    _save(_HP_CALLS_FILE, calls)
    return jsonify({"ok": True, "call": novo})

@app.route("/api/hp/gestoras2", methods=["GET","POST"])
def hp_gestoras2():
    """Armazena carteiras de gestoras externas com os 5 perfis. Estrutura: {gestora_id: {nome, referencia, perfis:{perfil:{classe:%}}}}"""
    if request.method == "GET":
        return jsonify(_load(_HP_GESTORAS2_FILE, {}))
    data = request.get_json()
    if data.get("_delete"):
        gestoras = _load(_HP_GESTORAS2_FILE, {})
        gid = data.get("id","")
        gestoras.pop(gid, None)
        _save(_HP_GESTORAS2_FILE, gestoras)
        return jsonify({"ok": True})
    CLASSES = ["pos_fixado","inflacao","pre_fixado","acoes","fiis","multimercado","internacional","alternativos","criptomoedas"]
    PERFIS  = ["super_conservadora","conservadora","moderada","agressiva","super_agressiva"]
    gestoras = _load(_HP_GESTORAS2_FILE, {})
    if len(gestoras) >= 5 and data.get("id","") not in gestoras:
        return jsonify({"ok": False, "error": "Máximo de 5 gestoras atingido."}), 400
    gid = data.get("id") or data.get("nome","gestora").lower().replace(" ","_")[:20]
    perfis_data = {}
    for p in PERFIS:
        alocacao = {}
        for c in CLASSES:
            v = data.get("perfis",{}).get(p,{}).get(c, 0)
            alocacao[c] = round(float(v or 0), 1)
        perfis_data[p] = alocacao
    gestoras[gid] = {"id": gid, "nome": data.get("nome","").strip(), "referencia": data.get("referencia","").strip(), "perfis": perfis_data}
    _save(_HP_GESTORAS2_FILE, gestoras)
    return jsonify({"ok": True, "id": gid})

@app.route("/api/hp/estruturadas", methods=["GET","POST"])
def hp_estruturadas():
    """Operações estruturadas cadastradas pelo Head."""
    if request.method == "GET":
        return jsonify(_load(_HP_ESTRUTURADAS_FILE, []))
    data = request.get_json()
    if data.get("_delete"):
        ops = _load(_HP_ESTRUTURADAS_FILE, [])
        ops = [o for o in ops if o.get("id") != data.get("id")]
        _save(_HP_ESTRUTURADAS_FILE, ops)
        return jsonify({"ok": True})
    nova = {
        "id": str(uuid.uuid4())[:8],
        "tipo": data.get("tipo","").strip(),
        "ativo": data.get("ativo","").strip(),
        "emissor": data.get("emissor","").strip(),
        "vencimento": data.get("vencimento","").strip(),
        "retorno": data.get("retorno","").strip(),
        "perfil_minimo": data.get("perfil_minimo","moderada"),
        "observacao": data.get("observacao","").strip(),
        "data": datetime.now().strftime("%d/%m/%Y"),
    }
    ops = _load(_HP_ESTRUTURADAS_FILE, [])
    ops.insert(0, nova)
    _save(_HP_ESTRUTURADAS_FILE, ops)
    return jsonify({"ok": True, "op": nova})

@app.route("/api/hp/gestores", methods=["GET","POST"])
def hp_gestores():
    if request.method == "GET":
        return jsonify(_load(_HP_GESTORES_FILE, []))
    data = request.get_json()
    if data.get("id") and data.get("_delete"):
        gestores = _load(_HP_GESTORES_FILE, [])
        gestores = [g for g in gestores if g.get("id") != data["id"]]
        _save(_HP_GESTORES_FILE, gestores)
        return jsonify({"ok": True})
    CLASSES = ["pos_fixado","inflacao","pre_fixado","acoes","fiis","multimercado","internacional","alternativos","criptomoedas"]
    alocacao = {}
    for c in CLASSES:
        v = data.get("alocacao", {}).get(c, 0)
        alocacao[c] = round(float(v or 0), 1)
    novo = {
        "id":         str(uuid.uuid4())[:8],
        "nome":       data.get("nome","").strip(),
        "gestora":    data.get("gestora","").strip(),
        "referencia": data.get("referencia","").strip(),
        "perfil":     data.get("perfil","moderada"),
        "alocacao":   alocacao,
        "observacao": data.get("observacao","").strip(),
        "data":       datetime.now().strftime("%d/%m/%Y"),
    }
    gestores = _load(_HP_GESTORES_FILE, [])
    # Atualiza se mesmo nome+gestora, senão insere (máx 5)
    existente = next((i for i,g in enumerate(gestores) if g.get("gestora")==novo["gestora"] and g.get("perfil")==novo["perfil"]), None)
    if existente is not None:
        gestores[existente] = novo
    else:
        if len(gestores) >= 5:
            return jsonify({"ok": False, "error": "Máximo de 5 carteiras atingido. Remova uma antes de adicionar."}), 400
        gestores.append(novo)
    _save(_HP_GESTORES_FILE, gestores)
    return jsonify({"ok": True, "gestora": novo})

@app.route("/api/hp/knowledge", methods=["GET"])
def hp_knowledge_list():
    return jsonify(_load(_HP_KNOW_FILE, []))

@app.route("/api/hp/knowledge/upload", methods=["POST"])
def hp_knowledge_upload():
    pdf_file = request.files.get("pdf")
    nome     = request.form.get("nome", "")
    tipo     = request.form.get("tipo", "carta")
    classes  = request.form.get("classes", "")   # comma-separated
    fonte    = request.form.get("fonte", "")
    if not pdf_file:
        return jsonify({"error": "PDF não enviado"}), 400
    try:
        import pdfplumber, io
        raw = pdf_file.read()
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            texto = "\n".join(p.extract_text() or "" for p in pdf.pages)
        texto = texto.strip()
        if len(texto) < 50:
            return jsonify({"error": "PDF sem texto extraível"}), 400
    except Exception as e:
        msg = str(e)
        if "Root" in msg or "EOF" in msg or "not a PDF" in msg.lower():
            msg = "PDF inválido ou protegido por senha. Tente exportar novamente como PDF."
        return jsonify({"error": msg}), 400

    docs = _load(_HP_KNOW_FILE, [])
    doc = {
        "id":       str(uuid.uuid4())[:8],
        "nome":     nome or pdf_file.filename,
        "tipo":     tipo,
        "classes":  [c.strip() for c in classes.split(",") if c.strip()],
        "fonte":    fonte,
        "chars":    len(texto),
        "texto":    texto[:12000],   # máx 12k chars para não explodir o JSON
        "data":     datetime.now().strftime("%d/%m/%Y %H:%M"),
        "publicado": False,
    }
    docs.insert(0, doc)
    docs = docs[:50]   # manter últimos 50 documentos
    _save(_HP_KNOW_FILE, docs)
    return jsonify({"ok": True, "id": doc["id"], "chars": len(texto), "texto": texto[:12000], "doc": doc})

@app.route("/api/hp/knowledge/delete", methods=["POST"])
def hp_knowledge_delete():
    id_del = request.get_json().get("id")
    docs = _load(_HP_KNOW_FILE, [])
    docs = [d for d in docs if d.get("id") != id_del]
    _save(_HP_KNOW_FILE, docs)
    return jsonify({"ok": True})

@app.route("/api/hp/carta-upload", methods=["POST"])
def hp_carta_upload():
    """Recebe PDF de carta de gestora, extrai texto e gera cenário macro via Claude. Não salva na base de conhecimento."""
    pdf_file = request.files.get("pdf")
    if not pdf_file:
        return jsonify({"error": "PDF não enviado"}), 400
    try:
        import pdfplumber, io
        raw = pdf_file.read()
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            texto = "\n".join(p.extract_text() or "" for p in pdf.pages)
        texto = texto.strip()
        if len(texto) < 50:
            return jsonify({"error": "PDF sem texto extraível (pode ser PDF de imagem ou protegido)"}), 400
    except Exception as e:
        msg = str(e)
        if "Root" in msg or "EOF" in msg or "not a pdf" in msg.lower():
            msg = "PDF inválido ou protegido por senha. Exporte novamente como PDF."
        return jsonify({"error": msg}), 400

    fonte = request.form.get("nome", pdf_file.filename or "Carta da Gestora").replace(".pdf","").strip()
    api_key = os.environ.get("ANTHROPIC_API_KEY","")
    resultado = {"ok": True, "fonte": fonte, "chars": len(texto), "ia": False,
                 "global": "", "brasil": "", "posicionamento": ""}

    if api_key and len(texto) > 100:
        try:
            import anthropic as _anthropic, re as _re
            client = _anthropic.Anthropic(api_key=api_key)
            prompt = f"""Você é analista de investimentos. Leia a carta da gestora abaixo e extraia 3 resumos objetivos em português, cada um com no máximo 3 frases:

1. CENÁRIO GLOBAL: principais pontos sobre economia mundial, Fed, inflação global, bolsas internacionais, geopolítica
2. CENÁRIO BRASIL: principais pontos sobre economia brasileira, Selic/COPOM, IPCA, Ibovespa, fiscal, câmbio
3. POSICIONAMENTO: o que a gestora está fazendo/recomendando (aumentando, reduzindo, neutro) por classe de ativo

Responda SOMENTE no formato JSON:
{{"global": "...", "brasil": "...", "posicionamento": "..."}}

CARTA:
{texto[:6000]}"""
            msg = client.messages.create(model="claude-haiku-4-5-20251001", max_tokens=600,
                                         messages=[{"role":"user","content": prompt}])
            raw = msg.content[0].text.strip()
            m = _re.search(r'\{.*\}', raw, _re.DOTALL)
            if m:
                parsed = json.loads(m.group())
                resultado.update({
                    "global": parsed.get("global","")[:1200],
                    "brasil": parsed.get("brasil","")[:1200],
                    "posicionamento": parsed.get("posicionamento","")[:1200],
                    "ia": True,
                })
        except Exception:
            pass  # fallback: retorna sem cenário extraído

    if not resultado["global"]:
        # Fallback: divide texto em terços
        linhas = [l for l in texto.split("\n") if l.strip()]
        n = max(len(linhas)//3, 1)
        resultado.update({
            "global": " ".join(linhas[:n])[:800],
            "brasil": " ".join(linhas[n:2*n])[:800],
            "posicionamento": " ".join(linhas[2*n:])[:800],
        })

    return jsonify(resultado)

@app.route("/api/hp/knowledge/apply", methods=["POST"])
def hp_knowledge_apply():
    """Extrai trechos do documento e aplica ao cenário macro ou retorna para o frontend usar."""
    body   = request.get_json()
    doc_id = body.get("id")
    destino= body.get("destino", "cenario")   # cenario | alerta | produto

    docs = _load(_HP_KNOW_FILE, [])
    doc  = next((d for d in docs if d.get("id") == doc_id), None)
    if not doc:
        return jsonify({"error": "Documento não encontrado"}), 404

    texto = doc.get("texto", "")

    if destino == "cenario":
        fonte = doc.get("fonte") or doc.get("nome","")

        # Usa Claude para extrair os 3 campos de forma inteligente
        api_key = os.environ.get("ANTHROPIC_API_KEY","")
        if api_key and len(texto) > 100:
            try:
                import anthropic as _anthropic
                client = _anthropic.Anthropic(api_key=api_key)
                prompt = f"""Você é analista de investimentos. Leia a carta da gestora abaixo e extraia 3 resumos objetivos em português, cada um com no máximo 3 frases:

1. CENÁRIO GLOBAL: principais pontos sobre economia mundial, Fed, inflação global, bolsas internacionais, geopolítica
2. CENÁRIO BRASIL: principais pontos sobre economia brasileira, Selic/COPOM, IPCA, Ibovespa, fiscal, câmbio
3. POSICIONAMENTO: o que a gestora está fazendo/recomendando (aumentando, reduzindo, neutro) por classe de ativo

Responda SOMENTE no formato JSON:
{{"global": "...", "brasil": "...", "posicionamento": "..."}}

CARTA:
{texto[:6000]}"""
                msg = client.messages.create(
                    model="claude-haiku-4-5-20251001",
                    max_tokens=600,
                    messages=[{"role":"user","content": prompt}]
                )
                raw = msg.content[0].text.strip()
                import re as _re
                m = _re.search(r'\{.*\}', raw, _re.DOTALL)
                if m:
                    parsed = json.loads(m.group())
                    return jsonify({
                        "ok": True,
                        "global":         parsed.get("global","")[:1200],
                        "brasil":         parsed.get("brasil","")[:1200],
                        "posicionamento": parsed.get("posicionamento","")[:1200],
                        "fonte": fonte,
                        "ia": True,
                    })
            except Exception:
                pass  # fallback abaixo

        # Fallback simples: divide em terços
        linhas = [l for l in texto.split("\n") if l.strip()]
        n = max(len(linhas), 1)
        return jsonify({
            "ok": True,
            "global":         " ".join(linhas[:n//3])[:1200].strip(),
            "brasil":         " ".join(linhas[n//3: 2*n//3])[:1200].strip(),
            "posicionamento": " ".join(linhas[2*n//3:])[:1200].strip(),
            "fonte": fonte,
            "ia": False,
        })

    return jsonify({"ok": True, "texto": texto[:3000]})

@app.route("/api/hp/knowledge/publicar", methods=["POST"])
def hp_knowledge_publicar():
    """Marca um documento como publicado (visível para assessores)."""
    body   = request.get_json()
    doc_id = body.get("id")
    docs   = _load(_HP_KNOW_FILE, [])
    for d in docs:
        if d.get("id") == doc_id:
            d["publicado"] = True
    _save(_HP_KNOW_FILE, docs)
    return jsonify({"ok": True})

@app.route("/api/analyze-xp", methods=["POST"])
def analyze_xp():
    """Lê PDF XPerformance, compara com modelo HP e sugere produtos cadastrados."""
    pdf_file = request.files.get("pdf")
    assessor  = request.form.get("assessor", "Assessor")
    if not pdf_file:
        return jsonify({"error": "PDF não enviado"}), 400

    try:
        dados = extrair_xperformance(pdf_file.read())
    except Exception as e:
        return jsonify({"error": f"Falha ao ler o PDF: {str(e)}"}), 400

    comp = dados["comp"]

    # Carrega modelo HP publicado (ou default Levante)
    hp_porto = _load(_HP_PORT_FILE, HP_PORTFOLIOS_DEFAULT)
    perfis_hp = hp_porto.get("perfis", {})

    # Detecta perfil mais próximo pela composição atual
    def distancia(comp, modelo_perfil):
        return sum(abs(comp.get(cls, 0) - modelo_perfil.get(cls, 0)) for cls in CATS)

    melhor_perfil = "moderada"
    menor_dist = float("inf")
    for pk, pv in perfis_hp.items():
        d = distancia(comp, pv)
        if d < menor_dist:
            menor_dist = d
            melhor_perfil = pk

    modelo_escolhido = perfis_hp.get(melhor_perfil, {})
    label_perfil = modelo_escolhido.get("label", melhor_perfil.capitalize())
    referencia_hp = hp_porto.get("referencia", "Levante Asset")

    # Desvios por classe
    desvios = []
    for cls in CATS:
        atual = comp.get(cls, 0.0)
        alvo  = float(modelo_escolhido.get(cls, 0))
        desvio = round(atual - alvo, 2)
        status = "ok" if abs(desvio) <= 2 else ("atencao" if abs(desvio) <= 5 else "fora")
        desvios.append({
            "cls": cls,
            "label": LABELS.get(cls, cls),
            "atual": round(atual, 2),
            "alvo": round(alvo, 2),
            "desvio": desvio,
            "status": status,
        })
    desvios.sort(key=lambda x: abs(x["desvio"]), reverse=True)

    # Cenário macro publicado pelo HP
    cenario = _load(_HP_CENARIO_FILE, HP_CENARIO_DEFAULT)

    # Produtos sugeridos para fechar os gaps (classes sub-alocadas)
    hp_prods = _load(_HP_PROD_FILE, {})
    sugestoes_prods = []
    for d in desvios:
        if d["desvio"] < -2:  # sub-alocado — sugerir produto
            cls_prods = hp_prods.get(d["cls"], [])
            for p in cls_prods[:2]:
                sugestoes_prods.append({
                    "classe": d["cls"],
                    "label_classe": d["label"],
                    "gap": d["desvio"],
                    "produto": p,
                })

    # ── Alertas para o assessor ────────────────────────────────────────────────
    # Filtro por assessor_destino: pega globais (vazio) + direcionados ao assessor deste cliente
    assessor_atual = (dados.get("assessor","") or assessor).lower().strip()
    alertas_hp_raw   = _load(_HP_ALERTS_FILE, [])
    alertas_hp       = [
        a for a in alertas_hp_raw
        if not a.get("assessor_destino","").strip()
        or a.get("assessor_destino","").lower().strip() == assessor_atual
        or assessor_atual in a.get("assessor_destino","").lower()
    ]
    docs_publicados  = [d for d in _load(_HP_KNOW_FILE, []) if d.get("publicado")]
    tickers_carteira = [a["ticker"] for a in dados.get("acoes", [])] + \
                       [f["ticker"] for f in dados.get("fiis", [])]
    nomes_fundos     = [a.get("nome","") for a in dados.get("acoes", [])] + \
                       [f.get("nome","") for f in dados.get("fiis", [])]

    alertas_relevantes = []

    # 1. Alertas manuais do HP por ticker ou classe com desvio
    classes_fora = [d["cls"] for d in desvios if d["status"] == "fora"]
    for a in alertas_hp:
        if any(t and t in a.get("produto","") for t in tickers_carteira) \
        or a.get("classe","") in classes_fora:
            alertas_relevantes.append({**a, "origem_tipo": "hp_manual"})

    # 2. Varredura nos documentos publicados da Base de Conhecimento
    #    Procura menções a tickers ou nomes de fundos da carteira
    for doc in docs_publicados:
        texto_doc = (doc.get("texto","") or "").upper()
        for tk in tickers_carteira:
            if tk and len(tk) >= 4 and tk.upper() in texto_doc:
                # Extrai trecho ao redor da menção
                idx = texto_doc.find(tk.upper())
                trecho = doc.get("texto","")[max(0, idx-80):idx+200].replace("\n"," ").strip()
                alertas_relevantes.append({
                    "id":       f"kb_{doc['id']}_{tk}",
                    "produto":  tk,
                    "classe":   "",
                    "tipo":     "atencao",
                    "mensagem": f"Mencionado em \"{doc.get('nome','')}\" ({doc.get('fonte','HP')}): ...{trecho}...",
                    "origem":   doc.get("nome","Base de Conhecimento"),
                    "origem_tipo": "knowledge_base",
                    "data":     doc.get("data",""),
                })
                break  # um alerta por doc por ticker

    # 3. Alertas automáticos por concentração e desvios graves
    comp_atual = dados.get("comp", {})
    LABELS_CLS = {"pos_fixado":"Pós Fixado","inflacao":"Inflação","pre_fixado":"Pré Fixado",
                  "acoes":"Ações","fiis":"FIIs","multimercado":"Multimercado",
                  "internacional":"Internacional","alternativos":"Alternativos","criptomoedas":"Criptomoedas"}
    for cls, pct in comp_atual.items():
        if pct >= 80:
            alertas_relevantes.append({
                "id": f"auto_conc_{cls}",
                "produto": LABELS_CLS.get(cls, cls),
                "classe": cls,
                "tipo": "urgente",
                "mensagem": f"Concentração muito alta: {pct:.1f}% em {LABELS_CLS.get(cls, cls)}. Discuta diversificação com o cliente.",
                "origem": "Análise automática",
                "origem_tipo": "auto",
                "data": datetime.now().strftime("%d/%m/%Y"),
            })
        elif pct >= 60:
            alertas_relevantes.append({
                "id": f"auto_conc_{cls}",
                "produto": LABELS_CLS.get(cls, cls),
                "classe": cls,
                "tipo": "atencao",
                "mensagem": f"Alocação elevada: {pct:.1f}% em {LABELS_CLS.get(cls, cls)}. Verifique se está alinhado ao objetivo do cliente.",
                "origem": "Análise automática",
                "origem_tipo": "auto",
                "data": datetime.now().strftime("%d/%m/%Y"),
            })

    for d_desvio in desvios:
        if abs(d_desvio.get("desvio_pp", 0)) >= 25:
            cls  = d_desvio["cls"]
            sinal = "+" if d_desvio["desvio_pp"] > 0 else ""
            alertas_relevantes.append({
                "id": f"auto_desvio_{cls}",
                "produto": LABELS_CLS.get(cls, cls),
                "classe": cls,
                "tipo": "atencao",
                "mensagem": f"Desvio significativo: {sinal}{d_desvio['desvio_pp']:.1f}pp em relação ao modelo. Pauta prioritária para rebalanceamento.",
                "origem": "Análise automática",
                "origem_tipo": "auto",
                "data": datetime.now().strftime("%d/%m/%Y"),
            })

    # Remove duplicatas por id
    seen = set()
    alertas_unicos = []
    for a in alertas_relevantes:
        k = a.get("id") or a.get("produto","")
        if k not in seen:
            seen.add(k)
            alertas_unicos.append(a)

    # Ordena: urgente primeiro, depois atenção, depois info
    _ord = {"urgente": 0, "atencao": 1, "info": 2}
    alertas_unicos.sort(key=lambda a: _ord.get(a.get("tipo","info"), 2))

    # Calls de ações relevantes para o cliente
    calls_hp = [c for c in _load(_HP_CALLS_FILE, []) if c.get("ativo", True)]
    calls_relevantes = []
    perfil_cliente = label_perfil.lower() if label_perfil else ""

    for call in calls_hp:
        # Aparece se: o cliente já tem o ticker OU o perfil está na lista de perfis do call
        ticker_ok  = any(call.get("ticker","") in a.get("ticker","") for a in dados.get("acoes",[]))
        perfil_ok  = not call.get("perfis") or any(p in perfil_cliente for p in call.get("perfis",[]))
        if ticker_ok or perfil_ok:
            calls_relevantes.append(call)

    # ── Salva carteira do cliente no Redis ────────────────────────────────────
    conta_key = (dados.get("conta") or dados.get("assessor") or "sem_conta").strip().replace(" ","_")
    clientes = _load(_CLIENTS_FILE, {})
    clientes[conta_key] = {
        "conta":       dados.get("conta",""),
        "assessor":    dados.get("assessor","") or assessor,
        "patrimonio":  dados.get("patrimonio", 0),
        "perfil":      label_perfil,
        "composicao":  comp,
        "acoes":       dados.get("acoes",[])[:30],
        "fiis":        dados.get("fiis",[])[:20],
        "rf_ativos":   dados.get("rf_ativos",[])[:40],
        "data_ref":    dados.get("data_ref",""),
        "atualizado_em": datetime.now().strftime("%d/%m/%Y %H:%M"),
    }
    _save(_CLIENTS_FILE, clientes)

    # ── Notícias online para os principais tickers ────────────────────────────
    alertas_mercado = []
    try:
        import requests as _req, xml.etree.ElementTree as ET, json as _json
        todos_ativos = dados.get("acoes",[]) + dados.get("fiis",[])
        top_tickers = [
            a.get("ticker","") for a in
            sorted(todos_ativos, key=lambda x: x.get("perc_carteira", x.get("perc", 0)), reverse=True)
            if a.get("ticker")
        ][:5]

        noticias_por_ticker = {}
        for tk in top_tickers:
            try:
                url = f"https://news.google.com/rss/search?q={tk}+bolsa&hl=pt-BR&gl=BR&ceid=BR:pt"
                r = _req.get(url, timeout=3, headers={"User-Agent": "Mozilla/5.0"})
                root = ET.fromstring(r.content)
                titulos = [item.findtext("title","") for item in root.findall(".//item")[:3]]
                titulos = [t for t in titulos if t]
                if titulos:
                    noticias_por_ticker[tk] = titulos
            except Exception:
                pass

        if noticias_por_ticker:
            noticias_txt = "\n".join(
                f"{tk}: " + " | ".join(ts) for tk, ts in noticias_por_ticker.items()
            )
            _ai = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY",""))
            resp_ai = _ai.messages.create(
                model="claude-haiku-4-5-20251001",
                max_tokens=700,
                messages=[{"role":"user","content":
                    f"Você é um assistente de assessor de investimentos. Analise estas notícias recentes "
                    f"sobre ativos na carteira de um cliente (perfil {label_perfil}, "
                    f"patrimônio R$ {dados.get('patrimonio',0):,.0f}).\n\n"
                    f"Notícias:\n{noticias_txt}\n\n"
                    f"Para cada ativo com notícia relevante (risco, oportunidade, evento corporativo, "
                    f"resultado, dividendo, split, rebaixamento), gere um alerta ACIONÁVEL para o assessor. "
                    f"Ignore notícias sem impacto prático. "
                    f"Responda APENAS com JSON array (sem markdown):\n"
                    f'[{{"ticker":"X","tipo":"urgente|atencao|info","mensagem":"Texto acionável para o assessor (max 160 chars)","fonte":"nome da notícia em 1 frase"}}]'
                }]
            )
            raw = resp_ai.content[0].text.strip()
            if raw.startswith("```"): raw = raw.split("```")[1].lstrip("json").strip()
            parsed = _json.loads(raw)
            for an in parsed:
                if not isinstance(an, dict) or not an.get("ticker"): continue
                alertas_mercado.append({
                    "id":       f"news_{an['ticker']}_{datetime.now().strftime('%H%M%S')}",
                    "produto":  an["ticker"],
                    "classe":   "",
                    "tipo":     an.get("tipo","atencao"),
                    "mensagem": an.get("mensagem",""),
                    "origem":   f"📰 {an.get('fonte','Notícias do mercado')}",
                    "origem_tipo": "mercado",
                    "data":     datetime.now().strftime("%d/%m/%Y"),
                })
    except Exception as _e:
        app.logger.warning(f"News fetch warning: {_e}")

    # Mescla: alertas internos primeiro, depois mercado
    alertas_final = alertas_unicos + alertas_mercado
    _ord2 = {"urgente": 0, "atencao": 1, "info": 2}
    alertas_final.sort(key=lambda a: _ord2.get(a.get("tipo","info"), 2))

    resultado = {
        "conta":       dados["conta"],
        "assessor":    dados["assessor"] or assessor,
        "data_ref":    dados["data_ref"],
        "patrimonio":  dados["patrimonio"],
        "perfil_detectado": label_perfil,
        "referencia_modelo": referencia_hp,
        "composicao":  comp,
        "caixa":       dados["caixa"],
        "rent":        dados["rent"],
        "desvios":     desvios,
        "acoes":       dados["acoes"],
        "fiis":        dados["fiis"],
        "sugestoes_produtos": sugestoes_prods,
        "cenario_macro": {
            "global":        cenario.get("global",""),
            "brasil":        cenario.get("brasil",""),
            "posicionamento":cenario.get("posicionamento",""),
            "vieses":        cenario.get("vieses",{}),
            "referencia":    cenario.get("referencia",""),
        },
        "alertas_relevantes": alertas_final,
        "alertas_mercado":    alertas_mercado,
        "calls_relevantes":   calls_relevantes,
        "carteira_salva":     True,
        "analisado_em": datetime.now().strftime("%d/%m/%Y %H:%M"),
    }
    return jsonify(resultado)


@app.route("/api/gerar-apresentacao", methods=["POST"])
def gerar_apresentacao():
    """Gera PDF de apresentação de reunião com dados HP + IA."""
    try:
        import sys, os
        _api_dir = os.path.dirname(os.path.abspath(__file__))
        if _api_dir not in sys.path:
            sys.path.insert(0, _api_dir)
        from apresentacao import gerar_apresentacao as _gerar
    except Exception as e:
        app.logger.error(f"Erro ao importar apresentacao: {e}")
        return jsonify({"error": f"Erro ao carregar módulo de apresentação: {e}"}), 500

    body = request.get_json()
    if not body:
        return jsonify({"error": "Dados não enviados"}), 400

    try:
        # Enriquece com dados HP publicados
        hp_cenario= _load(_HP_CENARIO_FILE, HP_CENARIO_DEFAULT)
        alertas_hp= _load(_HP_ALERTS_FILE, [])
        docs_pub  = [d for d in _load(_HP_KNOW_FILE, []) if d.get("publicado")]

        body["cenario_macro"] = {
            "global":        hp_cenario.get("global",""),
            "brasil":        hp_cenario.get("brasil",""),
            "posicionamento":hp_cenario.get("posicionamento",""),
            "vieses":        hp_cenario.get("vieses",{}),
            "fontes":        [{"nome": d.get("nome",""), "tipo": d.get("tipo",""), "fonte": d.get("fonte",""), "data": d.get("data","")} for d in docs_pub],
        }

        # Alertas relevantes para os ativos do cliente
        tickers = [a.get("ticker","") for a in body.get("acoes",[])] + \
                  [f.get("ticker","") for f in body.get("fiis",[])]
        body["alertas_relevantes"] = [
            a for a in alertas_hp
            if any(t in a.get("produto","") for t in tickers)
            or a.get("classe","") in [d.get("cls","") for d in body.get("desvios",[]) if d.get("status")=="fora"]
        ]

        # Garante campos de rentabilidade
        rent = body.get("rent", {})
        port_rent = rent.get("portfolio", {}) if isinstance(rent, dict) else {}
        cdi_rent  = rent.get("cdi", {}) if isinstance(rent, dict) else {}
        rent_12m  = port_rent.get("12m", 0)
        cdi_12m   = cdi_rent.get("12m", 0)
        body["rent_12m"] = rent_12m
        body["pct_cdi"]  = round(rent_12m / cdi_12m * 100, 1) if cdi_12m else 0

        pdf_bytes = _gerar(body)
    except Exception as e:
        import traceback
        app.logger.error(f"Erro ao gerar apresentação: {traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500

    nome = body.get("nome_cliente","cliente").replace(" ","-").lower()
    data = body.get("data_ref","").replace("/","-")
    filename = f"apresentacao_{nome}_{data}.pdf"

    return send_file(
        io.BytesIO(pdf_bytes),
        mimetype="application/pdf",
        as_attachment=True,
        download_name=filename,
    )


@app.route("/api/gerar-pptx", methods=["POST"])
def gerar_pptx():
    """Gera PPTX de apresentação de reunião com identidade Braúna."""
    try:
        import sys, os
        _api_dir = os.path.dirname(os.path.abspath(__file__))
        if _api_dir not in sys.path:
            sys.path.insert(0, _api_dir)
        from apresentacao_pptx import gerar_apresentacao_pptx as _gerar_pptx
    except Exception as e:
        app.logger.error(f"Erro ao importar apresentacao_pptx: {e}")
        return jsonify({"error": f"Erro ao carregar módulo PPTX: {e}"}), 500

    body = request.get_json()
    if not body:
        return jsonify({"error": "Dados não enviados"}), 400

    try:
        hp_cenario   = _load(_HP_CENARIO_FILE, HP_CENARIO_DEFAULT)
        macro_bcb    = buscar_macro_bcb()
        docs_pub     = [d for d in _load(_HP_KNOW_FILE, []) if d.get("publicado")]
        calls_hp     = [c for c in _load(_HP_CALLS_FILE, []) if c.get("ativo") is not False]
        estruturadas = _load(_HP_ESTRUTURADAS_FILE, [])
        produtos_hp  = _load(_HP_PROD_FILE, {})
        perfil_cli   = body.get("perfil", "moderada").lower()

        # Cenário macro completo do Head
        body["cenario_macro"] = {
            "global":         hp_cenario.get("global",""),
            "brasil":         hp_cenario.get("brasil",""),
            "posicionamento": hp_cenario.get("posicionamento",""),
            "vieses":         hp_cenario.get("vieses",{}),
            "sinais":         hp_cenario.get("sinais", []),
            "selic_meta":     macro_bcb.get("selic_meta"),
            "ipca_12m":       macro_bcb.get("ipca_12m"),
            "fontes":         [{"nome": d.get("nome",""), "tipo": d.get("tipo",""), "fonte": d.get("fonte","")} for d in docs_pub],
        }

        # Calls de ações do Head (se o frontend não enviou nenhum)
        if not body.get("calls"):
            body["calls"] = calls_hp[:8]

        # Operações estruturadas
        if not body.get("estruturadas"):
            body["estruturadas"] = estruturadas[:6]

        # Produtos do Head filtrados pelo perfil do cliente
        # (usa o que veio do frontend; complementa se vazio)
        if not body.get("sugestoes_produtos"):
            ORDEM_PERFIL = ["super_conservadora","conservadora","moderada","agressiva","super_agressiva"]
            idx = ORDEM_PERFIL.index(perfil_cli) if perfil_cli in ORDEM_PERFIL else 2
            sugs = []
            for cls, lista in produtos_hp.items():
                if not isinstance(lista, list): continue
                for p in lista:
                    perfis_prod = p.get("perfis") or []
                    if not perfis_prod:
                        sugs.append(p)
                    else:
                        perfis_norm = [x.lower().replace(" ","_") for x in perfis_prod]
                        if any(ORDEM_PERFIL.index(pf) <= idx for pf in perfis_norm if pf in ORDEM_PERFIL):
                            sugs.append(p)
            body["sugestoes_produtos"] = sugs[:8]

        rent = body.get("rent", {})
        port_rent = rent.get("portfolio", {}) if isinstance(rent, dict) else {}
        cdi_rent  = rent.get("cdi", {}) if isinstance(rent, dict) else {}
        rent_12m  = port_rent.get("12m", 0)
        cdi_12m   = cdi_rent.get("12m", 0)
        body["rent_12m"] = rent_12m
        body["pct_cdi"]  = round(rent_12m / cdi_12m * 100, 1) if cdi_12m else 0

        pptx_bytes = _gerar_pptx(body)
    except Exception as e:
        import traceback
        app.logger.error(f"Erro ao gerar PPTX: {traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500

    nome = body.get("nome_cliente","cliente").replace(" ","-").lower()
    data = body.get("data_ref","").replace("/","-")
    filename = f"apresentacao_{nome}_{data}.pptx"

    return send_file(
        io.BytesIO(pptx_bytes),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=filename,
    )


@app.route("/api/hp/produtos/salvar", methods=["POST"])
def hp_produtos_salvar():
    """Salva produtos com metadata de quem indicou e quando."""
    data = request.get_json()
    produtos_por_classe = data.get("produtos", {})
    indicado_por = data.get("indicado_por", "Head de Produtos")
    indicado_em  = datetime.now().strftime("%d/%m/%Y %H:%M")

    # Enriquecer cada produto com metadata
    for cls, lista in produtos_por_classe.items():
        for p in lista:
            if not p.get("indicado_por"):
                p["indicado_por"] = indicado_por
                p["indicado_em"]  = indicado_em

    _save(_HP_PROD_FILE, produtos_por_classe)

    total = sum(len(v) for v in produtos_por_classe.values())
    return jsonify({"ok": True, "total": total, "indicado_por": indicado_por, "indicado_em": indicado_em})

@app.route("/api/analyze", methods=["POST"])
def analyze():
    assessor = request.form.get("assessor","Assessor")
    nome   = request.form.get("nome","")
    perfil = request.form.get("perfil","conservadora")
    objetivo = request.form.get("objetivo","")
    data   = datetime.now().strftime("%d/%m/%Y")
    pdf    = request.files.get("pdf")
    checklist_raw  = request.form.get("checklist","{}")
    crosssell_raw  = request.form.get("cross_sell","{}")

    if not pdf:
        return jsonify({"error":"PDF não enviado"}), 400

    try: checklist_input = json.loads(checklist_raw)
    except: checklist_input = {}

    try: cross_input = json.loads(crosssell_raw)
    except: cross_input = {}

    # Usa carta carregada pelo admin (se houver)
    carta_info  = _load("/tmp/brauna_carta.json", {})
    carta_texto = carta_info.get("texto","")

    xp_parsed = extrair_xperformance(pdf.read())
    comp      = xp_parsed.get("comp", {c: 0.0 for c in CATS})
    caixa     = xp_parsed.get("caixa", 0.0)
    rent      = xp_parsed.get("rent", {})
    patrimonio= xp_parsed.get("patrimonio", 0.0)

    modelo  = MODELOS.get(perfil, MODELOS["conservadora"])
    desvios = calcular_desvios(comp, modelo)
    macro   = buscar_macro_bcb()
    ctx     = carregar_contexto()

    recomendacoes, alertas = gerar_recomendacoes(desvios, perfil, macro, ctx, carta_texto)
    checklist_servir, score_servir, pendentes_criticos, pendentes_altos = avaliar_modelo_servir(checklist_input)
    cross_sell_result, cross_tem, cross_nao_tem = avaliar_cross_sell(cross_input)

    # Determina status geral
    rent12 = 0
    if rent.get("portfolio") and rent.get("cdi") and rent["cdi"].get("12m",0)>0:
        rent12 = round(rent["portfolio"]["12m"] / rent["cdi"]["12m"] * 100, 1)
    status = "critico" if (pendentes_criticos>0 or rent12<70) else ("atencao" if (pendentes_altos>0 or len(alertas)>0) else "ok")

    # Salva resumo para o líder
    resumo = {
        "id": str(uuid.uuid4())[:8],
        "assessor": assessor,
        "nome": nome,
        "perfil": perfil,
        "data": data,
        "patrimonio": patrimonio,
        "score_servir": score_servir,
        "pendentes_criticos": pendentes_criticos,
        "pendentes_altos": pendentes_altos,
        "cross_ativos": len(cross_tem),
        "rent12_pct_cdi": rent12,
        "alertas_count": len(alertas),
        "status": status,
        "nota_lider": "",
        "analisado_em": datetime.now().strftime("%d/%m/%Y %H:%M"),
    }
    clientes = load_clientes()
    # Substitui se mesmo nome+assessor
    clientes = [c for c in clientes if not (c["nome"]==nome and c["assessor"]==assessor)]
    clientes.insert(0, resumo)
    save_clientes(clientes[:500])

    # Salva ficha do cliente (dados do formulário para pré-carregar na próxima vez)
    fichas = load_fichas()
    fkey = f"{assessor}|{nome}".lower().strip()
    ficha_salva = fichas.get(fkey, {})
    fichas[fkey] = {
        "nome": nome, "assessor": assessor, "perfil": perfil,
        "objetivo": objetivo,
        "checklist": checklist_input,
        "cross_ativos": list(cross_tem),
        "atualizado_em": datetime.now().strftime("%d/%m/%Y %H:%M"),
    }
    save_fichas(fichas)

    # Salva histórico de análises (até 4 por cliente)
    hist = load_hist()
    entrada_hist = {
        "data": datetime.now().strftime("%d/%m/%Y"),
        "hora": datetime.now().strftime("%H:%M"),
        "patrimonio": patrimonio,
        "perfil": perfil,
        "composicao": comp,
        "score_servir": score_servir,
        "cross_ativos": len(cross_tem),
        "status": status,
        "rent12_pct_cdi": rent12,
        "objetivo": objetivo,
    }
    reg_hist = hist.get(fkey, {"assessor": assessor, "nome": nome, "perfil": perfil, "objetivo": objetivo, "entradas": []})
    reg_hist["perfil"]   = perfil
    reg_hist["objetivo"] = objetivo
    reg_hist["entradas"] = ([entrada_hist] + reg_hist.get("entradas", []))[:4]
    hist[fkey] = reg_hist
    save_hist(hist)

    # Sugestões ativas filtradas para este cliente
    suge_data = load_suge()
    suge_ativa = next((s for s in suge_data.get("historico",[]) if s.get("ativa")), None)
    sugestoes_filtradas = filtrar_sugestoes(suge_ativa, perfil, comp)

    resultado = {
        "id": resumo["id"],
        "assessor": assessor,
        "nome":nome, "perfil":perfil, "data_ref":data, "objetivo":objetivo,
        "patrimonio":patrimonio, "caixa":caixa,
        "composicao":comp, "desvios":desvios, "rent":rent,
        "macro":macro, "recomendacoes":recomendacoes, "alertas":alertas,
        "visao_gestores": ctx.get("visao_gestores",[]),
        "ref_contexto": ctx.get("referencia",""),
        "checklist_servir": checklist_servir,
        "score_servir": score_servir,
        "pendentes_criticos": pendentes_criticos,
        "pendentes_altos": pendentes_altos,
        "cross_sell": cross_sell_result,
        "cross_tem": cross_tem,
        "cross_nao_tem": cross_nao_tem,
        "carta_info": {"nome": carta_info.get("nome",""), "atualizado": carta_info.get("atualizado","")},
        "sugestoes": sugestoes_filtradas,
    }
    return jsonify(resultado)


PRUNUS_MODELOS = {
    "conservadora": {
        "retorno_alvo": "CDI + 1,5% a.a.",
        "composicao": {"Renda Fixa": 72, "Multimercado": 15, "Internacional": 8, "Alternativos": 5},
        "descricao": "Carteira com foco em preservação de capital e geração de renda consistente através de gestão ativa profissional.",
    },
    "moderada": {
        "retorno_alvo": "IPCA + 4,5% a.a.",
        "composicao": {"Renda Fixa": 52, "Multimercado": 20, "Renda Variável": 15, "Internacional": 10, "Alternativos": 3},
        "descricao": "Portfólio equilibrado com diversificação ampla, buscando superar inflação com risco controlado.",
    },
    "arrojada": {
        "retorno_alvo": "IBOV + 2% a.a.",
        "composicao": {"Renda Variável": 35, "Renda Fixa": 28, "Multimercado": 17, "Internacional": 15, "Alternativos": 5},
        "descricao": "Carteira de maior crescimento com exposição diversificada à renda variável e mercados internacionais.",
    },
    "agressiva": {
        "retorno_alvo": "IBOV + 4% a.a.",
        "composicao": {"Renda Variável": 50, "Internacional": 20, "Multimercado": 15, "Renda Fixa": 10, "Alternativos": 5},
        "descricao": "Máxima exposição a crescimento patrimonial de longo prazo com gestão especializada.",
    },
}


def buscar_hist_macro():
    import urllib.request as _ur
    from datetime import date, timedelta
    hoje = date.today()
    ini  = (hoje - timedelta(days=220)).strftime("%d/%m/%Y")
    fim  = hoje.strftime("%d/%m/%Y")
    out  = {"cdi": [], "ipca": [], "focus": {}}
    for serie, chave in [("4391", "cdi"), ("433", "ipca")]:
        try:
            url = (f"https://api.bcb.gov.br/dados/serie/bcdata.sgs.{serie}/dados"
                   f"?formato=json&dataInicial={ini}&dataFinal={fim}")
            with _ur.urlopen(url, timeout=7) as r:
                dados = json.loads(r.read())
                out[chave] = [{"data": d["data"], "valor": float(d["valor"])} for d in dados[-6:]]
        except:
            pass
    try:
        url = ("https://olinda.bcb.gov.br/olinda/servico/Expectativas/versao/v1/odata/"
               "ExpectativaMercadoAnuais?%24top=80"
               "&%24filter=Indicador%20eq%20'IPCA'%20or%20Indicador%20eq%20'Selic'%20or%20Indicador%20eq%20'PIB%20Total'"
               "&%24format=json&%24select=Indicador,Data,Mediana")
        with _ur.urlopen(url, timeout=7) as r:
            for item in json.loads(r.read()).get("value", []):
                ind = item.get("Indicador", "")
                if ind and ind not in out["focus"]:
                    out["focus"][ind] = item.get("Mediana")
    except:
        pass
    return out


def gerar_ppt(dados):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    BG    = RGBColor(0x0D, 0x0D, 0x0D)
    DARK  = RGBColor(0x11, 0x11, 0x11)
    CARD  = RGBColor(0x1A, 0x1A, 0x1A)
    GOLD  = RGBColor(0xD6, 0xB2, 0x7A)
    WHITE = RGBColor(0xF0, 0xF0, 0xF0)
    GRAY  = RGBColor(0x55, 0x55, 0x55)
    GRAY2 = RGBColor(0x2A, 0x2A, 0x2A)
    GREEN = RGBColor(0x5D, 0xCA, 0xA5)
    RED   = RGBColor(0xFF, 0x6B, 0x6B)
    YEL   = RGBColor(0xFF, 0xD9, 0x66)
    BLUE  = RGBColor(0x8B, 0x9F, 0xE8)
    PURP  = RGBColor(0xB8, 0xA0, 0xE8)
    PINK  = RGBColor(0xFF, 0x8E, 0x8E)

    nome      = dados.get("nome", "Cliente")
    assessor  = dados.get("assessor", "Assessor")
    perfil    = dados.get("perfil", "conservadora")
    pat       = dados.get("patrimonio", 0)
    comp      = dados.get("composicao", {})
    desvios   = dados.get("desvios", [])
    rent      = dados.get("rent", {})
    macro     = dados.get("macro", {})
    suge      = dados.get("sugestoes") or {}
    cross_tem = dados.get("cross_tem", [])
    checklist = dados.get("checklist_servir", [])
    hist_m    = dados.get("hist_macro", {})
    objetivo  = dados.get("objetivo", "")
    pat_caindo = dados.get("patrimonio_caindo", False)

    PLABEL = {"conservadora": "Conservadora", "moderada": "Moderada",
              "arrojada": "Arrojada", "agressiva": "Agressiva"}
    fmt_brl = lambda v: f"R$ {int(v):,.0f}".replace(",", ".") if v else "—"
    data_hoje = datetime.now().strftime("%d/%m/%Y")
    mes_ano   = datetime.now().strftime("%B de %Y").capitalize()

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def nsl():
        return prs.slides.add_slide(blank)

    def rc(slide, x, y, w, h, fill, lc=None):
        s = slide.shapes.add_shape(1, x, y, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = fill
        if lc: s.line.color.rgb = lc
        else:  s.line.fill.background()
        return s

    def tb(slide, text, x, y, w, h, sz=11, bold=False, col=None,
           align=PP_ALIGN.LEFT, italic=False):
        col = col or WHITE
        box = slide.shapes.add_textbox(x, y, w, h)
        tf  = box.text_frame; tf.word_wrap = True
        p   = tf.paragraphs[0]; p.alignment = align
        r   = p.add_run(); r.text = str(text)
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = col
        r.font.name = "Calibri"
        return box

    def bar(slide, x, y, w, h, pct, col):
        rc(slide, x, y, w, h, GRAY2)
        filled = int(w * min(max(pct, 0), 100) / 100)
        if filled > 0:
            rc(slide, x, y, filled, h, col)

    def hdr(slide, title, page):
        rc(slide, 0, 0, W, H, BG)
        rc(slide, 0, 0, W, Inches(0.06), GOLD)
        rc(slide, 0, Inches(0.06), W, Inches(0.72), DARK)
        tb(slide, "BRAÚNA INVESTIMENTOS", Inches(0.25), Inches(0.12),
           Inches(4), Inches(0.55), sz=10, bold=True, col=GOLD)
        tb(slide, title.upper(), Inches(4), Inches(0.12),
           Inches(8), Inches(0.55), sz=13, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
        tb(slide, f"{page}/10", Inches(12.4), Inches(0.12),
           Inches(0.8), Inches(0.55), sz=10, col=GRAY, align=PP_ALIGN.RIGHT)
        rc(slide, 0, Inches(7.3), W, Inches(0.2), DARK)
        footer = f"{nome}  ·  Perfil {PLABEL.get(perfil, perfil)}  ·  Assessor: {assessor}  ·  {data_hoje}"
        tb(slide, footer, Inches(0.25), Inches(7.31), Inches(12.8), Inches(0.18), sz=8, col=GRAY)

    def metric(slide, x, y, w, h, label, value, val_col=None):
        val_col = val_col or GOLD
        rc(slide, x, y, w, h, CARD)
        tb(slide, label.upper(), x + Inches(0.12), y + Inches(0.1),
           w - Inches(0.24), Inches(0.25), sz=8, col=GRAY)
        tb(slide, value, x + Inches(0.12), y + Inches(0.38),
           w - Inches(0.24), h - Inches(0.48), sz=15, bold=True, col=val_col)

    def dv_col(d):
        if d < -1.5: return RED
        if d > 1.5:  return YEL
        return GREEN

    # ── SLIDE 1 — CAPA ───────────────────────────────────────────
    s1 = nsl()
    rc(s1, 0, 0, W, H, BG)
    rc(s1, 0, 0, Inches(4.2), H, DARK)
    rc(s1, 0, 0, Inches(0.08), H, GOLD)
    tb(s1, "BRAÚNA", Inches(0.25), Inches(1.2), Inches(3.8), Inches(1.1),
       sz=32, bold=True, col=GOLD)
    tb(s1, "INVESTIMENTOS", Inches(0.25), Inches(2.15), Inches(3.8), Inches(0.6),
       sz=15, bold=True, col=WHITE)
    rc(s1, Inches(0.25), Inches(2.95), Inches(3.5), Inches(0.04), GOLD)
    tb(s1, "Análise Personalizada de Carteira", Inches(0.25), Inches(3.1), Inches(3.8), Inches(0.6),
       sz=11, col=GRAY, italic=True)
    tb(s1, mes_ano, Inches(0.25), Inches(3.85), Inches(3.8), Inches(0.4),
       sz=10, col=GRAY)
    tb(s1, "Preparado exclusivamente para", Inches(4.6), Inches(1.0), Inches(8.5), Inches(0.4),
       sz=11, col=GRAY, italic=True)
    tb(s1, nome.upper(), Inches(4.6), Inches(1.45), Inches(8.5), Inches(1.4),
       sz=30, bold=True, col=WHITE)
    rc(s1, Inches(4.6), Inches(2.9), Inches(6.5), Inches(0.05), GOLD)
    infos = [
        ("Perfil de Investimento", PLABEL.get(perfil, perfil), GOLD),
        ("Patrimônio em Análise", fmt_brl(pat), WHITE),
        ("Assessor Responsável", assessor, WHITE),
    ]
    if objetivo:
        infos.append(("Objetivo do Cliente", objetivo, BLUE))
    for i, (lbl, val, col) in enumerate(infos):
        yy = Inches(3.1) + i * Inches(0.82)
        tb(s1, lbl.upper(), Inches(4.6), yy, Inches(8.5), Inches(0.25), sz=8, col=GRAY)
        tb(s1, val, Inches(4.6), yy + Inches(0.22), Inches(8.5), Inches(0.5),
           sz=18, bold=True, col=col)
    tb(s1, "DOCUMENTO CONFIDENCIAL — USO EXCLUSIVO DO CLIENTE",
       Inches(4.6), Inches(6.9), Inches(8.5), Inches(0.35), sz=8, col=GRAY2, italic=True)

    # ── SLIDE 2 — PANORAMA DE MERCADO ────────────────────────────
    s2 = nsl()
    hdr(s2, "Panorama de Mercado", 2)
    Y0 = Inches(0.95)
    cdi_hist  = hist_m.get("cdi", [])
    ipca_hist = hist_m.get("ipca", [])

    tb(s2, "CDI — Últimos 6 Meses (%)", Inches(0.25), Y0, Inches(6.2), Inches(0.3),
       sz=10, bold=True, col=GOLD)
    if cdi_hist:
        mv = max(d["valor"] for d in cdi_hist) or 1
        for j, d in enumerate(cdi_hist):
            yy = Y0 + Inches(0.35) + j * Inches(0.47)
            tb(s2, d["data"][:7], Inches(0.25), yy, Inches(1.1), Inches(0.35), sz=9, col=GRAY)
            bar(s2, Inches(1.4), yy + Inches(0.1), Inches(3.5), Inches(0.22), d["valor"] / mv * 100, BLUE)
            tb(s2, f"{d['valor']:.3f}%", Inches(5.0), yy, Inches(1.3), Inches(0.35),
               sz=9, bold=True, col=WHITE)
    else:
        tb(s2, f"CDI atual: {macro.get('selic_meta', 0):.2f}% a.a.", Inches(0.25),
           Y0 + Inches(0.35), Inches(6.2), Inches(0.3), sz=11, col=WHITE)

    tb(s2, "IPCA — Últimos 6 Meses (%)", Inches(6.8), Y0, Inches(6.2), Inches(0.3),
       sz=10, bold=True, col=GOLD)
    if ipca_hist:
        mv = max(d["valor"] for d in ipca_hist) or 1
        for j, d in enumerate(ipca_hist):
            yy = Y0 + Inches(0.35) + j * Inches(0.47)
            tb(s2, d["data"][:7], Inches(6.8), yy, Inches(1.1), Inches(0.35), sz=9, col=GRAY)
            bar(s2, Inches(8.0), yy + Inches(0.1), Inches(3.5), Inches(0.22), d["valor"] / mv * 100, YEL)
            tb(s2, f"{d['valor']:.2f}%", Inches(11.6), yy, Inches(1.3), Inches(0.35),
               sz=9, bold=True, col=WHITE)
    else:
        tb(s2, f"IPCA 12M: {macro.get('ipca_12m', 0):.2f}%", Inches(6.8),
           Y0 + Inches(0.35), Inches(6.2), Inches(0.3), sz=11, col=WHITE)

    rc(s2, Inches(6.55), Y0, Inches(0.03), Inches(3.6), GRAY2)

    rc(s2, Inches(0.25), Inches(4.35), W - Inches(0.5), Inches(0.04), GRAY2)
    tb(s2, "RELATÓRIO FOCUS — Expectativas do Mercado", Inches(0.25), Inches(4.5),
       Inches(12), Inches(0.3), sz=10, bold=True, col=GOLD)
    focus = hist_m.get("focus", {})
    focus_items = [
        ("IPCA 2026 (Exp.)",
         f"{focus['IPCA']}%" if focus.get("IPCA") else f"{macro.get('ipca_12m', 0):.2f}% (12M real)", YEL),
        ("Selic 2026 (Exp.)",
         f"{focus['Selic']}%" if focus.get("Selic") else f"{macro.get('selic_meta', 0):.2f}% (atual)", BLUE),
        ("PIB 2026 (Exp.)",
         f"{focus['PIB Total']}%" if focus.get("PIB Total") else "—", GREEN),
    ]
    for i, (lbl, val, col) in enumerate(focus_items):
        metric(s2, Inches(0.25) + i * Inches(4.3), Inches(4.9), Inches(4.1), Inches(1.15), lbl, val, col)
    tb(s2,
       "A conjuntura atual indica juros elevados com perspectiva de cortes graduais. "
       "A inflação acima do centro da meta favorece posições indexadas ao IPCA. "
       "Recomenda-se cautela na duration dos ativos de renda fixa e diversificação em ativos reais e internacionais.",
       Inches(0.25), Inches(6.2), Inches(12.8), Inches(0.8), sz=10, col=GRAY, italic=True)

    # ── SLIDES 3-8 — CARTEIRA POR SETOR ──────────────────────────
    SETORES = [
        ("Renda Fixa — Pós Fixado (CDI)", ["pos_fixado"], BLUE,
         "Proteção da liquidez com captura da taxa CDI. Indicado para reserva de curto prazo e parte estrutural da carteira."),
        ("Renda Fixa — Inflação (IPCA+)", ["inflacao"], YEL,
         "Preservação do poder de compra com retorno real garantido. Ideal para objetivos de médio e longo prazo."),
        ("Renda Fixa — Pré-fixado", ["pre_fixado"], PURP,
         "Trava de rentabilidade futura. Estratégico em cenário de expectativa de queda de juros."),
        ("Renda Variável — Ações & FIIs", ["acoes", "fiis"], GREEN,
         "Crescimento patrimonial de longo prazo com exposição ao mercado acionário brasileiro e fundos imobiliários."),
        ("Multimercado", ["multimercado"], GOLD,
         "Diversificação por meio de gestores especializados. Busca de retorno descorrelacionado do CDI."),
        ("Internacional & Alternativos", ["internacional", "alternativos", "criptomoedas"], PINK,
         "Proteção cambial e diversificação global. Acesso a classes de ativos complementares e descorrelacionadas."),
    ]

    for pg, (titulo, chaves, cor, descricao) in enumerate(SETORES, start=3):
        ss = nsl()
        hdr(ss, f"Carteira — {titulo}", pg)

        real = sum(comp.get(k, 0) for k in chaves)
        alvo = 0.0
        for k in chaves:
            alvo += next((d.get("alvo", 0) for d in desvios
                         if k in d.get("key", d.get("label", "")).lower()), 0.0)
        desvio = real - alvo
        dc = dv_col(desvio)
        val_r = pat * real / 100 if pat > 0 else 0

        metric(ss, Inches(0.25), Inches(0.92), Inches(3.0), Inches(1.1),
               "Alocação atual", f"{real:.1f}%", cor)
        metric(ss, Inches(3.4), Inches(0.92), Inches(3.0), Inches(1.1),
               "Meta do modelo", f"{alvo:.1f}%", GRAY)
        metric(ss, Inches(6.55), Inches(0.92), Inches(3.0), Inches(1.1),
               "Desvio", f"{desvio:+.1f}%", dc)
        metric(ss, Inches(9.7), Inches(0.92), Inches(3.4), Inches(1.1),
               "Valor estimado", fmt_brl(val_r), WHITE)

        Y1 = Inches(2.2)
        tb(ss, "Posição atual", Inches(0.25), Y1, Inches(2), Inches(0.3), sz=9, col=GRAY)
        bar(ss, Inches(0.25), Y1 + Inches(0.32), Inches(9.5), Inches(0.4), real, cor)
        tb(ss, f"{real:.1f}%", Inches(9.85), Y1 + Inches(0.32), Inches(1.2), Inches(0.4),
           sz=11, bold=True, col=cor)

        tb(ss, "Meta do perfil", Inches(0.25), Y1 + Inches(0.88), Inches(2), Inches(0.3), sz=9, col=GRAY)
        bar(ss, Inches(0.25), Y1 + Inches(1.2), Inches(9.5), Inches(0.4), alvo, GRAY2)
        if alvo > 0:
            mx = int(Inches(0.25) + Inches(9.5) * alvo / 100)
            rc(ss, mx - Inches(0.02), Y1 + Inches(1.2), Inches(0.04), Inches(0.4), GRAY)
        tb(ss, f"{alvo:.1f}%", Inches(9.85), Y1 + Inches(1.2), Inches(1.2), Inches(0.4),
           sz=11, bold=True, col=GRAY)

        # Rentabilidade se disponível
        pr = rent.get("portfolio", {})
        cr = rent.get("cdi", {})
        rent_items = [(per, lbl) for per, lbl in [("1m","1 Mês"),("3m","3 Meses"),("6m","6 Meses"),("12m","12 Meses")]
                      if pr.get(per) is not None]
        if rent_items:
            Y2 = Inches(3.9)
            rc(ss, Inches(0.25), Y2, W - Inches(0.5), Inches(0.04), GRAY2)
            tb(ss, "Rentabilidade da Carteira (Geral)", Inches(0.25), Y2 + Inches(0.1),
               Inches(6), Inches(0.3), sz=10, bold=True, col=GOLD)
            for i, (per, lbl) in enumerate(rent_items[:4]):
                v = pr[per]; pct_cdi = (v / cr[per] * 100) if cr.get(per) else 0
                xx = Inches(0.25) + i * Inches(3.25)
                metric(ss, xx, Y2 + Inches(0.5), Inches(3.1), Inches(1.0), lbl,
                       f"{v:.2f}%", GREEN if v > 0 else RED)
                tb(ss, f"{pct_cdi:.0f}% do CDI", xx + Inches(0.12), Y2 + Inches(1.3),
                   Inches(3.0), Inches(0.25), sz=8, col=GRAY)

        Y3 = Inches(5.55)
        rc(ss, Inches(0.25), Y3, W - Inches(0.5), Inches(0.04), GRAY2)
        status_txt = ("✓ Alocação alinhada ao modelo" if abs(desvio) <= 1.5
                      else f"{'▼ Abaixo' if desvio < 0 else '▲ Acima'} da meta em {abs(desvio):.1f}%")
        tb(ss, status_txt, Inches(0.25), Y3 + Inches(0.1), Inches(8), Inches(0.35),
           sz=11, bold=True, col=dc)
        tb(ss, descricao, Inches(0.25), Y3 + Inches(0.55), W - Inches(0.5), Inches(0.65),
           sz=10, col=GRAY, italic=True)

    # ── SLIDE 9 — AJUSTES E SUGESTÕES ────────────────────────────
    s9 = nsl()
    hdr(s9, "Ajustes & Sugestões Recomendadas", 9)
    Y_l, Y_r = Inches(0.95), Inches(0.95)

    # RF (coluna esquerda)
    rf_list = suge.get("renda_fixa", [])
    if rf_list:
        tb(s9, "RENDA FIXA", Inches(0.25), Y_l, Inches(6.2), Inches(0.28),
           sz=10, bold=True, col=GOLD)
        Y_l += Inches(0.35)
        for item in rf_list[:3]:
            urg_col = RED if item.get("urgencia") == "alta" else YEL if item.get("urgencia") == "media" else GREEN
            rc(s9, Inches(0.25), Y_l, Inches(6.2), Inches(0.75), CARD)
            rc(s9, Inches(0.25), Y_l, Inches(0.06), Inches(0.75), urg_col)
            acao = "↔ Substituir" if item.get("acao") == "substituir" else "+ Alocar"
            tb(s9, f"{acao}: {item.get('para', '')[:45]}", Inches(0.38), Y_l + Inches(0.05),
               Inches(5.95), Inches(0.3), sz=10, bold=True, col=WHITE)
            detail = f"Sair de: {item.get('de', '')}  " if item.get("de") else ""
            detail += f"Indexador: {item.get('indexador', '')}  " if item.get("indexador") else ""
            if item.get("motivo"):
                detail += item["motivo"][:50]
            tb(s9, detail, Inches(0.38), Y_l + Inches(0.38), Inches(5.95), Inches(0.3), sz=8, col=GRAY)
            Y_l += Inches(0.83)

    # RV (coluna direita)
    rv_list = suge.get("renda_variavel", [])
    if rv_list:
        tb(s9, "RENDA VARIÁVEL", Inches(6.8), Y_r, Inches(6.2), Inches(0.28),
           sz=10, bold=True, col=GOLD)
        Y_r += Inches(0.35)
        for item in rv_list[:3]:
            a = item.get("acao", "trocar")
            ac = GREEN if a == "comprar" else RED if a == "vender" else YEL
            rc(s9, Inches(6.8), Y_r, Inches(6.2), Inches(0.75), CARD)
            rc(s9, Inches(6.8), Y_r, Inches(0.06), Inches(0.75), ac)
            txt_rv = f"{a.upper()}: {item.get('de', '')}"
            if item.get("para"):
                txt_rv += f" → {item.get('para', '')}"
            tb(s9, txt_rv, Inches(6.93), Y_r + Inches(0.05), Inches(5.95), Inches(0.3),
               sz=10, bold=True, col=ac)
            detail2 = f"{item.get('motivo','')[:55]}  |  {item.get('fonte','')}"
            tb(s9, detail2, Inches(6.93), Y_r + Inches(0.38), Inches(5.95), Inches(0.3),
               sz=8, col=GRAY)
            Y_r += Inches(0.83)

    # Estruturadas
    estrut = suge.get("estruturadas", "")
    if estrut:
        Y_e = max(Y_l, Y_r) + Inches(0.1)
        rc(s9, Inches(0.25), Y_e, W - Inches(0.5), Inches(0.04), GRAY2)
        tb(s9, "OPERAÇÕES ESTRUTURADAS", Inches(0.25), Y_e + Inches(0.1),
           Inches(6), Inches(0.28), sz=10, bold=True, col=GOLD)
        linhas = [l for l in estrut.split("\n") if l.strip()][:2]
        for i, linha in enumerate(linhas):
            rc(s9, Inches(0.25), Y_e + Inches(0.45) + i * Inches(0.55), W - Inches(0.5), Inches(0.48), CARD)
            tb(s9, linha[:110], Inches(0.4), Y_e + Inches(0.5) + i * Inches(0.55),
               W - Inches(0.7), Inches(0.35), sz=10, col=WHITE)

    # Desvios resumo
    sig = [d for d in desvios if abs(d.get("desvio", 0)) > 1][:6]
    if sig:
        Y_bot = Inches(6.3)
        rc(s9, Inches(0.25), Y_bot, W - Inches(0.5), Inches(0.04), GRAY2)
        tb(s9, "DESVIOS DA CARTEIRA MODELO", Inches(0.25), Y_bot + Inches(0.08),
           Inches(6), Inches(0.25), sz=9, bold=True, col=GRAY)
        for i, d in enumerate(sig):
            xx = Inches(0.25) + i * Inches(2.15)
            tb(s9, f"{d.get('label', d.get('key', '')[:14])}: {d.get('desvio', 0):+.1f}%",
               xx, Y_bot + Inches(0.38), Inches(2.05), Inches(0.3),
               sz=9, bold=True, col=dv_col(d["desvio"]))

    # ── SLIDE 10 — OPORTUNIDADES BRAÚNA ──────────────────────────
    s10 = nsl()
    hdr(s10, "Oportunidades Braúna", 10)

    W_col = Inches(6.25)
    blocks = []

    # 1. Dividendos → Consórcio
    tem_rv_pct = comp.get("acoes", 0) + comp.get("fiis", 0)
    if tem_rv_pct > 3:
        blocks.append(("DIVERSIFICAÇÃO COM CONSÓRCIO",
            f"Sua carteira possui {tem_rv_pct:.1f}% em ativos geradores de dividendos (ações/FIIs). "
            "O consórcio Braúna é uma ferramenta complementar para aquisição de imóveis ou veículos "
            "sem juros — os dividendos recebidos podem cobrir as parcelas, diversificando o patrimônio.", GOLD))

    # 2. Prunus (≥ R$ 300k)
    if pat >= 300_000:
        prunus = PRUNUS_MODELOS.get(perfil, PRUNUS_MODELOS["moderada"])
        comp_str = ", ".join(f"{k} {v}%" for k, v in prunus["composicao"].items())
        blocks.append(("GESTÃO DISCRICIONÁRIA PRUNUS ASSET",
            f"Com {fmt_brl(pat)} sob gestão, você tem acesso à Prunus Asset. "
            f"Retorno alvo: {prunus['retorno_alvo']}. "
            f"Gestão profissional com zero conflito de interesse. "
            f"Carteira modelo ({PLABEL.get(perfil,'')}): {comp_str}.", GREEN))

    # 3. Planejamento patrimonial
    tem_pp = any(
        (c.get("id") == "planejamento_patrimonial" and c.get("status") == "ok")
        for c in (checklist if isinstance(checklist, list) else [])
    )
    if not tem_pp:
        blocks.append(("PLANEJAMENTO PATRIMONIAL",
            "Identificamos que você ainda não possui um planejamento patrimonial estruturado. "
            "Estratégias como holding familiar, doação com usufruto e seguro de vida permitem "
            "reduzir ITCMD, proteger o patrimônio e garantir transmissão ordenada para herdeiros. "
            "Agende uma análise gratuita com nosso especialista.", BLUE))

    # 4. Internacional < 14%
    intl_pct = comp.get("internacional", 0)
    if intl_pct < 14:
        falta = 14 - intl_pct
        val_falta = pat * falta / 100 if pat > 0 else 0
        blocks.append(("AUMENTO DA EXPOSIÇÃO INTERNACIONAL",
            f"Sua carteira possui {intl_pct:.1f}% em internacionalização (meta recomendada: 14%). "
            f"Sugerimos um plano de 12 meses para aportar {falta:.1f}% adicionais "
            f"(≈ {fmt_brl(val_falta)}), migrando gradualmente para ETFs e BDRs de qualidade "
            "para diversificação cambial e proteção patrimonial.", PINK))

    # 5. Patrimônio caindo → Planejamento financeiro
    if pat_caindo:
        blocks.append(("PLANEJAMENTO FINANCEIRO",
            "Identificamos redução patrimonial nas últimas análises. "
            "Um plano financeiro estruturado pode ajudar a controlar despesas, "
            "otimizar a taxa de poupança e definir metas claras de crescimento. "
            "Agende uma reunião de planejamento financeiro com seu assessor para revisão completa.", YEL))

    if not blocks:
        blocks.append(("CARTEIRA EM ORDEM",
            "Parabéns! Sua carteira está bem estruturada e alinhada com seu perfil de investimento. "
            "Continue acompanhando com seu assessor para manter o equilíbrio e aproveitar "
            "novas oportunidades de mercado conforme surgem.", GREEN))

    x_cols = [Inches(0.25), Inches(6.83)]
    y_cols = [Inches(0.95), Inches(0.95)]
    for i, (lbl, txt_b, cor) in enumerate(blocks[:6]):
        ci = i % 2
        xx, yy = x_cols[ci], y_cols[ci]
        h_blk = Inches(1.4)
        rc(s10, xx, yy, W_col, h_blk, CARD)
        rc(s10, xx, yy, Inches(0.07), h_blk, cor)
        tb(s10, lbl, xx + Inches(0.16), yy + Inches(0.08),
           W_col - Inches(0.26), Inches(0.28), sz=9, bold=True, col=cor)
        tb(s10, txt_b, xx + Inches(0.16), yy + Inches(0.42),
           W_col - Inches(0.26), Inches(0.9), sz=9, col=GRAY)
        y_cols[ci] += h_blk + Inches(0.12)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


@app.route("/api/ppt", methods=["POST"])
def ppt_endpoint():
    dados = request.get_json()
    if not dados:
        return jsonify({"error": "Dados necessários"}), 400
    try:
        dados["hist_macro"] = buscar_hist_macro()
        buf = gerar_ppt(dados)
        nome_arq = f"Brauna_{dados.get('nome','cliente').replace(' ','_')}_{datetime.now().strftime('%Y%m%d')}.pptx"
        return send_file(
            buf,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=nome_arq,
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/pdf", methods=["POST"])
def pdf_endpoint():
    d = request.get_json()
    buf = gerar_pdf(
        d.get("nome",""), d.get("perfil","conservadora"),
        d.get("desvios",[]), d.get("rent",{}),
        d.get("patrimonio",0), d.get("caixa",0), d.get("data_ref",""),
        d.get("recomendacoes",[]), d.get("alertas",[]),
        d.get("macro",{}), d.get("ref_contexto",""),
        d.get("checklist_servir"),
    )
    nome_arq = f"{d.get('nome','cliente').replace(' ','-').lower()}_{d.get('data_ref','')}_analise.pdf"
    return send_file(buf, mimetype="application/pdf", as_attachment=True, download_name=nome_arq)


HTML_LOGIN = r"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Braúna Investimentos — Acesso</title>
<style>
*{margin:0;padding:0;box-sizing:border-box}
html,body{height:100%;background:#081F18;font-family:'Segoe UI',system-ui,sans-serif;color:#F0F0F0}
body{display:flex;flex-direction:column;align-items:center;justify-content:center;min-height:100vh;padding:20px}

/* Logo */
.logo-area{text-align:center;margin-bottom:40px}
.logo-title{font-size:28px;font-weight:800;color:#C9A96E;letter-spacing:2px;text-transform:uppercase}
.logo-sub{font-size:13px;color:#2A5A3A;margin-top:6px;letter-spacing:1px}
.logo-line{width:60px;height:2px;background:linear-gradient(to right,transparent,#C9A96E,transparent);margin:14px auto 0}

/* Cards de papel */
.roles{display:flex;gap:16px;margin-bottom:36px;flex-wrap:wrap;justify-content:center}
.role-card{
  width:200px;padding:28px 20px;border-radius:16px;
  border:1.5px solid #1A4030;background:#111;
  cursor:pointer;text-align:center;transition:all .25s;
  position:relative;overflow:hidden;
}
.role-card::before{
  content:'';position:absolute;inset:0;opacity:0;transition:opacity .25s;
}
.role-card:hover{border-color:#C9A96E;transform:translateY(-3px);box-shadow:0 8px 32px rgba(214,178,122,.12)}
.role-card.selected{transform:translateY(-3px)}

/* Assessor */
.role-card.assessor.selected{border-color:#C9A96E;box-shadow:0 8px 32px rgba(214,178,122,.18);background:#1A2E1A}
.role-card.assessor:hover{border-color:#C9A96E}
/* Líder */
.role-card.lider.selected{border-color:#8B9FE8;box-shadow:0 8px 32px rgba(139,159,232,.18);background:#0D0D1F}
.role-card.lider:hover{border-color:#8B9FE8}
/* Head */
.role-card.head.selected{border-color:#D4B483;box-shadow:0 8px 32px rgba(232,201,107,.18);background:#1A1500}
.role-card.head:hover{border-color:#D4B483}
.head .role-check{background:#D4B483;color:#000}
/* Admin */
.role-card.admin.selected{border-color:#5DCAA5;box-shadow:0 8px 32px rgba(93,202,165,.18);background:#081410}
.role-card.admin:hover{border-color:#5DCAA5}

.role-icon{font-size:36px;margin-bottom:12px;display:block}
.role-name{font-size:15px;font-weight:700;margin-bottom:4px}
.role-desc{font-size:11px;color:#3A6A48;line-height:1.4}
.role-check{
  position:absolute;top:10px;right:10px;
  width:18px;height:18px;border-radius:50%;
  display:flex;align-items:center;justify-content:center;
  font-size:10px;font-weight:900;
  opacity:0;transition:opacity .2s;
}
.assessor .role-check{background:#C9A96E;color:#000}
.lider .role-check{background:#8B9FE8;color:#000}
.admin .role-check{background:#5DCAA5;color:#000}
.role-card.selected .role-check{opacity:1}

/* Caixa de senha */
.senha-box{
  width:100%;max-width:380px;
  background:#111;border-radius:14px;padding:28px;
  border:1px solid #1A4030;
  animation:fadeUp .25s ease;
}
@keyframes fadeUp{from{opacity:0;transform:translateY(10px)}to{opacity:1;transform:translateY(0)}}
.senha-label{font-size:11px;text-transform:uppercase;letter-spacing:.8px;margin-bottom:10px;display:flex;align-items:center;gap:8px}
.senha-label span{font-size:16px}
.senha-input-wrap{position:relative;margin-bottom:14px}
.senha-input{
  width:100%;background:#0A0A0A;border:1.5px solid #1C4A34;
  border-radius:10px;padding:13px 44px 13px 16px;
  color:#F0F0F0;font-size:16px;letter-spacing:4px;outline:none;
  transition:border .2s;
}
.senha-input::placeholder{letter-spacing:1px;font-size:13px;color:#1E4A30}
.senha-input:focus{border-color:var(--role-color)}
.toggle-pw{
  position:absolute;right:14px;top:50%;transform:translateY(-50%);
  background:none;border:none;color:#2A5A3A;cursor:pointer;font-size:16px;padding:0;
}
.btn-entrar{
  width:100%;padding:14px;border:none;border-radius:10px;
  font-size:14px;font-weight:800;cursor:pointer;
  letter-spacing:.5px;text-transform:uppercase;
  transition:all .2s;background:var(--role-color);color:#081F18;
}
.btn-entrar:hover{opacity:.88;transform:translateY(-1px)}
.btn-entrar:disabled{opacity:.3;cursor:not-allowed;transform:none}
.erro-msg{
  font-size:12px;color:#FF6B6B;text-align:center;margin-top:10px;
  height:18px;transition:opacity .2s;
}
.rodape{font-size:11px;color:#1C4A34;margin-top:32px;text-align:center}
</style>
</head>
<body>

<div class="logo-area">
  <div class="logo-title">Braúna Investimentos</div>
  <div class="logo-line"></div>
  <div class="logo-sub">Sistema de Análise de Carteiras</div>
</div>

<!-- Escolha de papel -->
<div class="roles">
  <div class="role-card assessor" onclick="selRole('assessor')" id="card-assessor">
    <div class="role-check">✓</div>
    <span class="role-icon">📊</span>
    <div class="role-name" style="color:#C9A96E">Assessor</div>
    <div class="role-desc">Análise de carteiras e atendimento ao cliente</div>
  </div>
  <div class="role-card lider" onclick="selRole('lider')" id="card-lider">
    <div class="role-check">✓</div>
    <span class="role-icon">👥</span>
    <div class="role-name" style="color:#8B9FE8">Líder</div>
    <div class="role-desc">Visão da equipe, OKRs e orientação aos assessores</div>
  </div>
  <div class="role-card head" onclick="selRole('head')" id="card-head">
    <div class="role-check">✓</div>
    <span class="role-icon">🏛️</span>
    <div class="role-name" style="color:#D4B483">Head de Produtos</div>
    <div class="role-desc">Portfólios modelo, cenário macro e produtos em destaque</div>
  </div>
  <div class="role-card admin" onclick="selRole('admin')" id="card-admin">
    <div class="role-check">✓</div>
    <span class="role-icon">⚙️</span>
    <div class="role-name" style="color:#5DCAA5">Administração</div>
    <div class="role-desc">Documentos estratégicos, cartas e configurações</div>
  </div>
</div>

<!-- Caixa de senha -->
<div class="senha-box" id="senha-box" style="display:none">
  <div class="senha-label" id="senha-label">
    <span id="senha-icon"></span>
    <span id="senha-role-name" style="color:var(--role-color)"></span>
  </div>
  <div class="senha-input-wrap">
    <input type="password" id="senha" class="senha-input" placeholder="Digite sua senha" onkeydown="if(event.key==='Enter')entrar()">
    <button class="toggle-pw" onclick="toggleSenha()" type="button">👁</button>
  </div>
  <button class="btn-entrar" id="btn-entrar" onclick="entrar()">Entrar</button>
  <div class="erro-msg" id="erro"></div>
</div>

<div class="rodape">Braúna Investimentos · Uso interno · Acesso restrito</div>

<script>
let roleAtual = null;
const ROLES = {
  assessor: {nome:"Assessor",        icon:"📊", color:"#C9A96E", dest:"/assessor"},
  lider:    {nome:"Líder",           icon:"👥", color:"#8B9FE8", dest:"/lider"},
  head:     {nome:"Head de Produtos",icon:"🏛️", color:"#D4B483", dest:"/head-produtos"},
  admin:    {nome:"Administração",   icon:"⚙️", color:"#5DCAA5", dest:"/admin"},
};

function selRole(role){
  roleAtual = role;
  document.querySelectorAll(".role-card").forEach(c=>c.classList.remove("selected"));
  document.getElementById("card-"+role).classList.add("selected");

  const r = ROLES[role];
  document.documentElement.style.setProperty("--role-color", r.color);
  document.getElementById("senha-icon").textContent = r.icon;
  document.getElementById("senha-role-name").textContent = r.nome;
  document.getElementById("senha").value = "";
  document.getElementById("erro").textContent = "";

  const inp = document.getElementById("senha");
  inp.placeholder = role === "assessor" ? "Código do assessor (ex: A74621)" : "Digite sua senha";
  inp.style.letterSpacing = role === "assessor" ? "2px" : "4px";

  document.getElementById("senha-box").style.display = "block";
  setTimeout(()=>inp.focus(), 50);
}

function toggleSenha(){
  const inp = document.getElementById("senha");
  inp.type = inp.type === "password" ? "text" : "password";
}

async function entrar(){
  if(!roleAtual) return;
  const senha = document.getElementById("senha").value;
  if(!senha){ document.getElementById("erro").textContent="Digite a senha."; return; }

  const btn = document.getElementById("btn-entrar");
  btn.disabled = true; btn.textContent = "Verificando...";
  document.getElementById("erro").textContent = "";

  try{
    const r = await fetch("/api/login",{
      method:"POST",
      headers:{"Content-Type":"application/json"},
      body: JSON.stringify({role: roleAtual, senha})
    });
    const d = await r.json();
    if(d.ok){
      localStorage.setItem("brauna_role", roleAtual);
      localStorage.setItem("brauna_role_ts", Date.now());
      if(d.nome)  localStorage.setItem("brauna_nome",   d.nome);
      if(d.codigo) localStorage.setItem("brauna_codigo", d.codigo);
      window.location.href = ROLES[roleAtual].dest;
    } else {
      document.getElementById("erro").textContent = d.msg || "Código inválido. Tente novamente.";
      document.getElementById("senha").value="";
      document.getElementById("senha").focus();
    }
  } catch(e){
    document.getElementById("erro").textContent = "Erro de conexão.";
  } finally {
    btn.disabled = false;
    btn.textContent = "Entrar";
  }
}

// Seleciona automaticamente se já existe role salvo
const saved = localStorage.getItem("brauna_role");
if(saved && ROLES[saved]) selRole(saved);
</script>
</body></html>"""

HTML_LIDER = r"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Braúna — Painel do Líder</title>
<style>
*{margin:0;padding:0;box-sizing:border-box}
body{background:#0A0A14;color:#F0F0F0;font-family:'Segoe UI',system-ui,sans-serif;min-height:100vh}
header{background:#0D0D1F;border-bottom:1px solid #1E1E3A;padding:14px 28px;display:flex;align-items:center;justify-content:space-between}
header h1{font-size:17px;color:#8B9FE8;font-weight:700}
header p{font-size:11px;color:#2A5A3A;margin-top:2px}
.nav{display:flex;gap:8px;align-items:center}
.nav a{font-size:12px;padding:5px 12px;border-radius:6px;border:1px solid #2A2A4A;color:#4A7055;text-decoration:none;transition:all .2s}
.nav a:hover{border-color:#8B9FE8;color:#8B9FE8}
.nav a.active{background:#8B9FE8;color:#000;border-color:#8B9FE8;font-weight:700}
.container{max-width:1200px;margin:0 auto;padding:24px 20px}
.card{background:#0F0F22;border:1px solid #1E1E3A;border-radius:12px;padding:20px;margin-bottom:16px}
.card h2{font-size:11px;color:#8B9FE8;font-weight:700;margin-bottom:16px;text-transform:uppercase;letter-spacing:.8px}
/* KPIs 6 colunas */
.kpi-row{display:grid;grid-template-columns:repeat(6,1fr);gap:8px;margin-bottom:16px}
.kpi{background:#0A0A1A;border-radius:10px;padding:12px 14px;text-align:center;border:1px solid #1E1E3A;transition:border-color .2s}
.kpi:hover{border-color:#8B9FE8}
.kpi .n{font-size:22px;font-weight:700;color:#8B9FE8}
.kpi .n.red{color:#FF6B6B}
.kpi .n.green{color:#5DCAA5}
.kpi .n.gold{color:#C9A96E}
.kpi .l{font-size:10px;color:#3A6A48;margin-top:3px}
/* Filtros */
.filtros-bar{display:flex;gap:10px;align-items:center;flex-wrap:wrap;margin-bottom:14px;padding:12px 16px;background:#0A0A1A;border:1px solid #1E1E3A;border-radius:10px}
.filtros-bar label{font-size:11px;color:#3A6A48;margin-right:4px}
.filtros-bar select{background:#0D0D22;border:1px solid #2A2A4A;color:#C0C0D0;font-size:11px;padding:4px 8px;border-radius:6px;outline:none;cursor:pointer}
.filtros-bar select:focus{border-color:#8B9FE8}
.btn-limpar{font-size:11px;padding:4px 12px;border-radius:6px;border:1px solid #2A2A4A;color:#4A7055;background:none;cursor:pointer;transition:all .2s;margin-left:auto}
.btn-limpar:hover{border-color:#FF6B6B;color:#FF6B6B}
/* Card risco */
.card-risco{background:#140A0A;border:1px solid #3A1A1A;border-radius:12px;padding:18px 20px;margin-bottom:16px}
.card-risco h2{font-size:11px;color:#FF6B6B;font-weight:700;margin-bottom:14px;text-transform:uppercase;letter-spacing:.8px}
.risco-item{display:flex;align-items:center;gap:12px;padding:10px 12px;border-radius:8px;background:#1A0A0A;border:1px solid #2A1010;margin-bottom:8px;flex-wrap:wrap}
.risco-nome{font-size:13px;font-weight:700;color:#F0F0F0;flex:1;min-width:120px}
.risco-assessor{font-size:11px;color:#3A6A48}
.risco-pat{font-size:11px;color:#8B9FE8}
.risco-motivo{font-size:11px;color:#FFD966;padding:2px 8px;background:#1A1500;border-radius:6px;border:1px solid #2A2000}
.btn-ver{font-size:11px;padding:4px 10px;border-radius:6px;border:1px solid #3A1010;color:#FF6B6B;background:none;cursor:pointer;transition:all .2s;margin-left:auto}
.btn-ver:hover{background:#2A1010;border-color:#FF6B6B}
.risco-ok{padding:12px 16px;color:#5DCAA5;font-size:13px;background:#0A2018;border-radius:8px;border:1px solid #0A3020}
/* Card alertas */
.card-alertas{background:#0D0A14;border:1px solid #2A1E3A;border-radius:12px;padding:18px 20px;margin-bottom:16px}
.card-alertas h2{font-size:11px;color:#C9A96E;font-weight:700;margin-bottom:14px;text-transform:uppercase;letter-spacing:.8px}
.alerta-acordeao{border:1px solid #2A2A4A;border-radius:8px;margin-bottom:6px;overflow:hidden}
.alerta-header{display:flex;align-items:center;gap:10px;padding:10px 14px;cursor:pointer;background:#0A0A1A;transition:background .2s}
.alerta-header:hover{background:#0D0D22}
.alerta-header-nome{font-size:13px;font-weight:700;color:#F0F0F0;flex:1}
.alerta-count{font-size:11px;color:#C9A96E;padding:2px 8px;background:#1A1200;border-radius:8px}
.alerta-body{display:none;padding:10px 14px;background:#080814}
.alerta-body.aberto{display:block}
.alerta-linha{display:flex;align-items:center;gap:10px;padding:7px 0;border-bottom:1px solid #1A1A2E;flex-wrap:wrap}
.alerta-linha:last-child{border-bottom:none}
.alerta-ticker{font-size:12px;font-weight:700;color:#8B9FE8;min-width:60px}
.alerta-titulo{font-size:12px;color:#D0D0E0;flex:1}
.alerta-tipo{font-size:10px;font-weight:700;padding:2px 8px;border-radius:8px}
.alerta-compra{background:#0A2018;color:#5DCAA5;border:1px solid #0A3020}
.alerta-venda{background:#2A1010;color:#FF6B6B;border:1px solid #3A1010}
.alerta-atencao{background:#1A1500;color:#FFD966;border:1px solid #2A2000}
/* Ranking assessores */
.rank-item{
  display:flex;align-items:center;gap:14px;
  padding:14px 16px;border-radius:10px;
  background:#0A0A1A;border:1px solid #1E1E3A;
  margin-bottom:8px;cursor:pointer;transition:all .2s;
}
.rank-item:hover,.rank-item.aberto{border-color:#8B9FE8;background:#0D0D22}
.rank-pos{font-size:20px;font-weight:900;color:#2A2A4A;min-width:28px;text-align:center}
.rank-pos.top1{color:#C9A96E}
.rank-pos.top2{color:#8B9FE8}
.rank-pos.top3{color:#5DCAA5}
.rank-nome{font-size:14px;font-weight:700;color:#F0F0F0;flex:1}
.rank-meta{font-size:11px;color:#3A6A48}
.rank-bars{flex:2;display:flex;flex-direction:column;gap:5px}
.rank-bar-row{display:flex;align-items:center;gap:8px;font-size:11px}
.rank-bar-row .lbl{width:90px;color:#3A6A48;text-align:right;flex-shrink:0}
.rank-bar-bg{flex:1;height:5px;background:#1A1A2E;border-radius:3px;overflow:hidden}
.rank-bar-fill{height:100%;border-radius:3px;transition:width .6s}
.rank-badges{display:flex;gap:6px;flex-direction:column;align-items:flex-end}
.rbadge{padding:3px 10px;border-radius:10px;font-size:11px;font-weight:700}
/* Painel detalhe assessor */
.detalhe-assessor{display:none;margin-bottom:8px;animation:fadeIn .2s}
.detalhe-assessor.aberto{display:block}
@keyframes fadeIn{from{opacity:0;transform:translateY(-6px)}to{opacity:1;transform:translateY(0)}}
/* Métricas assessor */
.metricas-assessor{display:grid;grid-template-columns:repeat(4,1fr);gap:8px;padding:10px;margin-bottom:8px;background:#080814;border-radius:8px;border:1px solid #1A1A2E}
.met-item{text-align:center;padding:8px}
.met-n{font-size:18px;font-weight:700;color:#8B9FE8}
.met-l{font-size:10px;color:#2A5A3A;margin-top:2px}
.meta-bar{height:4px;border-radius:2px;margin-top:6px}
.meta-ok{background:#5DCAA5}
.meta-nok{background:#FF6B6B}
/* Cards de cliente */
.cliente-card{
  background:#080816;border:1px solid #1A1A30;border-radius:10px;
  padding:14px 16px;margin-bottom:10px;
}
.cliente-header{display:flex;align-items:center;gap:12px;margin-bottom:10px;flex-wrap:wrap}
.cliente-nome{font-size:14px;font-weight:700;color:#F0F0F0}
.cliente-perfil{font-size:11px;padding:2px 10px;border-radius:10px;background:#1A1A3A;color:#8B9FE8;font-weight:700}
.cliente-objetivo{font-size:11px;color:#3A6A48;flex:1}
.c-badge{display:inline-block;padding:2px 8px;border-radius:8px;font-size:10px;font-weight:700}
.c-ok{background:#0A2018;color:#5DCAA5}
.c-atencao{background:#1A1500;color:#FFD966}
.c-critico{background:#2A1010;color:#FF6B6B}
/* Timeline de análises */
.timeline{position:relative;padding-left:20px}
.timeline::before{content:'';position:absolute;left:6px;top:0;bottom:0;width:2px;background:#1A1A30}
.tl-entry{position:relative;margin-bottom:12px}
.tl-dot{
  position:absolute;left:-17px;top:4px;
  width:10px;height:10px;border-radius:50%;border:2px solid #0A0A14;
}
.tl-dot.ok{background:#5DCAA5}
.tl-dot.atencao{background:#FFD966}
.tl-dot.critico{background:#FF6B6B}
.tl-dot.first{width:13px;height:13px;left:-18px;top:3px}
.tl-content{background:#0D0D20;border:1px solid #1E1E3A;border-radius:8px;padding:10px 12px}
.tl-data{font-size:10px;color:#2A5A3A;margin-bottom:6px;display:flex;align-items:center;gap:8px}
.tl-data b{color:#8B9FE8;font-size:11px}
.tl-metrics{display:flex;gap:16px;flex-wrap:wrap;margin-bottom:6px}
.tl-m{font-size:11px;color:#4A7055}
.tl-m span{color:#F0F0F0;font-weight:700}
/* Delta de carteira */
.delta-row{display:flex;gap:8px;flex-wrap:wrap;margin-top:6px}
.delta-item{
  display:flex;align-items:center;gap:4px;
  padding:3px 8px;border-radius:6px;font-size:11px;
  background:#0A0A1A;
}
.delta-up-ok{color:#5DCAA5;border:1px solid #0A2018}
.delta-up-bad{color:#FFD966;border:1px solid #1A1500}
.delta-dn-ok{color:#5DCAA5;border:1px solid #0A2018}
.delta-dn-bad{color:#FF6B6B;border:1px solid #2A1010}
.delta-neutral{color:#3A6A48;border:1px solid #1A1A2E}
/* Nota do líder */
.nota-input{
  width:100%;background:#0A0A1A;border:1px solid #1E1E3A;
  border-radius:6px;padding:6px 10px;color:#AAA;font-size:11px;
  margin-top:8px;outline:none;
}
.nota-input:focus{border-color:#8B9FE8}
/* Vazio */
.vazio{text-align:center;color:#1E4A30;padding:40px;font-size:13px}
.progress-anel{text-align:right}
@media(max-width:900px){.kpi-row{grid-template-columns:repeat(3,1fr)}}
@media(max-width:600px){.kpi-row{grid-template-columns:repeat(2,1fr)}.rank-bars{display:none}.metricas-assessor{grid-template-columns:repeat(2,1fr)}}
</style>
</head>
<body>
<header>
  <div><h1>👥 BRAÚNA — PAINEL DO LÍDER</h1><p>Ranking de assessores · Histórico de carteiras</p></div>
  <nav class="nav" id="nav-lider">
    <a href="/assessor">📊 Assessor</a>
    <a href="/lider" class="active">👥 Líder</a>
    <button onclick="sair()" style="font-size:12px;padding:5px 12px;border-radius:6px;border:1px solid #2A2A3A;color:#4A7055;background:none;cursor:pointer;transition:all .2s" onmouseover="this.style.borderColor='#FF6B6B';this.style.color='#FF6B6B'" onmouseout="this.style.borderColor='#2A2A3A';this.style.color='#4A7055'">Sair</button>
  </nav>
  <script>
  (function(){
    if(localStorage.getItem("brauna_role")==="admin"){
      const nav=document.getElementById("nav-lider");
      nav.innerHTML='<a href="/assessor">📊 Assessor</a><a href="/lider" class="active">👥 Líder</a><a href="/head-produtos">🏛️ Head</a><a href="/admin">⚙️ Admin</a>';
    }
  })();
  </script>
</header>
<div class="container">

<!-- KPIs 6 itens -->
<div class="kpi-row" id="kpi-row">
  <div class="kpi"><div class="n gold" id="kpi-patrimonio">—</div><div class="l">Total Patrimônio</div></div>
  <div class="kpi"><div class="n" id="kpi-assessores">—</div><div class="l">Assessores Ativos</div></div>
  <div class="kpi"><div class="n" id="kpi-clientes">—</div><div class="l">Clientes</div></div>
  <div class="kpi"><div class="n" id="kpi-relatorios">—</div><div class="l">Relatórios</div></div>
  <div class="kpi"><div class="n red" id="kpi-risco">—</div><div class="l">Clientes em Risco</div></div>
  <div class="kpi"><div class="n green" id="kpi-score">—</div><div class="l">Score Médio Servir</div></div>
</div>

<!-- Clientes em Risco -->
<div class="card-risco" id="sec-risco">
  <h2>🚨 Clientes em Risco</h2>
  <div id="risco-container"><p class="vazio">Carregando...</p></div>
</div>

<!-- Alertas de Mercado -->
<div class="card-alertas" id="sec-alertas" style="display:none">
  <h2>📡 Alertas de Mercado</h2>
  <div id="alertas-container"></div>
</div>

<!-- Filtros -->
<div class="filtros-bar">
  <label>Assessor</label>
  <select id="filtro-assessor" onchange="aplicarFiltros()">
    <option value="">Todos</option>
  </select>
  <label>Perfil</label>
  <select id="filtro-perfil" onchange="aplicarFiltros()">
    <option value="">Todos</option>
    <option value="conservadora">Conservadora</option>
    <option value="moderada">Moderada</option>
    <option value="arrojada">Arrojada</option>
    <option value="agressiva">Agressiva</option>
  </select>
  <label>Status</label>
  <select id="filtro-status" onchange="aplicarFiltros()">
    <option value="">Todos</option>
    <option value="saudavel">Saudável</option>
    <option value="atencao">Atenção</option>
    <option value="critico">Crítico</option>
  </select>
  <button class="btn-limpar" onclick="limparFiltros()">Limpar filtros</button>
</div>

<!-- Ranking assessores -->
<div class="card">
  <div style="display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:8px;margin-bottom:4px">
    <h2 style="margin:0">🏆 Ranking dos Assessores</h2>
    <label style="font-size:12px;color:#4A7055;cursor:pointer;display:flex;align-items:center;gap:6px">
      <input type="file" id="input-planilha" accept=".xlsx" style="display:none" onchange="importarPlanilha(this)">
      <span style="padding:5px 14px;border:1px solid #2A3A2A;border-radius:6px;background:#0A1A0A;color:#5DCAA5;font-size:12px" onclick="document.getElementById('input-planilha').click()">📥 Atualizar planilha</span>
    </label>
  </div>
  <div id="ranking-container">
    <p class="vazio">Carregando dados...</p>
  </div>
</div>

</div>
<script>
// Auth guard
(function(){
  const role = localStorage.getItem("brauna_role");
  if(role !== "lider" && role !== "admin") { localStorage.removeItem("brauna_role"); window.location.replace("/"); }
  const nome = localStorage.getItem("brauna_nome") || "Lider";
  fetch("/api/admin/activity",{method:"POST",headers:{"Content-Type":"application/json"},
    body:JSON.stringify({role:"lider",nome,acao:"acesso",detalhe:"Abriu a pagina do Lider"})});
})();
function sair(){ localStorage.removeItem("brauna_role"); window.location.replace("/"); }

function importarPlanilha(input){
  const file = input.files[0];
  if(!file) return;
  const formData = new FormData();
  formData.append("planilha", file);
  const btn = document.querySelector("[onclick=\"document.getElementById('input-planilha').click()\"]");
  if(btn) btn.textContent = "Importando...";
  fetch("/api/importar_assessores_xlsx", {method:"POST", body:formData})
    .then(r=>r.json())
    .then(function(res){
      if(res.ok){
        if(btn) btn.textContent = "✅ "+res.total+" importados";
        setTimeout(function(){ window.location.reload(); }, 1500);
      } else {
        if(btn) btn.textContent = "❌ Erro";
        alert(res.erro||"Erro ao importar planilha");
      }
    }).catch(function(){ if(btn) btn.textContent = "❌ Erro"; });
  input.value = "";
}

const LABEL_COMP = {
  pos_fixado:"Pos Fixado", inflacao:"Inflacao", pre_fixado:"Pre Fixado",
  acoes:"Acoes", fiis:"FIIs", multimercado:"Multimercado",
  internacional:"Internacional", alternativos:"Alternativos", criptomoedas:"Criptomoedas"
};

const MODELOS = {
  conservadora:{pos_fixado:70,inflacao:16,pre_fixado:7,acoes:0,fiis:0,multimercado:3,internacional:4,alternativos:0,criptomoedas:0},
  moderada:{pos_fixado:44,inflacao:23,pre_fixado:10,acoes:5,fiis:1.5,multimercado:6,internacional:9,alternativos:1,criptomoedas:0.5},
  arrojada:{pos_fixado:28,inflacao:28,pre_fixado:12,acoes:8,fiis:2.5,multimercado:9.5,internacional:10.25,alternativos:1,criptomoedas:0.75},
  agressiva:{pos_fixado:13,inflacao:31,pre_fixado:13,acoes:14,fiis:3.5,multimercado:10.5,internacional:13,alternativos:1,criptomoedas:1}
};

let historico = {};
let clientesData = [];
let notas = {};
let sugestoesHist = [];
let porAssessorGlobal = {};
let activityData = [];
let assessoresDados = {};  // mapa nome → dados financeiros da planilha

function fmtPat(v){
  if(!v || v<=0) return "—";
  if(v>=1e9) return "R$ " + (v/1e9).toFixed(1) + "B";
  if(v>=1e6) return "R$ " + (v/1e6).toFixed(1) + "M";
  return "R$ " + Number(v).toLocaleString("pt-BR",{maximumFractionDigits:0});
}

async function init(){
  const [hist, clientes, suge, assFinanc] = await Promise.all([
    fetch("/api/historico").then(r=>r.json()).catch(()=>({})),
    fetch("/api/clientes").then(r=>r.json()).catch(()=>[]),
    fetch("/api/sugestoes").then(r=>r.json()).catch(()=>({historico:[]})),
    fetch("/api/assessores_dados").then(r=>r.json()).catch(()=>({assessores:[]})),
  ]);
  historico = hist;
  clientesData = Array.isArray(clientes) ? clientes : [];
  sugestoesHist = suge.historico||[];
  clientesData.forEach(c=>{ notas[c.assessor+"|"+c.nome] = {id:c.id, nota:c.nota_lider||""}; });

  // Indexar dados financeiros por nome (busca aproximada)
  (assFinanc.assessores||[]).forEach(function(a){
    assessoresDados[a.nome] = a;
    // Indexar também pelo primeiro nome + sobrenome para match parcial
    const partes = a.nome.split(" ");
    if(partes.length >= 2) assessoresDados[partes[0]+" "+partes[partes.length-1]] = a;
  });

  renderKPIs(assFinanc);
  renderRisco();
  renderAlertas();
  popularFiltroAssessor();
  renderRanking();
}

function renderKPIs(assFinanc){
  const chaves = Object.keys(historico);
  const assessores = new Set(chaves.map(k=>historico[k].assessor||k.split("|")[0]));
  const totalRel = chaves.reduce((s,k)=>s+(historico[k].entradas&&historico[k].entradas.length||0),0);

  // Patrimônio: preferir soma da planilha quando disponível
  const assArr = (assFinanc && assFinanc.assessores) || [];
  const totalPatPlanilha = assArr.reduce((s,a)=>s+(a.patrimonio||0),0);
  const totalPatClientes = clientesData.reduce((s,c)=>s+(c.patrimonio||0),0);
  const totalPat = totalPatPlanilha > 0 ? totalPatPlanilha : totalPatClientes;

  let somaScore=0, cntScore=0;
  clientesData.forEach(c=>{
    const hist = historico[c.assessor+"|"+c.nome]||historico[Object.keys(historico).find(k=>k.includes(c.nome))||""]||null;
    const ult = hist && hist.entradas && hist.entradas[0];
    if(ult && ult.score_servir>0){ somaScore+=ult.score_servir; cntScore++; }
  });

  const emRisco = clientesData.filter(c=>{
    const hist = historico[c.assessor+"|"+c.nome]||null;
    const ult = hist && hist.entradas && hist.entradas[0];
    return ult && (ult.rent12_pct_cdi < 70 || ult.status === "critico");
  }).length;

  document.getElementById("kpi-patrimonio").textContent = fmtPat(totalPat);
  document.getElementById("kpi-assessores").textContent = assessores.size;
  document.getElementById("kpi-clientes").textContent   = clientesData.length;
  document.getElementById("kpi-relatorios").textContent = totalRel;
  document.getElementById("kpi-risco").textContent      = emRisco;
  document.getElementById("kpi-score").textContent      = cntScore>0 ? (somaScore/cntScore).toFixed(1) : "—";
}

function clientesEmRisco(filtroAssessor, filtroPerfil, filtroStatus){
  return clientesData.filter(c=>{
    if(filtroAssessor && c.assessor !== filtroAssessor) return false;
    if(filtroPerfil && c.perfil !== filtroPerfil) return false;
    const hist = historico[c.assessor+"|"+c.nome]||null;
    const ult = hist && hist.entradas && hist.entradas[0];
    if(!ult) return false;
    const isRisco = ult.rent12_pct_cdi < 70 || ult.status === "critico";
    if(filtroStatus==="critico" && ult.status!=="critico") return false;
    if(filtroStatus==="saudavel" || filtroStatus==="atencao") return false;
    return isRisco;
  });
}

function renderRisco(filtroAssessor, filtroPerfil, filtroStatus){
  const riscos = clientesEmRisco(filtroAssessor||"", filtroPerfil||"", filtroStatus||"");
  const cont = document.getElementById("risco-container");
  if(!riscos.length){
    cont.innerHTML = '<div class="risco-ok">Nenhum alerta critico com os filtros atuais</div>';
    return;
  }
  cont.innerHTML = riscos.map(c=>{
    const hist = historico[c.assessor+"|"+c.nome]||null;
    const ult = hist && hist.entradas && hist.entradas[0];
    const motivos = [];
    if(ult && ult.status==="critico") motivos.push("Status critico");
    if(ult && ult.rent12_pct_cdi < 70) motivos.push("Rent. 12M abaixo de 70% CDI (" + (ult.rent12_pct_cdi||0) + "%)");
    const motivoStr = motivos.join(" · ")||"Desvio detectado";
    const chaveHist = Object.keys(historico).find(k=>k.startsWith(c.assessor+"|"+c.nome)||k.endsWith("|"+c.nome));
    const assessorIdx = rankingOrder ? rankingOrder.findIndex(a=>a.nome===c.assessor) : -1;
    return '<div class="risco-item">'
      +'<div class="risco-nome">'+c.nome+'</div>'
      +'<span class="risco-assessor">'+c.assessor+'</span>'
      +'<span class="risco-pat">'+fmtPat(c.patrimonio)+'</span>'
      +'<span class="risco-motivo">'+motivoStr+'</span>'
      +'<button class="btn-ver" onclick="verDetalhesCliente(\''+c.assessor+'\',\''+c.nome+'\')">Ver detalhes</button>'
      +'</div>';
  }).join("");
}

function verDetalhesCliente(assessor, nome){
  // Encontra o assessor no ranking e abre o card dele
  const idx = rankingOrder ? rankingOrder.findIndex(a=>a.nome===assessor) : -1;
  if(idx>=0){
    const detalhe = document.getElementById("detalhe-"+idx);
    if(detalhe && !detalhe.classList.contains("aberto")) toggleAssessor(idx);
    setTimeout(function(){
      const cards = document.querySelectorAll(".cliente-card");
      for(var i=0;i<cards.length;i++){
        if(cards[i].getAttribute("data-nome")===nome){
          cards[i].scrollIntoView({behavior:"smooth",block:"center"});
          cards[i].style.borderColor="#FF6B6B";
          setTimeout(function(){ cards[i].style.borderColor=""; }, 3000);
          break;
        }
      }
    }, 300);
  }
}

function renderAlertas(){
  const porAssessor = {};
  clientesData.forEach(c=>{
    if(!c.alertas || !c.alertas.length) return;
    if(!porAssessor[c.assessor]) porAssessor[c.assessor] = [];
    c.alertas.forEach(function(al){
      porAssessor[c.assessor].push({ticker:al.ticker||"—", titulo:al.titulo||al.title||"Alerta", tipo:(al.tipo||al.type||"atencao").toLowerCase(), assessor:c.assessor});
    });
  });
  const assessoresComAlertas = Object.keys(porAssessor);
  if(!assessoresComAlertas.length){
    document.getElementById("sec-alertas").style.display = "none";
    return;
  }
  document.getElementById("sec-alertas").style.display = "block";
  let html = "";
  assessoresComAlertas.forEach(function(aN, idx){
    const als = porAssessor[aN];
    html += '<div class="alerta-acordeao">'
      +'<div class="alerta-header" onclick="toggleAlerta('+idx+')">'
      +'<span class="alerta-header-nome">'+aN+'</span>'
      +'<span class="alerta-count">'+als.length+' alerta(s)</span>'
      +'<span style="font-size:16px;color:#2A2A4A;transition:transform .2s" id="achev-'+idx+'">▾</span>'
      +'</div>'
      +'<div class="alerta-body" id="abody-'+idx+'">';
    als.forEach(function(al){
      const tipoCls = al.tipo.includes("compra")?"alerta-compra":al.tipo.includes("venda")?"alerta-venda":"alerta-atencao";
      html += '<div class="alerta-linha">'
        +'<span class="alerta-ticker">'+al.ticker+'</span>'
        +'<span class="alerta-titulo">'+al.titulo+'</span>'
        +'<span class="alerta-tipo '+tipoCls+'">'+al.tipo+'</span>'
        +'</div>';
    });
    html += '</div></div>';
  });
  document.getElementById("alertas-container").innerHTML = html;
}

function toggleAlerta(i){
  const body = document.getElementById("abody-"+i);
  const chev = document.getElementById("achev-"+i);
  const aberto = body.classList.toggle("aberto");
  chev.style.transform = aberto?"rotate(180deg)":"";
}

function popularFiltroAssessor(){
  const assessores = [...new Set(clientesData.map(c=>c.assessor).filter(Boolean))].sort();
  const sel = document.getElementById("filtro-assessor");
  assessores.forEach(function(a){
    const op = document.createElement("option");
    op.value = a; op.textContent = a;
    sel.appendChild(op);
  });
}

function limparFiltros(){
  document.getElementById("filtro-assessor").value="";
  document.getElementById("filtro-perfil").value="";
  document.getElementById("filtro-status").value="";
  aplicarFiltros();
}

function aplicarFiltros(){
  const fa = document.getElementById("filtro-assessor").value;
  const fp = document.getElementById("filtro-perfil").value;
  const fs = document.getElementById("filtro-status").value;
  renderRisco(fa, fp, fs);
  renderRankingFiltrado(fa, fp, fs);
}

let rankingOrder = null;

function buildPorAssessor(filtroAssessor, filtroPerfil, filtroStatus){
  const porAssessor = {};
  Object.entries(historico).forEach(function(entry){
    const key = entry[0]; const reg = entry[1];
    const assessor = reg.assessor || key.split("|")[0];
    const perfil = reg.perfil || "";
    const ult = reg.entradas && reg.entradas[0];
    const status = ult && ult.status || "saudavel";
    if(filtroAssessor && assessor !== filtroAssessor) return;
    if(filtroPerfil && perfil && perfil !== filtroPerfil) return;
    if(filtroStatus){
      const stMap = {saudavel:"saudavel",atencao:"atencao",critico:"critico"};
      const stNorm = status==="critico"?"critico":status==="atencao"?"atencao":"saudavel";
      if(filtroStatus !== stNorm) return;
    }
    if(!porAssessor[assessor]) porAssessor[assessor] = {nome:assessor, clientes:[], totalRel:0};
    porAssessor[assessor].clientes.push(Object.assign({key:key},reg));
    porAssessor[assessor].totalRel += (reg.entradas&&reg.entradas.length||0);
  });
  return porAssessor;
}

function renderRanking(){
  renderRankingFiltrado("","","");
}

function renderRankingFiltrado(filtroAssessor, filtroPerfil, filtroStatus){
  const container = document.getElementById("ranking-container");
  if(!Object.keys(historico).length){
    container.innerHTML = '<p class="vazio">Nenhum assessor usou o sistema ainda. Os relatorios aparecerao aqui apos a primeira analise.</p>';
    return;
  }

  const porAssessor = buildPorAssessor(filtroAssessor||"", filtroPerfil||"", filtroStatus||"");
  if(!Object.keys(porAssessor).length){
    container.innerHTML = '<p class="vazio">Nenhum resultado para os filtros selecionados.</p>';
    return;
  }

  const ranking = Object.values(porAssessor).sort(function(a,b){return b.totalRel-a.totalRel;});
  rankingOrder = ranking;
  const maxRel = ranking[0].totalRel || 1;

  // Contar PPTXs por assessor a partir de activity (se disponível)
  const pptxPorAssessor = {};
  if(activityData && activityData.length){
    activityData.forEach(function(ev){
      if(ev.acao==="pptx_gerado" && ev.nome){
        pptxPorAssessor[ev.nome] = (pptxPorAssessor[ev.nome]||0)+1;
      }
    });
  }

  // Função para buscar dados financeiros por nome (match aproximado)
  function findFinanc(nome){
    if(assessoresDados[nome]) return assessoresDados[nome];
    // Tentar match por primeiro nome
    const primeiro = nome.split(" ")[0].toLowerCase();
    for(var k in assessoresDados){
      if(k.toLowerCase().startsWith(primeiro)) return assessoresDados[k];
    }
    return null;
  }

  let html = "";
  ranking.forEach(function(a,i){
    const pos = i+1;
    const posClass = pos===1?"top1":pos===2?"top2":pos===3?"top3":"";
    const pctBar = Math.round(a.totalRel/maxRel*100);
    const emAjuste = a.clientes.filter(function(c){return (c.entradas&&c.entradas.length||0)>=2;}).length;
    const criticos = a.clientes.filter(function(c){return c.entradas&&c.entradas[0]&&c.entradas[0].status==="critico";}).length;
    const statusBadge = criticos>0
      ? '<span class="rbadge" style="background:#2A1010;color:#FF6B6B">🔴 '+criticos+' critico(s)</span>'
      : '<span class="rbadge" style="background:#0A2018;color:#5DCAA5">Saudavel</span>';

    // Dados financeiros da planilha
    const fin = findFinanc(a.nome);
    const patrimFin = fin ? fin.patrimonio : 0;
    const receitaFin = fin ? fin.receita_total : 0;
    const okrMensal  = fin ? fin.okr_mensal : 0;
    const pctOKR     = fin ? fin.pct_realizado_okr : 0;
    const roaAnual   = fin ? fin.roa_anual : 0;
    const okrColor   = pctOKR >= 100 ? "#5DCAA5" : pctOKR >= 70 ? "#FFD966" : "#FF6B6B";
    const okrBarW    = Math.min(pctOKR, 100);
    const fmtRec = function(v){ return v > 0 ? "R$ "+Number(v).toLocaleString("pt-BR",{maximumFractionDigits:0}) : "—"; };

    // Métricas de performance
    const mesAtual = new Date().toISOString().slice(0,7);
    const analisesMes = a.clientes.reduce(function(s,c){
      return s + (c.entradas||[]).filter(function(e){return e.data&&e.data.startsWith(mesAtual);}).length;
    },0);
    const pptxGerados = pptxPorAssessor[a.nome]||0;
    const clientesComScoreAlto = a.clientes.filter(function(c){return c.entradas&&c.entradas[0]&&c.entradas[0].score_servir>=4;}).length;
    const pctScore = a.clientes.length>0 ? Math.round(clientesComScoreAlto/a.clientes.length*100) : 0;
    const clientesComRentOk = a.clientes.filter(function(c){return c.entradas&&c.entradas[0]&&c.entradas[0].rent12_pct_cdi>=80;}).length;
    const pctRent = a.clientes.length>0 ? Math.round(clientesComRentOk/a.clientes.length*100) : 0;

    const finBlock = fin ? (
      '<div style="display:flex;gap:16px;align-items:center;padding:4px 0 6px;flex-wrap:wrap">'
      +'<span style="font-size:12px;color:#C9A96E;font-weight:700">'+fmtPat(patrimFin)+'</span>'
      +'<span style="font-size:11px;color:#8A9A8A">|</span>'
      +'<span style="font-size:11px;color:#B0B8B0">Rec: <b style="color:#5DCAA5">'+fmtRec(receitaFin)+'</b></span>'
      +'<span style="font-size:11px;color:#8A9A8A">|</span>'
      +'<span style="font-size:11px;color:#B0B8B0">OKR: '+fmtRec(okrMensal)+'</span>'
      +'<span style="font-size:11px;color:#8A9A8A">|</span>'
      +'<span style="font-size:11px;color:#B0B8B0">ROA: <b style="color:'+okrColor+'">'+roaAnual.toFixed(2)+'%a.a.</b></span>'
      +'</div>'
      +'<div class="rank-bar-row" style="margin-bottom:4px"><span class="lbl" style="min-width:68px">OKR '+Math.round(pctOKR)+'%</span><div class="rank-bar-bg" style="flex:1"><div class="rank-bar-fill" style="width:'+okrBarW+'%;background:'+okrColor+'"></div></div></div>'
    ) : "";

    html += '<div class="rank-item" id="rank-'+i+'" onclick="toggleAssessor('+i+')">'
      +'<div class="rank-pos '+posClass+'">#'+pos+'</div>'
      +'<div style="flex:1">'
      +'<div class="rank-nome">'+a.nome+'</div>'
      +finBlock
      +'<div class="rank-meta" style="margin-top:2px">'+a.clientes.length+' cliente(s) · '+a.totalRel+' relatorio(s) · '+emAjuste+' em ajuste ativo</div>'
      +'</div>'
      +'<div class="rank-bars">'
      +'<div class="rank-bar-row"><span class="lbl">Relatorios</span><div class="rank-bar-bg"><div class="rank-bar-fill" style="width:'+pctBar+'%;background:#8B9FE8"></div></div></div>'
      +'<div class="rank-bar-row"><span class="lbl">Em ajuste</span><div class="rank-bar-bg"><div class="rank-bar-fill" style="width:'+(a.clientes.length>0?Math.round(emAjuste/a.clientes.length*100):0)+'%;background:#5DCAA5"></div></div></div>'
      +'<div class="rank-bar-row"><span class="lbl">Score >=4</span><div class="rank-bar-bg"><div class="rank-bar-fill" style="width:'+pctScore+'%;background:'+(pctScore>=60?"#5DCAA5":"#FF6B6B")+'"></div></div><span style="font-size:10px;color:#3A6A48;margin-left:4px">'+pctScore+'%</span></div>'
      +'<div class="rank-bar-row"><span class="lbl">Rent >=80%</span><div class="rank-bar-bg"><div class="rank-bar-fill" style="width:'+pctRent+'%;background:'+(pctRent>=60?"#5DCAA5":"#FFD966")+'"></div></div><span style="font-size:10px;color:#3A6A48;margin-left:4px">'+pctRent+'%</span></div>'
      +'</div>'
      +'<div class="rank-badges">'+statusBadge+'<span style="font-size:18px;color:#2A2A4A;transition:transform .2s" id="chevron-'+i+'">▾</span></div>'
      +'</div>'
      +'<div class="detalhe-assessor" id="detalhe-'+i+'">'
      +'<div class="metricas-assessor">'
      +(fin?('<div class="met-item"><div class="met-n" style="color:'+okrColor+'">'+Math.round(pctOKR)+'%</div><div class="met-l">OKR realizado</div><div class="meta-bar '+(pctOKR>=100?"meta-ok":"meta-nok")+'" style="width:100%"></div></div>'):"")
      +'<div class="met-item"><div class="met-n">'+analisesMes+'</div><div class="met-l">Analises este mes</div><div class="meta-bar '+(analisesMes>=3?"meta-ok":"meta-nok")+'" style="width:100%"></div></div>'
      +'<div class="met-item"><div class="met-n">'+pptxGerados+'</div><div class="met-l">PPTX gerados</div><div class="meta-bar '+(pptxGerados>=1?"meta-ok":"meta-nok")+'" style="width:100%"></div></div>'
      +'<div class="met-item"><div class="met-n">'+pctScore+'%</div><div class="met-l">Clientes score >= 4</div><div class="meta-bar '+(pctScore>=60?"meta-ok":"meta-nok")+'" style="width:100%"></div></div>'
      +'<div class="met-item"><div class="met-n">'+pctRent+'%</div><div class="met-l">Clientes rent >= 80% CDI</div><div class="meta-bar '+(pctRent>=60?"meta-ok":"meta-nok")+'" style="width:100%"></div></div>'
      +'</div>'
      +renderClientesAssessor(a.clientes, a.nome)
      +'</div>';
  });

  container.innerHTML = html;
}

function toggleAssessor(i){
  const detalhe = document.getElementById("detalhe-"+i);
  const rank    = document.getElementById("rank-"+i);
  const chev    = document.getElementById("chevron-"+i);
  const aberto  = detalhe.classList.toggle("aberto");
  rank.classList.toggle("aberto", aberto);
  chev.style.transform = aberto?"rotate(180deg)":"";
}

function renderClientesAssessor(clientes, assessorNome){
  const sorted = clientes.slice().sort(function(a,b){
    const da = a.entradas&&a.entradas[0]&&a.entradas[0].data||"";
    const db = b.entradas&&b.entradas[0]&&b.entradas[0].data||"";
    return db.localeCompare(da);
  });

  return '<div style="padding:8px 0 4px">' + sorted.map(function(c){
    const entradas = c.entradas||[];
    const ultima   = entradas[0];
    const statusCls = ultima&&ultima.status==="critico"?"c-critico":ultima&&ultima.status==="atencao"?"c-atencao":"c-ok";
    const statusTxt = ultima&&ultima.status==="critico"?"🔴 Critico":ultima&&ultima.status==="atencao"?"🟡 Atencao":"🟢 Saudavel";
    const fmt = function(v){return v>0?"R$ "+Number(v).toLocaleString("pt-BR",{maximumFractionDigits:0}):"—";};
    const notaKey = assessorNome+"|"+c.nome;
    const notaObj = notas[notaKey]||{};
    const notaVal = (notaObj.nota||"").replace(/"/g,"&quot;");
    const clienteId = notaObj.id||"";

    let timelineHtml = '<div class="timeline">';
    entradas.forEach(function(e,idx){
      const isFirst = idx===0;
      const dotCls  = e.status==="critico"?"critico":e.status==="atencao"?"atencao":"ok";
      const prev    = entradas[idx+1];

      let deltaHtml = "";
      if(prev && e.composicao && prev.composicao){
        const modelo = MODELOS[c.perfil]||MODELOS.conservadora;
        const deltas = [];
        Object.keys(e.composicao).forEach(function(k){
          const v = e.composicao[k];
          const vPrev = prev.composicao[k]||0;
          const diff  = v - vPrev;
          if(Math.abs(diff) < 0.5) return;
          const meta  = modelo[k]||0;
          const distNow  = Math.abs(v - meta);
          const distPrev = Math.abs(vPrev - meta);
          const approx   = distNow < distPrev;
          const sinal    = diff>0?"▲":"▼";
          const cls      = diff>0 ? (approx?"delta-up-ok":"delta-up-bad") : (approx?"delta-dn-ok":"delta-dn-bad");
          deltas.push('<span class="delta-item '+cls+'">'+sinal+' '+(LABEL_COMP[k]||k)+' '+(diff>0?"+":"")+diff.toFixed(1)+'%</span>');
        });
        if(deltas.length){
          deltaHtml = '<div style="margin-top:4px;font-size:10px;color:#2A5A3A">Variacao vs analise anterior:</div><div class="delta-row">'+deltas.join("")+'</div>';
        } else {
          deltaHtml = '<div style="font-size:10px;color:#2A2A4A;margin-top:4px">Sem variacoes significativas vs analise anterior</div>';
        }

        if(sugestoesHist.length && deltas.length){
          const rfSugs = sugestoesHist[0]&&sugestoesHist[0].rf||[];
          const MAPA_IDX = {IPCA:"inflacao",CDI:"pos_fixado",SELIC:"pos_fixado","Pre-fixado":"pre_fixado"};
          let aderiuRF=false;
          rfSugs.forEach(function(s){
            if(!s.perfis||!s.perfis.includes(c.perfil)) return;
            const k = MAPA_IDX[s.indexador];
            if(!k) return;
            const diff2 = (e.composicao[k]||0)-(prev.composicao[k]||0);
            if(s.acao==="alocar"  && diff2>0.5) aderiuRF=true;
            if(s.acao==="substituir" && s.de && diff2>0.5) aderiuRF=true;
          });
          if(aderiuRF){
            deltaHtml += '<div style="margin-top:6px;font-size:11px;color:#5DCAA5;padding:5px 10px;background:#0A2018;border-radius:6px;border:1px solid #0A3020">Adesao a sugestao de Renda Fixa detectada nesta analise</div>';
          }
        }
      }

      const rentTxt = e.rent12_pct_cdi>0
        ? '<span style="color:'+(e.rent12_pct_cdi<70?"#FF6B6B":"#5DCAA5")+'">'+e.rent12_pct_cdi+'%</span> do CDI'
        : "—";

      timelineHtml += '<div class="tl-entry">'
        +'<div class="tl-dot '+dotCls+(isFirst?" first":"")+'"></div>'
        +'<div class="tl-content">'
        +'<div class="tl-data"><b>'+e.data+' '+(e.hora||"")+'</b>'
        +'<span class="c-badge '+statusCls+'">'+statusTxt+'</span>'
        +(isFirst?'<span style="font-size:10px;color:#8B9FE8;font-weight:700">MAIS RECENTE</span>':"")
        +'</div>'
        +'<div class="tl-metrics">'
        +'<div class="tl-m">Patrimonio <span>'+fmt(e.patrimonio)+'</span></div>'
        +'<div class="tl-m">Modelo Servir <span>'+e.score_servir+'/6</span></div>'
        +'<div class="tl-m">Cross Sell <span>'+e.cross_ativos+'/5</span></div>'
        +'<div class="tl-m">Rentab. CDI <span>'+rentTxt+'</span></div>'
        +'</div>'
        +deltaHtml
        +'</div></div>';
    });
    timelineHtml += "</div>";

    return '<div class="cliente-card" data-nome="'+c.nome+'">'
      +'<div class="cliente-header">'
      +'<div class="cliente-nome">'+c.nome+'</div>'
      +'<div class="cliente-perfil">'+(c.perfil||"—")+'</div>'
      +(c.objetivo?'<div class="cliente-objetivo">'+c.objetivo+'</div>':"")
      +'<span class="c-badge '+statusCls+'" style="margin-left:auto">'+statusTxt+'</span>'
      +'<span style="font-size:11px;color:#3A6A48">'+entradas.length+' analise(s)</span>'
      +'</div>'
      +timelineHtml
      +'<input class="nota-input" value="'+notaVal+'"'
      +' placeholder="Orientacao do lider para este assessor sobre este cliente..."'
      +' onchange="salvarNota(\''+clienteId+'\',this.value)">'
      +'</div>';
  }).join("") + "</div>";
}

async function salvarNota(id,nota){
  if(!id) return;
  await fetch("/api/clientes/"+id+"/nota",{method:"POST",headers:{"Content-Type":"application/json"},body:JSON.stringify({nota:nota})});
}

init();
</script>
</body></html>"""

HTML_ADMIN = r"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Braúna — Admin</title>
<script src="https://cdn.sheetjs.com/xlsx-0.20.3/package/dist/xlsx.full.min.js"></script>
<style>
*{margin:0;padding:0;box-sizing:border-box}
body{background:#060D08;color:#F0F0F0;font-family:'Segoe UI',system-ui,sans-serif;min-height:100vh}
header{background:#071209;border-bottom:2px solid #1A3A1A;padding:14px 28px;display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:10px}
header h1{font-size:16px;color:#5DCAA5;font-weight:700;letter-spacing:.5px}
header p{font-size:10px;color:#2A5A3A;margin-top:2px}
.nav{display:flex;gap:8px;flex-wrap:wrap}
.nav a{font-size:12px;padding:5px 12px;border-radius:6px;border:1px solid #1A2A1A;color:#4A7055;text-decoration:none;transition:all .2s}
.nav a:hover{border-color:#5DCAA5;color:#5DCAA5}
.nav a.active{background:#5DCAA5;color:#000;border-color:#5DCAA5;font-weight:700}
.container{max-width:1280px;margin:0 auto;padding:24px 20px}

/* Tabs principais */
.main-tabs{display:flex;gap:0;border-bottom:2px solid #1A3A1A;margin-bottom:24px}
.main-tab{padding:12px 24px;font-size:13px;font-weight:600;color:#3A6A48;background:none;border:none;border-bottom:3px solid transparent;cursor:pointer;font-family:inherit;transition:all .2s;margin-bottom:-2px}
.main-tab:hover{color:#5DCAA5}
.main-tab.active{color:#5DCAA5;border-bottom-color:#5DCAA5}
.main-panel{display:none}.main-panel.active{display:block}

/* Cards */
.card{background:#0A1A0C;border:1px solid #1A3A1A;border-radius:12px;padding:20px;margin-bottom:16px}
.card-title{font-size:11px;font-weight:700;text-transform:uppercase;letter-spacing:.8px;margin-bottom:14px;display:flex;align-items:center;gap:8px}
label{font-size:11px;color:#4A7055;display:block;margin-bottom:5px}
input[type=text],textarea,select{width:100%;background:#060E08;border:1px solid #1A2A1A;border-radius:7px;padding:8px 10px;color:#F0F0F0;font-size:13px;outline:none;font-family:inherit}
input[type=text]:focus,textarea:focus,select:focus{border-color:#5DCAA5}
.btn{display:inline-flex;align-items:center;gap:6px;background:#5DCAA5;color:#000;border:none;border-radius:7px;padding:9px 16px;font-size:13px;font-weight:700;cursor:pointer;transition:all .2s;font-family:inherit}
.btn:hover{opacity:.85}.btn:disabled{opacity:.35;cursor:not-allowed}
.btn-out{background:transparent;color:#5DCAA5;border:1px solid #5DCAA5}
.btn-sm{padding:5px 12px;font-size:11px;border-radius:6px}
.btn-ghost{background:transparent;border:1px solid #1A2A1A;color:#3A6A48;border-radius:6px;padding:5px 10px;font-size:11px;cursor:pointer;transition:all .2s;font-family:inherit}
.btn-ghost:hover{border-color:#5DCAA5;color:#5DCAA5}

/* KPI chips */
.kpi-row{display:flex;gap:12px;flex-wrap:wrap;margin-bottom:20px}
.kpi{background:#0D1C10;border:1px solid #1A3A1A;border-radius:10px;padding:14px 18px;min-width:130px;flex:1}
.kpi-val{font-size:28px;font-weight:900;line-height:1}
.kpi-lbl{font-size:10px;color:#3A6A48;text-transform:uppercase;letter-spacing:.5px;margin-top:4px}

/* Tabela de atividade */
.act-table{width:100%;border-collapse:collapse;font-size:12px}
.act-table th{text-align:left;padding:8px 10px;font-size:10px;color:#2A5A3A;text-transform:uppercase;letter-spacing:.5px;border-bottom:1px solid #1A3A1A;font-weight:700}
.act-table td{padding:9px 10px;border-bottom:1px solid #0D1A0E;vertical-align:top}
.act-table tr:hover td{background:#0C180E}

/* Badges */
.badge{display:inline-block;padding:2px 8px;border-radius:8px;font-size:10px;font-weight:700}
.badge-assessor{background:#0D2A1D;color:#5DCAA5;border:1px solid #1A4A2A}
.badge-lider{background:#0D1030;color:#8B9FE8;border:1px solid #1A2050}
.badge-head{background:#2A1A08;color:#C9A96E;border:1px solid #4A3010}
.badge-ok{background:#0A2018;color:#5DCAA5}
.badge-warn{background:#2A1808;color:#C9A96E}

/* Assessor cards */
.user-card{background:#0A1A0C;border:1px solid #1A3A1A;border-radius:10px;padding:14px 16px;margin-bottom:10px;display:flex;align-items:flex-start;gap:14px;transition:border-color .2s}
.user-card:hover{border-color:#2A5A3A}
.user-avatar{width:40px;height:40px;border-radius:50%;display:flex;align-items:center;justify-content:center;font-size:16px;font-weight:700;flex-shrink:0}
.user-info{flex:1;min-width:0}
.user-name{font-size:13px;font-weight:700;color:#F0F0F0}
.user-meta{font-size:10px;color:#3A6A48;margin-top:2px}
.user-log{margin-top:8px;display:flex;flex-direction:column;gap:4px;max-height:100px;overflow-y:auto}
.log-row{font-size:10px;color:#4A7055;display:flex;gap:8px;align-items:flex-start}
.log-ts{color:#1E4A30;flex-shrink:0;min-width:90px}

/* Head info cards */
.info-grid{display:grid;grid-template-columns:repeat(auto-fill,minmax(220px,1fr));gap:12px;margin-bottom:16px}
.info-card{background:#0A1A0C;border:1px solid #1A3A1A;border-radius:10px;padding:14px}
.info-card-title{font-size:10px;color:#3A6A48;text-transform:uppercase;letter-spacing:.5px;font-weight:700;margin-bottom:8px}
.info-card-body{font-size:12px;color:#CCC;line-height:1.6}

/* Doc list */
.doc-row{display:flex;align-items:center;gap:10px;padding:8px 10px;border-bottom:1px solid #0D1A0E;font-size:12px}
.doc-row:last-child{border-bottom:none}
.doc-icon{font-size:14px;flex-shrink:0}
.doc-name{flex:1;color:#C9A96E;font-weight:600;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.doc-meta{font-size:10px;color:#2A5A3A;flex-shrink:0}

/* Config tab */
.grid-2{display:grid;grid-template-columns:1fr 1fr;gap:16px}
.upload-area{border:1.5px dashed #1A2A1A;border-radius:10px;padding:22px;text-align:center;cursor:pointer;position:relative;background:#060E08;transition:all .2s}
.upload-area:hover{border-color:#5DCAA5;background:#050D05}
.upload-area input[type=file]{position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%}
.upload-area .ui{pointer-events:none}
.upload-area .icon{font-size:26px;margin-bottom:6px}
.upload-area p{font-size:12px;color:#3A6A48}
.upload-area .fname{color:#5DCAA5;font-weight:600;font-size:12px;margin-top:4px;min-height:18px}
.tab-toolbar{display:flex;align-items:center;gap:8px;flex-wrap:wrap;padding:10px 12px;background:#060C08;border:1px dashed #1A2A1A;border-radius:8px;margin-bottom:12px}
.tab-toolbar .lbl{font-size:11px;color:#3A4A3A;flex:1;min-width:120px}
.btn-file-wrap{position:relative;overflow:hidden;display:inline-flex}
.btn-file-wrap input[type=file]{position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%}
.pdf-preview{background:#060C08;border:1px solid #5DCAA5;border-radius:8px;padding:12px;margin-bottom:12px;display:none}
.pdf-preview .pdf-header{display:flex;align-items:center;justify-content:space-between;margin-bottom:8px}
.pdf-preview .pdf-title{font-size:11px;color:#5DCAA5;font-weight:700}
.pdf-preview textarea{height:140px;font-size:11px;font-family:monospace;color:#9A9;background:#040A04;resize:vertical}
.pdf-preview .pdf-actions{display:flex;gap:6px;margin-top:8px;flex-wrap:wrap}
.sg-tabs{display:flex;gap:6px;margin-bottom:14px;flex-wrap:wrap}
.sg-tab{padding:7px 16px;border-radius:8px;border:1px solid #1A2A1A;background:#060E08;color:#4A7055;font-size:12px;cursor:pointer;transition:all .2s;font-family:inherit}
.sg-tab.active{background:#5DCAA5;color:#000;border-color:#5DCAA5;font-weight:700}
.sg-panel{display:none}.sg-panel.active{display:block}
.sg-item{background:#060E08;border:1px solid #1A2A1A;border-radius:8px;padding:12px 14px;margin-bottom:8px;position:relative;animation:fadeIn .15s ease}
@keyframes fadeIn{from{opacity:0;transform:translateY(-4px)}to{opacity:1;transform:translateY(0)}}
.sg-del{position:absolute;top:8px;right:8px;background:none;border:none;color:#2A3A2A;cursor:pointer;font-size:16px;padding:2px 6px;transition:color .2s;border-radius:4px}
.sg-del:hover{color:#FF6B6B;background:#1A0505}
.sg-grid-2{display:grid;grid-template-columns:1fr 1fr;gap:8px;margin-bottom:8px}
.sg-grid-3{display:grid;grid-template-columns:1fr 1fr 1fr;gap:8px;margin-bottom:8px}
.sg-mini label{font-size:10px;color:#3A6A48;margin-bottom:3px}
.perfis-row{display:flex;gap:10px;flex-wrap:wrap;margin-top:4px}
.perfis-row label{display:flex;align-items:center;gap:4px;font-size:11px;color:#777;cursor:pointer;margin-bottom:0}
.hist-item{background:#060E08;border:1px solid #1A2A1A;border-radius:8px;padding:10px 14px;margin-bottom:6px;display:flex;align-items:center;gap:10px}
.hist-item.ativa{border-color:#5DCAA5}
.tag-ativa{display:inline-block;padding:2px 10px;border-radius:10px;font-size:10px;font-weight:700;background:#0A2018;color:#5DCAA5}
.info-box{background:#060E08;border-radius:8px;padding:12px 14px;font-size:12px;color:#4A7055;line-height:1.6;border:1px solid #1A2A1A}
.info-box b{color:#5DCAA5}
.status-ok{color:#5DCAA5}.status-err{color:#FF6B6B}
@media(max-width:640px){.grid-2,.sg-grid-2,.sg-grid-3{grid-template-columns:1fr}.kpi{min-width:100px}}
</style>
</head>
<body>
<header>
  <div><h1>⚙️ BRAÚNA — ADMINISTRAÇÃO</h1><p>Painel de controle · Atividade da equipe · Informações do sistema</p></div>
  <nav class="nav">
    <a href="/assessor">📊 Assessor</a>
    <a href="/lider">👥 Líder</a>
    <a href="/head-produtos">🏛️ Head</a>
    <a href="/admin" class="active">⚙️ Admin</a>
  </nav>
</header>

<div class="container">

<!-- Tabs principais -->
<div class="main-tabs">
  <button class="main-tab active" onclick="switchTab('assessores',this)">📊 Assessores</button>
  <button class="main-tab"        onclick="switchTab('lideres',this)">👥 Líderes</button>
  <button class="main-tab"        onclick="switchTab('head',this)">🏛️ Head — Informações</button>
  <button class="main-tab"        onclick="switchTab('config',this)">⚙️ Configurações</button>
</div>

<!-- ═══════ TAB: ASSESSORES ═══════ -->
<div class="main-panel active" id="tab-assessores">
  <div class="kpi-row" id="kpi-assessores">
    <div class="kpi"><div class="kpi-val" id="kpi-ativos" style="color:#5DCAA5">—</div><div class="kpi-lbl">Assessores Ativos</div></div>
    <div class="kpi"><div class="kpi-val" id="kpi-analises" style="color:#C9A96E">—</div><div class="kpi-lbl">Análises Realizadas</div></div>
    <div class="kpi"><div class="kpi-val" id="kpi-clientes" style="color:#8B9FE8">—</div><div class="kpi-lbl">Clientes Analisados</div></div>
    <div class="kpi"><div class="kpi-val" id="kpi-ultimo-acesso" style="color:#5DCAA5;font-size:14px">—</div><div class="kpi-lbl">Último Acesso</div></div>
  </div>

  <div class="card">
    <div class="card-title" style="color:#5DCAA5">👤 Assessores — Atividade por Usuário</div>
    <div id="lista-assessores"><div style="color:#1E4A30;font-size:12px;padding:20px;text-align:center">Carregando...</div></div>
  </div>

  <div class="card">
    <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:12px">
      <div class="card-title" style="color:#5DCAA5;margin-bottom:0">📋 Log de Atividade — Assessores</div>
      <button class="btn-ghost" onclick="carregarDashboard()">↻ Atualizar</button>
    </div>
    <div style="overflow-x:auto">
      <table class="act-table" id="log-assessores">
        <thead><tr><th>Hora</th><th>Assessor</th><th>Ação</th><th>Detalhe</th></tr></thead>
        <tbody id="log-assessores-body"><tr><td colspan="4" style="color:#1E4A30;text-align:center;padding:20px">Carregando...</td></tr></tbody>
      </table>
    </div>
  </div>
</div>

<!-- ═══════ TAB: LÍDERES ═══════ -->
<div class="main-panel" id="tab-lideres">
  <div class="kpi-row">
    <div class="kpi"><div class="kpi-val" id="kpi-lider-ativos" style="color:#8B9FE8">—</div><div class="kpi-lbl">Líderes Ativos</div></div>
    <div class="kpi"><div class="kpi-val" id="kpi-lider-acoes" style="color:#C9A96E">—</div><div class="kpi-lbl">Ações Realizadas</div></div>
    <div class="kpi"><div class="kpi-val" id="kpi-lider-ultimo" style="color:#8B9FE8;font-size:14px">—</div><div class="kpi-lbl">Último Acesso</div></div>
  </div>

  <div class="card">
    <div class="card-title" style="color:#8B9FE8">👥 Líderes — Atividade por Usuário</div>
    <div id="lista-lideres"><div style="color:#1E4A30;font-size:12px;padding:20px;text-align:center">Carregando...</div></div>
  </div>

  <div class="card">
    <div class="card-title" style="color:#8B9FE8;margin-bottom:12px">📋 Log de Atividade — Líderes</div>
    <div style="overflow-x:auto">
      <table class="act-table">
        <thead><tr><th>Hora</th><th>Líder</th><th>Ação</th><th>Detalhe</th></tr></thead>
        <tbody id="log-lideres-body"><tr><td colspan="4" style="color:#1E4A30;text-align:center;padding:20px">Carregando...</td></tr></tbody>
      </table>
    </div>
  </div>
</div>

<!-- ═══════ TAB: HEAD ═══════ -->
<div class="main-panel" id="tab-head">

  <!-- KPIs rápidos -->
  <div class="kpi-row" id="kpi-head-row">
    <div class="kpi"><div class="kpi-val" id="kpi-docs" style="color:#C9A96E">—</div><div class="kpi-lbl">Docs na Base</div></div>
    <div class="kpi"><div class="kpi-val" id="kpi-calls" style="color:#5DCAA5">—</div><div class="kpi-lbl">Calls Ativos</div></div>
    <div class="kpi"><div class="kpi-val" id="kpi-estr" style="color:#8B9FE8">—</div><div class="kpi-lbl">Op. Estruturadas</div></div>
    <div class="kpi"><div class="kpi-val" id="kpi-prod" style="color:#C9A96E">—</div><div class="kpi-lbl">Produtos Destaque</div></div>
    <div class="kpi"><div class="kpi-val" id="kpi-alertas" style="color:#FF6B6B">—</div><div class="kpi-lbl">Alertas Ativos</div></div>
    <div class="kpi"><div class="kpi-val" id="kpi-gestoras" style="color:#5DCAA5">—</div><div class="kpi-lbl">Gestoras Cadastradas</div></div>
    <div class="kpi"><div class="kpi-val" id="kpi-clientes" style="color:#8B9FE8">—</div><div class="kpi-lbl">Clientes Registrados</div></div>
  </div>

  <!-- Cenário Macro -->
  <div class="card" style="border-color:#2A3A20">
    <div class="card-title" style="color:#C9A96E">🌐 Cenário Macro Atual</div>
    <div id="head-cenario"><div style="color:#1E4A30;font-size:12px">Carregando...</div></div>
  </div>

  <!-- Grid: Docs + Calls -->
  <div style="display:grid;grid-template-columns:1fr 1fr;gap:16px;margin-bottom:16px">

    <!-- Base de Conhecimento -->
    <div class="card" style="border-color:#2A2A18">
      <div class="card-title" style="color:#C9A96E">📚 Base de Conhecimento</div>
      <div id="head-docs"><div style="color:#1E4A30;font-size:12px">Carregando...</div></div>
    </div>

    <!-- Calls de Ações -->
    <div class="card" style="border-color:#1A2E3A">
      <div class="card-title" style="color:#5DCAA5">📈 Calls de Ações</div>
      <div id="head-calls"><div style="color:#1E4A30;font-size:12px">Carregando...</div></div>
    </div>

  </div>

  <!-- Grid: Estruturadas + Gestoras -->
  <div style="display:grid;grid-template-columns:1fr 1fr;gap:16px;margin-bottom:16px">

    <!-- Operações Estruturadas -->
    <div class="card" style="border-color:#2A1A10">
      <div class="card-title" style="color:#C9A96E">💼 Operações Estruturadas</div>
      <div id="head-estruturadas"><div style="color:#1E4A30;font-size:12px">Carregando...</div></div>
    </div>

    <!-- Alocações de Referência -->
    <div class="card" style="border-color:#1A3A1A">
      <div class="card-title" style="color:#5DCAA5">📐 Alocações de Referência — Gestoras</div>
      <div id="head-gestoras"><div style="color:#1E4A30;font-size:12px">Carregando...</div></div>
    </div>

  </div>

  <!-- Alertas -->
  <div class="card" style="border-color:#2A1A08">
    <div class="card-title" style="color:#FF6B6B">🚨 Alertas Ativos</div>
    <div id="head-alertas"><div style="color:#1E4A30;font-size:12px">Carregando...</div></div>
  </div>

  <!-- Clientes Registrados -->
  <div class="card" style="border-color:#1A2A3A">
    <div class="card-title" style="color:#8B9FE8">👤 Carteiras de Clientes Registradas</div>
    <div id="head-clientes"><div style="color:#1E4A30;font-size:12px">Carregando...</div></div>
  </div>

</div>

<!-- ═══════ TAB: CONFIGURAÇÕES ═══════ -->
<div class="main-panel" id="tab-config">

<!-- Linha 1: Carta + Comunicado -->
<div class="grid-2" style="margin-bottom:16px">

  <!-- Carta da Gestão -->
  <div class="card" style="margin-bottom:0">
    <div class="card-title" style="color:#5DCAA5">📄 Carta da Gestão — Upload Mensal</div>
    <p style="font-size:11px;color:#2A5A3A;margin-bottom:12px;line-height:1.5">Suba o PDF da carta mensal (Levante, BTG, XP Research…). O sistema extrai os insights e usa em todas as análises de carteira dos assessores.</p>
    <div class="upload-area" id="drop-carta">
      <input type="file" id="pdf-carta" accept=".pdf" onchange="onCartaSelect()">
      <div class="ui">
        <div class="icon">📄</div>
        <p>Arraste o PDF ou clique aqui</p>
        <div class="fname" id="fname-carta"></div>
      </div>
    </div>
    <div style="margin-top:10px;display:flex;align-items:center;gap:10px;flex-wrap:wrap">
      <button class="btn btn-sm" id="btn-carta" onclick="uploadCarta()">Processar carta</button>
      <span id="st-carta" style="font-size:11px"></span>
    </div>
    <div id="carta-preview" style="display:none;margin-top:8px">
      <div class="info-box" id="carta-info" style="font-size:11px;color:#5DCAA5"></div>
    </div>
    <div style="margin-top:12px" id="status-contexto"><div class="info-box" style="color:#1E4A30;font-style:italic">Carregando status...</div></div>
  </div>

  <!-- Comunicado -->
  <div class="card" style="margin-bottom:0">
    <div class="card-title" style="color:#5DCAA5">📢 Comunicado para Assessores</div>
    <div style="margin-bottom:10px">
      <label>Estratégia do mês (visível na tela do assessor)</label>
      <textarea id="txt-estrategia" rows="3" placeholder="Ex: Foco em IPCA+ este mês. Clientes conservadores devem aumentar exposição em inflação..."></textarea>
    </div>
    <div style="margin-bottom:12px">
      <label>Alerta / mensagem rápida</label>
      <textarea id="txt-mensagem" rows="2" placeholder="Ex: ⚠️ Reunião de equipe na sexta às 14h."></textarea>
    </div>
    <div style="display:flex;align-items:center;gap:10px;flex-wrap:wrap">
      <button class="btn btn-sm" onclick="salvarMensagem()">📤 Publicar</button>
      <span id="st-msg" style="font-size:11px"></span>
    </div>
    <div style="margin-top:12px" id="status-msg"><div class="info-box" style="color:#1E4A30;font-style:italic">Carregando...</div></div>
  </div>

</div>

<!-- Sugestões de Alocação -->
<div class="card">
  <div class="card-title" style="color:#5DCAA5">💡 Sugestões de Alocação para Assessores</div>
  <p style="font-size:11px;color:#2A5A3A;margin-bottom:14px;line-height:1.5">Monte as sugestões por categoria. Ao publicar, ficam disponíveis para todos os assessores filtradas pelo perfil do cliente.</p>

  <div style="margin-bottom:14px">
    <label>Título da publicação</label>
    <input type="text" id="sg-titulo" placeholder="Ex: Carteira Recomendada — Julho 2026" style="max-width:420px">
  </div>

  <!-- Tabs -->
  <div class="sg-tabs">
    <button class="sg-tab active" onclick="sgTab('rf',this)">🏦 Renda Fixa</button>
    <button class="sg-tab"        onclick="sgTab('rv',this)">📈 Renda Variável</button>
    <button class="sg-tab"        onclick="sgTab('fii',this)">🏢 FIIs</button>
    <button class="sg-tab"        onclick="sgTab('intl',this)">🌎 Internacional</button>
    <button class="sg-tab"        onclick="sgTab('est',this)">⚙️ Estruturadas</button>
  </div>

  <!-- ── RENDA FIXA ── -->
  <div class="sg-panel active" id="sg-panel-rf">
    <div class="tab-toolbar">
      <span class="lbl">Importar de arquivo:</span>
      <div class="btn-file-wrap">
        <button class="btn btn-sm btn-out">📄 PDF</button>
        <input type="file" accept=".pdf" onchange="uploadPdfCategoria('rf',this)">
      </div>
      <div class="btn-file-wrap">
        <button class="btn btn-sm btn-out">📊 Planilha</button>
        <input type="file" accept=".csv,.xlsx,.xls" onchange="importarPlanilha('rf',this)">
      </div>
      <button class="btn-ghost" onclick="baixarModelo('rf')">⬇ Modelo CSV</button>
      <span id="import-st-rf" style="font-size:11px;color:#5DCAA5"></span>
    </div>
    <!-- Preview PDF -->
    <div class="pdf-preview" id="pdf-preview-rf">
      <div class="pdf-header">
        <span class="pdf-title">📄 Conteúdo extraído do PDF — Renda Fixa</span>
        <button class="btn-ghost" onclick="fecharPdfPreview('rf')" style="padding:2px 8px;font-size:11px">✕</button>
      </div>
      <textarea id="pdf-texto-rf" placeholder="Texto extraído aparece aqui..."></textarea>
      <div class="pdf-actions">
        <button class="btn btn-sm" onclick="usarPdfComoSugestao('rf')">→ Adicionar como sugestão RF</button>
        <button class="btn-ghost" onclick="copiarPdf('rf')">📋 Copiar</button>
      </div>
    </div>
    <button class="btn btn-sm btn-ghost" onclick="sgAdd('rf')" style="margin-bottom:10px">+ Adicionar manualmente</button>
    <div id="sg-list-rf"></div>
  </div>

  <!-- ── RENDA VARIÁVEL ── -->
  <div class="sg-panel" id="sg-panel-rv">
    <div class="tab-toolbar">
      <span class="lbl">Importar de arquivo:</span>
      <div class="btn-file-wrap">
        <button class="btn btn-sm btn-out">📄 PDF</button>
        <input type="file" accept=".pdf" onchange="uploadPdfCategoria('rv',this)">
      </div>
      <div class="btn-file-wrap">
        <button class="btn btn-sm btn-out">📊 Planilha</button>
        <input type="file" accept=".csv,.xlsx,.xls" onchange="importarPlanilha('rv',this)">
      </div>
      <button class="btn-ghost" onclick="baixarModelo('rv')">⬇ Modelo CSV</button>
      <span id="import-st-rv" style="font-size:11px;color:#5DCAA5"></span>
    </div>
    <div class="pdf-preview" id="pdf-preview-rv">
      <div class="pdf-header">
        <span class="pdf-title">📄 Conteúdo extraído do PDF — Renda Variável</span>
        <button class="btn-ghost" onclick="fecharPdfPreview('rv')" style="padding:2px 8px;font-size:11px">✕</button>
      </div>
      <textarea id="pdf-texto-rv" placeholder="Texto extraído aparece aqui..."></textarea>
      <div class="pdf-actions">
        <button class="btn btn-sm" onclick="usarPdfComoSugestao('rv')">→ Adicionar como sugestão RV</button>
        <button class="btn-ghost" onclick="copiarPdf('rv')">📋 Copiar</button>
      </div>
    </div>
    <button class="btn btn-sm btn-ghost" onclick="sgAdd('rv')" style="margin-bottom:10px">+ Adicionar manualmente</button>
    <div id="sg-list-rv"></div>
  </div>

  <!-- ── FIIs ── -->
  <div class="sg-panel" id="sg-panel-fii">
    <div class="tab-toolbar">
      <span class="lbl">Importar de arquivo:</span>
      <div class="btn-file-wrap">
        <button class="btn btn-sm btn-out">📄 PDF</button>
        <input type="file" accept=".pdf" onchange="uploadPdfCategoria('fii',this)">
      </div>
      <div class="btn-file-wrap">
        <button class="btn btn-sm btn-out">📊 Planilha</button>
        <input type="file" accept=".csv,.xlsx,.xls" onchange="importarFii(this)">
      </div>
      <button class="btn-ghost" onclick="baixarModelo('fii')">⬇ Modelo CSV</button>
      <span id="import-st-fii" style="font-size:11px;color:#5DCAA5"></span>
    </div>
    <div class="pdf-preview" id="pdf-preview-fii">
      <div class="pdf-header">
        <span class="pdf-title">📄 Conteúdo extraído do PDF — FIIs</span>
        <button class="btn-ghost" onclick="fecharPdfPreview('fii')" style="padding:2px 8px;font-size:11px">✕</button>
      </div>
      <textarea id="pdf-texto-fii" placeholder="Texto extraído aparece aqui..."></textarea>
      <div class="pdf-actions">
        <button class="btn btn-sm" onclick="usarPdfComoSugestao('fii')">→ Usar como referência FII</button>
        <button class="btn-ghost" onclick="copiarPdf('fii')">📋 Copiar</button>
      </div>
    </div>
    <p style="font-size:11px;color:#3A6A48;margin-bottom:6px">Uma linha por fundo: <b style="color:#5DCAA5">TICKER | Segmento | % sugerido | Motivo</b></p>
    <textarea id="sg-fii" rows="8" placeholder="MXRF11 | Papel | 0.5% | Alta distribuição, baixa volatilidade&#10;HGLG11 | Logística | 1.0% | Contratos longos, boa gestão&#10;XPML11 | Shopping | 0.8% | Recuperação do varejo&#10;BTLG11 | Logística | 0.7% | Portfólio diversificado"></textarea>
  </div>

  <!-- ── INTERNACIONAL ── -->
  <div class="sg-panel" id="sg-panel-intl">
    <div class="tab-toolbar">
      <span class="lbl">Importar de arquivo:</span>
      <div class="btn-file-wrap">
        <button class="btn btn-sm btn-out">📄 PDF</button>
        <input type="file" accept=".pdf" onchange="uploadPdfCategoria('intl',this)">
      </div>
      <div class="btn-file-wrap">
        <button class="btn btn-sm btn-out">📊 Planilha</button>
        <input type="file" accept=".csv,.xlsx,.xls" onchange="importarPlanilha('intl',this)">
      </div>
      <button class="btn-ghost" onclick="baixarModelo('intl')">⬇ Modelo CSV</button>
      <span id="import-st-intl" style="font-size:11px;color:#5DCAA5"></span>
    </div>
    <div class="pdf-preview" id="pdf-preview-intl">
      <div class="pdf-header">
        <span class="pdf-title">📄 Conteúdo extraído do PDF — Internacional</span>
        <button class="btn-ghost" onclick="fecharPdfPreview('intl')" style="padding:2px 8px;font-size:11px">✕</button>
      </div>
      <textarea id="pdf-texto-intl" placeholder="Texto extraído aparece aqui..."></textarea>
      <div class="pdf-actions">
        <button class="btn btn-sm" onclick="usarPdfComoSugestao('intl')">→ Adicionar como sugestão Internacional</button>
        <button class="btn-ghost" onclick="copiarPdf('intl')">📋 Copiar</button>
      </div>
    </div>
    <button class="btn btn-sm btn-ghost" onclick="sgAdd('intl')" style="margin-bottom:10px">+ Adicionar manualmente</button>
    <div id="sg-list-intl"></div>
  </div>

  <!-- ── ESTRUTURADAS ── -->
  <div class="sg-panel" id="sg-panel-est">
    <div class="tab-toolbar">
      <span class="lbl">Importar PDF com operações estruturadas disponíveis:</span>
      <div class="btn-file-wrap">
        <button class="btn btn-sm btn-out">📄 PDF</button>
        <input type="file" accept=".pdf" onchange="uploadPdfCategoria('est',this)">
      </div>
      <span id="import-st-est" style="font-size:11px;color:#5DCAA5"></span>
    </div>
    <div class="pdf-preview" id="pdf-preview-est">
      <div class="pdf-header">
        <span class="pdf-title">📄 Conteúdo extraído — Estruturadas</span>
        <button class="btn-ghost" onclick="fecharPdfPreview('est')" style="padding:2px 8px;font-size:11px">✕</button>
      </div>
      <textarea id="pdf-texto-est" placeholder="Texto extraído aparece aqui..."></textarea>
      <div class="pdf-actions">
        <button class="btn btn-sm" onclick="usarPdfComoSugestao('est')">→ Usar como Estruturadas</button>
        <button class="btn-ghost" onclick="copiarPdf('est')">📋 Copiar</button>
      </div>
    </div>
    <p style="font-size:11px;color:#3A6A48;margin-bottom:6px">Descreva as operações estruturadas disponíveis — exibido para clientes com ações na carteira.</p>
    <textarea id="sg-est" rows="8" placeholder="PETR4 — Financiamento com venda de call coberta OTM 10%, vencimento mensal.&#10;VALE3 — Trava de alta com compra de call ATM e venda de call OTM 15%.&#10;ITUB4 — Proteção com put ATM para carteiras acima de R$ 100k na ação."></textarea>
  </div>

  <!-- Ações -->
  <div style="display:flex;gap:10px;margin-top:16px;align-items:center;flex-wrap:wrap">
    <button class="btn" onclick="publicarSugestoes()">📤 Publicar sugestões</button>
    <button class="btn btn-out" onclick="limparForm()">Limpar</button>
    <span id="sg-status" style="font-size:12px"></span>
  </div>
</div>

<!-- Histórico -->
<div class="card">
  <div class="card-title">📋 Histórico de Publicações</div>
  <div id="sg-historico"><p style="color:#1E4A30;font-size:12px;font-style:italic">Carregando...</p></div>
</div>

<script>
// Admin sem autenticação — seta role para permitir navegar em todas as páginas
localStorage.setItem("brauna_role", "admin");

// ── Constantes ────────────────────────────────────────────────────────────────
const PERFIS     = ["conservadora","moderada","arrojada","agressiva"];
const PERFIS_LBL = {conservadora:"Conservadora",moderada:"Moderada",arrojada:"Arrojada",agressiva:"Agressiva"};
const INDEXADORES= ["IPCA","CDI","Pré-fixado","SELIC","Outro"];
const FONTES_RV  = ["BTG Pactual","XP Research","Safra","Itaú BBA","Goldman Sachs","Morgan Stanley","Levante","Outro"];

const MODELOS_CSV = {
  rf:   ["acao","de","para","indexador","urgencia","motivo","perfis"],
  rv:   ["acao","de","para","motivo","fonte","perfis"],
  intl: ["ativo","tipo","motivo","perfis"],
  fii:  ["ticker","segmento","alocacao","motivo"],
};
const MODELOS_EX = {
  rf:   [["substituir","LCI curto prazo","Tesouro IPCA+ 2035","IPCA","alta","Proteger contra inflação persistente","conservadora;moderada"],
         ["alocar","","CDB IPCA+5% 2027","IPCA","media","Duration adequada ao perfil","moderada;arrojada"]],
  rv:   [["trocar","PETR4","PRIO3","Melhor alocação de capital e gestão","BTG Pactual","arrojada;agressiva"],
         ["comprar","","WEGE3","Crescimento sólido, recomendação XP","XP Research","moderada;arrojada"]],
  intl: [["VTI","ETF","Diversificação em mercado americano amplo","moderada;arrojada;agressiva"],
         ["IVVB11","BDR","S&P 500 com liquidez em BRL","conservadora;moderada"]],
  fii:  [["MXRF11","Papel","0.5%","Alta distribuição, baixa volatilidade"],
         ["HGLG11","Logística","1.0%","Contratos longos, boa gestão"],
         ["XPML11","Shopping","0.8%","Recuperação do varejo físico"]],
};

// Estado local dos itens de sugestão
const sgData = {rf:[], rv:[], intl:[]};

// ── Inicialização ─────────────────────────────────────────────────────────────
async function init(){
  try{
    const info = await fetch("/api/admin/contexto-info").then(r=>r.json());

    // Status carta
    const ctxEl = document.getElementById("status-contexto");
    if(info.carta && info.carta.nome){
      ctxEl.innerHTML = `<div class="info-box"><b>Carta ativa:</b> ${info.carta.nome}<br><span style="color:#5DCAA5">Atualizado: ${info.carta.atualizado}</span></div>`;
    } else {
      ctxEl.innerHTML = `<div class="info-box" style="color:#1E4A30">Nenhuma carta carregada.</div>`;
    }

    // Status comunicado
    const msgEl = document.getElementById("status-msg");
    if(info.mensagem && info.mensagem.atualizado){
      document.getElementById("txt-estrategia").value = info.mensagem.estrategia || "";
      document.getElementById("txt-mensagem").value   = info.mensagem.mensagem   || "";
      msgEl.innerHTML = `<div class="info-box"><b>Publicado em:</b> ${info.mensagem.atualizado}</div>`;
    } else {
      msgEl.innerHTML = `<div class="info-box" style="color:#1E4A30">Nenhum comunicado publicado.</div>`;
    }
  } catch(e){
    document.getElementById("status-contexto").innerHTML = '<div class="info-box" style="color:#3A6A48">Erro ao carregar status.</div>';
  }
}

// ── Carta da Gestão ───────────────────────────────────────────────────────────
function onCartaSelect(){
  const f = document.getElementById("pdf-carta").files[0];
  document.getElementById("fname-carta").textContent = f ? f.name : "";
}

async function uploadCarta(){
  const f = document.getElementById("pdf-carta").files[0];
  if(!f){ alert("Selecione um PDF primeiro."); return; }
  const btn = document.getElementById("btn-carta");
  const st  = document.getElementById("st-carta");
  btn.disabled = true; btn.textContent = "Processando..."; st.textContent = "";
  try{
    const fd = new FormData(); fd.append("carta", f);
    const r = await fetch("/api/admin/upload-contexto", {method:"POST", body:fd});
    const d = await r.json();
    if(d.ok){
      st.innerHTML = `<span class="status-ok">✓ ${d.chars} caracteres extraídos. Carta ativa.</span>`;
      document.getElementById("carta-preview").style.display = "block";
      document.getElementById("carta-info").textContent = d.preview || "";
      init();
    } else {
      st.innerHTML = `<span class="status-err">Erro: ${d.error}</span>`;
    }
  } catch(e){ st.innerHTML = `<span class="status-err">Erro ao processar.</span>`; }
  finally{ btn.disabled=false; btn.textContent="Processar carta"; }
}

// ── Comunicado ────────────────────────────────────────────────────────────────
async function salvarMensagem(){
  const payload = {
    estrategia: document.getElementById("txt-estrategia").value,
    mensagem:   document.getElementById("txt-mensagem").value,
  };
  const st = document.getElementById("st-msg");
  try{
    await fetch("/api/admin/mensagem", {method:"POST", headers:{"Content-Type":"application/json"}, body:JSON.stringify(payload)});
    st.innerHTML = '<span class="status-ok">✓ Publicado!</span>';
    init();
    setTimeout(()=>st.textContent="", 3000);
  } catch(e){ st.innerHTML = '<span class="status-err">Erro ao salvar.</span>'; }
}

// ── Tabs ──────────────────────────────────────────────────────────────────────
function sgTab(t, btn){
  document.querySelectorAll(".sg-tab").forEach(b=>b.classList.remove("active"));
  if(btn) btn.classList.add("active");
  document.querySelectorAll(".sg-panel").forEach(p=>p.classList.remove("active"));
  document.getElementById("sg-panel-"+t).classList.add("active");
}

// ── PDF por categoria ─────────────────────────────────────────────────────────
async function uploadPdfCategoria(cat, input){
  const f = input.files[0]; if(!f) return;
  const st = document.getElementById("import-st-"+cat);
  st.textContent = "⏳ Extraindo PDF...";
  input.value = "";
  try{
    const fd = new FormData();
    fd.append("arquivo", f);
    fd.append("categoria", cat);
    const r = await fetch("/api/admin/upload-sugestao-doc", {method:"POST", body:fd});
    const d = await r.json();
    if(!d.ok){ st.innerHTML=`<span class="status-err">Erro: ${d.error||"falha"}</span>`; return; }
    const texto = d.texto || "";
    document.getElementById("pdf-texto-"+cat).value = texto;
    document.getElementById("pdf-preview-"+cat).style.display = "block";
    st.innerHTML = `<span class="status-ok">✓ ${d.chars||texto.length} caracteres extraídos</span>`;
    setTimeout(()=>st.textContent="", 4000);
  } catch(e){ st.innerHTML=`<span class="status-err">Erro: ${e.message}</span>`; }
}

function fecharPdfPreview(cat){
  document.getElementById("pdf-preview-"+cat).style.display = "none";
}

function copiarPdf(cat){
  const txt = document.getElementById("pdf-texto-"+cat).value;
  navigator.clipboard.writeText(txt).catch(()=>{});
}

function usarPdfComoSugestao(cat){
  const txt = document.getElementById("pdf-texto-"+cat).value.trim();
  if(!txt) return;
  if(cat === "rf"){
    // Adiciona item RF com motivo preenchido do PDF
    sgAdd("rf");
    const id = sgData.rf[sgData.rf.length-1];
    if(id) setVal("rf-mot", id, txt.substring(0,300));
  } else if(cat === "rv"){
    sgAdd("rv");
    const id = sgData.rv[sgData.rv.length-1];
    if(id) setVal("rv-mot", id, txt.substring(0,300));
  } else if(cat === "intl"){
    sgAdd("intl");
    const id = sgData.intl[sgData.intl.length-1];
    if(id) setVal("intl-mot", id, txt.substring(0,300));
  } else if(cat === "fii"){
    const cur = document.getElementById("sg-fii").value.trim();
    document.getElementById("sg-fii").value = cur ? cur+"\n"+txt : txt;
  } else if(cat === "est"){
    const cur = document.getElementById("sg-est").value.trim();
    document.getElementById("sg-est").value = cur ? cur+"\n\n"+txt : txt;
  }
  fecharPdfPreview(cat);
}

// ── Itens de sugestão ─────────────────────────────────────────────────────────
function perfisHtml(nome){
  return `<div class="perfis-row">`+PERFIS.map(p=>
    `<label><input type="checkbox" name="${nome}" value="${p}" checked> ${PERFIS_LBL[p]}</label>`
  ).join("")+`</div>`;
}

function sgAdd(tipo){
  const id = Date.now();
  let html = "";
  if(tipo==="rf"){
    html=`<div class="sg-item" id="sgitem-rf-${id}">
      <button class="sg-del" onclick="sgDel('rf','${id}')">✕</button>
      <div class="sg-grid-2">
        <div class="sg-mini"><label>Ação</label><select id="rf-acao-${id}"><option value="alocar">Alocar novo</option><option value="substituir">Substituir</option></select></div>
        <div class="sg-mini"><label>Urgência</label><select id="rf-urg-${id}"><option value="alta">🔴 Alta</option><option value="media">🟡 Média</option><option value="baixa">🟢 Baixa</option></select></div>
      </div>
      <div class="sg-grid-2">
        <div class="sg-mini"><label>Sair de (opcional)</label><input type="text" id="rf-de-${id}" placeholder="Ex: LCI curto prazo"></div>
        <div class="sg-mini"><label>Produto recomendado</label><input type="text" id="rf-para-${id}" placeholder="Ex: Tesouro IPCA+ 2035"></div>
      </div>
      <div class="sg-grid-2" style="margin-bottom:8px">
        <div class="sg-mini"><label>Indexador</label><select id="rf-idx-${id}">${INDEXADORES.map(x=>`<option>${x}</option>`).join("")}</select></div>
        <div class="sg-mini"><label>Motivo</label><input type="text" id="rf-mot-${id}" placeholder="Ex: Proteger contra inflação persistente"></div>
      </div>
      <div class="sg-mini"><label>Perfis que recebem esta sugestão</label>${perfisHtml("rf-perf-"+id)}</div>
    </div>`;
    sgData.rf.push(id);
  } else if(tipo==="rv"){
    html=`<div class="sg-item" id="sgitem-rv-${id}">
      <button class="sg-del" onclick="sgDel('rv','${id}')">✕</button>
      <div class="sg-grid-3">
        <div class="sg-mini"><label>Ação</label><select id="rv-acao-${id}"><option value="comprar">Comprar</option><option value="vender">Vender</option><option value="trocar">Trocar</option></select></div>
        <div class="sg-mini"><label>Ticker / Sair de</label><input type="text" id="rv-de-${id}" placeholder="Ex: PETR4"></div>
        <div class="sg-mini"><label>Entrar em (opcional)</label><input type="text" id="rv-para-${id}" placeholder="Ex: PRIO3"></div>
      </div>
      <div class="sg-grid-2" style="margin-bottom:8px">
        <div class="sg-mini"><label>Motivo / Tese</label><input type="text" id="rv-mot-${id}" placeholder="Ex: Melhor gestão de capital, recomendação BTG"></div>
        <div class="sg-mini"><label>Fonte</label><select id="rv-fonte-${id}">${FONTES_RV.map(x=>`<option>${x}</option>`).join("")}</select></div>
      </div>
      <div class="sg-mini"><label>Perfis</label>${perfisHtml("rv-perf-"+id)}</div>
    </div>`;
    sgData.rv.push(id);
  } else if(tipo==="intl"){
    html=`<div class="sg-item" id="sgitem-intl-${id}">
      <button class="sg-del" onclick="sgDel('intl','${id}')">✕</button>
      <div class="sg-grid-3" style="margin-bottom:8px">
        <div class="sg-mini"><label>Ativo (ticker)</label><input type="text" id="intl-ativo-${id}" placeholder="Ex: VTI, IVVB11"></div>
        <div class="sg-mini"><label>Tipo</label><select id="intl-tipo-${id}"><option>ETF</option><option>BDR</option><option>Ação</option><option>Fundo</option><option>Outro</option></select></div>
        <div class="sg-mini"><label>Motivo</label><input type="text" id="intl-mot-${id}" placeholder="Ex: Diversificação em tech americana"></div>
      </div>
      <div class="sg-mini"><label>Perfis</label>${perfisHtml("intl-perf-"+id)}</div>
    </div>`;
    sgData.intl.push(id);
  }
  document.getElementById("sg-list-"+tipo).insertAdjacentHTML("beforeend", html);
}

function sgDel(tipo, id){
  const el = document.getElementById("sgitem-"+tipo+"-"+id);
  if(el) el.remove();
  sgData[tipo] = sgData[tipo].filter(x=>String(x)!==String(id));
}

// ── Coleta e publicação ───────────────────────────────────────────────────────
function getPerfis(name){
  return [...document.querySelectorAll(`input[name="${name}"]:checked`)].map(i=>i.value);
}
function setVal(prefix, id, val){
  const el = document.getElementById(prefix+"-"+id); if(el) el.value = val;
}
function setSelect(prefix, id, val){
  const el = document.getElementById(prefix+"-"+id);
  if(!el) return;
  [...el.options].forEach(o=>{ o.selected = o.value===val||o.text===val; });
}
function setPerfis(name, lista){
  document.querySelectorAll(`input[name="${name}"]`).forEach(cb=>{
    cb.checked = lista.length===0 || lista.includes(cb.value);
  });
}

async function publicarSugestoes(){
  const titulo = document.getElementById("sg-titulo").value.trim();
  if(!titulo){ alert("Digite um título para a publicação."); return; }

  const rf = sgData.rf.map(id=>({
    acao:     document.getElementById(`rf-acao-${id}`)?.value,
    urgencia: document.getElementById(`rf-urg-${id}`)?.value,
    de:       document.getElementById(`rf-de-${id}`)?.value,
    para:     document.getElementById(`rf-para-${id}`)?.value,
    indexador:document.getElementById(`rf-idx-${id}`)?.value,
    motivo:   document.getElementById(`rf-mot-${id}`)?.value,
    perfis:   getPerfis(`rf-perf-${id}`),
  })).filter(i=>i.para);

  const rv = sgData.rv.map(id=>({
    acao:   document.getElementById(`rv-acao-${id}`)?.value,
    de:     document.getElementById(`rv-de-${id}`)?.value,
    para:   document.getElementById(`rv-para-${id}`)?.value,
    motivo: document.getElementById(`rv-mot-${id}`)?.value,
    fonte:  document.getElementById(`rv-fonte-${id}`)?.value,
    perfis: getPerfis(`rv-perf-${id}`),
  })).filter(i=>i.de||i.para);

  const intl = sgData.intl.map(id=>({
    ativo:  document.getElementById(`intl-ativo-${id}`)?.value,
    tipo:   document.getElementById(`intl-tipo-${id}`)?.value,
    motivo: document.getElementById(`intl-mot-${id}`)?.value,
    perfis: getPerfis(`intl-perf-${id}`),
  })).filter(i=>i.ativo);

  const fiiText = document.getElementById("sg-fii").value;
  const fiis = fiiText.split("\n").map(l=>l.trim()).filter(Boolean).map(l=>{
    const p = l.split("|").map(x=>x.trim());
    return {ticker:p[0]||"", segmento:p[1]||"", alocacao:p[2]||"", motivo:p[3]||""};
  }).filter(i=>i.ticker);

  const payload = {
    titulo, rf, rv, intl,
    estruturadas: document.getElementById("sg-est").value,
    internacional: intl,
    fiis,
  };

  const st = document.getElementById("sg-status");
  st.textContent = "Publicando...";
  try{
    const r = await fetch("/api/sugestoes", {method:"POST", headers:{"Content-Type":"application/json"}, body:JSON.stringify(payload)});
    const d = await r.json();
    if(d.ok){
      st.innerHTML = '<span class="status-ok">✓ Publicado! Visível para todos os assessores.</span>';
      carregarHistorico();
    } else { st.innerHTML='<span class="status-err">Erro ao publicar.</span>'; }
  } catch(e){ st.innerHTML='<span class="status-err">Erro de conexão.</span>'; }
  setTimeout(()=>st.textContent="", 5000);
}

function limparForm(){
  if(!confirm("Limpar todos os campos de sugestão?")) return;
  document.getElementById("sg-titulo").value = "";
  ["rf","rv","intl"].forEach(t=>{ document.getElementById("sg-list-"+t).innerHTML=""; sgData[t]=[]; });
  document.getElementById("sg-fii").value = "";
  document.getElementById("sg-est").value = "";
  ["rf","rv","fii","intl","est"].forEach(c=>fecharPdfPreview(c));
}

// ── Histórico ─────────────────────────────────────────────────────────────────
async function carregarHistorico(){
  const data = await fetch("/api/sugestoes").then(r=>r.json()).catch(()=>({historico:[]}));
  const el   = document.getElementById("sg-historico");
  const hist = data.historico || [];
  if(!hist.length){ el.innerHTML='<p style="color:#1E4A30;font-size:12px;font-style:italic">Nenhuma publicação ainda.</p>'; return; }
  el.innerHTML = hist.map(s=>{
    const total = (s.rf?.length||0)+(s.rv?.length||0)+(s.intl?.length||0)+(s.internacional?.length||0)+(s.fiis?.length||0);
    const ativaHtml = s.ativa
      ? `<span class="tag-ativa">● ATIVA</span>`
      : `<span style="font-size:10px;color:#1C4A34">inativa</span>`;
    return `<div class="hist-item ${s.ativa?'ativa':''}">
      <div style="flex:1">
        <div style="font-size:13px;font-weight:700;color:${s.ativa?'#5DCAA5':'#3A6A48'}">${s.titulo||"—"}</div>
        <div style="font-size:10px;color:#2A5A3A;margin-top:3px">${s.criado_em||""} · RF:${s.rf?.length||0} RV:${s.rv?.length||0} FIIs:${s.fiis?.length||0} Intl:${s.internacional?.length||s.intl?.length||0} · ${total} sugestões</div>
      </div>
      ${ativaHtml}
    </div>`;
  }).join("");
}

// ── Importação planilhas ──────────────────────────────────────────────────────
function lerPlanilha(file, callback){
  const reader = new FileReader();
  const isXls  = /\.(xlsx|xls)$/i.test(file.name);
  reader.onload = function(e){
    let rows = [];
    if(isXls){
      if(typeof XLSX==="undefined"){ alert("Biblioteca XLSX ainda carregando. Aguarde 2s e tente novamente."); return; }
      const wb = XLSX.read(new Uint8Array(e.target.result), {type:"array"});
      const ws = wb.Sheets[wb.SheetNames[0]];
      rows = XLSX.utils.sheet_to_json(ws, {header:1, defval:""});
    } else {
      const txt = e.target.result;
      const sep = (txt.split(";").length > txt.split(",").length) ? ";" : ",";
      rows = txt.split("\n").map(l=>l.split(sep).map(c=>c.replace(/^"|"$/g,"").trim()));
    }
    if(rows.length>1) rows = rows.slice(1);
    callback(rows.filter(r=>r.some(c=>c)));
  };
  isXls ? reader.readAsArrayBuffer(file) : reader.readAsText(file,"UTF-8");
}

function importarPlanilha(tipo, input){
  const file = input.files[0]; if(!file) return;
  const st = document.getElementById("import-st-"+tipo);
  lerPlanilha(file, function(rows){
    let count = 0;
    rows.forEach(function(r){
      if(!r[0]) return;
      sgAdd(tipo);
      const id = sgData[tipo][sgData[tipo].length-1];
      if(!id) return;
      count++;
      if(tipo==="rf"){
        setSelect("rf-acao",id,r[0]||"alocar");
        setVal("rf-de",  id, r[1]||"");
        setVal("rf-para",id, r[2]||"");
        setSelect("rf-idx",id,r[3]||"IPCA");
        setSelect("rf-urg",id,r[4]||"media");
        setVal("rf-mot", id, r[5]||"");
        setPerfis("rf-perf-"+id, (r[6]||"").split(";").map(s=>s.trim()).filter(Boolean));
      } else if(tipo==="rv"){
        setSelect("rv-acao", id, r[0]||"trocar");
        setVal("rv-de",   id, r[1]||"");
        setVal("rv-para", id, r[2]||"");
        setVal("rv-mot",  id, r[3]||"");
        setSelect("rv-fonte",id,r[4]||"BTG Pactual");
        setPerfis("rv-perf-"+id, (r[5]||"").split(";").map(s=>s.trim()).filter(Boolean));
      } else if(tipo==="intl"){
        setVal("intl-ativo",id,r[0]||"");
        setSelect("intl-tipo",id,r[1]||"ETF");
        setVal("intl-mot",  id,r[2]||"");
        setPerfis("intl-perf-"+id,(r[3]||"").split(";").map(s=>s.trim()).filter(Boolean));
      }
    });
    if(st) st.innerHTML = `<span class="status-ok">✓ ${count} importados</span>`;
    setTimeout(()=>{ if(st) st.textContent=""; },4000);
    input.value="";
  });
}

function importarFii(input){
  const file = input.files[0]; if(!file) return;
  const st = document.getElementById("import-st-fii");
  lerPlanilha(file, function(rows){
    const linhas = rows.filter(r=>r[0]).map(r=>r.slice(0,4).join(" | "));
    const cur = document.getElementById("sg-fii").value.trim();
    document.getElementById("sg-fii").value = cur ? cur+"\n"+linhas.join("\n") : linhas.join("\n");
    if(st) st.innerHTML = `<span class="status-ok">✓ ${rows.length} fundos importados</span>`;
    setTimeout(()=>{ if(st) st.textContent=""; },4000);
    input.value="";
  });
}

function baixarModelo(tipo){
  const cols  = MODELOS_CSV[tipo];
  const exs   = MODELOS_EX[tipo];
  const linhas= [cols.join(","), ...exs.map(r=>r.map(c=>/,/.test(c)?`"${c}"`:c).join(","))];
  const blob  = new Blob(["﻿"+linhas.join("\r\n")], {type:"text/csv;charset=utf-8"});
  const a = document.createElement("a");
  a.href = URL.createObjectURL(blob);
  a.download = `modelo_brauna_${tipo}.csv`;
  a.click();
}

init();
carregarHistorico();
</script>

</div><!-- /tab-config -->

</div><!-- /container -->

<script>
// ─── Tab switching (main tabs) ────────────────────────────────────────────────
function switchTab(name, btn){
  document.querySelectorAll(".main-panel").forEach(p=>p.classList.remove("active"));
  document.querySelectorAll(".main-tab").forEach(b=>b.classList.remove("active"));
  const panel = document.getElementById("tab-"+name);
  if(panel) panel.classList.add("active");
  if(btn)   btn.classList.add("active");
  if(name==="assessores"||name==="lideres"||name==="head") carregarDashboard();
}

// ─── Dashboard data ───────────────────────────────────────────────────────────
let _dashData = null;

async function carregarDashboard(){
  try{
    const r = await fetch("/api/admin/dashboard");
    _dashData = await r.json();
    renderAssessores(_dashData.assessores||[]);
    renderLideres(_dashData.lideres||[]);
    renderHead(_dashData.head||{});
  } catch(e){ console.warn("Dashboard error", e); }
}

// ─── Assessores ───────────────────────────────────────────────────────────────
function renderAssessores(assessores){
  const totalAnalises = assessores.reduce((s,a)=>s+(a.acoes||[]).filter(x=>x.acao==="análise").length,0);
  const clientes = new Set();
  assessores.forEach(a=>(a.acoes||[]).filter(x=>x.acao==="análise").forEach(x=>{
    const m = (x.detalhe||"").match(/Analisou carteira de (.+?) —/);
    if(m) clientes.add(m[1]);
  }));
  const ultimo = assessores.length ? (assessores[0].ultimo||"—") : "—";
  document.getElementById("kpi-ativos").textContent    = assessores.length;
  document.getElementById("kpi-analises").textContent  = totalAnalises;
  document.getElementById("kpi-clientes").textContent  = clientes.size;
  document.getElementById("kpi-ultimo-acesso").textContent = ultimo;

  // Cards
  const lista = document.getElementById("lista-assessores");
  if(!assessores.length){
    lista.innerHTML = '<div style="color:#1E4A30;font-size:12px;padding:20px;text-align:center">Nenhum assessor registrado ainda.</div>';
  } else {
    const colors = ["#5DCAA5","#8B9FE8","#C9A96E","#FF8B6B","#A0D8A0"];
    lista.innerHTML = assessores.map((a,i)=>{
      const analises = (a.acoes||[]).filter(x=>x.acao==="análise");
      const logHtml  = (a.acoes||[]).slice(0,5).map(x=>
        `<div class="log-row"><span class="log-ts">${x.ts||""}</span><span>${x.acao}: ${x.detalhe||""}</span></div>`
      ).join("");
      return `<div class="user-card">
        <div class="user-avatar" style="background:${colors[i%colors.length]}20;color:${colors[i%colors.length]}">${(a.nome||"?")[0].toUpperCase()}</div>
        <div class="user-info">
          <div class="user-name">${a.nome||"Desconhecido"}</div>
          <div class="user-meta">Último acesso: ${a.ultimo||"—"} · ${(a.acoes||[]).length} ações · ${analises.length} análises</div>
          <div class="user-log">${logHtml||'<div class="log-row"><span style="color:#1E4A30">Sem atividade registrada</span></div>'}</div>
        </div>
        <span class="badge badge-assessor">Assessor</span>
      </div>`;
    }).join("");
  }

  // Log table
  const activity = _dashData?._raw_activity || [];
  const rows = activity.filter(e=>e.role==="assessor").slice(0,50);
  document.getElementById("log-assessores-body").innerHTML = rows.length
    ? rows.map(e=>`<tr><td style="color:#2A5A3A">${e.ts||""}</td><td style="color:#F0F0F0;font-weight:600">${e.nome||""}</td><td><span class="badge badge-${e.acao==="análise"?"warn":"ok"}">${e.acao||""}</span></td><td style="color:#888">${e.detalhe||""}</td></tr>`).join("")
    : `<tr><td colspan="4" style="color:#1E4A30;text-align:center;padding:20px">Sem atividade registrada.</td></tr>`;
}

// ─── Líderes ──────────────────────────────────────────────────────────────────
function renderLideres(lideres){
  const totalAcoes = lideres.reduce((s,l)=>s+(l.acoes||[]).length,0);
  const ultimo = lideres.length ? (lideres[0].ultimo||"—") : "—";
  document.getElementById("kpi-lider-ativos").textContent = lideres.length;
  document.getElementById("kpi-lider-acoes").textContent  = totalAcoes;
  document.getElementById("kpi-lider-ultimo").textContent = ultimo;

  const lista = document.getElementById("lista-lideres");
  if(!lideres.length){
    lista.innerHTML = '<div style="color:#1E4A30;font-size:12px;padding:20px;text-align:center">Nenhum líder registrado ainda.</div>';
  } else {
    lista.innerHTML = lideres.map(l=>{
      const logHtml = (l.acoes||[]).slice(0,5).map(x=>
        `<div class="log-row"><span class="log-ts">${x.ts||""}</span><span>${x.acao}: ${x.detalhe||""}</span></div>`
      ).join("");
      return `<div class="user-card">
        <div class="user-avatar" style="background:#8B9FE820;color:#8B9FE8">${(l.nome||"?")[0].toUpperCase()}</div>
        <div class="user-info">
          <div class="user-name">${l.nome||"Desconhecido"}</div>
          <div class="user-meta">Último acesso: ${l.ultimo||"—"} · ${(l.acoes||[]).length} ações</div>
          <div class="user-log">${logHtml||'<div class="log-row"><span style="color:#1E4A30">Sem atividade</span></div>'}</div>
        </div>
        <span class="badge badge-lider">Líder</span>
      </div>`;
    }).join("");
  }

  const activity = _dashData?._raw_activity || [];
  const rows = activity.filter(e=>e.role==="lider"||e.role==="líder").slice(0,50);
  document.getElementById("log-lideres-body").innerHTML = rows.length
    ? rows.map(e=>`<tr><td style="color:#2A5A3A">${e.ts||""}</td><td style="color:#F0F0F0;font-weight:600">${e.nome||""}</td><td><span class="badge badge-lider">${e.acao||""}</span></td><td style="color:#888">${e.detalhe||""}</td></tr>`).join("")
    : `<tr><td colspan="4" style="color:#1E4A30;text-align:center;padding:20px">Sem atividade registrada.</td></tr>`;
}

// ─── Head ─────────────────────────────────────────────────────────────────────
function renderHead(h){
  document.getElementById("kpi-docs").textContent    = h.total_docs_base   ?? "—";
  document.getElementById("kpi-calls").textContent   = h.total_calls        ?? "—";
  document.getElementById("kpi-estr").textContent    = h.total_estruturadas ?? "—";
  document.getElementById("kpi-prod").textContent    = h.total_produtos     ?? "—";
  document.getElementById("kpi-alertas").textContent = (h.alertas||[]).length;
  document.getElementById("kpi-gestoras").textContent= (h.gestoras2||[]).length;
  document.getElementById("kpi-clientes").textContent= h.total_clientes ?? "—";

  // Cenário
  const c = h.cenario || {};
  const cenRef = c.referencia ? `<div style="font-size:10px;color:#3A6A48;margin-bottom:8px">📅 Referência: <b style="color:#C9A96E">${c.referencia}</b></div>` : "";
  const cenHtml = c.global || c.brasil || c.posicionamento
    ? `${cenRef}
      <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:10px">
        <div><div style="font-size:10px;color:#3A6A48;margin-bottom:4px;font-weight:700">🌐 GLOBAL</div><div style="font-size:11px;color:#CCC;line-height:1.6">${(c.global||"").substring(0,200)}${c.global_chars>200?"…":""}</div></div>
        <div><div style="font-size:10px;color:#3A6A48;margin-bottom:4px;font-weight:700">🇧🇷 BRASIL</div><div style="font-size:11px;color:#CCC;line-height:1.6">${(c.brasil||"").substring(0,200)}${c.brasil_chars>200?"…":""}</div></div>
        <div><div style="font-size:10px;color:#3A6A48;margin-bottom:4px;font-weight:700">🎯 POSICIONAMENTO</div><div style="font-size:11px;color:#CCC;line-height:1.6">${(c.posicionamento||"").substring(0,200)}${c.pos_chars>200?"…":""}</div></div>
      </div>`
    : `<div style="color:#1E4A30;font-size:12px;font-style:italic">Nenhum cenário publicado ainda.</div>`;
  document.getElementById("head-cenario").innerHTML = cenHtml;

  // Base de Conhecimento
  const docs = h.knowledge || [];
  document.getElementById("head-docs").innerHTML = docs.length
    ? docs.map(d=>`<div class="doc-row"><span class="doc-icon">📄</span><span class="doc-name" title="${d.nome||""}">${d.nome||"Documento"}</span><span class="doc-meta">${d.data||""}</span></div>`).join("")
    : `<div style="color:#1E4A30;font-size:12px;font-style:italic;padding:8px">Nenhum documento na base.</div>`;

  // Calls
  const calls = h.calls || [];
  document.getElementById("head-calls").innerHTML = calls.length
    ? calls.map(c=>`<div class="doc-row"><span class="doc-icon">📈</span><span class="doc-name">${c.ativo||c.nome||"—"} <span style="font-size:10px;color:#888">${c.acao||""}</span></span><span class="doc-meta">${c.data||""}</span></div>`).join("")
    : `<div style="color:#1E4A30;font-size:12px;font-style:italic;padding:8px">Nenhuma call ativa.</div>`;

  // Estruturadas
  const estr = h.estruturadas || [];
  document.getElementById("head-estruturadas").innerHTML = estr.length
    ? estr.map(e=>`<div class="doc-row"><span class="doc-icon">💼</span><span class="doc-name">${e.ativo||"—"} <span style="font-size:10px;color:#888">${e.tipo||""}</span></span><span class="doc-meta" style="color:#C9A96E">${e.retorno||""}</span></div>`).join("")
    : `<div style="color:#1E4A30;font-size:12px;font-style:italic;padding:8px">Nenhuma operação cadastrada.</div>`;

  // Gestoras
  const gest = h.gestoras2 || [];
  document.getElementById("head-gestoras").innerHTML = gest.length
    ? gest.map(g=>`<div class="doc-row"><span class="doc-icon">🏦</span><span class="doc-name">${g.nome||"—"}</span><span class="doc-meta">${g.referencia||""}</span></div>`).join("")
    : `<div style="color:#1E4A30;font-size:12px;font-style:italic;padding:8px">Nenhuma gestora cadastrada.</div>`;

  // Alertas
  const alertas = h.alertas || [];
  document.getElementById("head-alertas").innerHTML = alertas.length
    ? alertas.map(a=>`<div class="doc-row"><span class="doc-icon">🚨</span><span class="doc-name">${a.titulo||a.texto||"Alerta"}</span><span class="doc-meta">${a.data||""}</span></div>`).join("")
    : `<div style="color:#1E4A30;font-size:12px;font-style:italic;padding:8px">Nenhum alerta ativo.</div>`;

  // Clientes registrados
  const clientes = h.clientes || [];
  const fmtPat = v => v >= 1e6 ? `R$ ${(v/1e6).toFixed(1)}M` : v >= 1e3 ? `R$ ${(v/1e3).toFixed(0)}k` : `R$ ${v}`;
  document.getElementById("head-clientes").innerHTML = clientes.length
    ? `<div style="overflow-x:auto"><table class="act-table">
        <thead><tr><th>Conta</th><th>Assessor</th><th>Perfil</th><th>Patrimônio</th><th>Principais ativos</th><th>Atualizado</th></tr></thead>
        <tbody>${clientes.map(c=>{
          const ativos = [...(c.acoes||[]).slice(0,3).map(a=>a.ticker), ...(c.fiis||[]).slice(0,2).map(f=>f.ticker)].filter(Boolean).join(", ");
          return `<tr>
            <td style="color:#C9A96E;font-weight:600">${c.conta||"—"}</td>
            <td style="color:#AAA">${c.assessor||"—"}</td>
            <td><span class="badge badge-assessor">${c.perfil||"—"}</span></td>
            <td style="color:#5DCAA5;font-weight:700">${fmtPat(c.patrimonio||0)}</td>
            <td style="color:#888;font-size:11px">${ativos||"—"}</td>
            <td style="color:#2A5A3A;font-size:10px">${c.atualizado_em||"—"}</td>
          </tr>`;
        }).join("")}</tbody>
      </table></div>`
    : `<div style="color:#1E4A30;font-size:12px;font-style:italic;padding:8px">Nenhuma carteira registrada ainda.</div>`;
}

// ─── Load on start ────────────────────────────────────────────────────────────
carregarDashboard();
</script>

</body></html>"""

# ═══════════════════════════════════════════════════════════════════════════════
#  HEAD DE PRODUTOS
# ═══════════════════════════════════════════════════════════════════════════════
HTML_HEAD = r"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Braúna — Head de Produtos</title>
<style>
*{margin:0;padding:0;box-sizing:border-box}
body{background:#071E17;color:#F0F0F0;font-family:'Segoe UI',system-ui,sans-serif;min-height:100vh}
header{background:#0F0F0A;border-bottom:1px solid #2A2A18;padding:14px 28px;display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:10px}
header h1{font-size:17px;color:#D4B483;font-weight:700}
header p{font-size:11px;color:#2A5A3A;margin-top:2px}
.nav{display:flex;gap:8px;flex-wrap:wrap}
.nav a,.nav button{font-size:12px;padding:5px 12px;border-radius:6px;border:1px solid #2A2A18;color:#4A7055;text-decoration:none;background:none;cursor:pointer;transition:all .2s;font-family:inherit}
.nav a:hover,.nav button:hover{border-color:#D4B483;color:#D4B483}
.nav a.active{background:#D4B483;color:#000;border-color:#D4B483;font-weight:700}
.container{max-width:1200px;margin:0 auto;padding:24px 20px}
.card{background:#0F0F0A;border:1px solid #2A2A18;border-radius:12px;padding:22px;margin-bottom:18px}
.card-title{font-size:11px;color:#D4B483;font-weight:700;text-transform:uppercase;letter-spacing:.8px;margin-bottom:16px;display:flex;align-items:center;gap:8px}
.card-title span{font-size:16px}
/* Tabs */
.tabs{display:flex;gap:6px;margin-bottom:16px;flex-wrap:wrap}
.tab{padding:7px 16px;border-radius:8px;border:1px solid #2A2A18;background:#060F0B;color:#4A7055;font-size:12px;cursor:pointer;transition:all .2s;font-family:inherit}
.tab.active{background:#D4B483;color:#000;border-color:#D4B483;font-weight:700}
.tab-panel{display:none}.tab-panel.active{display:block}
/* Portfólio modelo — tabela */
.porto-table{width:100%;border-collapse:collapse;font-size:12px}
.porto-table th{padding:8px 10px;text-align:center;font-size:10px;font-weight:700;text-transform:uppercase;letter-spacing:.5px;border-bottom:1px solid #2A2A18;color:#4A7055}
.porto-table th.classe{text-align:left;color:#D4B483;width:140px}
.porto-table td{padding:6px 8px;border-bottom:1px solid #1A1A10;text-align:center}
.porto-table td.classe-label{text-align:left;font-size:11px;color:#CCC;padding-left:4px}
.porto-table td input[type=number]{width:70px;background:#060F0B;border:1px solid #2A2A18;border-radius:5px;padding:5px 6px;color:#F0F0F0;font-size:12px;text-align:center;outline:none}
.porto-table td input[type=number]:focus{border-color:#D4B483}
.porto-table .total-row td{font-weight:700;font-size:12px;border-top:2px solid #2A2A18;padding-top:8px}
.total-val{color:#D4B483}
.total-ok{color:#5DCAA5}.total-err{color:#FF6B6B}
/* Viés badges */
.vies-select{background:#060F0B;border:1px solid #2A2A18;border-radius:5px;padding:4px 6px;color:#CCC;font-size:11px;outline:none;cursor:pointer}
.vies-select:focus{border-color:#D4B483}
/* Inputs gerais */
label{font-size:11px;color:#4A7055;display:block;margin-bottom:5px}
input[type=text],textarea,select:not(.vies-select){width:100%;background:#060F0B;border:1px solid #2A2A18;border-radius:7px;padding:8px 10px;color:#F0F0F0;font-size:13px;outline:none;font-family:inherit}
input[type=text]:focus,textarea:focus{border-color:#D4B483}
textarea{resize:vertical}
/* Produto item */
.prod-item{background:#060F0B;border:1px solid #2A2A18;border-radius:8px;padding:12px 14px;margin-bottom:8px;position:relative;animation:fadeIn .15s ease}
@keyframes fadeIn{from{opacity:0;transform:translateY(-4px)}to{opacity:1;transform:translateY(0)}}
.prod-del{position:absolute;top:8px;right:8px;background:none;border:none;color:#2A2A18;cursor:pointer;font-size:16px;padding:2px 6px;transition:color .2s;border-radius:4px}
.prod-del:hover{color:#FF6B6B}
.prod-grid-3{display:grid;grid-template-columns:1fr 1fr 1fr;gap:8px;margin-bottom:8px}
.prod-grid-2{display:grid;grid-template-columns:1fr 1fr;gap:8px}
.prod-mini label{font-size:10px;color:#3A6A48;margin-bottom:3px}
/* Botões */
.btn{display:inline-flex;align-items:center;gap:6px;background:#D4B483;color:#000;border:none;border-radius:7px;padding:10px 20px;font-size:13px;font-weight:700;cursor:pointer;transition:all .2s;font-family:inherit}
.btn:hover{opacity:.88;transform:translateY(-1px)}
.btn:disabled{opacity:.35;cursor:not-allowed;transform:none}
.btn-out{background:transparent;color:#D4B483;border:1px solid #D4B483}
.btn-out:hover{background:#D4B483;color:#000}
.btn-sm{padding:6px 14px;font-size:12px}
.btn-ghost{background:transparent;border:1px solid #2A2A18;color:#3A6A48;border-radius:6px;padding:6px 12px;font-size:11px;cursor:pointer;transition:all .2s;font-family:inherit}
.btn-ghost:hover{border-color:#D4B483;color:#D4B483}
.btn-add{background:#1A1A08;color:#D4B483;border:1px solid #D4B483;border-radius:6px;padding:6px 14px;font-size:12px;cursor:pointer;transition:all .2s;margin-bottom:10px;font-family:inherit}
.btn-add:hover{background:#D4B483;color:#000}
/* Status */
.status-ok{color:#5DCAA5}.status-err{color:#FF6B6B}
.info-box{background:#060F0B;border-radius:8px;padding:12px 14px;font-size:12px;color:#3A6A48;line-height:1.6;border:1px solid #2A2A18}
.info-box b{color:#D4B483}
/* Publish bar */
.publish-bar{background:#1A1A08;border:2px solid #D4B483;border-radius:12px;padding:18px 22px;display:flex;align-items:center;gap:16px;flex-wrap:wrap;transition:border-color .3s,background .3s;position:relative;overflow:hidden}
.publish-bar.pendente{border-color:#FF4444 !important;background:#1A0808 !important;animation:piscar-pub 1.2s ease-in-out infinite}
@keyframes piscar-pub{0%,100%{border-color:#FF4444;box-shadow:0 0 0 0 #FF444400}50%{border-color:#FF8888;box-shadow:0 0 16px 4px #FF444455}}
.publish-bar .pub-info{flex:1}
.publish-bar .pub-info h3{font-size:14px;color:#D4B483;font-weight:700;transition:color .3s}
.publish-bar.pendente .pub-info h3{color:#FF6B6B}
.publish-bar .pub-info p{font-size:11px;color:#4A7055;margin-top:3px}
/* Referência */
.ref-row{display:flex;align-items:center;gap:12px;margin-bottom:16px;flex-wrap:wrap}
.ref-row input{max-width:320px;font-size:12px;padding:6px 10px}
.ref-badge{font-size:10px;padding:3px 10px;background:#1A1A08;color:#D4B483;border:1px solid #D4B483;border-radius:10px;white-space:nowrap}
/* Viés indicator */
.vies-pos{color:#5DCAA5}.vies-neg{color:#FF6B6B}.vies-neu{color:#4A7055}
@media(max-width:640px){.prod-grid-3,.prod-grid-2{grid-template-columns:1fr}}
</style>
</head>
<body>
<header>
  <div>
    <h1>🏛️ HEAD DE PRODUTOS — BRAÚNA</h1>
    <p>Portfólios modelo · Cenário macro · Produtos em destaque · Base de conhecimento</p>
  </div>
  <nav class="nav" id="nav-head">
    <a href="/head-produtos" class="active">🏛️ Head</a>
    <button onclick="sair()">Sair</button>
  </nav>
  <script>
  (function(){
    if(localStorage.getItem("brauna_role")==="admin"){
      const nav=document.getElementById("nav-head");
      nav.innerHTML=`<a href="/assessor">📊 Assessor</a><a href="/lider">👥 Líder</a><a href="/head-produtos" class="active">🏛️ Head</a><a href="/admin">⚙️ Admin</a>`;
    }
  })();
  </script>
</header>

<div class="container">

<!-- ── Status da publicação ── -->
<div id="status-pub" class="info-box" style="margin-bottom:18px;display:none"></div>

<!-- ══ 1. CENÁRIO MACRO ═══════════════════════════════════════════════════════ -->
<div class="card">
  <div class="card-title"><span>🌐</span> Cenário Macro</div>
  <div style="background:#060C08;border:1px solid #1A2E1A;border-radius:10px;padding:14px 16px;margin-bottom:16px;display:flex;align-items:center;gap:14px;flex-wrap:wrap">
    <span style="font-size:12px;color:#5DCAA5;font-weight:700">📩 Carta de Gestora</span>
    <span style="font-size:11px;color:#2A5A3A;flex:1">Suba a carta mensal — a IA extrai e preenche Global, Brasil e Posicionamento automaticamente</span>
    <input type="file" id="carta-input" accept=".pdf" style="display:none" onchange="uploadCartaCenario(this)">
    <button class="btn btn-sm" onclick="document.getElementById('carta-input').click()" style="white-space:nowrap">📂 Selecionar carta (PDF)</button>
    <span id="carta-st" style="font-size:11px;min-width:120px"></span>
  </div>
  <p style="font-size:11px;color:#2A5A3A;margin-bottom:14px;line-height:1.5">Ou edite manualmente os campos abaixo. Este cenário embasa as sugestões e scripts dos assessores.</p>
  <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:12px;margin-bottom:12px">
    <div><label>🌍 Cenário Global</label><textarea id="cenario-global" rows="5" placeholder="Contexto macro internacional..." oninput="autosaveCenario()"></textarea></div>
    <div><label>🇧🇷 Cenário Brasil</label><textarea id="cenario-brasil" rows="5" placeholder="Selic, IPCA, risco fiscal..." oninput="autosaveCenario()"></textarea></div>
    <div><label>📌 Posicionamento da Gestora</label><textarea id="cenario-pos" rows="5" placeholder="O que está sendo reduzido / aumentado..." oninput="autosaveCenario()"></textarea></div>
  </div>
  <div style="display:flex;align-items:center;gap:14px;flex-wrap:wrap">
    <div><label>Fonte / Referência do cenário</label><input type="text" id="cenario-ref" value="Levante Asset — Junho 2026" style="max-width:320px" oninput="autosaveCenario()"></div>
    <span id="cenario-autosave-st" style="font-size:10px;color:#2A5A3A;margin-top:18px"></span>
  </div>
</div>

<!-- ══ 2. UPLOAD RÁPIDO ═══════════════════════════════════════════════════════ -->
<div class="card" style="border-color:#3A3A20">
  <input type="file" id="pdf-rapido" accept=".pdf" multiple style="display:none" onchange="uploadRapidoSelecionar(this)">
  <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:14px;flex-wrap:wrap;gap:10px">
    <div style="display:flex;align-items:center;gap:10px">
      <span style="font-size:13px;color:#D4B483;font-weight:700;text-transform:uppercase;letter-spacing:.5px">⚡ Upload Rápido</span>
      <span style="font-size:11px;color:#2A5A3A">Carteiras recomendadas, produtos e material sobre ativos → <b style="color:#D4B483">Base de Conhecimento</b></span>
    </div>
    <button class="btn-ghost" onclick="fecharRapido()" id="btn-limpar-rapido" style="display:none;font-size:11px;padding:4px 10px">✕ Limpar</button>
  </div>
  <div id="rapido-zona-vazia" onclick="document.getElementById('pdf-rapido').click()"
    style="border:2px dashed #2A4A2A;border-radius:10px;padding:28px 20px;text-align:center;cursor:pointer;transition:border-color .2s"
    onmouseover="this.style.borderColor='#D4B483'" onmouseout="this.style.borderColor='#2A4A2A'">
    <div style="font-size:26px;margin-bottom:8px">📂</div>
    <div style="font-size:13px;color:#4A7055;font-weight:600;margin-bottom:4px">Clique para selecionar PDFs</div>
    <div style="font-size:11px;color:#1E3A2A">Até 5 arquivos · Carteiras recomendadas Levante / XP, relatórios, material sobre ativos</div>
  </div>
  <div id="pdf-rapido-itens" style="display:none;margin-top:12px">
    <div style="font-size:11px;color:#3A6A48;font-weight:700;text-transform:uppercase;letter-spacing:.5px;margin-bottom:8px">Arquivos selecionados — <span id="pdf-rapido-contador">0</span></div>
    <div id="pdf-rapido-lista"></div>
    <div style="margin-top:14px;padding-top:14px;border-top:1px solid #1A2A1A;display:flex;align-items:center;gap:12px;flex-wrap:wrap">
      <button class="btn" id="btn-extrair-todos" onclick="uploadRapidoTodos()" style="font-size:13px;padding:10px 22px;background:#C9A96E;color:#071E17;font-weight:700;border-color:#C9A96E">📤 Subir para a Memória</button>
      <button class="btn-ghost" onclick="document.getElementById('pdf-rapido').click()" style="font-size:11px;padding:6px 14px">+ Adicionar mais</button>
      <span id="rapido-upload-st" style="font-size:11px;color:#5DCAA5;margin-left:auto"></span>
    </div>
  </div>
</div>

<!-- ══ 3. ALOCAÇÕES DE REFERÊNCIA ════════════════════════════════════════════ -->
<div class="card">
  <div class="card-title"><span>📐</span> Alocações de Referência — Carteiras por Gestora</div>
  <p style="font-size:11px;color:#2A5A3A;margin-bottom:16px;line-height:1.5">
    Cadastre até <b style="color:#D4B483">5 gestoras</b> com as 5 categorias de perfil cada. Os % são usados como benchmark nas análises e sempre citados com a fonte.
  </p>

  <!-- Abas de gestoras (dinâmicas) -->
  <div style="display:flex;gap:0;margin-bottom:0;border-bottom:1px solid #1A2A1A;flex-wrap:wrap" id="gest2-tabs"></div>

  <!-- Painel: lista de gestoras / form -->
  <div id="gest2-painel" style="margin-top:16px"></div>

  <!-- Botão adicionar gestora -->
  <div style="margin-top:14px;display:flex;align-items:center;gap:12px;flex-wrap:wrap" id="gest2-add-area">
    <button class="btn btn-sm" onclick="gest2MostrarForm()" id="btn-gest2-add">➕ Adicionar Gestora</button>
    <span style="font-size:11px;color:#2A5A3A" id="gest2-count-label"></span>
  </div>

  <!-- Formulário oculto para adicionar/editar gestora -->
  <div id="gest2-form" style="display:none;margin-top:16px;background:#060C10;border:1px solid #1A2E3A;border-radius:10px;padding:18px">
    <div style="font-size:12px;color:#8BcFEF;font-weight:700;margin-bottom:14px" id="gest2-form-titulo">➕ Nova Gestora</div>
    <input type="hidden" id="gest2-id">
    <div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;margin-bottom:14px">
      <div><label>Nome da Gestora</label><input type="text" id="gest2-nome" placeholder="ex: Levante Asset"></div>
      <div><label>Referência / Mês</label><input type="text" id="gest2-ref" placeholder="ex: Junho 2026"></div>
    </div>
    <!-- 5 perfis com alocações -->
    <div id="gest2-perfis-tabs" style="display:flex;gap:4px;margin-bottom:12px;flex-wrap:wrap">
      <button type="button" class="gest2-ptab active" data-p="super_conservadora" onclick="gest2SwitchPerfil('super_conservadora',this)">Super Conserv.</button>
      <button type="button" class="gest2-ptab" data-p="conservadora" onclick="gest2SwitchPerfil('conservadora',this)">Conservadora</button>
      <button type="button" class="gest2-ptab" data-p="moderada" onclick="gest2SwitchPerfil('moderada',this)">Moderada</button>
      <button type="button" class="gest2-ptab" data-p="agressiva" onclick="gest2SwitchPerfil('agressiva',this)">Agressiva</button>
      <button type="button" class="gest2-ptab" data-p="super_agressiva" onclick="gest2SwitchPerfil('super_agressiva',this)">Super Agressiva</button>
    </div>
    <div id="gest2-aloc-panel" style="display:grid;grid-template-columns:repeat(3,1fr);gap:8px;margin-bottom:12px">
      <div><label style="color:#4A7055">Pós Fixado %</label><input type="number" id="g2aloc-pos_fixado" min="0" max="100" step="0.5" value="0" oninput="g2Total()" style="width:100%;background:#060F0B;border:1px solid #1A2E3A;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:13px;outline:none"></div>
      <div><label style="color:#4A7055">Inflação %</label><input type="number" id="g2aloc-inflacao" min="0" max="100" step="0.5" value="0" oninput="g2Total()" style="width:100%;background:#060F0B;border:1px solid #1A2E3A;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:13px;outline:none"></div>
      <div><label style="color:#4A7055">Pré Fixado %</label><input type="number" id="g2aloc-pre_fixado" min="0" max="100" step="0.5" value="0" oninput="g2Total()" style="width:100%;background:#060F0B;border:1px solid #1A2E3A;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:13px;outline:none"></div>
      <div><label style="color:#4A7055">Ações %</label><input type="number" id="g2aloc-acoes" min="0" max="100" step="0.5" value="0" oninput="g2Total()" style="width:100%;background:#060F0B;border:1px solid #1A2E3A;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:13px;outline:none"></div>
      <div><label style="color:#4A7055">FIIs %</label><input type="number" id="g2aloc-fiis" min="0" max="100" step="0.5" value="0" oninput="g2Total()" style="width:100%;background:#060F0B;border:1px solid #1A2E3A;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:13px;outline:none"></div>
      <div><label style="color:#4A7055">Multimercado %</label><input type="number" id="g2aloc-multimercado" min="0" max="100" step="0.5" value="0" oninput="g2Total()" style="width:100%;background:#060F0B;border:1px solid #1A2E3A;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:13px;outline:none"></div>
      <div><label style="color:#4A7055">Internacional %</label><input type="number" id="g2aloc-internacional" min="0" max="100" step="0.5" value="0" oninput="g2Total()" style="width:100%;background:#060F0B;border:1px solid #1A2E3A;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:13px;outline:none"></div>
      <div><label style="color:#4A7055">Alternativos %</label><input type="number" id="g2aloc-alternativos" min="0" max="100" step="0.5" value="0" oninput="g2Total()" style="width:100%;background:#060F0B;border:1px solid #1A2E3A;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:13px;outline:none"></div>
      <div><label style="color:#4A7055">Criptomoedas %</label><input type="number" id="g2aloc-criptomoedas" min="0" max="100" step="0.5" value="0" oninput="g2Total()" style="width:100%;background:#060F0B;border:1px solid #1A2E3A;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:13px;outline:none"></div>
    </div>
    <div style="display:flex;align-items:center;gap:12px;flex-wrap:wrap">
      <button class="btn" onclick="gest2Salvar()">💾 Salvar Gestora</button>
      <span style="font-size:12px;color:#4A7055">Total: <b id="g2-total">0</b>%</span>
      <span id="gest2-st" style="font-size:12px"></span>
      <button class="btn-ghost" onclick="document.getElementById('gest2-form').style.display='none'" style="margin-left:auto;font-size:11px">✕ Fechar</button>
    </div>
  </div>
</div>
<style>
.gest2-ptab{padding:5px 12px;border-radius:6px;border:1px solid #1A2E3A;background:#060C10;color:#4A7055;font-size:11px;cursor:pointer;font-family:inherit;transition:all .15s}
.gest2-ptab.active{background:#1A2E3A;color:#8BcFEF;border-color:#8BcFEF;font-weight:700}
</style>

<!-- ══ 3. PRODUTOS EM DESTAQUE ════════════════════════════════════════════════ -->
<div class="card">
  <div class="card-title"><span>⭐</span> Produtos em Destaque do Mês</div>
  <p style="font-size:11px;color:#2A5A3A;margin-bottom:14px;line-height:1.5">
    Para cada classe/indexador, cadastre os produtos disponíveis e recomendados. O agente usa estes produtos para sugerir ao assessor como fechar o gap da carteira do cliente.
  </p>

  <div class="tabs">
    <button class="tab active" onclick="prodTab('pos_fixado',this)">Pós Fixado</button>
    <button class="tab" onclick="prodTab('inflacao',this)">Inflação</button>
    <button class="tab" onclick="prodTab('pre_fixado',this)">Pré Fixado</button>
    <button class="tab" onclick="prodTab('acoes',this)">Ações</button>
    <button class="tab" onclick="prodTab('fiis',this)">FIIs</button>
    <button class="tab" onclick="prodTab('multimercado',this)">Multimercado</button>
    <button class="tab" onclick="prodTab('internacional',this)">Internacional</button>
    <button class="tab" onclick="prodTab('alternativos',this)">Alternativos</button>
  </div>

  <div id="prod-pos_fixado"   class="tab-panel active"></div>
  <div id="prod-inflacao"     class="tab-panel"></div>
  <div id="prod-pre_fixado"   class="tab-panel"></div>
  <div id="prod-acoes"        class="tab-panel"></div>
  <div id="prod-fiis"         class="tab-panel"></div>
  <div id="prod-multimercado" class="tab-panel"></div>
  <div id="prod-internacional"class="tab-panel"></div>
  <div id="prod-alternativos" class="tab-panel"></div>

  <div style="display:flex;align-items:center;gap:12px;margin-top:16px;padding-top:16px;border-top:1px solid #2A2A18;flex-wrap:wrap">
    <button class="btn" id="btn-salvar-prod" onclick="salvarProdutos()">💾 Salvar produtos</button>
    <span style="font-size:11px;color:#2A5A3A">Salva apenas os produtos em destaque. Use "Publicar" no final para enviar tudo de uma vez.</span>
    <span id="prod-save-st" style="font-size:12px;margin-left:auto"></span>
  </div>

  <!-- Alertas de produtos inline -->
  <div style="margin-top:16px;padding-top:16px;border-top:1px solid #1A1A28">
    <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:10px;cursor:pointer" onclick="this.nextElementSibling.style.display=this.nextElementSibling.style.display==='none'?'':'none'">
      <span style="font-size:11px;color:#B08FCF;font-weight:700;text-transform:uppercase;letter-spacing:.5px">🔔 Alertas de Produtos</span>
      <span style="font-size:10px;color:#4A7055">▾ expandir</span>
    </div>
    <div style="display:none">
      <div style="background:#060F0B;border:1px solid #2A1A30;border-radius:10px;padding:14px;margin-bottom:12px">
        <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:10px;margin-bottom:10px">
          <div><label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Produto / Ticker</label>
            <input id="alert-produto" type="text" placeholder="ex: NTN-B 2035, VALE3" style="width:100%;background:#081F18;border:1px solid #2A1A30;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none"></div>
          <div><label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Classe</label>
            <select id="alert-classe" style="width:100%;background:#081F18;border:1px solid #2A1A30;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
              <option value="">— Selecionar —</option>
              <option value="pos_fixado">Pós Fixado</option><option value="inflacao">Inflação</option>
              <option value="pre_fixado">Pré Fixado</option><option value="acoes">Ações</option>
              <option value="fiis">FIIs</option><option value="multimercado">Multimercado</option>
              <option value="internacional">Internacional</option><option value="alternativos">Alternativos</option>
            </select></div>
          <div><label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Tipo</label>
            <select id="alert-tipo" style="width:100%;background:#081F18;border:1px solid #2A1A30;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
              <option value="info">ℹ️ Informativo</option>
              <option value="atencao">⚠️ Atenção</option>
              <option value="urgente">🚨 Urgente</option>
            </select></div>
        </div>
        <div style="margin-bottom:10px">
          <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Direcionar para assessor (opcional)</label>
          <input id="alert-assessor" type="text" placeholder="Deixe em branco para todos" style="width:100%;background:#081F18;border:1px solid #2A1A30;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
        </div>
        <div style="margin-bottom:10px">
          <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Mensagem</label>
          <textarea id="alert-msg" rows="3" placeholder="ex: COPOM elevou Selic +50bps. Revisar alocação em pré-fixado." style="width:100%;background:#081F18;border:1px solid #2A1A30;border-radius:6px;padding:10px;color:#CCC;font-size:12px;resize:vertical;outline:none"></textarea>
        </div>
        <div style="display:flex;align-items:center;gap:10px">
          <button class="btn" onclick="registrarAlerta()">🔔 Registrar alerta</button>
          <span id="alert-st" style="font-size:11px"></span>
        </div>
      </div>
      <div style="font-size:11px;color:#3A6A48;margin-bottom:8px;display:flex;align-items:center;justify-content:space-between">
        <span>Histórico de alertas</span>
        <button onclick="carregarAlertas()" class="btn-ghost">↻ Atualizar</button>
      </div>
      <div id="alertas-feed" style="max-height:280px;overflow-y:auto;display:flex;flex-direction:column;gap:8px">
        <div style="color:#1E4A30;font-size:11px;text-align:center;padding:20px">Carregando alertas...</div>
      </div>
    </div>
  </div>
</div>

<!-- ══ 5. CALLS DE AÇÕES + OPERAÇÕES ESTRUTURADAS ════════════════════════════ -->
<div class="card" style="border-color:#1A2E3A">
  <div class="card-title"><span>📈</span> Calls de Ações</div>
  <p style="font-size:11px;color:#2A5A3A;margin-bottom:14px;line-height:1.5">
    Registre calls de compra, venda ou neutro por ticker. O agente entrega o call ao assessor quando o cliente já tem o ativo na carteira ou quando o perfil está na lista de destinatários.
  </p>

  <!-- Formulário novo call -->
  <div style="background:#060F0B;border:1px solid #1A2E3A;border-radius:10px;padding:16px;margin-bottom:16px">
    <div style="font-size:12px;color:#8BcFEF;font-weight:700;margin-bottom:12px">➕ Novo Call</div>

    <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:10px;margin-bottom:10px">
      <div>
        <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Ticker</label>
        <input id="call-ticker" type="text" placeholder="ex: VALE3" oninput="this.value=this.value.toUpperCase()" style="width:100%;background:#081F18;border:1px solid #1A2E3A;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
      </div>
      <div>
        <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Nome da empresa</label>
        <input id="call-nome" type="text" placeholder="ex: Vale S.A." style="width:100%;background:#081F18;border:1px solid #1A2E3A;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
      </div>
      <div>
        <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Direção</label>
        <div style="display:flex;gap:6px">
          <button id="call-dir-compra" onclick="setCallDir('compra')" style="flex:1;padding:7px;border-radius:6px;border:1.5px solid #5DCAA5;background:#0A2018;color:#5DCAA5;font-size:12px;font-weight:700;cursor:pointer">▲ COMPRA</button>
          <button id="call-dir-neutro" onclick="setCallDir('neutro')" style="flex:1;padding:7px;border-radius:6px;border:1.5px solid #2A5A3A;background:#111;color:#888;font-size:12px;font-weight:700;cursor:pointer">→ NEUTRO</button>
          <button id="call-dir-venda" onclick="setCallDir('venda')" style="flex:1;padding:7px;border-radius:6px;border:1.5px solid #2A1010;background:#1A0808;color:#FF6B6B;font-size:12px;font-weight:700;cursor:pointer">▼ VENDA</button>
        </div>
        <input type="hidden" id="call-direcao" value="compra">
      </div>
    </div>

    <div style="display:grid;grid-template-columns:1fr 1fr 1fr 1fr;gap:10px;margin-bottom:10px">
      <div>
        <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Preço de entrada (R$)</label>
        <input id="call-entrada" type="number" step="0.01" placeholder="62.50" oninput="calcCallUpside()" style="width:100%;background:#081F18;border:1px solid #1A2E3A;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
      </div>
      <div>
        <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Preço alvo (R$)</label>
        <input id="call-alvo" type="number" step="0.01" placeholder="74.00" oninput="calcCallUpside()" style="width:100%;background:#081F18;border:1px solid #1A2E3A;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
      </div>
      <div>
        <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Upside calculado</label>
        <div id="call-upside-display" style="padding:8px;border-radius:6px;background:#0A2A18;border:1px solid #1A3A28;font-size:14px;font-weight:700;color:#5DCAA5;text-align:center">—</div>
      </div>
      <div>
        <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Stop (R$)</label>
        <input id="call-stop" type="number" step="0.01" placeholder="57.00" style="width:100%;background:#081F18;border:1px solid #1A2E3A;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
      </div>
    </div>

    <div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;margin-bottom:10px">
      <div>
        <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Prazo</label>
        <input id="call-prazo" type="text" placeholder="ex: 6 meses, 12 meses" style="width:100%;background:#081F18;border:1px solid #1A2E3A;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
      </div>
      <div>
        <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Fonte</label>
        <input id="call-fonte" type="text" placeholder="ex: Levante Asset — Jun 2026" style="width:100%;background:#081F18;border:1px solid #1A2E3A;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
      </div>
    </div>

    <div style="margin-bottom:10px">
      <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Tese de investimento</label>
      <textarea id="call-tese" rows="3" placeholder="Descreva a tese resumidamente..." style="width:100%;background:#081F18;border:1px solid #1A2E3A;border-radius:6px;padding:8px;color:#CCC;font-size:12px;resize:vertical;outline:none"></textarea>
    </div>

    <div style="margin-bottom:12px">
      <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:6px">Perfis alvo (quem deve receber)</label>
      <div style="display:flex;gap:10px">
        <label style="display:flex;align-items:center;gap:6px;font-size:12px;cursor:pointer">
          <input type="checkbox" id="call-perf-arrojada" style="accent-color:#C9A96E"> Arrojada
        </label>
        <label style="display:flex;align-items:center;gap:6px;font-size:12px;cursor:pointer">
          <input type="checkbox" id="call-perf-agressiva" style="accent-color:#C9A96E"> Agressiva
        </label>
        <label style="display:flex;align-items:center;gap:6px;font-size:12px;cursor:pointer">
          <input type="checkbox" id="call-perf-moderada" style="accent-color:#C9A96E"> Moderada
        </label>
        <label style="display:flex;align-items:center;gap:6px;font-size:12px;cursor:pointer">
          <input type="checkbox" id="call-perf-conservadora" style="accent-color:#C9A96E"> Conservadora
        </label>
      </div>
    </div>

    <div style="display:flex;align-items:center;gap:10px">
      <button class="btn" onclick="adicionarCall()" style="max-width:200px">📈 Adicionar Call</button>
      <span id="call-st" style="font-size:11px"></span>
    </div>
  </div>

  <!-- Lista de calls ativos -->
  <div style="font-size:11px;color:#3A6A48;margin-bottom:8px;display:flex;align-items:center;justify-content:space-between">
    <span>Calls ativos</span>
    <button onclick="carregarCalls()" class="btn-ghost">↻ Atualizar</button>
  </div>
  <div id="calls-feed" style="display:flex;flex-direction:column;gap:8px">
    <div style="color:#1E4A30;font-size:11px;text-align:center;padding:20px">Carregando calls...</div>
  </div>

  <!-- ── Operações Estruturadas ─────────────────────────────────────────────── -->
  <div style="margin-top:24px;border-top:1px solid #1A2E3A;padding-top:18px">
    <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:14px;cursor:pointer" onclick="toggleEstruturadas()">
      <div style="display:flex;align-items:center;gap:10px">
        <span style="font-size:14px;font-weight:700;color:#C9A96E">💼 Operações Estruturadas</span>
        <span style="font-size:11px;color:#2A5A3A">COE, Termo, Opção e outras oportunidades manuais</span>
      </div>
      <span id="icon-estr" style="color:#4A7055;font-size:14px">▼</span>
    </div>

    <div id="painel-estruturadas">
      <!-- Formulário novo -->
      <div style="background:#060C08;border:1px solid #2A1A10;border-radius:10px;padding:16px;margin-bottom:16px">
        <div style="font-size:12px;color:#C9A96E;font-weight:700;margin-bottom:12px">➕ Nova Operação</div>
        <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:10px;margin-bottom:10px">
          <div>
            <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Tipo</label>
            <select id="estr-tipo" style="width:100%;background:#081208;border:1px solid #2A1A10;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
              <option>COE</option><option>Termo</option><option>Opção de Compra</option><option>Opção de Venda</option>
              <option>Debenture Incentivada</option><option>CRI</option><option>CRA</option><option>Outro</option>
            </select>
          </div>
          <div>
            <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Ativo / Ativo Base</label>
            <input id="estr-ativo" type="text" placeholder="ex: PETR4, IPCA+6,5%" style="width:100%;background:#081208;border:1px solid #2A1A10;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
          </div>
          <div>
            <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Emissor / Banco</label>
            <input id="estr-emissor" type="text" placeholder="ex: BTG Pactual" style="width:100%;background:#081208;border:1px solid #2A1A10;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
          </div>
        </div>
        <div style="display:grid;grid-template-columns:1fr 1fr 1fr 1fr;gap:10px;margin-bottom:10px">
          <div>
            <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Vencimento</label>
            <input id="estr-venc" type="text" placeholder="ex: Jan/2028" style="width:100%;background:#081208;border:1px solid #2A1A10;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
          </div>
          <div>
            <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Retorno Esperado</label>
            <input id="estr-retorno" type="text" placeholder="ex: IPCA+7% ou 120% CDI" style="width:100%;background:#081208;border:1px solid #2A1A10;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
          </div>
          <div>
            <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Perfil Mínimo</label>
            <select id="estr-perfil" style="width:100%;background:#081208;border:1px solid #2A1A10;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
              <option value="super_conservadora">Super Conservadora</option>
              <option value="conservadora">Conservadora</option>
              <option value="moderada" selected>Moderada</option>
              <option value="agressiva">Agressiva</option>
              <option value="super_agressiva">Super Agressiva</option>
            </select>
          </div>
          <div>
            <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Aplicação Mínima</label>
            <input id="estr-minimo" type="text" placeholder="ex: R$ 1.000" style="width:100%;background:#081208;border:1px solid #2A1A10;border-radius:6px;padding:8px;color:#CCC;font-size:12px;outline:none">
          </div>
        </div>
        <div style="margin-bottom:10px">
          <label style="font-size:10px;color:#3A6A48;display:block;margin-bottom:4px">Observação / Tese</label>
          <textarea id="estr-obs" rows="2" placeholder="Descreva a oportunidade, riscos e diferenciais..." style="width:100%;background:#081208;border:1px solid #2A1A10;border-radius:6px;padding:8px;color:#CCC;font-size:12px;resize:vertical;outline:none"></textarea>
        </div>
        <div style="display:flex;align-items:center;gap:10px">
          <button class="btn" onclick="adicionarEstruturada()" style="background:#C9A96E;color:#071E17;border-color:#C9A96E;font-weight:700">💼 Adicionar Operação</button>
          <span id="estr-st" style="font-size:11px"></span>
        </div>
      </div>

      <!-- Lista -->
      <div style="font-size:11px;color:#C9A96E;font-weight:700;text-transform:uppercase;letter-spacing:.5px;margin-bottom:10px">Operações Ativas</div>
      <div id="estr-lista"><div style="color:#2A2A18;font-size:11px;text-align:center;padding:16px">Carregando...</div></div>
    </div>
  </div>
</div>

<!-- ══ 6. BASE DE CONHECIMENTO ═══════════════════════════════════════════════ -->
<div class="card" style="border-color:#2A2A18">
  <div class="card-title"><span>📚</span> Base de Conhecimento — Cartas de Gestores</div>
  <p style="font-size:11px;color:#2A5A3A;margin-bottom:16px;line-height:1.6">
    Repositório de <b style="color:#D4B483">cartas mensais de gestoras</b>. Cada carta fica salva e pode ser consultada pelos assessores.
  </p>
  <div style="display:grid;grid-template-columns:220px 1fr;gap:16px;margin-bottom:20px">
    <div>
      <label>PDF da carta</label>
      <div id="know-drop" style="border:1.5px dashed #3A3A20;border-radius:10px;padding:22px;text-align:center;cursor:pointer;background:#060F0B;position:relative;transition:all .2s" onmouseover="this.style.borderColor='#D4B483'" onmouseout="this.style.borderColor='#3A3A20'">
        <input type="file" id="know-pdf" accept=".pdf" multiple style="position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%" onchange="knowUpload(this)">
        <div style="font-size:28px;margin-bottom:8px">📄</div>
        <p style="font-size:12px;color:#3A6A48">Arraste ou clique — múltiplos PDFs</p>
        <p style="font-size:11px;color:#D4B483;margin-top:6px;min-height:16px;word-break:break-all" id="know-fname"></p>
      </div>
      <div style="margin-top:8px;text-align:center"><span id="know-st" style="font-size:11px"></span></div>
    </div>
    <div style="display:flex;flex-direction:column;gap:10px">
      <div><label>Nome do documento</label><input type="text" id="know-nome" placeholder="ex: Carta Levante Junho 2026"></div>
      <div style="display:grid;grid-template-columns:1fr 1fr;gap:10px">
        <div><label>Tipo</label>
          <select id="know-tipo" style="width:100%;background:#060F0B;border:1px solid #2A2A18;border-radius:7px;padding:8px 10px;color:#F0F0F0;font-size:13px;outline:none">
            <option value="carta">Carta de Gestora</option>
            <option value="research">Nota de Research</option>
            <option value="cenario">Cenário Macro</option>
            <option value="relatorio">Relatório Mensal</option>
            <option value="outro">Outro</option>
          </select>
        </div>
        <div><label>Fonte / Gestora</label><input type="text" id="know-fonte" placeholder="ex: Levante Asset"></div>
      </div>
      <div>
        <label style="margin-bottom:6px">Classes relacionadas</label>
        <div style="display:flex;gap:6px;flex-wrap:wrap" id="know-classes-btns">
          <button type="button" class="know-cls-btn" data-cls="pos_fixado">Pós Fixado</button>
          <button type="button" class="know-cls-btn" data-cls="inflacao">Inflação</button>
          <button type="button" class="know-cls-btn" data-cls="pre_fixado">Pré Fixado</button>
          <button type="button" class="know-cls-btn" data-cls="acoes">Ações</button>
          <button type="button" class="know-cls-btn" data-cls="fiis">FIIs</button>
          <button type="button" class="know-cls-btn" data-cls="multimercado">Multimercado</button>
          <button type="button" class="know-cls-btn" data-cls="internacional">Internacional</button>
          <button type="button" class="know-cls-btn" data-cls="alternativos">Alternativos</button>
          <button type="button" class="know-cls-btn" data-cls="geral" style="background:#2A2A18;color:#D4B483;border-color:#D4B483">Geral</button>
        </div>
      </div>
      <div style="display:flex;align-items:center;gap:10px;margin-top:4px">
        <button class="btn" id="know-btn-upload" onclick="knowSalvar()" disabled>⬆️ Adicionar à Base</button>
        <span id="know-save-st" style="font-size:12px"></span>
      </div>
    </div>
  </div>
  <div>
    <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:10px">
      <span style="font-size:11px;color:#D4B483;font-weight:700;text-transform:uppercase;letter-spacing:.8px">📂 Documentos na Base</span>
      <span id="know-count" style="font-size:10px;color:#2A5A3A"></span>
    </div>
    <div id="know-lista"><div style="color:#1E4A30;font-size:11px;font-style:italic;text-align:center;padding:20px">Nenhum documento ainda.</div></div>
  </div>
</div>

<style>
.know-cls-btn{padding:4px 10px;border-radius:12px;border:1px solid #2A2A18;background:#0A0A0A;color:#4A7055;font-size:11px;cursor:pointer;transition:all .15s;font-family:inherit}
.know-cls-btn.sel{background:#D4B483;color:#000;border-color:#D4B483;font-weight:700}
.know-doc-card{background:#060F0B;border:1px solid #1A1A10;border-radius:10px;padding:14px 16px;margin-bottom:10px;transition:border-color .2s}
.know-doc-card:hover{border-color:#2A2A18}
.know-tipo-badge{font-size:10px;padding:2px 8px;border-radius:10px;font-weight:700}
</style>

<!-- ══ PUBLICAR ══════════════════════════════════════════════════════════════ -->
<div class="publish-bar" id="publish-bar">
  <!-- Barra de progresso (oculta por padrão) -->
  <div id="pub-progress-wrap" style="display:none;position:absolute;bottom:0;left:0;right:0;height:4px;border-radius:0 0 12px 12px;overflow:hidden;background:#1A0A0A">
    <div id="pub-progress-bar" style="height:100%;width:0%;background:linear-gradient(90deg,#FF4444,#FF8C42,#FFD700);border-radius:0 0 12px 12px;transition:width .4s ease;box-shadow:0 0 8px #FF444488"></div>
  </div>
  <div class="pub-info">
    <h3 id="pub-titulo">📤 Publicar para os Assessores</h3>
    <p id="pub-desc">Ao publicar, o cenário macro, alocações, produtos e calls ficam disponíveis para todos os assessores.</p>
  </div>
  <div style="display:flex;gap:10px;align-items:center;flex-wrap:wrap">
    <span id="pub-pendente-badge" style="display:none;font-size:11px;font-weight:700;color:#FF6B6B;background:#2A0808;border:1px solid #FF444455;border-radius:8px;padding:3px 10px;animation:piscar-pub 1.2s ease-in-out infinite">⚠️ Alterações não publicadas</span>
    <button class="btn" id="btn-pub" onclick="publicar()" style="background:#C9A96E;color:#071E17;border-color:#C9A96E;font-weight:700">📤 Publicar agora</button>
    <span id="pub-st" style="font-size:12px"></span>
  </div>
</div>

</div><!-- /container -->

<script>
// ── Auth ──────────────────────────────────────────────────────────────────────
(function(){
  const _r = localStorage.getItem("brauna_role");
  if(_r !== "head" && _r !== "admin"){
    localStorage.removeItem("brauna_role");
    window.location.replace("/");
  }
})();
function sair(){ localStorage.removeItem("brauna_role"); window.location.replace("/"); }

// ── Publicar pendente ─────────────────────────────────────────────────────────
function marcarPendente(){
  const bar   = document.getElementById("publish-bar");
  const badge = document.getElementById("pub-pendente-badge");
  const titulo= document.getElementById("pub-titulo");
  if(bar)   bar.classList.add("pendente");
  if(badge) badge.style.display = "";
  if(titulo) titulo.textContent = "🔴 Publicar para os Assessores";
  // scroll suave até a publish-bar para chamar atenção
  if(bar) bar.scrollIntoView({behavior:"smooth", block:"nearest"});
}
// ── Auto-save Cenário Macro ───────────────────────────────────────────────────
let _cenarioTimer = null;
function autosaveCenario(){
  clearTimeout(_cenarioTimer);
  const st = document.getElementById("cenario-autosave-st");
  if(st) st.textContent = "✏️ digitando...";
  _cenarioTimer = setTimeout(async ()=>{
    const payload = {
      global:        document.getElementById("cenario-global").value,
      brasil:        document.getElementById("cenario-brasil").value,
      posicionamento:document.getElementById("cenario-pos").value,
      referencia:    document.getElementById("cenario-ref").value,
    };
    try{
      await fetch("/api/hp/cenario", {method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify(payload)});
      if(st){ st.textContent = "✓ salvo automaticamente"; setTimeout(()=>{ if(st) st.textContent=""; }, 3000); }
      marcarPendente();
    } catch(e){ if(st) st.textContent = ""; }
  }, 1500);
}

function limparPendente(){
  const bar   = document.getElementById("publish-bar");
  const badge = document.getElementById("pub-pendente-badge");
  const titulo= document.getElementById("pub-titulo");
  if(bar)   bar.classList.remove("pendente");
  if(badge) badge.style.display = "none";
  if(titulo) titulo.textContent = "✅ Publicar para os Assessores";
  setTimeout(()=>{ if(titulo) titulo.textContent="📤 Publicar para os Assessores"; }, 4000);
}

// ── Estado ────────────────────────────────────────────────────────────────────
let _portfolios = null;
let _cenario    = null;
let _produtos   = {};
let _base       = [];

const CLASSES = [
  {key:"pos_fixado",   label:"Pós Fixado"},
  {key:"inflacao",     label:"Inflação"},
  {key:"pre_fixado",   label:"Pré Fixado"},
  {key:"acoes",        label:"Ações"},
  {key:"fiis",         label:"FIIs"},
  {key:"multimercado", label:"Multimercado"},
  {key:"alternativos", label:"Alternativos"},
  {key:"internacional",label:"Internacional"},
  {key:"criptomoedas", label:"Criptomoedas"},
];
const PERFIS = ["super_conservadora","conservadora","moderada","arrojada","agressiva"];
const VIESES = ["positivo","neutro","negativo"];
const TIPOS_RF = ["CDB","LCI","LCA","LIG","CRA","CRI","Debênture","Tesouro Direto","Fundo RF","Outro"];
const TIPOS_RV = ["Ação","ETF","BDR","Fundo RV","Outro"];
const TIPOS_FII= ["Papel","Logística","Shopping","Escritório","Híbrido","Fundo de Fundos","Outro"];
const TIPOS_MM = ["Long & Short","Macro","Multiestratégia","Quantitativo","Outro"];
const TIPOS_IN = ["ETF","BDR","Ação","Fundo","COE","Outro"];

// ── Portfólio Modelo — renderizar tabela ──────────────────────────────────────
function renderPortoTable(data){
  _portfolios = data;
  document.getElementById("porto-ref").value = data.referencia || "";
  document.getElementById("porto-badge").textContent = data.referencia || "";

  const tbody = document.getElementById("porto-body");
  tbody.innerHTML = "";

  CLASSES.forEach(cls=>{
    const tr = document.createElement("tr");
    const viresAtual = data.vieses?.[cls.key] || "neutro";
    let cells = `<td class="classe-label">${cls.label}</td>`;
    PERFIS.forEach(p=>{
      const val = data.perfis?.[p]?.[cls.key] ?? 0;
      cells += `<td><input type="number" min="0" max="100" step="0.25" value="${val}" id="p-${p}-${cls.key}" onchange="atualizarTotais()"></td>`;
    });
    cells += `<td>
      <select class="vies-select" id="vies-${cls.key}" onchange="colorirVies('${cls.key}')">
        <option value="positivo" ${viresAtual==="positivo"?"selected":""}>▲ Positivo</option>
        <option value="neutro"   ${viresAtual==="neutro"?"selected":""}  >→ Neutro</option>
        <option value="negativo" ${viresAtual==="negativo"?"selected":""}>▼ Negativo</option>
      </select>
    </td>`;
    tr.innerHTML = cells;
    tbody.appendChild(tr);
  });

  atualizarTotais();
}

function atualizarTotais(){
  PERFIS.forEach(p=>{
    let soma = 0;
    CLASSES.forEach(cls=>{
      const el = document.getElementById(`p-${p}-${cls.key}`);
      if(el) soma += parseFloat(el.value)||0;
    });
    const el = document.getElementById(`tot-${p}`);
    if(el){
      el.textContent = soma.toFixed(2) + "%";
      el.className = Math.abs(soma - 100) < 0.01 ? "total-ok" : (soma > 100 ? "total-err" : "total-val");
    }
  });
}

function colorirVies(cls){
  const el = document.getElementById("vies-"+cls);
  if(!el) return;
  const v = el.value;
  el.style.color = v==="positivo" ? "#5DCAA5" : v==="negativo" ? "#FF6B6B" : "#4A7055";
}

function coletarPortfolios(){
  const ref = document.getElementById("porto-ref").value;
  const perfis = {};
  PERFIS.forEach(p=>{
    perfis[p] = {label: p.replace(/_/g," ")};
    CLASSES.forEach(cls=>{
      const el = document.getElementById(`p-${p}-${cls.key}`);
      perfis[p][cls.key] = parseFloat(el?.value||0);
    });
  });
  const vieses = {};
  CLASSES.forEach(cls=>{
    const el = document.getElementById("vies-"+cls.key);
    vieses[cls.key] = el?.value || "neutro";
  });
  return {referencia: ref, publicado_em: new Date().toLocaleDateString("pt-BR"), perfis, vieses,
          classes: CLASSES.map(c=>c.key), labels: Object.fromEntries(CLASSES.map(c=>[c.key,c.label]))};
}

// ── Cenário Macro ─────────────────────────────────────────────────────────────
function carregarCenario(data){
  _cenario = data;
  document.getElementById("cenario-global").value = data.global || "";
  document.getElementById("cenario-brasil").value = data.brasil || "";
  document.getElementById("cenario-pos").value    = data.posicionamento || "";
  document.getElementById("cenario-ref").value    = data.referencia || "";
}

function coletarCenario(){
  return {
    referencia:    document.getElementById("cenario-ref").value,
    global:        document.getElementById("cenario-global").value,
    brasil:        document.getElementById("cenario-brasil").value,
    posicionamento:document.getElementById("cenario-pos").value,
    publicado_em:  new Date().toLocaleDateString("pt-BR"),
  };
}

// ── Produtos ──────────────────────────────────────────────────────────────────
function prodTab(cat, btn){
  document.querySelectorAll(".tab").forEach(b=>b.classList.remove("active"));
  if(btn) btn.classList.add("active");
  document.querySelectorAll(".tab-panel").forEach(p=>p.classList.remove("active"));
  document.getElementById("prod-"+cat).classList.add("active");
}

function tiposPorClasse(cls){
  if(["pos_fixado","inflacao","pre_fixado"].includes(cls)) return TIPOS_RF;
  if(cls==="acoes") return TIPOS_RV;
  if(cls==="fiis")  return TIPOS_FII;
  if(cls==="multimercado") return TIPOS_MM;
  if(cls==="internacional") return TIPOS_IN;
  return ["Outro"];
}

function buildProdHtml(cls, id, p){
  p = p || {};
  const tipos = tiposPorClasse(cls).map(t=>`<option ${t===(p.tipo||"")?"selected":""}>${t}</option>`).join("");
  const v = k => (p[k]||"").toString().replace(/"/g,"&quot;");

  if(cls === "acoes"){
    const setores = ["Bancário","Energia","Mineração","Varejo","Tecnologia","Saúde","Agronegócio","Construção","Utilities","Transporte","Financeiro","Indústria","Outro"].map(s=>`<option ${s===p.setor?"selected":""}>${s}</option>`).join("");
    const coberturas = ["Compra","Compra Forte","Neutro","Vender","Sem cobertura"].map(s=>`<option ${s===p.cobertura?"selected":""}>${s}</option>`).join("");
    return `<div class="prod-item" id="prod-item-${cls}-${id}" style="border-color:#1A2A18">
      <button class="prod-del" onclick="delProduto('${cls}',${id})">✕</button>
      <div class="prod-grid-3">
        <div class="prod-mini"><label>Ticker</label><input type="text" id="pn-${cls}-${id}" value="${v("nome")}" placeholder="Ex: VALE3, WEGE3, PETR4" style="text-transform:uppercase"></div>
        <div class="prod-mini"><label>Nome da empresa</label><input type="text" id="pe-${cls}-${id}" value="${v("emissor")}" placeholder="Ex: Vale S.A."></div>
        <div class="prod-mini"><label>Setor</label><select id="pt-${cls}-${id}" style="background:#060F0B;border:1px solid #1A2A18;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:12px;outline:none;width:100%">${setores}</select></div>
      </div>
      <div style="display:grid;grid-template-columns:1fr 1fr 1fr 1fr;gap:8px;margin-top:8px">
        <div class="prod-mini"><label>Preço atual (R$)</label><input type="text" id="pr-${cls}-${id}" value="${v("taxa")}" placeholder="Ex: 68,40"></div>
        <div class="prod-mini"><label>Preço-alvo (R$)</label><input type="text" id="pv-${cls}-${id}" value="${v("vencimento")}" placeholder="Ex: 85,00"></div>
        <div class="prod-mini"><label>Upside estimado</label><input type="text" id="pup-${cls}-${id}" value="${v("upside")}" placeholder="Ex: +24%"></div>
        <div class="prod-mini"><label>Recomendação</label><select id="pcob-${cls}-${id}" style="background:#060F0B;border:1px solid #1A2A18;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:12px;outline:none;width:100%">${coberturas}</select></div>
      </div>
      <div style="display:grid;grid-template-columns:1fr 1fr;gap:8px;margin-top:8px">
        <div class="prod-mini"><label>Dividend Yield esperado</label><input type="text" id="pdy-${cls}-${id}" value="${v("dy")}" placeholder="Ex: 8,5% a.a."></div>
        <div class="prod-mini"><label>Prazo da tese</label><input type="text" id="ppz-${cls}-${id}" value="${v("prazo")}" placeholder="Ex: 12-18 meses"></div>
      </div>
      <div class="prod-mini" style="margin-top:8px"><label>Tese de investimento</label><textarea id="pm-${cls}-${id}" rows="2" style="width:100%;background:#060F0B;border:1px solid #1A2A18;border-radius:6px;padding:8px;color:#F0F0F0;font-size:12px;resize:vertical;outline:none" placeholder="Por que recomendar? Quais os gatilhos de valorização?">${v("motivo")}</textarea></div>
      <div class="prod-mini" style="margin-top:6px"><label>Riscos / Pontos de atenção</label><input type="text" id="prisk-${cls}-${id}" value="${v("riscos")}" placeholder="Ex: Exposição ao câmbio, risco regulatório, ciclo de commodities"></div>
    </div>`;
  }

  if(cls === "fiis"){
    const segmentos = ["Papel (CRI)","Lajes Corporativas","Logística","Shopping","Residencial","Híbrido","Fundo de Fundos","Agro","Outro"].map(s=>`<option ${s===p.setor?"selected":""}>${s}</option>`).join("");
    const recs = ["Compra","Compra Forte","Neutro","Vender"].map(s=>`<option ${s===p.cobertura?"selected":""}>${s}</option>`).join("");
    return `<div class="prod-item" id="prod-item-${cls}-${id}" style="border-color:#1A1A2A">
      <button class="prod-del" onclick="delProduto('${cls}',${id})">✕</button>
      <div class="prod-grid-3">
        <div class="prod-mini"><label>Ticker</label><input type="text" id="pn-${cls}-${id}" value="${v("nome")}" placeholder="Ex: MXRF11, HGLG11" style="text-transform:uppercase"></div>
        <div class="prod-mini"><label>Nome do fundo</label><input type="text" id="pe-${cls}-${id}" value="${v("emissor")}" placeholder="Ex: Maxi Renda FII"></div>
        <div class="prod-mini"><label>Segmento</label><select id="pt-${cls}-${id}" style="background:#060F0B;border:1px solid #1A1A2A;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:12px;outline:none;width:100%">${segmentos}</select></div>
      </div>
      <div style="display:grid;grid-template-columns:1fr 1fr 1fr 1fr;gap:8px;margin-top:8px">
        <div class="prod-mini"><label>Preço atual (R$)</label><input type="text" id="pr-${cls}-${id}" value="${v("taxa")}" placeholder="Ex: 9,85"></div>
        <div class="prod-mini"><label>P/VP</label><input type="text" id="pv-${cls}-${id}" value="${v("vencimento")}" placeholder="Ex: 0,92"></div>
        <div class="prod-mini"><label>Dividend Yield (a.a.)</label><input type="text" id="pdy-${cls}-${id}" value="${v("dy")}" placeholder="Ex: 12,8%"></div>
        <div class="prod-mini"><label>Recomendação</label><select id="pcob-${cls}-${id}" style="background:#060F0B;border:1px solid #1A1A2A;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:12px;outline:none;width:100%">${recs}</select></div>
      </div>
      <div style="display:grid;grid-template-columns:1fr 1fr;gap:8px;margin-top:8px">
        <div class="prod-mini"><label>Rendimento mensal esperado (R$/cota)</label><input type="text" id="pup-${cls}-${id}" value="${v("upside")}" placeholder="Ex: R$ 0,105/cota"></div>
        <div class="prod-mini"><label>Gestora</label><input type="text" id="ppz-${cls}-${id}" value="${v("prazo")}" placeholder="Ex: CSHG, Kinea, BTG Pactual"></div>
      </div>
      <div class="prod-mini" style="margin-top:8px"><label>Tese de investimento</label><textarea id="pm-${cls}-${id}" rows="2" style="width:100%;background:#060F0B;border:1px solid #1A1A2A;border-radius:6px;padding:8px;color:#F0F0F0;font-size:12px;resize:vertical;outline:none" placeholder="Portfólio, qualidade dos contratos, vacância, gestão...">${v("motivo")}</textarea></div>
      <div class="prod-mini" style="margin-top:6px"><label>Riscos</label><input type="text" id="prisk-${cls}-${id}" value="${v("riscos")}" placeholder="Ex: Alta de juros reduz atratividade, vacância elevada, risco de crédito dos CRIs"></div>
    </div>`;
  }

  // Formulário genérico para demais classes
  return `<div class="prod-item" id="prod-item-${cls}-${id}">
    <button class="prod-del" onclick="delProduto('${cls}',${id})">✕</button>
    <div class="prod-grid-3">
      <div class="prod-mini"><label>Nome do produto / ticker</label><input type="text" id="pn-${cls}-${id}" value="${v("nome")}" placeholder="Ex: CDB Banco BTG 108% CDI"></div>
      <div class="prod-mini"><label>Tipo</label><select id="pt-${cls}-${id}" style="background:#060F0B;border:1px solid #2A2A18;border-radius:6px;padding:6px 8px;color:#F0F0F0;font-size:12px;outline:none;width:100%">${tipos}</select></div>
      <div class="prod-mini"><label>Taxa / retorno esperado</label><input type="text" id="pr-${cls}-${id}" value="${v("taxa")}" placeholder="Ex: 108% CDI / IPCA+7% / 13,5% a.a."></div>
    </div>
    <div class="prod-grid-2">
      <div class="prod-mini"><label>Emissor / Gestora</label><input type="text" id="pe-${cls}-${id}" value="${v("emissor")}" placeholder="Ex: BTG Pactual, XP Asset, Kinea"></div>
      <div class="prod-mini"><label>Vencimento / Liquidez</label><input type="text" id="pv-${cls}-${id}" value="${v("vencimento")}" placeholder="Ex: 2027-06 / D+30 / 2 anos"></div>
    </div>
    <div class="prod-mini" style="margin-top:8px"><label>Motivo / Tese de recomendação</label><input type="text" id="pm-${cls}-${id}" value="${v("motivo")}" placeholder="Por que recomendar este produto?"></div>
  </div>`;
}

function addProduto(cls){
  if(!_produtos[cls]) _produtos[cls] = [];
  const id = Date.now();
  _produtos[cls].push(id);
  document.getElementById("prod-"+cls).insertAdjacentHTML("beforeend", buildProdHtml(cls, id, {}));
}

function delProduto(cls, id){
  const el = document.getElementById(`prod-item-${cls}-${id}`);
  if(el) el.remove();
  _produtos[cls] = (_produtos[cls]||[]).filter(x=>x!==id);
}

function renderProdutos(data){
  _produtos = {};
  Object.keys(data).forEach(cls=>{
    const lista = data[cls];
    if(!Array.isArray(lista)) return;
    _produtos[cls] = [];
    lista.forEach(p=>{
      const id = Date.now() + Math.random();
      _produtos[cls].push(id);
      const container = document.getElementById("prod-"+cls);
      if(container) container.insertAdjacentHTML("beforeend", buildProdHtml(cls, id, p));
    });
  });
}

// Cria toolbar + botão "+" em cada tab-panel
function initProdButtons(){
  CLASSES.forEach(cls=>{
    const panel = document.getElementById("prod-"+cls.key);
    if(!panel) return;

    // Barra de importação
    const toolbar = document.createElement("div");
    toolbar.style.cssText = "display:flex;align-items:center;gap:8px;flex-wrap:wrap;padding:10px 12px;background:#070C07;border:1px dashed #2A2A18;border-radius:8px;margin-bottom:12px";
    toolbar.innerHTML = `
      <span style="font-size:11px;color:#3A4A3A;flex:1">📂 Importar produtos via arquivo:</span>
      <label style="display:inline-flex;align-items:center;gap:6px;background:#1A1A08;color:#D4B483;border:1px solid #D4B483;border-radius:6px;padding:5px 12px;font-size:11px;cursor:pointer;transition:all .2s;position:relative;overflow:hidden" onmouseover="this.style.background='#D4B483';this.style.color='#000'" onmouseout="this.style.background='#1A1A08';this.style.color='#D4B483'">
        📄 PDF
        <input type="file" accept=".pdf" style="position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%" onchange="uploadPdfProduto('${cls.key}',this)">
      </label>
      <label style="display:inline-flex;align-items:center;gap:6px;background:#1A1A08;color:#D4B483;border:1px solid #D4B483;border-radius:6px;padding:5px 12px;font-size:11px;cursor:pointer;transition:all .2s;position:relative;overflow:hidden" onmouseover="this.style.background='#D4B483';this.style.color='#000'" onmouseout="this.style.background='#1A1A08';this.style.color='#D4B483'">
        📊 Planilha CSV/Excel
        <input type="file" accept=".csv,.xlsx,.xls" style="position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%" onchange="importarProdPlanilha('${cls.key}',this)">
      </label>
      <button class="btn-ghost" onclick="baixarModeloProd('${cls.key}')">⬇ Modelo CSV</button>
      <span id="prod-import-st-${cls.key}" style="font-size:11px;color:#5DCAA5"></span>
    `;
    panel.appendChild(toolbar);

    // Preview PDF por aba
    const preview = document.createElement("div");
    preview.id = `prod-pdf-preview-${cls.key}`;
    preview.style.cssText = "display:none;background:#070C07;border:1px solid #D4B483;border-radius:8px;padding:12px;margin-bottom:12px";
    preview.innerHTML = `
      <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:6px">
        <span style="font-size:11px;color:#D4B483;font-weight:700">📄 Conteúdo extraído — ${cls.label}</span>
        <button onclick="document.getElementById('prod-pdf-preview-${cls.key}').style.display='none'" style="background:none;border:none;color:#3A6A48;cursor:pointer;font-size:14px">✕</button>
      </div>
      <textarea id="prod-pdf-texto-${cls.key}" rows="6" style="width:100%;background:#050A05;border:1px solid #2A2A18;border-radius:7px;padding:10px;color:#AAA;font-size:11px;font-family:monospace;resize:vertical;outline:none"></textarea>
      <div style="display:flex;gap:8px;margin-top:8px;flex-wrap:wrap">
        <button class="btn btn-sm" onclick="usarPdfComoProduto('${cls.key}')">→ Criar produto a partir deste texto</button>
        <button class="btn-ghost" onclick="navigator.clipboard.writeText(document.getElementById('prod-pdf-texto-${cls.key}').value)">📋 Copiar</button>
      </div>
    `;
    panel.appendChild(preview);

    // Botão adicionar manual
    const btn = document.createElement("button");
    btn.className = "btn-add";
    btn.textContent = `+ Adicionar produto de ${cls.label}`;
    btn.onclick = ()=>addProduto(cls.key);
    panel.appendChild(btn);

    const list = document.createElement("div");
    list.id = "prod-list-"+cls.key;
    panel.appendChild(list);
  });
}

// Upload PDF por categoria de produto
async function uploadPdfProduto(cls, input){
  const f = input.files[0]; if(!f) return;
  const st = document.getElementById("prod-import-st-"+cls);
  if(st) st.textContent = "⏳ Extraindo...";
  input.value = "";
  try{
    const fd = new FormData();
    fd.append("arquivo", f);
    fd.append("categoria", cls);
    const r = await fetch("/api/admin/upload-sugestao-doc", {method:"POST", body:fd});
    const d = await r.json();
    if(d.ok){
      const texto = d.texto || "";
      document.getElementById(`prod-pdf-texto-${cls}`).value = texto;
      document.getElementById(`prod-pdf-preview-${cls}`).style.display = "block";
      if(st) st.innerHTML = `<span class="status-ok">✓ ${d.chars||texto.length} caracteres extraídos</span>`;
      setTimeout(()=>{ if(st) st.textContent=""; }, 4000);
      _base.push({nome: f.name, chars: d.chars||0, data: new Date().toLocaleDateString("pt-BR")});
      renderBase();
    } else { if(st) st.innerHTML = `<span class="status-err">Erro: ${d.error}</span>`; }
  } catch(e){ if(st) st.innerHTML = `<span class="status-err">Erro: ${e.message}</span>`; }
}

function usarPdfComoProduto(cls){
  const txt = document.getElementById(`prod-pdf-texto-${cls}`).value.trim();
  if(!txt) return;
  addProduto(cls);
  const ids = _produtos[cls] || [];
  const id = ids[ids.length-1];
  if(id){
    const el = document.getElementById(`pm-${cls}-${id}`);
    if(el) el.value = txt.substring(0, 400);
  }
  document.getElementById(`prod-pdf-preview-${cls}`).style.display = "none";
}

// Importar planilha de produtos
const MODELOS_PROD_CSV  = ["nome","tipo","taxa","emissor","vencimento","motivo"];
const MODELOS_PROD_EX   = {
  pos_fixado:    [["CDB Banco Inter 108% CDI","CDB","108% CDI","Banco Inter","2027-06","Boa liquidez e retorno acima da média"],["LCI Bradesco 97% CDI","LCI","97% CDI","Bradesco","2026-12","Isento IR, ideal para conservadores"]],
  inflacao:      [["Tesouro IPCA+ 2035","Tesouro Direto","IPCA+7,35%","Tesouro Nacional","2035","Proteção inflacionária longo prazo"],["CRI IPCA+8%","CRI","IPCA+8%","Securitizadora","2028","Isento IR, crédito privado"]],
  pre_fixado:    [["CDB Pré 14% a.a.","CDB","14% a.a.","Banco XP","2027","Trava taxa antes de queda Selic"]],
  acoes:         [["VALE3","Ação","—","Vale S.A.","—","Commodities, geração de caixa, dividendos"],["WEGE3","Ação","—","WEG S.A.","—","Crescimento sólido, exportação"]],
  fiis:          [["MXRF11","Papel","Yield ~13% a.a.","Maxi Renda","Aberto","Alta distribuição, baixa volatilidade"],["HGLG11","Logística","Yield ~12% a.a.","CSHG","Aberto","Contratos longos, boa gestão"]],
  multimercado:  [["Kinea Chronos","Multiestratégia","CDI+4% alvo","Kinea","D+30","Track record sólido, baixa correlação"]],
  internacional: [["IVVB11","BDR","S&P 500","BlackRock","Aberto","Exposição EUA em BRL, liquidez diária"],["VTI","ETF","Mercado amplo EUA","Vanguard","Aberto","Diversificação total EUA"]],
  alternativos:  [["COE S&P 500 Capital Protegido","COE","Participação 80% alta S&P","BTG","2027","Proteção de capital + upside"]],
};

function baixarModeloProd(cls){
  const cols = MODELOS_PROD_CSV;
  const exs  = MODELOS_PROD_EX[cls] || [[cls+" Exemplo","Tipo","Taxa","Emissor","Vencimento","Motivo"]];
  const linhas = [cols.join(","), ...exs.map(r=>r.map(c=>/,/.test(c)?`"${c}"`:c).join(","))];
  const blob = new Blob(["﻿"+linhas.join("\r\n")], {type:"text/csv;charset=utf-8"});
  const a = document.createElement("a"); a.href=URL.createObjectURL(blob);
  a.download=`modelo_prod_${cls}.csv`; a.click();
}

function importarProdPlanilha(cls, input){
  const file = input.files[0]; if(!file) return;
  const st = document.getElementById("prod-import-st-"+cls);
  lerPlanilha(file, function(rows){
    let count = 0;
    rows.forEach(function(r){
      if(!r[0]) return;
      addProduto(cls);
      const ids = _produtos[cls]||[];
      const id = ids[ids.length-1];
      if(!id) return;
      count++;
      const el = (suf)=>document.getElementById(`p${suf}-${cls}-${id}`);
      const n=document.getElementById(`pn-${cls}-${id}`); if(n) n.value=r[0]||"";
      const t=document.getElementById(`pt-${cls}-${id}`); if(t) t.value=r[1]||"";
      const x=document.getElementById(`pr-${cls}-${id}`); if(x) x.value=r[2]||"";
      const e=document.getElementById(`pe-${cls}-${id}`); if(e) e.value=r[3]||"";
      const v=document.getElementById(`pv-${cls}-${id}`); if(v) v.value=r[4]||"";
      const m=document.getElementById(`pm-${cls}-${id}`); if(m) m.value=r[5]||"";
    });
    if(st) st.innerHTML = `<span class="status-ok">✓ ${count} produtos importados</span>`;
    setTimeout(()=>{ if(st) st.textContent=""; }, 4000);
    input.value="";
  });
}

function lerPlanilha(file, callback){
  const reader = new FileReader();
  const isXls  = /\.(xlsx|xls)$/i.test(file.name);
  reader.onload = function(e){
    let rows = [];
    if(isXls){
      alert("Para importar Excel, use a versão CSV por enquanto.");
      return;
    }
    const txt = e.target.result;
    const sep = (txt.split(";").length > txt.split(",").length) ? ";" : ",";
    rows = txt.split("\n").map(l=>l.split(sep).map(c=>c.replace(/^"|"$/g,"").trim()));
    if(rows.length>1) rows=rows.slice(1);
    callback(rows.filter(r=>r.some(c=>c)));
  };
  isXls ? reader.readAsArrayBuffer(file) : reader.readAsText(file,"UTF-8");
}

function coletarProdutos(){
  const result = {};
  const g = (id) => document.getElementById(id)?.value || "";
  CLASSES.forEach(cls=>{
    const ids = _produtos[cls.key] || [];
    const k = cls.key;
    result[k] = ids.map(id=>{
      const base = {
        nome:       g(`pn-${k}-${id}`),
        tipo:       g(`pt-${k}-${id}`),
        taxa:       g(`pr-${k}-${id}`),
        emissor:    g(`pe-${k}-${id}`),
        vencimento: g(`pv-${k}-${id}`),
        motivo:     g(`pm-${k}-${id}`),
      };
      if(k === "acoes" || k === "fiis"){
        base.upside  = g(`pup-${k}-${id}`);
        base.cobertura = g(`pcob-${k}-${id}`);
        base.dy      = g(`pdy-${k}-${id}`);
        base.prazo   = g(`ppz-${k}-${id}`);
        base.setor   = g(`pt-${k}-${id}`);
        base.riscos  = g(`prisk-${k}-${id}`);
      }
      return base;
    }).filter(p=>p.nome);
  });
  return result;
}

// ── Base de Conhecimento ──────────────────────────────────────────────────────
const KNOW_TIPO_COLORS = {carta:"#5DCAA5", research:"#8B9FE8", cenario:"#D4B483", relatorio:"#B08FCF", outro:"#888"};
const KNOW_TIPO_LABELS = {carta:"Carta de Gestora", research:"Research", cenario:"Cenário Macro", relatorio:"Relatório", outro:"Outro"};
let _knowFile = null;

// seleção de classes
document.addEventListener("DOMContentLoaded", ()=>{
  document.querySelectorAll(".know-cls-btn").forEach(btn=>{
    btn.addEventListener("click", ()=>btn.classList.toggle("sel"));
  });
  carregarKnowledge();
});

let _knowFiles = [];  // suporta múltiplos

function knowUpload(input){
  const files = [...input.files];
  if(!files.length) return;
  _knowFiles = files;
  _knowFile  = files[0];  // compatibilidade
  document.getElementById("know-fname").textContent =
    files.length === 1 ? files[0].name : `${files.length} arquivos selecionados`;
  if(files.length === 1 && !document.getElementById("know-nome").value)
    document.getElementById("know-nome").value = files[0].name.replace(/\.pdf$/i,"");
  document.getElementById("know-btn-upload").disabled = false;
  document.getElementById("know-st").innerHTML =
    `<span style="color:#5DCAA5">✓ ${files.length} arquivo${files.length>1?"s":""} prontos para subir</span>`;
}

async function knowSalvar(){
  if(!_knowFiles.length){ document.getElementById("know-save-st").innerHTML=`<span class="status-err">Selecione ao menos um PDF.</span>`; return; }
  const btn = document.getElementById("know-btn-upload");
  const st  = document.getElementById("know-save-st");
  btn.disabled = true;
  const classes = [...document.querySelectorAll(".know-cls-btn.sel")].map(b=>b.dataset.cls).join(",");
  const tipo    = document.getElementById("know-tipo").value;
  const fonte   = document.getElementById("know-fonte").value.trim();
  const total   = _knowFiles.length;
  let ok = 0, erros = [];

  for(let i = 0; i < total; i++){
    const f = _knowFiles[i];
    st.innerHTML = `⏳ Processando ${i+1}/${total}: <b>${f.name}</b>`;
    try{
      const fd = new FormData();
      fd.append("pdf",     f);
      fd.append("nome",    total === 1 ? (document.getElementById("know-nome").value.trim() || f.name.replace(/\.pdf$/i,"")) : f.name.replace(/\.pdf$/i,""));
      fd.append("tipo",    tipo);
      fd.append("fonte",   fonte);
      fd.append("classes", classes);
      const r = await fetch("/api/hp/knowledge/upload", {method:"POST", body:fd});
      const d = await r.json();
      if(d.ok) ok++;
      else erros.push(f.name);
    } catch(e){ erros.push(f.name); }
  }

  if(ok > 0){
    st.innerHTML = `<span class="status-ok">✓ ${ok} arquivo${ok>1?"s":""} adicionado${ok>1?"s":""} à base${erros.length ? ` · ${erros.length} erro(s)` : ""}</span>`; marcarPendente();
    _knowFiles = []; _knowFile = null;
    document.getElementById("know-pdf").value = "";
    document.getElementById("know-fname").textContent = "";
    document.getElementById("know-nome").value = "";
    document.getElementById("know-fonte").value = "";
    document.getElementById("know-st").textContent = "";
    document.querySelectorAll(".know-cls-btn.sel").forEach(b=>b.classList.remove("sel"));
    btn.disabled = true;
    carregarKnowledge();
    setTimeout(()=>st.textContent="", 5000);
  } else {
    st.innerHTML = `<span class="status-err">${erros.length} erro(s) ao enviar.</span>`;
    btn.disabled = false;
  }
}

async function carregarKnowledge(){
  try{
    const r = await fetch("/api/hp/knowledge");
    const docs = await r.json();
    const ct = document.getElementById("know-count");
    if(ct) ct.textContent = docs.length ? `${docs.length} documento${docs.length>1?"s":""}` : "";
    renderKnowledge(docs);
  } catch(e){ console.warn("knowledge:", e); }
}

function renderKnowledge(docs){
  const el = document.getElementById("know-lista");
  if(!el) return;
  if(!docs.length){
    el.innerHTML = `<div style="color:#1E4A30;font-size:11px;font-style:italic;text-align:center;padding:20px">Nenhum documento ainda. Suba o primeiro PDF acima.</div>`;
    return;
  }
  el.innerHTML = docs.map(d=>{
    const cor = KNOW_TIPO_COLORS[d.tipo]||"#888";
    const label = KNOW_TIPO_LABELS[d.tipo]||d.tipo;
    const clsTags = (d.classes||[]).map(c=>`<span style="font-size:10px;background:#1A1A08;color:#888;padding:2px 7px;border-radius:10px">${c}</span>`).join(" ");
    const pubBadge = d.publicado ? `<span style="font-size:10px;background:#0A2A0A;color:#5DCAA5;padding:2px 8px;border-radius:10px;border:1px solid #1A4A1A">✓ publicado</span>` : `<span style="font-size:10px;color:#2A5A3A">rascunho</span>`;
    return `<div class="know-doc-card">
      <div style="display:flex;align-items:flex-start;gap:12px">
        <div style="font-size:24px;flex-shrink:0">📄</div>
        <div style="flex:1;min-width:0">
          <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;margin-bottom:5px">
            <span style="font-size:13px;font-weight:700;color:#F0F0F0">${d.nome||"Documento"}</span>
            <span class="know-tipo-badge" style="background:${cor}22;color:${cor};border:1px solid ${cor}44">${label}</span>
            ${pubBadge}
          </div>
          <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;margin-bottom:8px">
            ${d.fonte ? `<span style="font-size:11px;color:#D4B483">📎 ${d.fonte}</span>` : ""}
            <span style="font-size:10px;color:#2A5A3A">${d.chars?.toLocaleString()||0} chars · ${d.data||""}</span>
            ${clsTags}
          </div>
          <!-- Ações -->
          <div style="display:flex;gap:8px;flex-wrap:wrap">
            <button class="btn btn-sm" onclick="knowAplicarCenario('${d.id}')">→ Cenário Macro</button>
            <button class="btn btn-sm btn-out" onclick="knowVerTexto('${d.id}')">👁 Ver texto</button>
            ${!d.publicado ? `<button class="btn-ghost" onclick="knowPublicar('${d.id}')">📢 Publicar</button>` : ""}
            <button class="btn-ghost" style="color:#FF6B6B;border-color:#FF6B6B" onclick="knowDeletar('${d.id}')">🗑 Remover</button>
          </div>
        </div>
      </div>
      <!-- Preview expansível -->
      <div id="know-preview-${d.id}" style="display:none;margin-top:10px;background:#050A05;border:1px solid #1A2A1A;border-radius:8px;padding:10px;font-size:11px;color:#AAA;font-family:monospace;white-space:pre-wrap;max-height:200px;overflow-y:auto"></div>
    </div>`;
  }).join("");
}

async function knowAplicarCenario(id){
  const st = document.getElementById("know-save-st");
  st.innerHTML = "⏳ Aplicando ao cenário...";
  try{
    const r = await fetch("/api/hp/knowledge/apply", {method:"POST", headers:{"Content-Type":"application/json"}, body:JSON.stringify({id, destino:"cenario"})});
    const d = await r.json();
    if(d.ok){
      if(d.global) document.getElementById("cenario-global").value = d.global;
      if(d.brasil) document.getElementById("cenario-brasil").value = d.brasil;
      if(d.posicionamento) document.getElementById("cenario-pos").value = d.posicionamento;
      if(d.fonte) document.getElementById("cenario-ref").value = d.fonte;
      st.innerHTML = `<span class="status-ok">✓ Cenário preenchido — revise e publique</span>`; marcarPendente();
      // Scroll até o cenário
      document.querySelector('[id="cenario-global"]')?.scrollIntoView({behavior:"smooth", block:"center"});
      setTimeout(()=>st.textContent="", 6000);
    } else { st.innerHTML=`<span class="status-err">Erro: ${d.error}</span>`; }
  } catch(e){ st.innerHTML=`<span class="status-err">Erro: ${e.message}</span>`; }
}

async function knowPublicar(id){
  try{
    const r = await fetch("/api/hp/knowledge/publicar", {method:"POST", headers:{"Content-Type":"application/json"}, body:JSON.stringify({id})});
    const d = await r.json();
    if(d.ok) carregarKnowledge();
  } catch(e){ console.warn(e); }
}

async function knowDeletar(id){
  if(!confirm("Remover este documento da base?")) return;
  try{
    await fetch("/api/hp/knowledge/delete", {method:"POST", headers:{"Content-Type":"application/json"}, body:JSON.stringify({id})});
    carregarKnowledge();
  } catch(e){ console.warn(e); }
}

function knowVerTexto(id){
  const el = document.getElementById(`know-preview-${id}`);
  if(!el) return;
  if(el.style.display === "none"){
    // busca o texto da lista em cache
    fetch("/api/hp/knowledge").then(r=>r.json()).then(docs=>{
      const doc = docs.find(d=>d.id===id);
      if(doc){ el.textContent = doc.texto||"(sem texto)"; el.style.display = "block"; }
    });
  } else { el.style.display = "none"; }
}

// ── Alertas de Monitoramento ──────────────────────────────────────────────────
const ALERT_COLORS = { info:"#5DCAA5", atencao:"#D4B483", urgente:"#FF6B6B" };
const ALERT_LABELS = { info:"ℹ️ Informativo", atencao:"⚠️ Atenção", urgente:"🚨 Urgente" };

async function registrarAlerta(){
  const produto          = document.getElementById("alert-produto").value.trim();
  const classe           = document.getElementById("alert-classe").value;
  const tipo             = document.getElementById("alert-tipo").value;
  const mensagem         = document.getElementById("alert-msg").value.trim();
  const assessor_destino = (document.getElementById("alert-assessor")?.value||"").trim().toLowerCase();
  const st               = document.getElementById("alert-st");
  if(!produto || !mensagem){ st.innerHTML=`<span class="status-err">Preencha produto e mensagem.</span>`; return; }
  try{
    const r = await fetch("/api/hp/alertas", {method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify({produto, classe, tipo, mensagem, origem:"Head de Produtos", assessor_destino})});
    const d = await r.json();
    if(d.ok){
      st.innerHTML = `<span class="status-ok">✓ Alerta registrado${assessor_destino?" para "+assessor_destino:""}</span>`; marcarPendente();
      document.getElementById("alert-produto").value = "";
      document.getElementById("alert-msg").value = "";
      document.getElementById("alert-tipo").value = "info";
      document.getElementById("alert-classe").value = "";
      const elAss = document.getElementById("alert-assessor");
      if(elAss) elAss.value = "";
      carregarAlertas();
    }
  } catch(e){ st.innerHTML=`<span class="status-err">Erro: ${e.message}</span>`; }
  setTimeout(()=>st.textContent="",4000);
}

async function carregarAlertas(){
  const feed = document.getElementById("alertas-feed");
  try{
    const r = await fetch("/api/hp/alertas");
    const lista = await r.json();
    if(!lista.length){
      feed.innerHTML = `<div style="color:#1E4A30;font-size:11px;text-align:center;padding:20px">Nenhum alerta registrado ainda.</div>`;
      return;
    }
    feed.innerHTML = lista.map(a=>`
      <div style="background:#0A0A0A;border-left:3px solid ${ALERT_COLORS[a.tipo]||"#2A5A3A"};border-radius:0 8px 8px 0;padding:12px 14px;display:flex;align-items:flex-start;gap:12px">
        <div style="flex:1">
          <div style="display:flex;align-items:center;gap:8px;margin-bottom:4px;flex-wrap:wrap">
            <span style="font-size:12px;font-weight:700;color:${ALERT_COLORS[a.tipo]||"#CCC"}">${ALERT_LABELS[a.tipo]||a.tipo}</span>
            <span style="font-size:11px;color:#D4B483;font-weight:700">${a.produto||""}</span>
            ${a.classe ? `<span style="font-size:10px;background:#1A1A08;color:#888;padding:2px 7px;border-radius:10px">${a.classe}</span>` : ""}
            ${a.assessor_destino ? `<span style="background:#1A2E1A;color:#C9A96E;font-size:10px;padding:2px 7px;border-radius:10px">→ ${a.assessor_destino}</span>` : ""}
            <span style="font-size:10px;color:#1E4A30;margin-left:auto">${a.data||""} · ${a.origem||"HP"}</span>
          </div>
          <p style="font-size:12px;color:#AAA;margin:0;line-height:1.5">${a.mensagem||""}</p>
        </div>
        <button onclick="deletarAlerta('${a.id}')" title="Remover" style="background:none;border:none;color:#1E4A30;cursor:pointer;font-size:14px;flex-shrink:0;padding:2px 4px" onmouseover="this.style.color='#FF6B6B'" onmouseout="this.style.color='#1E4A30'">✕</button>
      </div>
    `).join("");
  } catch(e){ feed.innerHTML=`<div style="color:#FF6B6B;font-size:11px;padding:12px">Erro ao carregar alertas.</div>`; }
}

async function deletarAlerta(id){
  await fetch("/api/hp/alertas", {method:"DELETE", headers:{"Content-Type":"application/json"}, body: JSON.stringify({id})});
  carregarAlertas();
}

// ── Calls de Ações ────────────────────────────────────────────────────────────
let _callDir = "compra";
function setCallDir(dir){
  _callDir = dir;
  document.getElementById("call-direcao").value = dir;
  const styles = {
    compra: {border:"#5DCAA5", bg:"#0A2018", color:"#5DCAA5"},
    neutro: {border:"#888",    bg:"#1A1A1A", color:"#AAA"},
    venda:  {border:"#FF6B6B", bg:"#1A0808", color:"#FF6B6B"},
  };
  ["compra","neutro","venda"].forEach(d=>{
    const btn = document.getElementById("call-dir-"+d);
    if(!btn) return;
    const s = styles[d];
    const active = d === dir;
    btn.style.borderColor = active ? s.border : "#1A1A1A";
    btn.style.background  = active ? s.bg     : "#111";
    btn.style.color       = active ? s.color  : "#444";
  });
}
setCallDir("compra");

function calcCallUpside(){
  const entrada = parseFloat(document.getElementById("call-entrada").value) || 0;
  const alvo    = parseFloat(document.getElementById("call-alvo").value)    || 0;
  const el      = document.getElementById("call-upside-display");
  if(entrada > 0 && alvo > 0){
    const up = ((alvo / entrada - 1) * 100).toFixed(1);
    const cor = up >= 0 ? "#5DCAA5" : "#FF6B6B";
    el.textContent = (up >= 0 ? "+" : "") + up + "%";
    el.style.color = cor;
  } else {
    el.textContent = "—"; el.style.color = "#5DCAA5";
  }
}

async function adicionarCall(){
  const ticker  = document.getElementById("call-ticker").value.trim().toUpperCase();
  const nome    = document.getElementById("call-nome").value.trim();
  const direcao = document.getElementById("call-direcao").value;
  const entrada = parseFloat(document.getElementById("call-entrada").value) || 0;
  const alvo    = parseFloat(document.getElementById("call-alvo").value)    || 0;
  const stop    = parseFloat(document.getElementById("call-stop").value)    || 0;
  const prazo   = document.getElementById("call-prazo").value.trim();
  const tese    = document.getElementById("call-tese").value.trim();
  const fonte   = document.getElementById("call-fonte").value.trim();
  const st      = document.getElementById("call-st");
  const perfis  = [];
  ["arrojada","agressiva","moderada","conservadora"].forEach(p=>{
    if(document.getElementById("call-perf-"+p)?.checked) perfis.push(p);
  });
  if(!ticker){ st.innerHTML=`<span class="status-err">Informe o ticker.</span>`; return; }
  if(entrada <= 0){ st.innerHTML=`<span class="status-err">Informe o preço de entrada.</span>`; return; }
  try{
    const r = await fetch("/api/hp/calls", {
      method:"POST", headers:{"Content-Type":"application/json"},
      body: JSON.stringify({ticker, nome, direcao, preco_entrada:entrada, preco_alvo:alvo, stop, prazo, tese, perfis, fonte})
    });
    const d = await r.json();
    if(d.ok){
      st.innerHTML = `<span class="status-ok">✓ Call ${ticker} adicionado</span>`; marcarPendente();
      ["call-ticker","call-nome","call-entrada","call-alvo","call-stop","call-prazo","call-tese","call-fonte"].forEach(id=>{ const el=document.getElementById(id); if(el) el.value=""; });
      ["arrojada","agressiva","moderada","conservadora"].forEach(p=>{ const el=document.getElementById("call-perf-"+p); if(el) el.checked=false; });
      document.getElementById("call-upside-display").textContent="—";
      setCallDir("compra");
      carregarCalls();
    }
  } catch(e){ st.innerHTML=`<span class="status-err">Erro: ${e.message}</span>`; }
  setTimeout(()=>st.textContent="", 4000);
}

async function carregarCalls(){
  const feed = document.getElementById("calls-feed");
  if(!feed) return;
  try{
    const r = await fetch("/api/hp/calls");
    const lista = await r.json();
    if(!lista.length){
      feed.innerHTML = `<div style="color:#1E4A30;font-size:11px;text-align:center;padding:20px">Nenhum call cadastrado ainda.</div>`;
      return;
    }
    const DIR_COR   = {compra:"#5DCAA5", neutro:"#AAA", venda:"#FF6B6B"};
    const DIR_LABEL = {compra:"▲ COMPRA", neutro:"→ NEUTRO", venda:"▼ VENDA"};
    feed.innerHTML = lista.filter(c=>c.ativo!==false).map(c=>{
      const cor = DIR_COR[c.direcao] || "#AAA";
      const upCor = (c.upside||0) >= 0 ? "#5DCAA5" : "#FF6B6B";
      const perfBadges = (c.perfis||[]).map(p=>`<span style="font-size:10px;background:#1A2E1A;color:#C9A96E;border:1px solid #C9A96E44;padding:1px 7px;border-radius:8px">${p}</span>`).join(" ");
      return `<div style="background:#0A0F0B;border:1px solid #1A2E3A;border-left:3px solid ${cor};border-radius:0 10px 10px 0;padding:12px 14px;display:flex;align-items:flex-start;gap:12px">
        <div style="flex:1">
          <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;margin-bottom:6px">
            <span style="font-size:15px;font-weight:900;color:${cor}">${c.ticker}</span>
            <span style="font-size:10px;font-weight:700;padding:2px 8px;border-radius:8px;background:${cor}22;color:${cor}">${DIR_LABEL[c.direcao]||c.direcao}</span>
            ${c.nome ? `<span style="font-size:12px;color:#AAA">${c.nome}</span>` : ""}
            ${perfBadges}
            <span style="font-size:10px;color:#1E4A30;margin-left:auto">${c.data||""}</span>
          </div>
          <div style="display:flex;gap:14px;flex-wrap:wrap;font-size:12px;margin-bottom:6px">
            ${c.preco_entrada ? `<span style="color:#888">Entrada: <b style="color:#F0F0F0">R$ ${c.preco_entrada.toLocaleString("pt-BR",{minimumFractionDigits:2})}</b></span>` : ""}
            ${c.preco_alvo ? `<span style="color:#888">→ Alvo: <b style="color:#F0F0F0">R$ ${c.preco_alvo.toLocaleString("pt-BR",{minimumFractionDigits:2})}</b></span>` : ""}
            ${c.upside !== undefined ? `<span style="color:${upCor};font-weight:700">Upside ${c.upside >= 0 ? "+" : ""}${c.upside}%</span>` : ""}
            ${c.stop ? `<span style="color:#888">Stop: <b style="color:#FF6B6B">R$ ${c.stop.toLocaleString("pt-BR",{minimumFractionDigits:2})}</b></span>` : ""}
            ${c.prazo ? `<span style="color:#888">Prazo: <b style="color:#AAA">${c.prazo}</b></span>` : ""}
          </div>
          ${c.tese ? `<p style="font-size:11px;color:#888;margin:0 0 4px;line-height:1.5;font-style:italic">"${c.tese.length>120?c.tese.slice(0,120)+"...":c.tese}"</p>` : ""}
          ${c.fonte ? `<div style="font-size:10px;color:#2A5A3A">${c.fonte}</div>` : ""}
        </div>
        <button onclick="deletarCall('${c.id}')" title="Remover" style="background:none;border:none;color:#1E4A30;cursor:pointer;font-size:14px;flex-shrink:0;padding:2px 4px" onmouseover="this.style.color='#FF6B6B'" onmouseout="this.style.color='#1E4A30'">✕</button>
      </div>`;
    }).join("");
  } catch(e){ feed.innerHTML=`<div style="color:#FF6B6B;font-size:11px;padding:12px">Erro ao carregar calls.</div>`; }
}

async function deletarCall(id){
  await fetch("/api/hp/calls", {method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify({id, _delete:true})});
  carregarCalls();
}

// ── Operações Estruturadas ────────────────────────────────────────────────────
function toggleEstruturadas(){
  const p = document.getElementById("painel-estruturadas");
  const i = document.getElementById("icon-estr");
  const aberto = p.style.display !== "none";
  p.style.display = aberto ? "none" : "";
  i.textContent = aberto ? "▶" : "▼";
}

async function adicionarEstruturada(){
  const st = document.getElementById("estr-st");
  const payload = {
    tipo:         document.getElementById("estr-tipo").value,
    ativo:        document.getElementById("estr-ativo").value.trim(),
    emissor:      document.getElementById("estr-emissor").value.trim(),
    vencimento:   document.getElementById("estr-venc").value.trim(),
    retorno:      document.getElementById("estr-retorno").value.trim(),
    perfil_minimo:document.getElementById("estr-perfil").value,
    observacao:   document.getElementById("estr-obs").value.trim(),
    minimo:       document.getElementById("estr-minimo").value.trim(),
  };
  if(!payload.ativo){ st.innerHTML=`<span class="status-err">Informe o ativo/produto.</span>`; return; }
  st.innerHTML = `<span style="color:#4A7055">Salvando...</span>`;
  try{
    const r = await fetch("/api/hp/estruturadas", {method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify(payload)});
    const d = await r.json();
    if(d.ok){
      st.innerHTML = `<span class="status-ok">✓ Operação adicionada.</span>`; marcarPendente();
      ["estr-ativo","estr-emissor","estr-venc","estr-retorno","estr-obs","estr-minimo"].forEach(id=>{
        const el=document.getElementById(id); if(el) el.value="";
      });
      carregarEstruturadas();
    } else { st.innerHTML=`<span class="status-err">Erro: ${d.error||"desconhecido"}</span>`; }
  } catch(e){ st.innerHTML=`<span class="status-err">Erro: ${e.message}</span>`; }
  setTimeout(()=>st.textContent="", 5000);
}

async function carregarEstruturadas(){
  const lista = document.getElementById("estr-lista");
  if(!lista) return;
  try{
    const r = await fetch("/api/hp/estruturadas");
    const ops = await r.json();
    if(!ops.length){
      lista.innerHTML=`<div style="color:#2A2A18;font-size:11px;text-align:center;padding:16px">Nenhuma operação cadastrada ainda.</div>`;
      return;
    }
    const PERFIL_LABEL = {super_conservadora:"Super Conserv.",conservadora:"Conservadora",moderada:"Moderada",agressiva:"Agressiva",super_agressiva:"Super Agressiva"};
    lista.innerHTML = ops.map(op=>`
      <div style="background:#080C06;border:1px solid #2A1A10;border-left:3px solid #C9A96E;border-radius:0 10px 10px 0;padding:12px 14px;margin-bottom:8px;display:flex;gap:12px;align-items:flex-start">
        <div style="flex:1">
          <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;margin-bottom:6px">
            <span style="font-size:12px;font-weight:700;background:#C9A96E22;color:#C9A96E;padding:2px 8px;border-radius:8px">${op.tipo}</span>
            <span style="font-size:14px;font-weight:700;color:#F0F0F0">${op.ativo}</span>
            ${op.emissor ? `<span style="font-size:11px;color:#888">${op.emissor}</span>` : ""}
            <span style="font-size:10px;color:#1E4A30;margin-left:auto">${op.data||""}</span>
          </div>
          <div style="display:flex;gap:14px;flex-wrap:wrap;font-size:11px;margin-bottom:4px">
            ${op.vencimento ? `<span style="color:#888">Venc: <b style="color:#DDD">${op.vencimento}</b></span>` : ""}
            ${op.retorno ? `<span style="color:#888">Retorno: <b style="color:#5DCAA5">${op.retorno}</b></span>` : ""}
            ${op.perfil_minimo ? `<span style="color:#888">Mínimo: <b style="color:#C9A96E">${PERFIL_LABEL[op.perfil_minimo]||op.perfil_minimo}</b></span>` : ""}
            ${op.minimo ? `<span style="color:#888">Apl. mín: <b style="color:#DDD">${op.minimo}</b></span>` : ""}
          </div>
          ${op.observacao ? `<p style="font-size:11px;color:#888;margin:0;line-height:1.5;font-style:italic">"${op.observacao}"</p>` : ""}
        </div>
        <button onclick="deletarEstruturada('${op.id}')" title="Remover" style="background:none;border:none;color:#2A2A18;cursor:pointer;font-size:14px;flex-shrink:0;padding:2px 4px" onmouseover="this.style.color='#FF6B6B'" onmouseout="this.style.color='#2A2A18'">✕</button>
      </div>`).join("");
  } catch(e){ lista.innerHTML=`<div style="color:#FF6B6B;font-size:11px;padding:12px">Erro ao carregar operações.</div>`; }
}

async function deletarEstruturada(id){
  await fetch("/api/hp/estruturadas", {method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify({id, _delete:true})});
  carregarEstruturadas();
}

// ── Alocações de Referência — Gestoras 2 ─────────────────────────────────────
let _gest2_data = {};
let _gest2_perfil_ativo = "super_conservadora";
let _gest2_perfis_temp  = {};

const G2_CLASSES = ["pos_fixado","inflacao","pre_fixado","acoes","fiis","multimercado","internacional","alternativos","criptomoedas"];
const G2_PERFIS  = ["super_conservadora","conservadora","moderada","agressiva","super_agressiva"];
const G2_PERFIL_LABEL = {super_conservadora:"Super Conserv.",conservadora:"Conservadora",moderada:"Moderada",agressiva:"Agressiva",super_agressiva:"Super Agressiva"};

async function carregarGest2(){
  try{
    const r = await fetch("/api/hp/gestoras2");
    _gest2_data = await r.json();
    renderGest2Tabs();
    renderGest2Painel();
    atualizarGest2Count();
  } catch(e){}
}

function atualizarGest2Count(){
  const n = Object.keys(_gest2_data).length;
  const el = document.getElementById("gest2-count-label");
  if(el) el.textContent = `${n}/5 gestoras cadastradas`;
  const btnAdd = document.getElementById("btn-gest2-add");
  if(btnAdd) btnAdd.style.display = n >= 5 ? "none" : "";
}

function renderGest2Tabs(){
  const tabs = document.getElementById("gest2-tabs");
  if(!tabs) return;
  const gestoras = Object.values(_gest2_data);
  tabs.innerHTML = gestoras.map(g=>`
    <button onclick="gest2ShowGestora('${g.id}')" style="font-size:12px;padding:8px 14px;border:none;border-bottom:2px solid transparent;background:none;color:#4A7055;cursor:pointer;transition:all .15s" id="gest2tab-${g.id}">
      🏦 ${g.nome||g.id}
    </button>`).join("") + (gestoras.length ? "" : `<span style="font-size:11px;color:#1E3A2A;padding:8px 14px">Nenhuma gestora ainda.</span>`);
}

let _gest2_ativa = null;
function gest2ShowGestora(gid){
  _gest2_ativa = gid;
  const g = _gest2_data[gid];
  // highlight tab
  document.querySelectorAll("[id^='gest2tab-']").forEach(b=>{
    b.style.borderBottomColor = "transparent";
    b.style.color = "#4A7055";
    b.style.fontWeight = "400";
  });
  const tab = document.getElementById("gest2tab-"+gid);
  if(tab){ tab.style.borderBottomColor="#D4B483"; tab.style.color="#D4B483"; tab.style.fontWeight="700"; }

  const painel = document.getElementById("gest2-painel");
  if(!g){ painel.innerHTML=""; return; }

  const CLASSE_LABEL = {pos_fixado:"Pós Fixado",inflacao:"Inflação",pre_fixado:"Pré Fixado",acoes:"Ações",fiis:"FIIs",multimercado:"Multimercado",internacional:"Internacional",alternativos:"Alternativos",criptomoedas:"Criptomoedas"};
  painel.innerHTML = `
    <div style="overflow-x:auto;margin-bottom:10px">
      <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:10px">
        <div>
          <span style="font-size:14px;font-weight:700;color:#D4B483">${g.nome}</span>
          ${g.referencia ? `<span style="font-size:11px;color:#3A6A48;margin-left:10px">${g.referencia}</span>` : ""}
        </div>
        <div style="display:flex;gap:8px">
          <button class="btn-ghost" onclick="gest2EditarGestora('${gid}')" style="font-size:11px">✏️ Editar</button>
          <button class="btn-ghost" onclick="gest2DeletarGestora('${gid}')" style="font-size:11px;color:#FF6B6B">✕ Remover</button>
        </div>
      </div>
      <table style="width:100%;border-collapse:collapse;font-size:11px">
        <thead>
          <tr style="background:#060F0B">
            <th style="padding:8px;text-align:left;color:#3A6A48;border-bottom:1px solid #1A2A1A">Classe</th>
            ${G2_PERFIS.map(p=>`<th style="padding:8px;text-align:center;color:#3A6A48;border-bottom:1px solid #1A2A1A">${G2_PERFIL_LABEL[p]}</th>`).join("")}
          </tr>
        </thead>
        <tbody>
          ${G2_CLASSES.map(c=>{
            const vals = G2_PERFIS.map(p=>g.perfis?.[p]?.[c]||0);
            const allZero = vals.every(v=>v===0);
            return `<tr style="border-bottom:1px solid #0A1A0A${allZero?";opacity:.4":""}">
              <td style="padding:7px 8px;color:#4A7055;font-weight:600">${CLASSE_LABEL[c]||c}</td>
              ${G2_PERFIS.map(p=>{
                const v = g.perfis?.[p]?.[c]||0;
                return `<td style="padding:7px 8px;text-align:center;color:${v>0?"#F0F0F0":"#1E3A2A"};font-weight:${v>0?"700":"400"}">${v>0?v+"%":"—"}</td>`;
              }).join("")}
            </tr>`;
          }).join("")}
        </tbody>
        <tfoot>
          <tr style="border-top:2px solid #1A2A1A;background:#060F0B">
            <td style="padding:8px;color:#D4B483;font-weight:700">TOTAL</td>
            ${G2_PERFIS.map(p=>{
              const soma = G2_CLASSES.reduce((s,c)=>s+(g.perfis?.[p]?.[c]||0),0);
              const ok = Math.abs(soma-100)<0.1;
              return `<td style="padding:8px;text-align:center;font-weight:700;color:${ok?"#5DCAA5":"#FF6B6B"}">${soma.toFixed(1)}%</td>`;
            }).join("")}
          </tr>
        </tfoot>
      </table>
    </div>`;
}

function renderGest2Painel(){
  const gestoras = Object.values(_gest2_data);
  if(!gestoras.length){
    document.getElementById("gest2-painel").innerHTML = `<div style="color:#1E3A2A;font-size:11px;text-align:center;padding:20px">Adicione a primeira gestora clicando em ➕ Adicionar Gestora.</div>`;
    return;
  }
  if(!_gest2_ativa && gestoras.length) gest2ShowGestora(gestoras[0].id);
}

function gest2MostrarForm(g){
  const form = document.getElementById("gest2-form");
  form.style.display = "";
  document.getElementById("gest2-form-titulo").textContent = g ? "✏️ Editar Gestora" : "➕ Nova Gestora";
  document.getElementById("gest2-id").value    = g ? g.id    : "";
  document.getElementById("gest2-nome").value  = g ? g.nome  : "";
  document.getElementById("gest2-ref").value   = g ? g.referencia : "";
  _gest2_perfis_temp = g ? JSON.parse(JSON.stringify(g.perfis || {})) : {};
  _gest2_perfil_ativo = "super_conservadora";
  gest2SwitchPerfil("super_conservadora", document.querySelector(".gest2-ptab"));
}

function gest2EditarGestora(gid){
  gest2MostrarForm(_gest2_data[gid]);
}

async function gest2DeletarGestora(gid){
  if(!confirm("Remover esta gestora?")) return;
  await fetch("/api/hp/gestoras2", {method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify({_delete:true, id:gid})});
  _gest2_ativa = null;
  carregarGest2();
}

function gest2SwitchPerfil(p, btn){
  // salvar valores do perfil anterior
  const prev = _gest2_perfil_ativo;
  const prevData = {};
  G2_CLASSES.forEach(c=>{ const el=document.getElementById("g2aloc-"+c); if(el) prevData[c]=parseFloat(el.value)||0; });
  _gest2_perfis_temp[prev] = prevData;
  // trocar aba
  document.querySelectorAll(".gest2-ptab").forEach(b=>b.classList.remove("active"));
  if(btn) btn.classList.add("active");
  _gest2_perfil_ativo = p;
  // carregar valores do novo perfil
  const perfData = _gest2_perfis_temp[p] || {};
  G2_CLASSES.forEach(c=>{ const el=document.getElementById("g2aloc-"+c); if(el) el.value = perfData[c]||0; });
  g2Total();
}

function g2Total(){
  let s=0;
  G2_CLASSES.forEach(c=>{ const el=document.getElementById("g2aloc-"+c); if(el) s+=parseFloat(el.value)||0; });
  const el=document.getElementById("g2-total");
  if(el){ el.textContent=s.toFixed(1); el.style.color=Math.abs(s-100)<0.1?"#5DCAA5":"#FF6B6B"; }
}

async function gest2Salvar(){
  const st = document.getElementById("gest2-st");
  // save current tab values
  const cur = {}; G2_CLASSES.forEach(c=>{ const el=document.getElementById("g2aloc-"+c); if(el) cur[c]=parseFloat(el.value)||0; });
  _gest2_perfis_temp[_gest2_perfil_ativo] = cur;
  const payload = {
    id:         document.getElementById("gest2-id").value || undefined,
    nome:       document.getElementById("gest2-nome").value.trim(),
    referencia: document.getElementById("gest2-ref").value.trim(),
    perfis:     _gest2_perfis_temp,
  };
  if(!payload.nome){ st.innerHTML=`<span class="status-err">Informe o nome.</span>`; return; }
  st.innerHTML=`<span style="color:#4A7055">Salvando...</span>`;
  try{
    const r = await fetch("/api/hp/gestoras2", {method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify(payload)});
    const d = await r.json();
    if(d.ok){
      st.innerHTML=`<span class="status-ok">✓ Salvo.</span>`; marcarPendente();
      document.getElementById("gest2-form").style.display="none";
      _gest2_ativa = d.id;
      carregarGest2();
    } else { st.innerHTML=`<span class="status-err">Erro: ${d.error||"desconhecido"}</span>`; }
  } catch(e){ st.innerHTML=`<span class="status-err">Erro: ${e.message}</span>`; }
  setTimeout(()=>st.textContent="",5000);
}

// ── Salvar só produtos ────────────────────────────────────────────────────────
async function salvarProdutos(){
  const btn = document.getElementById("btn-salvar-prod");
  const st  = document.getElementById("prod-save-st");
  btn.disabled = true; btn.textContent = "Salvando...";
  const prods = coletarProdutos();
  const total = Object.values(prods).reduce((s,a)=>s+a.length, 0);
  try{
    const r = await fetch("/api/hp/produtos/salvar", {
      method:"POST",
      headers:{"Content-Type":"application/json"},
      body: JSON.stringify({ produtos: prods, indicado_por: "Head de Produtos" })
    });
    const d = await r.json();
    if(d.ok){
      st.innerHTML = `<span class="status-ok">✓ ${d.total} produto(s) salvos · indicado por <strong>${d.indicado_por}</strong> em ${d.indicado_em}</span>`; marcarPendente();
    } else { st.innerHTML = `<span class="status-err">Erro ao salvar.</span>`; }
  } catch(e){ st.innerHTML = `<span class="status-err">Erro: ${e.message}</span>`; }
  finally{ btn.disabled=false; btn.textContent="💾 Salvar produtos"; }
  setTimeout(()=>st.textContent="", 6000);
}

// ── Publicar ──────────────────────────────────────────────────────────────────
function _pubProgress(pct, cor){
  const wrap = document.getElementById("pub-progress-wrap");
  const bar  = document.getElementById("pub-progress-bar");
  if(!wrap || !bar) return;
  wrap.style.display = "";
  bar.style.width = pct + "%";
  if(cor) bar.style.background = cor;
}
function _pubProgressHide(){
  const wrap = document.getElementById("pub-progress-wrap");
  const bar  = document.getElementById("pub-progress-bar");
  if(wrap && bar){
    bar.style.width = "100%";
    bar.style.background = "linear-gradient(90deg,#5DCAA5,#C9A96E)";
    setTimeout(()=>{ if(wrap) wrap.style.display="none"; bar.style.width="0%"; }, 800);
  }
}

async function publicar(){
  const btn = document.getElementById("btn-pub");
  const st  = document.getElementById("pub-st");
  btn.disabled = true;
  st.innerHTML = `<span style="color:#D4B483;font-size:11px">⏳ Publicando...</span>`;

  // Progresso simulado: avança até 85% enquanto aguarda resposta do servidor
  // Timing baseado no cold start típico do Vercel (~6-8s)
  _pubProgress(5);
  const etapas = [
    {pct:15, delay:400,  label:"Preparando dados..."},
    {pct:30, delay:1000, label:"Enviando cenário macro..."},
    {pct:50, delay:2000, label:"Salvando portfólios..."},
    {pct:65, delay:3500, label:"Salvando produtos..."},
    {pct:78, delay:5000, label:"Aguardando servidor..."},
    {pct:85, delay:7000, label:"Quase pronto..."},
  ];
  const timers = etapas.map(e => setTimeout(()=>{
    _pubProgress(e.pct);
    st.innerHTML = `<span style="color:#D4B483;font-size:11px">⏳ ${e.label}</span>`;
  }, e.delay));

  const payload = {
    portfolios: coletarPortfolios(),
    cenario:    coletarCenario(),
    produtos:   coletarProdutos(),
  };

  try{
    const r = await fetch("/api/hp/publicar", {method:"POST", headers:{"Content-Type":"application/json"}, body: JSON.stringify(payload)});
    const d = await r.json();
    timers.forEach(t=>clearTimeout(t));
    if(d.ok){
      _pubProgress(100, "linear-gradient(90deg,#5DCAA5,#C9A96E)");
      st.innerHTML = `<span class="status-ok">✓ Publicado em ${d.publicado_em}</span>`;
      const bar = document.getElementById("status-pub");
      if(bar){ bar.style.display="block"; bar.innerHTML=`<b>✓ Última publicação:</b> ${d.publicado_em}`; }
      limparPendente();
      setTimeout(()=>_pubProgressHide(), 600);
    } else {
      timers.forEach(t=>clearTimeout(t));
      _pubProgress(100, "#FF4444");
      st.innerHTML = `<span class="status-err">Erro ao publicar.</span>`;
      setTimeout(()=>_pubProgressHide(), 1000);
    }
  } catch(e){
    timers.forEach(t=>clearTimeout(t));
    _pubProgress(100, "#FF4444");
    st.innerHTML = `<span class="status-err">Erro de conexão. Tente novamente.</span>`;
    setTimeout(()=>_pubProgressHide(), 1000);
  }
  finally{ btn.disabled=false; btn.textContent="📤 Publicar agora"; }
  setTimeout(()=>st.textContent="", 8000);
}

// ── Upload Rápido ─────────────────────────────────────────────────────────────
let _arquivosRapido = [];

function uploadRapidoSelecionar(input){
  const files = Array.from(input.files);
  if(!files.length) return;
  const MAX = 5;
  // Merge novos arquivos com existentes (sem duplicar por nome)
  const novos = files.filter(f => !_arquivosRapido.some(a => a.file.name === f.name));
  const total = _arquivosRapido.length + novos.length;
  if(total > MAX){ alert(`Máximo de ${MAX} arquivos. Você já tem ${_arquivosRapido.length}.`); input.value=""; return; }
  novos.forEach(f => _arquivosRapido.push({file: f, status: "pendente", texto: "", chars: 0}));
  renderListaRapido();
  // Mostra seção de lista, esconde zona vazia
  document.getElementById("rapido-zona-vazia").style.display = "none";
  document.getElementById("pdf-rapido-itens").style.display = "";
  document.getElementById("btn-limpar-rapido").style.display = "";
  document.getElementById("pdf-rapido-contador").textContent =
    `${_arquivosRapido.length} arquivo${_arquivosRapido.length>1?"s":""}`;
  input.value = ""; // permite re-selecionar mesmos arquivos
}

function renderListaRapido(){
  const el = document.getElementById("pdf-rapido-lista");
  el.innerHTML = _arquivosRapido.map((item, i) => {
    const cor = item.status==="ok" ? "#5DCAA5" : item.status==="erro" ? "#FF6B6B" : item.status==="processando" ? "#D4B483" : "#4A7055";
    const icone = item.status==="ok" ? "✅" : item.status==="erro" ? "❌" : item.status==="processando" ? "⏳" : "📄";
    const tamanho = `${(item.file.size/1024).toFixed(0)} KB`;
    const statusTxt = item.status==="ok" ? `Salvo · ${item.chars} caracteres` : item.status==="erro" ? `Erro: ${item.erro}` : item.status==="processando" ? "Enviando..." : `Aguardando upload · ${tamanho}`;
    return `<div style="display:flex;align-items:center;gap:12px;padding:12px 16px;background:#0A140A;border:1px solid ${item.status==='ok'?'#1A3A1A':'#1A2A1A'};border-radius:8px;margin-bottom:6px">
      <span style="font-size:20px;flex-shrink:0">${icone}</span>
      <div style="flex:1;min-width:0">
        <div style="font-size:13px;color:#D4B483;font-weight:600;white-space:nowrap;overflow:hidden;text-overflow:ellipsis">${item.file.name}</div>
        <div style="font-size:10px;color:${cor};margin-top:3px">${statusTxt}</div>
      </div>
      <div style="display:flex;gap:6px;flex-shrink:0">
        <button class="btn-ghost" onclick="_arquivosRapido.splice(${i},1);if(!_arquivosRapido.length){fecharRapido();}else{renderListaRapido();document.getElementById('pdf-rapido-contador').textContent=_arquivosRapido.length+' arquivo'+(_arquivosRapido.length>1?'s':'');}" style="font-size:11px;padding:4px 8px;color:#FF6B6B" title="Remover">✕</button>
      </div>
    </div>`;
  }).join("");

  // Atualiza botão principal
  const btn = document.getElementById("btn-extrair-todos");
  const pendentes = _arquivosRapido.filter(a=>a.status==="pendente").length;
  const todos_ok = _arquivosRapido.every(a=>a.status==="ok");
  btn.textContent = todos_ok ? "✅ Todos enviados" : `📤 Subir para a Memória${pendentes>0?" ("+pendentes+")":""}`;
  btn.disabled = todos_ok || pendentes===0;
  btn.style.opacity = (todos_ok || pendentes===0) ? ".5" : "1";

  // Status geral
  const st = document.getElementById("rapido-upload-st");
  const ok = _arquivosRapido.filter(a=>a.status==="ok").length;
  if(ok>0) st.textContent = `${ok} de ${_arquivosRapido.length} salvo${ok>1?"s":""}`;
}

async function uploadRapidoUm(i){
  const item = _arquivosRapido[i];
  if(item.file.size > 4.3 * 1024 * 1024){
    item.status = "erro";
    item.erro = `Arquivo muito grande (${(item.file.size/1024/1024).toFixed(1)} MB). Limite: 4,3 MB.`;
    renderListaRapido();
    return;
  }
  item.status = "processando";
  renderListaRapido();
  try{
    const fd = new FormData();
    fd.append("pdf", item.file);
    fd.append("nome", item.file.name.replace(/\.pdf$/i,""));
    fd.append("tipo", "outro");
    fd.append("classes", "geral");
    fd.append("fonte", "");
    const r = await fetch("/api/hp/knowledge/upload", {method:"POST", body:fd});
    const txt = await r.text();
    let d;
    try { d = JSON.parse(txt); }
    catch(_){ item.status="erro"; item.erro = r.status===413?"Arquivo muito grande.": `Erro ${r.status} do servidor.`; renderListaRapido(); return; }
    if(d.ok){
      item.status = "ok"; item.chars = d.chars||0; item.doc_id = d.id||"";
      carregarKnowledge();
    } else { item.status="erro"; item.erro = d.error||"falha no servidor"; }
  } catch(e){ item.status="erro"; item.erro = e.message; }
  renderListaRapido();
}

function rapidoAplicarCenario(i){
  const c = _arquivosRapido[i]?.cenario;
  if(!c) return;
  aplicarCenarioNaPagina(c);
}

async function uploadCartaCenario(input){
  const file = input.files[0];
  if(!file) return;
  input.value = "";
  const st = document.getElementById("carta-st");
  st.innerHTML = `<span style="color:#D4B483">⏳ Processando carta...</span>`;
  if(file.size > 4.3*1024*1024){ st.innerHTML=`<span style="color:#FF6B6B">Arquivo muito grande (máx 4,3 MB)</span>`; return; }
  try{
    const fd = new FormData();
    fd.append("pdf", file);
    fd.append("nome", file.name.replace(/\.pdf$/i,""));
    const r = await fetch("/api/hp/carta-upload", {method:"POST", body:fd});
    const txt = await r.text();
    let d;
    try{ d = JSON.parse(txt); } catch(_){ st.innerHTML=`<span style="color:#FF6B6B">Erro ${r.status} do servidor</span>`; return; }
    if(d.ok){
      aplicarCenarioNaPagina({global:d.global, brasil:d.brasil, posicionamento:d.posicionamento, referencia:d.fonte});
      st.innerHTML = `<span style="color:#5DCAA5">✅ Cenário atualizado${d.ia?" (IA)":""}</span>`;
      setTimeout(()=>st.textContent="", 5000);
    } else { st.innerHTML=`<span style="color:#FF6B6B">${d.error}</span>`; }
  } catch(e){ st.innerHTML=`<span style="color:#FF6B6B">${e.message}</span>`; }
}

function aplicarCenarioNaPagina(c){
  if(!c) return;
  if(c.global)         { const el=document.getElementById("cenario-global");        if(el) el.value=c.global; }
  if(c.brasil)         { const el=document.getElementById("cenario-brasil");        if(el) el.value=c.brasil; }
  if(c.posicionamento) { const el=document.getElementById("cenario-pos");           if(el) el.value=c.posicionamento; }
  if(c.referencia)     { const el=document.getElementById("cenario-ref");           if(el) el.value=c.referencia; }
  fetch("/api/hp/cenario", {method:"POST", headers:{"Content-Type":"application/json"},
    body: JSON.stringify({global:c.global||"", brasil:c.brasil||"", posicionamento:c.posicionamento||"", referencia:c.referencia||""})
  });
  document.getElementById("cenario-global")?.scrollIntoView({behavior:"smooth", block:"center"});
}

async function uploadRapidoTodos(){
  const btn = document.getElementById("btn-extrair-todos");
  btn.disabled = true;
  for(let i=0; i<_arquivosRapido.length; i++){
    if(_arquivosRapido[i].status === "pendente") await uploadRapidoUm(i);
  }
  btn.disabled = false;
  const algumOk = _arquivosRapido.some(a=>a.status==="ok");
  if(algumOk) marcarPendente();
}

function rapidoUsarCenario(i){
  const txt = _arquivosRapido[i].texto;
  if(!txt){ alert("Nenhum texto extraído."); return; }
  document.getElementById("cenario-global").value = txt.slice(0,1200);
  document.getElementById("cenario-global").scrollIntoView({behavior:"smooth", block:"center"});
}

function rapidoUsarBase(i){
  const item = _arquivosRapido[i];
  document.getElementById("know-nome").value = item.file.name.replace(".pdf","");
  document.getElementById("know-fonte").value = "";
  document.getElementById("cenario-global").scrollIntoView({behavior:"smooth", block:"start"});
  document.querySelector('[data-cls="geral"]')?.click();
  alert(`Preencha os metadados na Base de Conhecimento abaixo e clique em "Adicionar à Base".`);
}

function rapidoCopiar(i){
  const txt = _arquivosRapido[i].texto;
  if(txt) navigator.clipboard.writeText(txt).then(()=>alert("Texto copiado!"));
}

async function uploadRapido(){ await uploadRapidoTodos(); }

function fecharRapido(){
  _arquivosRapido = [];
  document.getElementById("pdf-rapido-lista").innerHTML = "";
  document.getElementById("pdf-rapido-contador").textContent = "0";
  document.getElementById("pdf-rapido-itens").style.display = "none";
  document.getElementById("rapido-zona-vazia").style.display = "";
  document.getElementById("btn-limpar-rapido").style.display = "none";
  document.getElementById("rapido-upload-st").textContent = "";
  document.getElementById("pdf-rapido").value = "";
}

function usarComoCenario(){
  const txt = document.getElementById("pdf-rapido-texto").value.trim();
  if(!txt) return;
  // Preenche o campo de posicionamento do cenário macro
  const cur = document.getElementById("cenario-pos").value.trim();
  document.getElementById("cenario-pos").value = cur ? cur + "\n\n" + txt : txt;
  // Rola até o cenário
  document.getElementById("cenario-global").scrollIntoView({behavior:"smooth", block:"start"});
  fecharRapido();
}

function usarComoBase(){
  // Já foi adicionado à base no momento do upload — só fecha o preview
  fecharRapido();
}

function copiarRapido(){
  const txt = document.getElementById("pdf-rapido-texto").value;
  navigator.clipboard.writeText(txt).then(()=>{
    const btn = event.target;
    btn.textContent = "✓ Copiado!";
    setTimeout(()=>btn.textContent="📋 Copiar", 2000);
  }).catch(()=>{});
}

// ── Init ──────────────────────────────────────────────────────────────────────
async function init(){
  initProdButtons();
  // Colorir vieses no load
  CLASSES.forEach(cls=>colorirVies(cls.key));

  try{
    const [porto, cenario, prods] = await Promise.all([
      fetch("/api/hp/portfolios").then(r=>r.json()),
      fetch("/api/hp/cenario").then(r=>r.json()),
      fetch("/api/hp/produtos").then(r=>r.json()),
    ]);
    renderPortoTable(porto);
    carregarCenario(cenario);
    if(prods.produtos && Object.keys(prods.produtos).length) renderProdutos(prods.produtos);
    else if(typeof prods === "object" && !prods.produtos) renderProdutos(prods);
  } catch(e){
    console.error("Erro ao carregar dados:", e);
  }
  CLASSES.forEach(cls=>colorirVies(cls.key));
  carregarAlertas();
  carregarCalls();
  carregarGestores();
  carregarEstruturadas();
  carregarGest2();
}

// ── Carteiras de Gestores ─────────────────────────────────────────────────────
const GEST_CLASSES = [
  {key:"pos_fixado",lbl:"Pós Fixado"}, {key:"inflacao",lbl:"Inflação"},
  {key:"pre_fixado",lbl:"Pré Fixado"}, {key:"acoes",lbl:"Ações"},
  {key:"fiis",lbl:"FIIs"}, {key:"multimercado",lbl:"Multimercado"},
  {key:"internacional",lbl:"Internacional"}, {key:"alternativos",lbl:"Alternativos"},
  {key:"criptomoedas",lbl:"Cripto"},
];
const GEST_CORES = ["#8BcFEF","#5DCAA5","#D4B483","#FF9966","#BB86FC","#4FC3F7","#A5D6A7","#FFD180","#F48FB1"];

let _gestores = [];

function atualizarTotalGest(){
  const tot = GEST_CLASSES.reduce((s,c)=>s+parseFloat(document.getElementById(`galoc-${c.key}`)?.value||0),0);
  const el = document.getElementById("gest-total");
  if(el){ el.textContent = tot.toFixed(1); el.style.color = Math.abs(tot-100)<0.1?"#5DCAA5":"#FF6B6B"; }
}

async function carregarGestores(){
  try{
    const r = await fetch("/api/hp/gestores"); _gestores = await r.json();
    renderGestores();
  }catch(e){}
}

function renderGestores(){
  const lista = document.getElementById("gest-lista");
  const count = document.getElementById("gest-count");
  if(count) count.textContent = _gestores.length;
  if(!lista) return;
  if(!_gestores.length){
    lista.innerHTML = '<div style="color:#1E3A4A;font-size:11px;font-style:italic;text-align:center;padding:20px">Nenhuma carteira salva ainda.</div>';
    const comp = document.getElementById("gest-comparativo");
    if(comp) comp.style.display = "none";
    return;
  }
  lista.innerHTML = _gestores.map((g,i)=>`
    <div style="background:#060C10;border:1px solid #1A2E3A;border-left:3px solid ${GEST_CORES[i%GEST_CORES.length]};border-radius:10px;padding:12px 14px;margin-bottom:8px;display:flex;align-items:flex-start;gap:12px;flex-wrap:wrap">
      <div style="flex:1;min-width:200px">
        <div style="display:flex;align-items:center;gap:8px;margin-bottom:4px">
          <span style="font-size:13px;color:${GEST_CORES[i%GEST_CORES.length]};font-weight:700">${g.nome||g.gestora}</span>
          <span style="font-size:10px;background:#1A2E3A;border-radius:8px;padding:2px 8px;color:#8BcFEF">${g.perfil}</span>
          <span style="font-size:10px;color:#2A5A6A">📅 ${g.referencia||g.data}</span>
        </div>
        <div style="font-size:10px;color:#3A7A8A;margin-bottom:6px">Fonte: <b>${g.gestora}</b>${g.observacao?" · "+g.observacao:""}</div>
        <div style="display:flex;flex-wrap:wrap;gap:4px">
          ${GEST_CLASSES.filter(c=>(g.alocacao[c.key]||0)>0).map(c=>
            `<span style="font-size:10px;background:#0A1A20;border:1px solid #1A3A4A;border-radius:8px;padding:2px 8px;color:#8BcFEF">${c.lbl}: <b>${g.alocacao[c.key]}%</b></span>`
          ).join("")}
        </div>
      </div>
      <div style="display:flex;gap:6px;align-items:center">
        <button class="btn-ghost" onclick="editarGestora(${i})" style="font-size:10px;padding:4px 10px">✏️ Editar</button>
        <button class="btn-ghost" onclick="deletarGestora('${g.id}')" style="font-size:10px;padding:4px 10px;color:#FF6B6B;border-color:#FF6B6B44">🗑️</button>
      </div>
    </div>`).join("");
  renderComparativo();
}

function renderComparativo(){
  const comp = document.getElementById("gest-comparativo");
  const tbl  = document.getElementById("gest-table");
  if(!comp||!tbl||_gestores.length < 2){ if(comp) comp.style.display="none"; return; }
  comp.style.display = "";
  const header = `<thead><tr>
    <th style="text-align:left;padding:6px 10px;color:#3A7A8A;font-size:10px;border-bottom:1px solid #1A2E3A">Classe</th>
    ${_gestores.map((g,i)=>`<th style="text-align:center;padding:6px 10px;color:${GEST_CORES[i%GEST_CORES.length]};font-size:10px;border-bottom:1px solid #1A2E3A">${g.gestora}<br><span style="font-weight:400;color:#2A5A6A">${g.referencia||""}</span></th>`).join("")}
    ${_gestores.length>1?`<th style="text-align:center;padding:6px 10px;color:#FFD180;font-size:10px;border-bottom:1px solid #1A2E3A">Média</th>`:""}
  </tr></thead>`;
  const rows = GEST_CLASSES.map(c=>{
    const vals = _gestores.map(g=>g.alocacao[c.key]||0);
    const media = vals.reduce((a,b)=>a+b,0)/vals.length;
    const max = Math.max(...vals);
    return `<tr style="border-bottom:1px solid #0D1E28">
      <td style="padding:5px 10px;color:#6A9AAA;font-size:11px">${c.lbl}</td>
      ${vals.map(v=>`<td style="text-align:center;padding:5px 10px;font-size:11px;color:${v===max&&max>0?"#5DCAA5":"#AAA"};font-weight:${v===max&&max>0?"700":"400"}">${v>0?v+"%":"—"}</td>`).join("")}
      ${_gestores.length>1?`<td style="text-align:center;padding:5px 10px;font-size:11px;color:#FFD180;font-weight:600">${media>0?media.toFixed(1)+"%":"—"}</td>`:""}
    </tr>`;
  });
  tbl.innerHTML = header + "<tbody>" + rows.join("") + "</tbody>";
}

function editarGestora(i){
  const g = _gestores[i];
  document.getElementById("gest-nome").value = g.nome||"";
  document.getElementById("gest-gestora").value = g.gestora||"";
  document.getElementById("gest-ref").value = g.referencia||"";
  document.getElementById("gest-perfil").value = g.perfil||"moderada";
  document.getElementById("gest-obs").value = g.observacao||"";
  GEST_CLASSES.forEach(c=>{
    const el = document.getElementById(`galoc-${c.key}`);
    if(el) el.value = g.alocacao[c.key]||0;
  });
  atualizarTotalGest();
  document.getElementById("gest-nome").scrollIntoView({behavior:"smooth",block:"center"});
}

async function deletarGestora(id){
  if(!confirm("Remover esta carteira?")) return;
  await fetch("/api/hp/gestores", {method:"POST", headers:{"Content-Type":"application/json"}, body:JSON.stringify({id, _delete:true})});
  carregarGestores();
}

async function salvarGestora(){
  const nome    = document.getElementById("gest-nome").value.trim();
  const gestora = document.getElementById("gest-gestora").value.trim();
  if(!nome||!gestora){ alert("Preencha Nome da Carteira e Gestora."); return; }
  const alocacao = {};
  GEST_CLASSES.forEach(c=>{ alocacao[c.key] = parseFloat(document.getElementById(`galoc-${c.key}`)?.value||0); });
  const total = Object.values(alocacao).reduce((a,b)=>a+b,0);
  if(Math.abs(total-100)>5){ if(!confirm(`Total = ${total.toFixed(1)}% (diferente de 100%). Salvar assim?`)) return; }
  const st = document.getElementById("gest-save-st");
  st.textContent = "⏳ Salvando...";
  try{
    const r = await fetch("/api/hp/gestores",{method:"POST",headers:{"Content-Type":"application/json"},body:JSON.stringify({
      nome, gestora, referencia: document.getElementById("gest-ref").value.trim(),
      perfil: document.getElementById("gest-perfil").value,
      observacao: document.getElementById("gest-obs").value.trim(), alocacao,
    })});
    const d = await r.json();
    if(d.ok){ st.innerHTML='<span class="status-ok">✓ Salvo!</span>'; marcarPendente(); carregarGestores(); }
    else { st.innerHTML=`<span class="status-err">${d.error}</span>`; }
  }catch(e){ st.innerHTML=`<span class="status-err">Erro: ${e.message}</span>`; }
  setTimeout(()=>st.textContent="",3000);
}

init();
</script>
</body></html>"""

if __name__ == "__main__":
    app.run(debug=True, port=5000)
