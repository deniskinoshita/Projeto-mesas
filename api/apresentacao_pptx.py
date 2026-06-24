"""
Gerador de Apresentação PPTX — Braúna Investimentos
Identidade visual teal premium, gráficos nativos python-pptx, 15 slides.
Modelo: Reunião_Trimestral_15slides_Brauna.pptx
"""
import io, math, datetime
from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Paleta (referência Reunião_Trimestral_15slides_Brauna) ───────────────────
C_BG     = RGBColor(0x07, 0x1E, 0x17)   # fundo escuro
C_CARD   = RGBColor(0x0D, 0x2B, 0x28)   # card teal escuro
C_CARD2  = RGBColor(0x11, 0x35, 0x32)   # card teal médio
C_GOLD   = RGBColor(0xC8, 0xA7, 0x33)   # ouro/âmbar
C_GOLD2  = RGBColor(0xE8, 0xA8, 0x38)   # âmbar médio
C_TEAL   = RGBColor(0x3E, 0xB8, 0xA8)   # teal principal (KPIs, acentos)
C_TEAL_D = RGBColor(0x14, 0x51, 0x4F)   # teal escuro (badges, destaques)
C_TEAL_M = RGBColor(0x1C, 0x6B, 0x67)   # teal médio (subtítulos)
C_TEAL_DP= RGBColor(0x29, 0x89, 0x7F)   # teal profundo (#AntecipeSeusSonhos)
C_LTEAL  = RGBColor(0x9E, 0xCB, 0xC6)   # teal claro (subtítulos suaves)
C_VLTEAL = RGBColor(0xD0, 0xE8, 0xE5)   # teal muito claro (body text suave)
C_GREEN  = RGBColor(0x1E, 0x84, 0x49)   # positivo
C_RED    = RGBColor(0xC0, 0x39, 0x2B)   # negativo
C_AMBER  = RGBColor(0xE8, 0xA8, 0x38)   # atenção/neutro
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0x6B, 0x7B, 0x7A)   # cinza rodapé/subtítulos
C_TEXT   = RGBColor(0x37, 0x41, 0x40)   # texto escuro sobre fundo claro
C_DKGRAY = RGBColor(0xAA, 0xCC, 0xCA)   # teal acinzentado claro

CLS_COR = {
    "pos_fixado":    C_GOLD,
    "inflacao":      C_AMBER,
    "pre_fixado":    RGBColor(0xF0, 0x78, 0x50),
    "acoes":         C_TEAL,
    "fiis":          RGBColor(0x5B, 0x9B, 0xD5),
    "multimercado":  RGBColor(0xB0, 0x8F, 0xCF),
    "alternativos":  RGBColor(0x4E, 0xC9, 0xB0),
    "internacional": RGBColor(0x5B, 0x9B, 0xD5),
    "criptomoedas":  C_LGRAY,
}
CLS_LABEL = {
    "pos_fixado":    "Pós Fixado",
    "inflacao":      "Inflação",
    "pre_fixado":    "Pré Fixado",
    "acoes":         "Ações / RV Brasil",
    "fiis":          "FIIs / Fundos Listados",
    "multimercado":  "Multimercado",
    "alternativos":  "Alternativos",
    "internacional": "Internacional",
    "criptomoedas":  "Criptomoedas",
}
CATS = ["pos_fixado","inflacao","pre_fixado","acoes","fiis","multimercado","alternativos","internacional","criptomoedas"]

SW     = Cm(33.87)
SH     = Cm(19.05)
MARGIN = Cm(0.8)
CONTENT_TOP = Cm(2.7)    # onde começa o conteúdo (abaixo do header)
CONTENT_BOT = Cm(18.4)   # onde termina o conteúdo (acima do footer)

# ── Primitivas ────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, border=None, bw=0.4):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bw)
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, text, x, y, w, h, size=11, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def bg(slide):
    add_rect(slide, 0, 0, SW, SH, C_BG)

def gold_line(slide, y=0, x=0, w=None, thickness=Cm(0.065)):
    if w is None:
        w = SW
    add_rect(slide, x, y, w, thickness, C_GOLD)

def header(slide, num, titulo, sub="", data_ref="", conta="", tipo_conta="XP Unique"):
    """Padrão de header: linha ouro + 'NN | Título' 27pt + footer URL."""
    add_rect(slide, 0, 0, SW, Cm(0.12), C_GOLD)
    title_txt = f"{num:02d}  |  {titulo}" if num > 0 else titulo
    add_text(slide, title_txt, MARGIN, Cm(0.22), Cm(30), Cm(1.5),
             size=27, bold=True, color=C_WHITE)
    if sub:
        add_text(slide, sub, MARGIN, Cm(1.85), Cm(30), Cm(0.7),
                 size=12, color=C_TEAL_M)
    footer(slide, data_ref, conta, tipo_conta)

def footer(slide, data_ref="", conta="", tipo_conta="XP Unique"):
    add_rect(slide, 0, SH - Cm(0.12), SW, Cm(0.12), C_GOLD)
    add_text(slide, "braunainvestimentos.com.br", MARGIN, SH - Cm(0.56),
             Cm(12), Cm(0.45), size=9, color=C_LGRAY)
    parts = []
    if data_ref:
        parts.append(f"Ref.: {data_ref}")
    if conta:
        parts.append(f"Conta {conta}")
    if tipo_conta:
        parts.append(tipo_conta)
    ref_txt = "  |  ".join(parts)
    if ref_txt:
        add_text(slide, ref_txt, SW - Cm(14), SH - Cm(0.56),
                 Cm(13.5), Cm(0.45), size=9, color=C_LGRAY, align=PP_ALIGN.RIGHT)

def fmt_brl(v):
    try:
        v = float(v)
        if v >= 1_000_000: return f"R$ {v/1_000_000:.2f} MM"
        if v >= 1_000:     return f"R$ {v/1_000:.0f}K"
        return f"R$ {v:.0f}"
    except Exception:
        return "—"

def fmt_brl_short(v):
    try:
        v = float(v)
        if v >= 1_000_000: return f"R$ {v/1_000_000:.2f}M"
        if v >= 1_000:     return f"R$ {v/1_000:.1f}K"
        return f"R$ {v:.0f}"
    except Exception:
        return "—"

def pct_cor(v):
    try:
        s = str(v).replace("%","").replace(",",".").strip()
        return C_GREEN if float(s) >= 0 else C_RED
    except Exception:
        return C_LGRAY

def _pct_float(v):
    try:
        return float(str(v).replace("%","").replace(",",".").strip())
    except Exception:
        return 0.0

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _get_base(d):
    """Extrai campos comuns do dict d."""
    return {
        "data_ref":   d.get("data_ref", ""),
        "conta":      str(d.get("conta", d.get("codigo_cliente", ""))),
        "tipo_conta": d.get("tipo_conta", "XP Unique"),
    }

# ── Chart helpers ─────────────────────────────────────────────────────────────
def _rgb_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

def _set_series_color(series, rgb: RGBColor):
    try:
        spPr = series.format.fill._element
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", _rgb_hex(rgb))
    except Exception:
        pass

def _style_chart_plot_area(chart, bg_rgb: RGBColor = C_BG):
    try:
        plotArea = chart._element.find(qn("c:plotArea"))
        if plotArea is None: return
        spPr = plotArea.find(qn("c:spPr"))
        if spPr is None:
            spPr = etree.SubElement(plotArea, qn("c:spPr"))
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", _rgb_hex(bg_rgb))
        ln = etree.SubElement(spPr, qn("a:ln"))
        etree.SubElement(ln, qn("a:noFill"))
    except Exception:
        pass

def _style_chart_bg(chart, bg_rgb: RGBColor = C_BG):
    try:
        chartSpace = chart._element
        spPr = chartSpace.find(qn("c:spPr"))
        if spPr is None:
            spPr = etree.SubElement(chartSpace, qn("c:spPr"))
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", _rgb_hex(bg_rgb))
        ln = etree.SubElement(spPr, qn("a:ln"))
        etree.SubElement(ln, qn("a:noFill"))
    except Exception:
        pass

def _hide_gridlines(chart):
    try:
        for axis_tag in [qn("c:valAx"), qn("c:catAx")]:
            for ax in chart._element.iter(axis_tag):
                mj = ax.find(qn("c:majorGridlines"))
                if mj is not None:
                    ax.remove(mj)
    except Exception:
        pass

def _set_axis_label_color(chart, rgb: RGBColor):
    try:
        for axis_tag in [qn("c:valAx"), qn("c:catAx")]:
            for ax in chart._element.iter(axis_tag):
                txPr = ax.find(qn("c:txPr"))
                if txPr is None:
                    txPr = etree.SubElement(ax, qn("c:txPr"))
                bodyPr   = etree.SubElement(txPr, qn("a:bodyPr"))
                lstStyle = etree.SubElement(txPr, qn("a:lstStyle"))
                p        = etree.SubElement(txPr, qn("a:p"))
                pPr      = etree.SubElement(p, qn("a:pPr"))
                defRPr   = etree.SubElement(pPr, qn("a:defRPr"))
                solidFill = etree.SubElement(defRPr, qn("a:solidFill"))
                srgbClr  = etree.SubElement(solidFill, qn("a:srgbClr"))
                srgbClr.set("val", _rgb_hex(rgb))
    except Exception:
        pass

# ── SLIDE 1: CAPA ─────────────────────────────────────────────────────────────
def s_capa(prs, d):
    sl = new_slide(prs)
    bg(sl)

    # Faixa vertical esquerda teal escura
    add_rect(sl, 0, 0, Cm(1.2), SH, C_CARD)
    add_rect(sl, 0, 0, Cm(0.12), SH, C_GOLD)

    b = _get_base(d)
    nome     = d.get("nome_cliente", "Cliente")
    patrimonio = d.get("patrimonio", 0)
    data_r   = b["data_ref"]
    conta    = b["conta"]
    tipo_conta = b["tipo_conta"]
    assessor = d.get("assessor", "Assessor Braúna")

    # Obter mês/ano da data de referência
    try:
        parts = data_r.split("/")
        if len(parts) == 3:
            mes_num = int(parts[1])
            ano_num = int(parts[2])
            MESES = ["","Janeiro","Fevereiro","Março","Abril","Maio","Junho",
                     "Julho","Agosto","Setembro","Outubro","Novembro","Dezembro"]
            mes_txt = MESES[mes_num] if 1 <= mes_num <= 12 else str(mes_num)
            data_titulo = f"{mes_txt} {ano_num}"
        else:
            data_titulo = data_r
    except Exception:
        data_titulo = data_r

    # Título da reunião (topo)
    add_text(sl, "Reunião Trimestral de Investimentos",
             Cm(1.7), Cm(1.8), Cm(26), Cm(0.7),
             size=14, bold=True, color=C_GOLD)

    # Mês / Ano grande
    add_text(sl, data_titulo, Cm(1.7), Cm(2.7), Cm(26), Cm(3.2),
             size=60, bold=True, color=C_WHITE)

    # Tagline
    add_text(sl, "Estratégia  ·  Resultados  ·  Análise  ·  Próximos Passos",
             Cm(1.7), Cm(6.0), Cm(26), Cm(0.7),
             size=15, color=C_LTEAL)

    gold_line(sl, y=Cm(6.9), x=Cm(1.7), w=Cm(22), thickness=Cm(0.04))

    # Assessor + empresa
    add_text(sl, f"{assessor}  |  Assessor de Investimentos XP Unique",
             Cm(1.7), Cm(7.1), Cm(26), Cm(0.65),
             size=13, color=C_DKGRAY)
    add_text(sl, "Braúna Ecossistema de Investimentos",
             Cm(1.7), Cm(7.85), Cm(26), Cm(0.55),
             size=12, color=RGBColor(0x6A, 0x9E, 0x9A))

    # Card de patrimônio
    card_x = Cm(1.7)
    card_y = Cm(9.0)
    card_w = Cm(13)
    card_h = Cm(4.2)
    add_rect(sl, card_x, card_y, card_w, card_h, C_CARD, C_TEAL_D, 0.6)
    add_rect(sl, card_x, card_y, Cm(0.08), card_h, C_TEAL)

    add_text(sl, "Patrimônio Total",
             card_x + Cm(0.5), card_y + Cm(0.3), Cm(12), Cm(0.55),
             size=11, bold=True, color=C_TEAL)

    pat_str = fmt_brl(patrimonio)
    add_text(sl, pat_str,
             card_x + Cm(0.5), card_y + Cm(0.95), Cm(12), Cm(1.2),
             size=21, bold=True, color=C_WHITE)

    if data_r:
        add_text(sl, f"Ref.: {data_r}",
                 card_x + Cm(0.5), card_y + Cm(2.2), Cm(5), Cm(0.45),
                 size=10, color=RGBColor(0x7A, 0xAC, 0xAA))
    if tipo_conta:
        add_text(sl, tipo_conta,
                 card_x + Cm(0.5), card_y + Cm(2.7), Cm(5), Cm(0.45),
                 size=10, color=RGBColor(0x7A, 0xAC, 0xAA))
    if conta:
        add_text(sl, f"Conta {conta}",
                 card_x + Cm(0.5), card_y + Cm(3.2), Cm(5), Cm(0.45),
                 size=10, color=RGBColor(0x7A, 0xAC, 0xAA))

    # Hashtag
    add_text(sl, "#AntecipeSeusSonhos",
             Cm(1.7), Cm(14.0), Cm(18), Cm(0.7),
             size=13, bold=True, color=C_TEAL_DP)

    # Nome do cliente (destaque)
    if nome and nome != "Cliente":
        add_text(sl, nome,
                 Cm(1.7), Cm(15.0), Cm(26), Cm(1.0),
                 size=18, bold=False, color=C_LTEAL)

    # Footer
    footer(sl, data_r, conta, tipo_conta)


# ── SLIDE 2: AGENDA ───────────────────────────────────────────────────────────
def s_agenda(prs, d, agenda_items):
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)

    add_rect(sl, 0, 0, SW, Cm(0.12), C_GOLD)
    add_text(sl, "Agenda da Reunião",
             MARGIN, Cm(0.22), Cm(30), Cm(1.5),
             size=27, bold=True, color=C_WHITE)
    footer(sl, b["data_ref"], b["conta"], b["tipo_conta"])

    # Items em duas colunas
    col_w = Cm(15.5)
    col1_x = MARGIN
    col2_x = Cm(17.8)
    item_h = Cm(1.55)
    start_y = Cm(2.4)

    for idx, item in enumerate(agenda_items):
        if isinstance(item, (list, tuple)) and len(item) >= 2:
            num_item = item[0]
            titulo_item = item[1]
            sub_item = item[2] if len(item) > 2 else ""
        else:
            num_item = idx + 1
            titulo_item = str(item)
            sub_item = ""

        col = idx % 2
        row = idx // 2
        x = col1_x if col == 0 else col2_x
        y = start_y + row * item_h

        # Fundo do item
        add_rect(sl, x, y, col_w, item_h - Cm(0.1), C_CARD, C_TEAL_D, 0.3)

        # Badge numérico
        badge_w = Cm(1.0)
        add_rect(sl, x, y, badge_w, item_h - Cm(0.1), C_TEAL_D)
        add_text(sl, f"{num_item:02d}",
                 x + Cm(0.08), y + Cm(0.2), badge_w - Cm(0.16), Cm(0.9),
                 size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Título do item
        add_text(sl, titulo_item,
                 x + badge_w + Cm(0.2), y + Cm(0.15), col_w - badge_w - Cm(0.4), Cm(0.65),
                 size=12, bold=True, color=C_TEXT)

        # Sub-item
        if sub_item:
            add_text(sl, str(sub_item)[:55],
                     x + badge_w + Cm(0.2), y + Cm(0.82), col_w - badge_w - Cm(0.4), Cm(0.55),
                     size=10, color=C_LGRAY)


# ── SLIDE 3: CONTEXTO MACRO ───────────────────────────────────────────────────
def s_macro(prs, d, num=1):
    """01 | Contexto Macro & Mercado"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)

    # Determinar mês/ano do título
    data_r = b["data_ref"]
    try:
        parts = data_r.split("/")
        MESES = ["","Janeiro","Fevereiro","Março","Abril","Maio","Junho",
                 "Julho","Agosto","Setembro","Outubro","Novembro","Dezembro"]
        mes_txt = MESES[int(parts[1])] if len(parts) == 3 else ""
        ano_txt = parts[2] if len(parts) == 3 else ""
        titulo_macro = f"01  |  Contexto Macro & Mercado — {mes_txt} {ano_txt}"
    except Exception:
        titulo_macro = "01  |  Contexto Macro & Mercado"

    add_rect(sl, 0, 0, SW, Cm(0.12), C_GOLD)
    add_text(sl, titulo_macro, MARGIN, Cm(0.22), Cm(31), Cm(1.5),
             size=22, bold=True, color=C_WHITE)
    footer(sl, data_r, b["conta"], b["tipo_conta"])

    # Subtítulo
    cenario_txt = d.get("cenario_global", "") or d.get("cenario_brasil", "")
    sub = "Cenário que moldou a carteira no semestre e perspectivas"
    if cenario_txt and len(cenario_txt) > 10:
        sub = cenario_txt[:65]
    add_text(sl, sub, MARGIN, Cm(1.85), Cm(30), Cm(0.65),
             size=13, color=C_GOLD)

    # 4 KPIs macro
    selic = d.get("selic", "—")
    ipca  = d.get("ipca",  "—")
    dolar = d.get("dolar", "—")
    ibov  = d.get("ibov",  "—")

    macros = [
        ("SELIC Meta",    "▼", selic, "a.a.", "Ciclo de queda em curso",     C_TEAL),
        ("IPCA 12M",      "▲", ipca,  "",     "Acima do centro da meta (3%)", C_GOLD2),
        ("USD / BRL",     "▶", dolar, "",     "Câmbio de referência",         C_TEAL_DP),
        ("Ibovespa 12M",  "▲", ibov,  "",     "Mercado de renda variável",    C_GREEN),
    ]

    kpi_w = Cm(7.5)
    kpi_h = Cm(4.2)
    gap   = Cm(0.5)
    row_y = Cm(2.75)

    for i, (label, seta, valor, suf, desc, cor_seta) in enumerate(macros):
        x = MARGIN + i * (kpi_w + gap)
        add_rect(sl, x, row_y, kpi_w, kpi_h, C_CARD, C_TEAL_D, 0.4)
        add_text(sl, label, x + Cm(0.3), row_y + Cm(0.2), kpi_w - Cm(0.4), Cm(0.55),
                 size=10, bold=True, color=C_TEAL)
        add_text(sl, seta, x + Cm(0.2), row_y + Cm(0.75), Cm(0.9), Cm(1.0),
                 size=18, color=cor_seta)
        val_txt = f"{valor} {suf}".strip() if suf else str(valor)
        add_text(sl, val_txt, x + Cm(1.1), row_y + Cm(0.8), kpi_w - Cm(1.3), Cm(1.0),
                 size=20, bold=True, color=C_WHITE)
        add_text(sl, desc, x + Cm(0.3), row_y + Cm(1.9), kpi_w - Cm(0.6), Cm(0.55),
                 size=10, color=RGBColor(0x8E, 0xCE, 0xCA))

    # Leituras-chave
    leituras_y = Cm(7.3)
    add_text(sl, "Leituras-chave para a carteira",
             MARGIN, leituras_y, Cm(28), Cm(0.65),
             size=13, bold=True, color=C_TEAL)

    leituras = d.get("leituras_macro", [])
    if not leituras:
        # Gerar leituras baseadas nos dados disponíveis
        bullets_raw = []
        cenario_b = d.get("cenario_brasil", "")
        cenario_g = d.get("cenario_global", "")
        if cenario_b:
            bullets_raw.append(("Brasil:", cenario_b[:90], C_GOLD2))
        if cenario_g:
            bullets_raw.append(("Global:", cenario_g[:90], C_TEAL))
        # defaults genéricos se não houver texto
        if not bullets_raw:
            bullets_raw = [
                ("Renda Fixa:", f"Com SELIC em {selic}, juros ainda elevados favorecem estratégias pós-fixadas.", C_TEAL),
                ("Renda Variável:", f"Ibovespa {ibov} no período — acompanhar seletividade.", C_GOLD),
                ("Câmbio:", f"Dólar a {dolar} — impacto em BDRs e ativos internacionais.", C_GOLD2),
                ("Oportunidade:", "Queda de juros favorece diversificação e alongamento de duration.", C_GREEN),
            ]

    for i, item in enumerate(bullets_raw[:4]):
        titulo_b, texto_b, cor_b = item
        add_rect(sl, MARGIN, leituras_y + Cm(0.8) + i * Cm(1.55),
                 Cm(28.5), Cm(1.4), C_CARD, None)
        add_rect(sl, MARGIN, leituras_y + Cm(0.8) + i * Cm(1.55),
                 Cm(0.07), Cm(1.4), cor_b)
        add_text(sl, f"{titulo_b} {texto_b}",
                 MARGIN + Cm(0.3), leituras_y + Cm(0.97) + i * Cm(1.55),
                 Cm(28), Cm(1.0), size=11, bold=False, color=cor_b)


# ── SLIDE 4: VISÃO GERAL DO PATRIMÔNIO ───────────────────────────────────────
def s_patrimonio(prs, d, num=2):
    """02 | Visão Geral do Patrimônio"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)

    header(sl, num, "Visão Geral do Patrimônio",
           sub="Evolução patrimonial e estatísticas da carteira",
           **b)

    nome   = d.get("nome_cliente", "")
    pat    = d.get("patrimonio", 0)
    r_mes  = d.get("rentabilidade_mes", "—")
    r_ano  = d.get("rentabilidade_ano", "—")
    r_12m  = d.get("rentabilidade_12m", "—")
    r_24m  = d.get("rentabilidade_24m", "—")
    g_mes  = d.get("ganho_mes", 0)
    _pct_cdi_raw = d.get("pct_cdi", {})
    if isinstance(_pct_cdi_raw, dict):
        pct_cdi = _pct_cdi_raw
    else:
        pct_cdi = {"mes": "—", "12m": _pct_cdi_raw, "24m": "—"}
    cdi_mes = pct_cdi.get("mes", pct_cdi.get("pct_cdi_mes", "—"))
    cdi_12m = pct_cdi.get("12m", pct_cdi.get("pct_cdi_12m", "—"))
    cdi_24m = pct_cdi.get("24m", pct_cdi.get("pct_cdi_24m", "—"))

    # 4 KPIs principais no topo
    kpis_top = [
        ("Patrimônio Total",   fmt_brl(pat),  f"Bruto em {b['data_ref']}"),
        ("Rent. Mês",          r_mes,         f"CDI mês: {cdi_mes}%"),
        ("Ganho no Mês",       fmt_brl_short(g_mes), "Variação absoluta"),
        ("Rent. 12 Meses",     r_12m,         f"{cdi_12m}% do CDI"),
    ]

    kpi_w = Cm(7.5)
    kpi_h = Cm(2.7)
    gap   = Cm(0.5)
    row_y = Cm(2.6)

    for i, (lab, val, nota) in enumerate(kpis_top):
        x = MARGIN + i * (kpi_w + gap)
        add_rect(sl, x, row_y, kpi_w, kpi_h, C_CARD, C_TEAL_D, 0.4)
        add_text(sl, lab, x + Cm(0.3), row_y + Cm(0.15), kpi_w - Cm(0.5), Cm(0.55),
                 size=10, bold=True, color=C_LGRAY)
        cor_val = pct_cor(val) if "%" in str(val) else C_TEXT
        add_text(sl, str(val), x + Cm(0.3), row_y + Cm(0.75), kpi_w - Cm(0.5), Cm(1.3),
                 size=22, bold=True, color=cor_val)
        add_text(sl, nota, x + Cm(0.3), row_y + Cm(2.1), kpi_w - Cm(0.5), Cm(0.5),
                 size=10, color=C_LGRAY)

    # Linha divisória
    add_rect(sl, MARGIN, Cm(5.55), SW - Cm(1.6), Cm(0.04), C_TEAL_D)

    # Segunda linha: estatísticas adicionais
    stats = [
        ("Rent. 24 Meses", r_24m, f"{cdi_24m}% do CDI"),
        ("Rent. no Ano",   r_ano, "Ano corrente"),
    ]
    extras = d.get("estatisticas_extras", {})
    if extras:
        if extras.get("meses_positivos"):
            stats.append(("Meses Positivos", extras["meses_positivos"], "últimos 24 meses"))
        if extras.get("volatilidade_12m"):
            stats.append(("Volatilidade 12M", extras["volatilidade_12m"], "risco de mercado"))
        if extras.get("retorno_max"):
            stats.append(("Retorno Máx. Mensal", extras["retorno_max"], extras.get("retorno_max_mes", "")))
        if extras.get("retorno_min"):
            stats.append(("Retorno Mín. Mensal", extras["retorno_min"], extras.get("retorno_min_mes", "")))

    stat_w = Cm(7.0)
    stat_h = Cm(2.5)
    row2_y = Cm(5.85)
    for i, (lab, val, nota) in enumerate(stats[:4]):
        x = MARGIN + i * (stat_w + Cm(0.65))
        add_rect(sl, x, row2_y, stat_w, stat_h, C_CARD2, C_TEAL_D, 0.3)
        add_text(sl, lab, x + Cm(0.3), row2_y + Cm(0.12), stat_w - Cm(0.4), Cm(0.5),
                 size=10, color=C_LGRAY)
        cor_v = pct_cor(val) if "%" in str(val) else C_TEXT
        add_text(sl, str(val), x + Cm(0.3), row2_y + Cm(0.65), stat_w - Cm(0.4), Cm(1.0),
                 size=14, bold=True, color=cor_v)
        if nota:
            add_text(sl, nota, x + Cm(0.3), row2_y + Cm(1.75), stat_w - Cm(0.4), Cm(0.45),
                     size=9, color=C_LGRAY)

    # Histórico anual se disponível
    historico_anos = d.get("historico_anos", [])
    hist_y = Cm(9.0)
    if historico_anos:
        add_text(sl, "Rentabilidade anual",
                 MARGIN, hist_y - Cm(0.5), Cm(28), Cm(0.5),
                 size=11, bold=True, color=C_TEAL)
        hist_txt = "  |  ".join([f"{h.get('ano','')}: {h.get('retorno','—')}" for h in historico_anos[:5]])
        add_rect(sl, MARGIN, hist_y, SW - Cm(1.6), Cm(0.65), C_TEAL_D)
        add_text(sl, hist_txt, MARGIN + Cm(0.3), hist_y + Cm(0.1),
                 SW - Cm(2.2), Cm(0.5), size=10, bold=True, color=C_TEAL_M)
    else:
        # mostrar linha de rentabilidade com anos disponíveis
        anos_txt = ""
        if r_ano and r_ano != "—":
            anos_txt += f"Ano corrente: {r_ano}"
        if r_12m and r_12m != "—":
            anos_txt += f"  |  12M: {r_12m}"
        if r_24m and r_24m != "—":
            anos_txt += f"  |  24M: {r_24m}"
        if anos_txt:
            add_rect(sl, MARGIN, hist_y, SW - Cm(1.6), Cm(0.65), C_TEAL_D)
            add_text(sl, anos_txt, MARGIN + Cm(0.3), hist_y + Cm(0.1),
                     SW - Cm(2.2), Cm(0.5), size=10, bold=True, color=C_TEAL_M)

    # Tabela de rentabilidade por período vs benchmarks
    bench_y = Cm(10.1)
    add_rect(sl, MARGIN, bench_y - Cm(0.05), SW - Cm(1.6), Cm(0.5), C_TEAL_D)
    cols_h = ["Referência", "Mês", "% CDI", "Ano", "% CDI", "12M", "% CDI", "24M", "% CDI"]
    col_xs = [MARGIN + Cm(0.1), Cm(4.8), Cm(7.3), Cm(9.4), Cm(12.0), Cm(14.3), Cm(17.0), Cm(19.5), Cm(22.3)]
    for j, (cx, ch) in enumerate(zip(col_xs, cols_h)):
        add_text(sl, ch, cx, bench_y, Cm(2.4), Cm(0.5),
                 size=10, bold=True, color=C_WHITE)

    bench_rows = [
        ("Portfólio", r_mes, str(cdi_mes), r_ano, "—", r_12m, str(cdi_12m), r_24m, str(cdi_24m)),
    ]
    # CDI, Ibovespa se tiverem nos dados macro
    if d.get("selic"):
        bench_rows.append(("CDI", "—", "—", "—", "—", "—", "—", "—", "—"))

    row_colors = [C_CARD, C_CARD2]
    for ri, row in enumerate(bench_rows):
        ry = bench_y + Cm(0.55) + ri * Cm(0.6)
        add_rect(sl, MARGIN, ry - Cm(0.02), SW - Cm(1.6), Cm(0.6), row_colors[ri % 2])
        for j, (cx, val) in enumerate(zip(col_xs, row)):
            is_pct_col = j > 0
            cor_v = pct_cor(val) if is_pct_col else C_WHITE
            bold_v = (j == 0)
            add_text(sl, str(val), cx, ry, Cm(2.4), Cm(0.55),
                     size=10, bold=bold_v, color=cor_v if is_pct_col else C_WHITE)


# ── SLIDE 5: COMPOSIÇÃO DA CARTEIRA ──────────────────────────────────────────
def s_composicao(prs, d, num=4):
    """04 | Composição da Carteira — Estratégia Atual"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)

    pat = d.get("patrimonio", 0)
    pat_str = fmt_brl(pat)
    header(sl, num, "Composição da Carteira — Estratégia Atual",
           sub=f"Distribuição patrimonial por classe de ativos  —  {pat_str}",
           **b)

    composicao = d.get("composicao", {})
    saldos     = d.get("saldos_por_classe", {})
    rents      = d.get("rents_por_classe", {})

    classes_ord = [(k, composicao[k]) for k in CATS if k in composicao and composicao[k] > 0]
    classes_ord.sort(key=lambda x: -x[1])

    card_w = Cm(7.5)
    card_h = Cm(3.6)
    gap_x  = Cm(0.45)
    gap_y  = Cm(0.3)
    start_y = Cm(2.65)
    per_row = 4

    for idx, (cls, pct) in enumerate(classes_ord[:8]):
        col = idx % per_row
        row = idx // per_row
        x = MARGIN + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        cor = CLS_COR.get(cls, C_TEAL)
        label = CLS_LABEL.get(cls, cls)
        saldo = saldos.get(cls, 0)
        rent_d = rents.get(cls, {}) if isinstance(rents, dict) else {}

        add_rect(sl, x, y, card_w, card_h, C_CARD, C_TEAL_D, 0.3)
        add_rect(sl, x, y, Cm(0.1), card_h, cor)

        add_text(sl, label, x + Cm(0.3), y + Cm(0.15), card_w - Cm(0.5), Cm(0.6),
                 size=12, bold=True, color=C_WHITE)
        add_text(sl, f"{pct:.2f}% do patrimônio", x + Cm(0.3), y + Cm(0.78), card_w - Cm(0.5), Cm(0.5),
                 size=11, color=C_WHITE)

        if saldo:
            add_text(sl, "Saldo", x + Cm(0.3), y + Cm(1.35), Cm(1.8), Cm(0.4),
                     size=10, color=C_LGRAY)
            add_text(sl, fmt_brl_short(saldo), x + Cm(0.3), y + Cm(1.75), card_w - Cm(0.5), Cm(0.7),
                     size=16, bold=True, color=C_TEXT)

        r_mes_c = rent_d.get("mes", "")
        r_12m_c = rent_d.get("12m", rent_d.get("12M", ""))
        if r_mes_c or r_12m_c:
            perf_txt = ""
            if r_mes_c:
                perf_txt += f"Mês: {r_mes_c}"
            if r_12m_c:
                perf_txt += f"  |  12M: {r_12m_c}" if r_mes_c else f"12M: {r_12m_c}"
            add_text(sl, perf_txt, x + Cm(0.3), y + Cm(2.9), card_w - Cm(0.5), Cm(0.5),
                     size=10, color=pct_cor(r_mes_c or r_12m_c))

    # Gráfico de rosca (donut) se houver dados suficientes
    if len(classes_ord) >= 2:
        chart_x = MARGIN
        chart_y = start_y + 2 * (card_h + gap_y) + Cm(0.2)
        chart_w = SW - Cm(1.6)
        chart_h_c = Cm(4.5)

        if chart_y + chart_h_c < SH - Cm(0.8):
            try:
                cd = ChartData()
                labels = [CLS_LABEL.get(k, k) for k, _ in classes_ord[:8]]
                values = [v for _, v in classes_ord[:8]]
                cd.categories = labels
                cd.add_series("Alocação", values)
                chart = sl.shapes.add_chart(
                    XL_CHART_TYPE.DOUGHNUT, chart_x, chart_y, chart_w, chart_h_c, cd
                ).chart
                _style_chart_bg(chart, C_BG)
                _style_chart_plot_area(chart, C_BG)
                chart.has_legend = True
                try:
                    chart.legend.font.color.rgb = C_LGRAY
                    chart.legend.font.size = Pt(8)
                except Exception:
                    pass
            except Exception:
                pass


# ── SLIDE 6: CARTEIRA vs. MODELO ─────────────────────────────────────────────
def s_vs_modelo(prs, d, num=3):
    """03 | Carteira vs. Modelo — Desvios do Perfil"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)

    perfil  = str(d.get("perfil", "Moderado")).capitalize()
    header(sl, num, f"Carteira vs. Modelo — Perfil {perfil}",
           sub="Comparativo entre alocação atual e modelo de referência",
           **b)

    composicao = d.get("composicao", {})
    modelo_hp  = d.get("modelo_hp", {})
    desvios    = d.get("desvios", [])

    classes_ativas = [k for k in CATS if composicao.get(k, 0) > 0 or modelo_hp.get(k, 0) > 0]
    if not classes_ativas:
        add_text(sl, "Dados de composição não disponíveis.",
                 MARGIN, Cm(5), Cm(28), Cm(1), size=12, color=C_LGRAY)
        return

    # Tabela comparativa
    COLS = ["Classe", "Atual %", "Modelo %", "Desvio pp", "Status"]
    col_xs = [MARGIN + Cm(0.1), Cm(11.5), Cm(16.0), Cm(20.5), Cm(25.5)]
    col_ws = [Cm(9.0), Cm(4.0), Cm(4.0), Cm(4.5), Cm(6.0)]

    tab_y = Cm(2.65)
    add_rect(sl, MARGIN, tab_y, SW - Cm(1.6), Cm(0.65), C_TEAL_D)
    for j, (cx, cw, ch) in enumerate(zip(col_xs, col_ws, COLS)):
        add_text(sl, ch, cx, tab_y + Cm(0.05), cw, Cm(0.55),
                 size=10, bold=True, color=C_WHITE)

    # Mapear desvios por cls
    desvios_map = {}
    for dv in desvios:
        ck = dv.get("cls") or dv.get("classe", "")
        if ck:
            desvios_map[ck] = dv

    row_colors = [C_CARD, C_CARD2]
    row_h = Cm(0.72)
    for ri, cls in enumerate(classes_ativas[:12]):
        ry = tab_y + Cm(0.65) + ri * row_h
        add_rect(sl, MARGIN, ry, SW - Cm(1.6), row_h, row_colors[ri % 2])

        dv   = desvios_map.get(cls, {})
        lbl  = dv.get("label") or CLS_LABEL.get(cls, cls)
        atual = composicao.get(cls, dv.get("atual", 0))
        alvo  = modelo_hp.get(cls, dv.get("alvo", dv.get("modelo", 0)))
        desvio = atual - alvo
        status = dv.get("status", "Adequado")

        cor_dev = C_GREEN if abs(desvio) <= 2 else (C_RED if abs(desvio) > 5 else C_AMBER)
        STATUS_LABEL = {
            "ok": "✅ OK",
            "Adequado": "✅ OK",
            "atencao": "⚠️ Atenção",
            "Atencao": "⚠️ Atenção",
            "fora": "🔴 Fora",
            "Fora do perfil": "🔴 Fora",
        }
        status_txt = STATUS_LABEL.get(status, status)

        row_data = [
            (lbl, C_WHITE, False),
            (f"{atual:.2f}%", pct_cor(atual), False),
            (f"{alvo:.2f}%", C_LGRAY, False),
            (f"{desvio:+.2f}pp", cor_dev, True),
            (status_txt, cor_dev, False),
        ]
        for j, (cx, cw, (val, cor, bld)) in enumerate(zip(col_xs, col_ws, row_data)):
            add_text(sl, val, cx, ry + Cm(0.1), cw, Cm(0.55),
                     size=10, bold=bld, color=cor)

    # Rodapé informativo
    last_row_y = tab_y + Cm(0.65) + len(classes_ativas[:12]) * row_h + Cm(0.2)
    if last_row_y < SH - Cm(1.5):
        nota_txt = "Desvio > ±5pp: fora do perfil  |  Desvio entre ±2pp e ±5pp: atenção  |  Desvio < ±2pp: adequado"
        add_text(sl, nota_txt, MARGIN, last_row_y, SW - Cm(1.6), Cm(0.5),
                 size=9, color=C_LGRAY)


# ── SLIDE 7: DESVIOS ──────────────────────────────────────────────────────────
def s_desvios(prs, d, num=5):
    """05 | Análise de Risco & Diagnóstico"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)
    header(sl, num, "Análise de Risco & Diagnóstico da Carteira",
           sub="Desvios do perfil e pontos de atenção identificados",
           **b)

    desvios = d.get("desvios", [])
    if not desvios:
        add_text(sl, "Sem desvios significativos identificados.",
                 MARGIN, Cm(4), Cm(28), Cm(1), size=12, color=C_LGRAY)
        return

    # KPIs de risco no topo
    risk_kpis_y = Cm(2.65)
    extras = d.get("estatisticas_extras", {})
    vol_12m = extras.get("volatilidade_12m", d.get("volatilidade_12m", "—"))
    vol_24m = extras.get("volatilidade_24m", d.get("volatilidade_24m", "—"))
    drawdown = extras.get("max_drawdown", d.get("max_drawdown", "—"))
    ret_max  = extras.get("retorno_max", d.get("retorno_max", "—"))

    risk_kpis = [
        ("Volatilidade 12M",  str(vol_12m), "risco de mercado", C_RED if _pct_float(vol_12m) > 15 else C_AMBER),
        ("Volatilidade 24M",  str(vol_24m), "histórico",        C_RED if _pct_float(vol_24m) > 15 else C_AMBER),
        ("Máx. Drawdown Mês", str(drawdown), "mês mais negativo", C_RED),
        ("Máx. Retorno Mês",  str(ret_max),  "mês mais positivo", C_GREEN),
    ]
    kpi_w = Cm(7.5)
    kpi_h = Cm(2.5)
    gap   = Cm(0.5)
    for i, (lab, val, nota, cor) in enumerate(risk_kpis):
        x = MARGIN + i * (kpi_w + gap)
        add_rect(sl, x, risk_kpis_y, kpi_w, kpi_h, C_CARD, C_TEAL_D, 0.4)
        add_text(sl, lab, x + Cm(0.3), risk_kpis_y + Cm(0.15), kpi_w - Cm(0.5), Cm(0.55),
                 size=10, bold=True, color=C_TEAL)
        add_text(sl, val, x + Cm(0.3), risk_kpis_y + Cm(0.75), kpi_w - Cm(0.5), Cm(1.0),
                 size=21, bold=True, color=cor)
        add_text(sl, nota, x + Cm(0.3), risk_kpis_y + Cm(1.85), kpi_w - Cm(0.5), Cm(0.45),
                 size=9, color=RGBColor(0x8E, 0xCE, 0xCA))

    # Diagnóstico: desvios
    diag_y = Cm(5.45)
    add_text(sl, "⚠️  Diagnóstico: Desvios de Alocação",
             MARGIN, diag_y, Cm(20), Cm(0.65),
             size=13, bold=True, color=C_GOLD)

    add_text(sl, "💡  Pontos de Melhoria Identificados",
             Cm(18.5), diag_y, Cm(14), Cm(0.65),
             size=13, bold=True, color=C_TEAL)

    col_split = Cm(16.5)
    diag_item_y = diag_y + Cm(0.75)

    # Itens de desvio (lado esquerdo)
    fora = [dv for dv in desvios if abs(_pct_float(dv.get("desvio", dv.get("desvio_pp", 0)))) > 5]
    atencao = [dv for dv in desvios if 2 < abs(_pct_float(dv.get("desvio", dv.get("desvio_pp", 0)))) <= 5]
    ok = [dv for dv in desvios if abs(_pct_float(dv.get("desvio", dv.get("desvio_pp", 0)))) <= 2]

    item_h_diag = Cm(1.4)
    for ri, dv in enumerate(fora[:3] + atencao[:2] + ok[:1]):
        dev_val = _pct_float(dv.get("desvio", dv.get("desvio_pp", 0)))
        lbl = dv.get("label") or CLS_LABEL.get(dv.get("cls",""), dv.get("cls",""))
        atual_v = dv.get("atual", 0)
        alvo_v  = dv.get("alvo", dv.get("modelo", 0))

        if abs(dev_val) > 5:
            icone, cor = "🔴", C_RED
        elif abs(dev_val) > 2:
            icone, cor = "🟡", C_AMBER
        else:
            icone, cor = "🟢", C_GREEN

        ry = diag_item_y + ri * item_h_diag
        add_rect(sl, MARGIN, ry, col_split - Cm(0.3), item_h_diag - Cm(0.1), C_CARD, None)
        add_rect(sl, MARGIN, ry, Cm(0.06), item_h_diag - Cm(0.1), cor)
        add_text(sl, icone, MARGIN + Cm(0.2), ry + Cm(0.15), Cm(0.8), Cm(0.9), size=16)
        add_text(sl, f"{lbl}: atual {atual_v:.1f}% | modelo {alvo_v:.1f}% | desvio {dev_val:+.1f}pp",
                 MARGIN + Cm(1.1), ry + Cm(0.15), col_split - Cm(1.5), Cm(0.9),
                 size=10, bold=True, color=cor)

    # Pontos de melhoria (lado direito)
    melhorias = d.get("melhorias", d.get("sugestoes_melhoria", []))
    if not melhorias:
        # Gerar automaticamente a partir dos desvios
        for dv in (fora + atencao)[:4]:
            dev_val = _pct_float(dv.get("desvio", dv.get("desvio_pp", 0)))
            lbl = dv.get("label") or CLS_LABEL.get(dv.get("cls",""), dv.get("cls",""))
            if dev_val > 0:
                melhorias.append(f"Reduzir {lbl}: {dev_val:+.1f}pp acima do modelo")
            else:
                melhorias.append(f"Aumentar {lbl}: {dev_val:+.1f}pp abaixo do modelo")

    for ri, txt in enumerate(melhorias[:5]):
        ry = diag_item_y + ri * item_h_diag
        add_rect(sl, col_split + Cm(0.3), ry, SW - col_split - Cm(1.1), item_h_diag - Cm(0.1),
                 C_CARD, None)
        add_rect(sl, col_split + Cm(0.3), ry, Cm(0.06), item_h_diag - Cm(0.1), C_TEAL)
        add_text(sl, str(txt),
                 col_split + Cm(0.6), ry + Cm(0.25), SW - col_split - Cm(1.4), Cm(0.9),
                 size=10, bold=False, color=C_VLTEAL, wrap=True)


# ── SLIDE 8: CALLS & OPORTUNIDADES ───────────────────────────────────────────
def s_calls(prs, d, num=6):
    """06 | Calls & Oportunidades"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)
    header(sl, num, "Calls & Oportunidades",
           sub="Recomendações e operações estruturadas do Head de Produtos",
           **b)

    calls       = d.get("calls", [])
    estruturadas = d.get("estruturadas", [])
    todos       = calls + estruturadas

    if not todos:
        add_text(sl, "Sem calls ou operações estruturadas para esta reunião.",
                 MARGIN, Cm(5), Cm(28), Cm(1), size=12, color=C_LGRAY)
        return

    card_h = Cm(1.8)
    gap    = Cm(0.25)
    col_w  = Cm(15.2)
    start_y = Cm(2.65)

    for idx, call in enumerate(todos[:10]):
        col = idx % 2
        row = idx // 2
        x = MARGIN + col * (col_w + Cm(0.5) + Cm(2.0))
        y = start_y + row * (card_h + gap)

        if isinstance(call, dict):
            titulo_c = call.get("titulo") or call.get("ativo") or call.get("nome", "Call")
            tipo_c   = call.get("tipo", "")
            desc_c   = call.get("descricao") or call.get("obs") or call.get("texto", "")
        else:
            titulo_c = str(call)
            tipo_c   = ""
            desc_c   = ""

        cor_tipo = {"compra": C_GREEN, "venda": C_RED, "estruturada": C_GOLD}.get(
            tipo_c.lower() if tipo_c else "", C_TEAL)

        add_rect(sl, x, y, col_w, card_h, C_CARD, C_TEAL_D, 0.3)
        add_rect(sl, x, y, Cm(0.08), card_h, cor_tipo)

        add_text(sl, titulo_c, x + Cm(0.3), y + Cm(0.15), col_w - Cm(0.5), Cm(0.65),
                 size=12, bold=True, color=C_WHITE)
        if tipo_c:
            add_text(sl, tipo_c.upper(), x + col_w - Cm(3.5), y + Cm(0.15), Cm(3.2), Cm(0.45),
                     size=9, bold=True, color=cor_tipo, align=PP_ALIGN.RIGHT)
        if desc_c:
            add_text(sl, str(desc_c)[:70], x + Cm(0.3), y + Cm(0.9), col_w - Cm(0.5), Cm(0.65),
                     size=10, color=C_LGRAY)


# ── SLIDE 9: RENDA VARIÁVEL ───────────────────────────────────────────────────
def s_rv(prs, d, num=7):
    """07 | Renda Variável — Destaques da Carteira"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)
    header(sl, num, "Renda Variável — Destaques da Carteira",
           sub="Principais posições e performance | Ações e FIIs",
           **b)

    acoes = d.get("acoes", [])
    fiis  = d.get("fiis", [])
    todos = acoes + fiis
    estruturadas = d.get("estruturadas", [])
    sugerir_estruturada = d.get("sugerir_estruturada_rv", False) or bool(estruturadas)

    # Decidir se mostra tabela RV compacta (para caber box de estruturada) ou completa
    max_rows_rv = 9 if sugerir_estruturada else 14

    if not todos and not sugerir_estruturada:
        add_text(sl, "Sem posições de renda variável registradas.",
                 MARGIN, Cm(5), Cm(28), Cm(1), size=12, color=C_LGRAY)
        return

    # Tabela de posições RV (metade esquerda quando há estruturada)
    tab_w   = Cm(18.5) if sugerir_estruturada else SW - Cm(1.6)
    tab_y   = Cm(2.65)
    COLS    = ["Ticker", "Saldo", "% Cart.", "Mês", "Ano", "12M"]
    col_xs  = [MARGIN + Cm(0.1), Cm(4.2), Cm(8.5), Cm(11.8), Cm(14.8), Cm(17.5)]
    col_ws  = [Cm(3.5), Cm(3.8), Cm(3.0), Cm(3.0), Cm(3.0), Cm(3.0)]

    if todos:
        add_rect(sl, MARGIN, tab_y, tab_w, Cm(0.6), C_TEAL_D)
        for j, (cx, cw, ch) in enumerate(zip(col_xs, col_ws, COLS)):
            if cx + cw > MARGIN + tab_w + Cm(0.3): break
            add_text(sl, ch, cx, tab_y + Cm(0.05), cw, Cm(0.5),
                     size=10, bold=True, color=C_LGRAY)

        row_colors = [C_CARD, C_CARD2]
        row_h = Cm(0.72)

        for ri, ativo in enumerate(todos[:max_rows_rv]):
            ry = tab_y + Cm(0.6) + ri * row_h
            add_rect(sl, MARGIN, ry, tab_w, row_h, row_colors[ri % 2])

            if isinstance(ativo, dict):
                ticker = ativo.get("ticker") or ativo.get("ativo") or ativo.get("nome", "—")
                saldo  = ativo.get("saldo", ativo.get("valor", 0))
                pct_c  = ativo.get("pct_carteira", ativo.get("percentual", 0))
                r_mes  = ativo.get("rentabilidade_mes", ativo.get("mes", "—"))
                r_ano  = ativo.get("rentabilidade_ano", ativo.get("ano", "—"))
                r_12m  = ativo.get("rentabilidade_12m", ativo.get("12m", "—"))
            else:
                ticker = str(ativo)
                saldo = pct_c = r_mes = r_ano = r_12m = "—"

            row_vals = [
                (str(ticker), C_TEAL_D, True),
                (fmt_brl_short(saldo) if saldo and saldo != "—" else "—", C_TEXT, True),
                (f"{float(pct_c):.2f}%" if pct_c and pct_c != "—" else "—", C_TEXT, True),
                (str(r_mes), pct_cor(r_mes), False),
                (str(r_ano), pct_cor(r_ano), False),
                (str(r_12m), pct_cor(r_12m), False),
            ]
            for j, (cx, cw, (val, cor, bld)) in enumerate(zip(col_xs, col_ws, row_vals)):
                if cx + cw > MARGIN + tab_w + Cm(0.3): break
                add_text(sl, val, cx, ry + Cm(0.1), cw, Cm(0.55),
                         size=10, bold=bld, color=cor)

    # ── Box de Operação Estruturada Sugerida (lado direito) ──────────────────
    if sugerir_estruturada:
        box_x = MARGIN + tab_w + Cm(0.4)
        box_w = SW - box_x - Cm(0.8)
        box_y = tab_y

        add_rect(sl, box_x, box_y, box_w, Cm(0.6), C_GOLD)
        add_text(sl, "📐 Operação Estruturada",
                 box_x + Cm(0.2), box_y + Cm(0.05), box_w - Cm(0.3), Cm(0.55),
                 size=10, bold=True, color=C_BG)

        item_y = box_y + Cm(0.7)
        if estruturadas:
            for ei, est in enumerate(estruturadas[:5]):
                if isinstance(est, dict):
                    nome_e = est.get("titulo") or est.get("nome") or est.get("ativo", "Estruturada")
                    tipo_e = est.get("tipo", "")
                    desc_e = est.get("descricao") or est.get("obs", "")
                    rent_e = est.get("rentabilidade") or est.get("taxa", "")
                else:
                    nome_e = str(est)
                    tipo_e = desc_e = rent_e = ""

                h_item = Cm(2.4)
                add_rect(sl, box_x, item_y, box_w, h_item, C_CARD, C_GOLD, 0.3)
                add_rect(sl, box_x, item_y, Cm(0.08), h_item, C_GOLD)
                add_text(sl, nome_e[:35], box_x + Cm(0.25), item_y + Cm(0.1),
                         box_w - Cm(0.4), Cm(0.6), size=11, bold=True, color=C_WHITE)
                if tipo_e:
                    add_text(sl, tipo_e, box_x + Cm(0.25), item_y + Cm(0.75),
                             box_w - Cm(0.4), Cm(0.4), size=9, color=C_GOLD)
                if desc_e:
                    add_text(sl, str(desc_e)[:55], box_x + Cm(0.25), item_y + Cm(1.18),
                             box_w - Cm(0.4), Cm(0.75), size=9, color=C_LGRAY, wrap=True)
                if rent_e:
                    add_text(sl, str(rent_e), box_x + Cm(0.25), item_y + Cm(1.92),
                             box_w - Cm(0.4), Cm(0.4), size=10, bold=True, color=C_GOLD)
                item_y += h_item + Cm(0.2)
        else:
            # Sugestão genérica quando não há estruturada do HP mas tem RV expressiva
            pct_rv = float((d.get("composicao") or {}).get("acoes", 0) or 0)
            pat = d.get("patrimonio", 0) or 0
            val_rv_str = fmt_brl(pat * pct_rv / 100) if pat else "—"
            sugestoes_gen = [
                ("Proteção de Carteira (Put)", "Hedge",
                 f"Proteção sobre posição RV ({pct_rv:.1f}% / {val_rv_str})\nIdeal para cenário de incerteza"),
                ("Trava de Alta (Call Spread)", "Renda",
                 "Potencializar ganho com risco limitado\nEntrada em queda com prêmio reduzido"),
                ("Certificado de Operação Estruturada", "COE",
                 "Proteção de capital + participação em ativos\nHorizonte 2-3 anos"),
            ]
            for ei, (nome_e, tipo_e, desc_e) in enumerate(sugestoes_gen):
                h_item = Cm(2.8)
                add_rect(sl, box_x, item_y, box_w, h_item, C_CARD, C_GOLD, 0.3)
                add_rect(sl, box_x, item_y, Cm(0.08), h_item, C_GOLD)
                add_text(sl, nome_e, box_x + Cm(0.25), item_y + Cm(0.1),
                         box_w - Cm(0.4), Cm(0.6), size=11, bold=True, color=C_WHITE)
                add_text(sl, tipo_e, box_x + Cm(0.25), item_y + Cm(0.78),
                         box_w - Cm(0.4), Cm(0.4), size=9, color=C_GOLD)
                add_text(sl, desc_e, box_x + Cm(0.25), item_y + Cm(1.25),
                         box_w - Cm(0.4), Cm(1.3), size=9, color=C_LGRAY, wrap=True)
                item_y += h_item + Cm(0.2)

        # Nota sobre estruturadas
        nota_y = item_y + Cm(0.1)
        if nota_y < SH - Cm(1.0):
            add_text(sl, "💡 Fale com seu assessor para estruturar a operação ideal",
                     box_x, nota_y, box_w, Cm(0.5), size=8, color=C_TEAL, wrap=True)


# ── SLIDE: INTERNACIONAL — 15% DÓLAR ─────────────────────────────────────────
def s_internacional(prs, d, num=8):
    """Internacional — Estudo Braúna: mínimo 15% em dólar"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)

    alerta = d.get("alerta_internacional", {})
    intl_atual  = alerta.get("atual_pct", 0)
    meta_ef     = alerta.get("meta_efetiva", 15.0)
    gap_pp      = alerta.get("gap_pp", 0)
    gap_val     = alerta.get("gap_valor", 0)
    produtos    = alerta.get("produtos_sugeridos", [])
    modelo_pct  = alerta.get("modelo_pct", 9.0)

    header(sl, num, "Internacional — Proteção Cambial & Diversificação Global",
           sub="Estudo Braúna: portfólio equilibrado deve ter mínimo 15% em dólar",
           **b)

    # KPIs topo
    kpis = [
        ("Alocação Atual",   f"{intl_atual:.2f}%",   "em ativos internacionais",  C_RED if intl_atual < 10 else C_AMBER),
        ("Meta do Modelo",   f"{modelo_pct:.1f}%",   "recomendação Levante Asset", C_TEAL),
        ("Mínimo Braúna",    "15,0%",                 "estudo dólar no patrimônio", C_GOLD),
        ("Gap a Alocar",     f"{gap_pp:+.1f}pp",      fmt_brl(gap_val),             C_RED if gap_pp > 0 else C_GREEN),
    ]

    kpi_w = Cm(7.5)
    kpi_h = Cm(2.5)
    gap   = Cm(0.5)
    row_y = Cm(2.65)
    for i, (lab, val, nota, cor) in enumerate(kpis):
        x = MARGIN + i * (kpi_w + gap)
        add_rect(sl, x, row_y, kpi_w, kpi_h, C_CARD, C_TEAL_D, 0.4)
        add_text(sl, lab, x + Cm(0.3), row_y + Cm(0.12), kpi_w - Cm(0.5), Cm(0.55),
                 size=10, bold=True, color=C_LGRAY)
        add_text(sl, val, x + Cm(0.3), row_y + Cm(0.72), kpi_w - Cm(0.5), Cm(1.0),
                 size=21, bold=True, color=cor)
        add_text(sl, nota, x + Cm(0.3), row_y + Cm(1.83), kpi_w - Cm(0.5), Cm(0.45),
                 size=9, color=RGBColor(0x8E, 0xCE, 0xCA))

    # Justificativa
    just_y = Cm(5.4)
    add_rect(sl, MARGIN, just_y, SW - Cm(1.6), Cm(1.3), C_CARD2, C_GOLD, 0.4)
    add_rect(sl, MARGIN, just_y, Cm(0.1), Cm(1.3), C_GOLD)
    add_text(sl, "📊 Por que 15% em dólar?",
             MARGIN + Cm(0.3), just_y + Cm(0.1), Cm(14), Cm(0.55),
             size=12, bold=True, color=C_GOLD)
    add_text(sl, "Portfólios com exposição cambial entre 15-20% apresentam menor drawdown máximo e maior Sharpe no longo prazo. "
                 "O dólar funciona como seguro nos ciclos de desvalorização do real e acesso ao crescimento global.",
             MARGIN + Cm(0.3), just_y + Cm(0.7), SW - Cm(2.0), Cm(0.55),
             size=10, color=C_VLTEAL, wrap=True)

    # Grid de produtos sugeridos
    prod_y = Cm(7.0)
    add_text(sl, "🎯 Produtos Sugeridos para a Carteira",
             MARGIN, prod_y - Cm(0.45), Cm(28), Cm(0.45),
             size=12, bold=True, color=C_TEAL)

    card_w = Cm(9.5)
    card_h = Cm(3.3)
    gap_x  = Cm(0.5)
    per_row = 3
    for idx, prod in enumerate(produtos[:6]):
        col = idx % per_row
        row = idx // per_row
        x = MARGIN + col * (card_w + gap_x)
        y = prod_y + row * (card_h + Cm(0.3))

        ticker = prod.get("ticker", "")
        nome_p = prod.get("nome", "")
        tipo_p = prod.get("tipo", "")
        desc_p = prod.get("descricao", "")

        cor_tipo = {"ETF": C_GOLD, "BDR": C_TEAL, "Fundo": C_LTEAL}.get(tipo_p, C_TEAL)

        add_rect(sl, x, y, card_w, card_h, C_CARD, C_TEAL_D, 0.4)
        add_rect(sl, x, y, Cm(0.1), card_h, cor_tipo)

        # Badge tipo
        add_rect(sl, x + card_w - Cm(2.2), y + Cm(0.2), Cm(2.0), Cm(0.5), cor_tipo)
        add_text(sl, tipo_p, x + card_w - Cm(2.2), y + Cm(0.2), Cm(2.0), Cm(0.5),
                 size=9, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

        add_text(sl, ticker, x + Cm(0.3), y + Cm(0.15), Cm(5), Cm(0.7),
                 size=16, bold=True, color=C_WHITE)
        add_text(sl, nome_p, x + Cm(0.3), y + Cm(0.9), card_w - Cm(0.5), Cm(0.55),
                 size=11, bold=False, color=C_LTEAL)
        add_text(sl, desc_p, x + Cm(0.3), y + Cm(1.55), card_w - Cm(0.5), Cm(1.2),
                 size=10, color=C_LGRAY, wrap=True)

    # Cálculo de quanto aportar
    if gap_val > 0:
        nota_y = prod_y + 2 * (card_h + Cm(0.3)) + Cm(0.1)
        if nota_y < SH - Cm(0.9):
            add_rect(sl, MARGIN, nota_y, SW - Cm(1.6), Cm(0.7), C_TEAL_D)
            add_text(sl,
                     f"💡 Para atingir {meta_ef:.0f}% de exposição internacional, "
                     f"seria necessário aportar aprox. {fmt_brl(gap_val)} em ativos dolarizados.",
                     MARGIN + Cm(0.3), nota_y + Cm(0.12),
                     SW - Cm(1.8), Cm(0.5), size=11, bold=False, color=C_WHITE)


# ── SLIDE 10: ECOSSISTEMA BRAÚNA ──────────────────────────────────────────────
def s_ecossistema(prs, d, num=11):
    """11 | Ecossistema Braúna — Muito Além dos Investimentos"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)

    pat_str = fmt_brl(d.get("patrimonio", 0))
    header(sl, num, "Ecossistema Braúna — Muito Além dos Investimentos",
           sub=f"Soluções que agregam valor ao seu patrimônio de {pat_str}",
           **b)

    cross = d.get("cross_ativos", [])
    cross_nomes = d.get("cross_ativos_nomes", [])

    # 6 serviços padrão + itens marcados do cross
    servicos = [
        ("📊", "Assessoria de Investimentos",
         "Carteiras personalizadas, rebalanceamento e acompanhamento contínuo"),
        ("🛡️", "Seguros & Proteção",
         "Vida, saúde e patrimônio — apólices revisadas e atualizadas"),
        ("🏢", "Consultoria PJ",
         "Gestão de caixa, conta PJ e planejamento tributário empresarial"),
        ("💳", "Crédito Estruturado",
         "Crédito com garantia em ativos — taxas competitivas ao patrimônio"),
        ("🌍", "Câmbio & Global",
         "Remessas, BDRs, ETFs globais e diversificação internacional"),
        ("📋", "Planejamento Sucessório",
         "Holding familiar, testamento, inventário e estruturação do legado"),
    ]

    card_w = Cm(9.5)
    card_h = Cm(3.5)
    gap_x  = Cm(0.4)
    gap_y  = Cm(0.35)
    start_y = Cm(2.65)
    per_row = 3

    for idx, (icone, titulo_s, descricao_s) in enumerate(servicos):
        col = idx % per_row
        row = idx // per_row
        x = MARGIN + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        in_cross = any(titulo_s.lower() in (str(c).lower()) for c in cross + cross_nomes)
        borda_cor = C_TEAL if in_cross else C_TEAL_D
        borda_w   = 1.0 if in_cross else 0.4

        add_rect(sl, x, y, card_w, card_h, C_CARD, borda_cor, borda_w)

        add_text(sl, icone, x + Cm(0.3), y + Cm(0.25), Cm(1.1), Cm(1.0), size=22)
        add_text(sl, titulo_s, x + Cm(1.7), y + Cm(0.3), card_w - Cm(2.0), Cm(0.65),
                 size=11, bold=True, color=C_TEAL)
        add_text(sl, descricao_s, x + Cm(1.7), y + Cm(1.05), card_w - Cm(2.0), Cm(1.6),
                 size=10, color=RGBColor(0x9E, 0xCF, 0xCC), wrap=True)

        if in_cross:
            add_rect(sl, x + card_w - Cm(2.0), y + Cm(0.15), Cm(1.8), Cm(0.5), C_TEAL)
            add_text(sl, "Em foco", x + card_w - Cm(2.0), y + Cm(0.15), Cm(1.8), Cm(0.5),
                     size=9, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # Rodapé destacado
    nota_y = start_y + 2 * (card_h + gap_y) + Cm(0.1)
    if nota_y < SH - Cm(0.9):
        add_rect(sl, MARGIN, nota_y, SW - Cm(1.6), Cm(0.65), C_CARD, C_TEAL_D, 0.3)
        add_text(sl, "💡 Com um patrimônio estruturado, a sinergia entre assessoria, proteção e planejamento potencializa seus resultados.",
                 MARGIN + Cm(0.3), nota_y + Cm(0.1), SW - Cm(2.2), Cm(0.5),
                 size=11, color=C_VLTEAL)


# ── SLIDE 11: SUGESTÕES DE PRODUTOS ──────────────────────────────────────────
def s_sugestoes(prs, d, num=9):
    """08 | Sugestões de Produtos"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)
    header(sl, num, "Sugestões de Produtos",
           sub="Produtos recomendados para otimizar a alocação da carteira",
           **b)

    sugs = d.get("sugestoes_produtos", [])
    if not sugs:
        add_text(sl, "Sem sugestões de produtos para esta reunião.",
                 MARGIN, Cm(5), Cm(28), Cm(1), size=12, color=C_LGRAY)
        return

    card_w = Cm(15.0)
    card_h = Cm(2.4)
    gap    = Cm(0.3)
    start_y = Cm(2.65)
    per_row = 2

    for idx, sug in enumerate(sugs[:8]):
        col = idx % per_row
        row = idx // per_row
        x = MARGIN + col * (card_w + Cm(0.6))
        y = start_y + row * (card_h + gap)

        if isinstance(sug, dict):
            if "produto" in sug and isinstance(sug["produto"], dict):
                prod   = sug["produto"]
                classe = sug.get("classe", "")
                gap_v  = sug.get("gap", 0)
                label_c = sug.get("label_classe", CLS_LABEL.get(classe, classe))
            else:
                prod   = sug
                classe = sug.get("classe", sug.get("cls", ""))
                gap_v  = sug.get("gap", 0)
                label_c = CLS_LABEL.get(classe, classe)

            nome_p = prod.get("nome") or prod.get("titulo") or prod.get("produto", "Produto")
            tipo_p = prod.get("tipo", prod.get("subtipo", ""))
            desc_p = prod.get("descricao") or prod.get("obs", "")
            rent_p = prod.get("rentabilidade") or prod.get("taxa", "")
        else:
            nome_p = str(sug)
            tipo_p = label_c = desc_p = rent_p = ""
            gap_v  = 0

        cor = CLS_COR.get(classe, C_TEAL) if classe else C_TEAL
        add_rect(sl, x, y, card_w, card_h, C_CARD, C_TEAL_D, 0.3)
        add_rect(sl, x, y, Cm(0.1), card_h, cor)

        add_text(sl, nome_p, x + Cm(0.3), y + Cm(0.15), card_w - Cm(0.5), Cm(0.65),
                 size=12, bold=True, color=C_WHITE)
        if label_c or tipo_p:
            tag = f"{label_c}  ·  {tipo_p}" if label_c and tipo_p else (label_c or tipo_p)
            add_text(sl, tag, x + Cm(0.3), y + Cm(0.85), card_w - Cm(0.5), Cm(0.45),
                     size=10, color=C_TEAL)
        if desc_p:
            add_text(sl, str(desc_p)[:60], x + Cm(0.3), y + Cm(1.4), card_w - Cm(0.5), Cm(0.6),
                     size=10, color=C_LGRAY)
        if rent_p:
            add_text(sl, str(rent_p), x + card_w - Cm(4.0), y + Cm(1.7), Cm(3.7), Cm(0.45),
                     size=10, bold=True, color=C_GOLD, align=PP_ALIGN.RIGHT)


# ── SLIDE 12: ESTRATÉGIAS & RECOMENDAÇÕES ────────────────────────────────────
def s_recomendacoes(prs, d, num=10):
    """10 | Estratégias & Recomendações — usa metas do assessor + regra 15% dólar"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)

    data_r = b["data_ref"]
    try:
        ano_txt = "2º Semestre " + data_r.split("/")[2]
    except Exception:
        ano_txt = ""

    header(sl, num, f"Estratégias & Recomendações — {ano_txt}",
           sub="Baseadas na análise da carteira, metas do assessor e cenário macroeconômico",
           **b)

    recomendacoes = list(d.get("recomendacoes", []))  # cópia para não mutar
    desvios  = d.get("desvios", [])
    sugs     = d.get("sugestoes_produtos", [])
    alerta_i = d.get("alerta_internacional", {})
    metas_a  = d.get("assessor_metas", {}) or {}
    pat      = d.get("patrimonio", 0) or 0
    composicao = d.get("composicao", {}) or {}

    # ── Bloco 1: Rebalanceamento — desvios do modelo ─────────────────────────
    if not recomendacoes:
        fora = [dv for dv in desvios if abs(_pct_float(dv.get("desvio", dv.get("desvio_pp", 0)))) > 5]
        atencao = [dv for dv in desvios if 2 < abs(_pct_float(dv.get("desvio", dv.get("desvio_pp", 0)))) <= 5]

        if fora or atencao:
            bullets_reb = []
            for dv in (fora + atencao)[:4]:
                lbl   = dv.get("label") or CLS_LABEL.get(dv.get("cls",""), "")
                dev_v = _pct_float(dv.get("desvio", dv.get("desvio_pp", 0)))
                atual = dv.get("atual", 0)
                alvo  = dv.get("alvo", dv.get("modelo", 0))
                val_mov = abs(pat * dev_v / 100)
                mov_str = fmt_brl(val_mov) if val_mov > 1000 else ""
                if dev_v > 0:
                    bullets_reb.append(f"Reduzir {lbl}: {dev_v:+.1f}pp ({atual:.1f}% → {alvo:.1f}%){' · ' + mov_str if mov_str else ''}")
                else:
                    bullets_reb.append(f"Ampliar {lbl}: {dev_v:+.1f}pp ({atual:.1f}% → {alvo:.1f}%){' · ' + mov_str if mov_str else ''}")
            recomendacoes.append(("🔄 Rebalanceamento Prioritário", C_TEAL, bullets_reb))

    # ── Bloco 2: Internacional — regra 15% dólar ─────────────────────────────
    if alerta_i.get("perfil_elegivel") and alerta_i.get("gap_pp", 0) > 0:
        gap_pp_i  = alerta_i.get("gap_pp", 0)
        gap_val_i = alerta_i.get("gap_valor", 0)
        meta_i    = alerta_i.get("meta_efetiva", 15)
        atual_i   = alerta_i.get("atual_pct", 0)
        prods_i   = alerta_i.get("produtos_sugeridos", [])[:3]
        tickers   = " · ".join(p.get("ticker","") for p in prods_i if p.get("ticker"))

        bullets_intl = [
            f"Exposição atual: {atual_i:.1f}% — meta Braúna: {meta_i:.0f}% (gap: {gap_pp_i:+.1f}pp)",
            f"Valor a alocar para atingir meta: {fmt_brl(gap_val_i)}",
            f"Produtos recomendados: {tickers}" if tickers else "BDRs e ETFs globais (ver slide Internacional)",
            "Dólar como hedge cambial e acesso ao crescimento global",
        ]
        recomendacoes.append(("🌍 Internacional — Mínimo 15% em Dólar", C_GOLD, bullets_intl))

    # ── Bloco 3: Estruturadas — se assessor tem gap nessa área ───────────────
    estrt_meta = metas_a.get("estruturadas", {})
    pct_estrt_atual = float(estrt_meta.get("pct_max", 0) or 0)
    pct_rv_cli = float(composicao.get("acoes", 0) or 0)

    if pct_rv_cli >= 3.0:
        bullets_estrt = [
            f"Carteira com {pct_rv_cli:.1f}% em RV — elegível para proteção estruturada",
            "Trava de alta (call spread) para potencializar ganho com risco limitado",
            "Proteção de portfólio (put) para cenários de volatilidade",
            "COE: proteção de capital + participação em ativos internacionais",
        ]
        recomendacoes.append(("📐 Operações Estruturadas", C_AMBER, bullets_estrt))

    # ── Bloco 4: Produto sugerido / complementar ──────────────────────────────
    # Prioriza classes onde o assessor tem gap de receita
    meta_rv_asses  = metas_a.get("rv", {})
    meta_fii_asses = metas_a.get("fii", {})
    meta_rf_asses  = metas_a.get("rf", {})
    pct_falta_meta = float(metas_a.get("pct_falta_meta", 0) or 0)

    bullets_prod = []
    if sugs:
        for s in sugs[:3]:
            if isinstance(s, dict) and "produto" in s:
                pnome = s["produto"].get("nome") or s["produto"].get("titulo","")
            elif isinstance(s, dict):
                pnome = s.get("nome") or s.get("titulo","")
            else:
                pnome = str(s)
            if pnome:
                bullets_prod.append(pnome)

    # Adiciona hint de FII se assessor está abaixo em fii e cliente tem pouco
    pct_fii_cli = float(composicao.get("fiis", 0) or 0)
    modelo_hp_fii = float((d.get("modelo_hp") or {}).get("fiis", 0) or 0)
    if pct_fii_cli < modelo_hp_fii - 1 and meta_fii_asses.get("pct", 0) < 10:
        bullets_prod.append(f"FIIs: alocação atual {pct_fii_cli:.1f}% abaixo do modelo {modelo_hp_fii:.1f}% — oportunidade de renda passiva")

    if bullets_prod:
        recomendacoes.append(("📈 Produtos & Ampliações", RGBColor(0x1E, 0x84, 0x49), bullets_prod))

    # Fallback se ainda vazio
    if not recomendacoes:
        recomendacoes = [
            ("🔄 Rebalanceamento", C_TEAL,
             ["Revisar alocação em relação ao modelo de perfil",
              "Adequar exposições fora do range recomendado"]),
            ("📋 Planejamento", C_GOLD,
             ["Revisar objetivos financeiros e horizonte de investimento",
              "Avaliar necessidades de proteção e previdência"]),
        ]

    bloco_w = Cm(14.0)
    bloco_y = Cm(2.65)
    bloco_gap = Cm(0.4)

    for idx, bloco in enumerate(recomendacoes[:4]):
        if isinstance(bloco, (list, tuple)) and len(bloco) >= 2:
            titulo_b = bloco[0]
            cor_b    = bloco[1] if len(bloco) > 1 else C_TEAL
            bullets  = bloco[2] if len(bloco) > 2 else []
        else:
            titulo_b = str(bloco)
            cor_b    = C_TEAL
            bullets  = []

        col = idx % 2
        row = idx // 2
        x = MARGIN + col * (bloco_w + Cm(1.0))
        y = bloco_y + row * (Cm(4.0) + bloco_gap)

        bullet_h = max(Cm(3.0), Cm(0.6) * len(bullets)) + Cm(1.0)
        add_rect(sl, x, y, bloco_w, bullet_h, C_CARD, C_TEAL_D, 0.4)
        add_rect(sl, x, y, Cm(0.1), bullet_h, cor_b)

        add_text(sl, titulo_b, x + Cm(0.3), y + Cm(0.15), bloco_w - Cm(0.5), Cm(0.75),
                 size=13, bold=True, color=C_WHITE)

        for bi, bullet in enumerate(bullets[:4]):
            by = y + Cm(1.05) + bi * Cm(0.65)
            add_rect(sl, x + Cm(0.3), by + Cm(0.15), Cm(0.25), Cm(0.25), cor_b)
            add_text(sl, str(bullet)[:70], x + Cm(0.75), by, bloco_w - Cm(1.0), Cm(0.6),
                     size=10, color=C_VLTEAL)


# ── SLIDE 13: PRÓXIMOS PASSOS ─────────────────────────────────────────────────
def s_proximos(prs, d, num=12):
    """12 | Próximos Passos — Plano de Ação"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)
    header(sl, num, "Próximos Passos — Plano de Ação",
           sub="Ações acordadas nesta reunião — responsabilidades e prazos",
           **b)

    proximos = d.get("proximos_passos", [])

    # Separar por horizonte temporal (imediato / 30 dias / 90 dias)
    imediato = []
    trinta   = []
    noventa  = []

    for item in proximos:
        if isinstance(item, dict):
            txt = item.get("descricao") or item.get("texto") or item.get("acao", str(item))
            prazo = str(item.get("prazo", "")).lower()
        else:
            txt = str(item)
            prazo = ""

        if "imediato" in prazo or "semana" in prazo or "urgente" in prazo:
            imediato.append(txt)
        elif "30" in prazo or "trinta" in prazo or "julho" in prazo.lower():
            trinta.append(txt)
        elif "90" in prazo or "noventa" in prazo or "setembro" in prazo.lower():
            noventa.append(txt)
        else:
            # Distribuir proporcionalmente
            if len(imediato) <= len(trinta) and len(imediato) <= len(noventa):
                imediato.append(txt)
            elif len(trinta) <= len(noventa):
                trinta.append(txt)
            else:
                noventa.append(txt)

    # Defaults se vazio
    if not imediato and not trinta and not noventa:
        imediato = ["Agendar próxima reunião de acompanhamento",
                    "Revisar cadastro e suitability",
                    "Confirmar ações acordadas nesta reunião"]
        trinta   = ["Implementar ajustes de rebalanceamento identificados",
                    "Apresentar novas sugestões de produtos"]
        noventa  = ["Revisão trimestral completa da carteira",
                    "Relatório consolidado de planejamento financeiro"]

    # 3 colunas de horizonte
    col_w   = Cm(9.7)
    col_gap = Cm(0.55)
    col_ys  = [Cm(2.65)] * 3
    cols    = [
        ("Imediato",   "Essa semana",     C_RED,   imediato),
        ("30 dias",    "Até próximo mês", C_AMBER, trinta),
        ("90 dias",    "Próximo trimestre", C_GREEN, noventa),
    ]

    for ci, (titulo_h, subtit, cor_h, items) in enumerate(cols):
        x = MARGIN + ci * (col_w + col_gap)
        y = col_ys[ci]

        # Header da coluna
        add_rect(sl, x, y, col_w, Cm(1.5), cor_h)
        add_text(sl, titulo_h, x + Cm(0.3), y + Cm(0.05), col_w - Cm(0.4), Cm(0.95),
                 size=22, bold=True, color=C_WHITE)
        add_text(sl, subtit, x + Cm(0.3), y + Cm(1.1), col_w - Cm(0.4), Cm(0.4),
                 size=11, color=C_WHITE)

        # Items
        for ii, item_txt in enumerate(items[:6]):
            iy = y + Cm(1.65) + ii * Cm(1.5)
            add_rect(sl, x, iy, col_w, Cm(1.4), C_CARD, C_TEAL_D, 0.3)

            # Badge numérico
            add_rect(sl, x, iy, Cm(0.55), Cm(1.4), cor_h)
            add_text(sl, str(ii + 1), x + Cm(0.05), iy + Cm(0.3), Cm(0.45), Cm(0.6),
                     size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

            add_text(sl, str(item_txt)[:55], x + Cm(0.7), iy + Cm(0.2),
                     col_w - Cm(0.85), Cm(1.0), size=10, color=C_TEXT, wrap=True)


# ── SLIDE 14: ENCERRAMENTO ────────────────────────────────────────────────────
def s_encerramento(prs, d):
    """Encerramento — Antecipe seus sonhos."""
    sl = new_slide(prs)
    bg(sl)

    b = _get_base(d)
    assessor = d.get("assessor", "Assessor Braúna")
    nome     = d.get("nome_cliente", "")
    pat      = d.get("patrimonio", 0)
    r_mes    = d.get("rentabilidade_mes", "—")
    r_12m    = d.get("rentabilidade_12m", "—")
    r_24m    = d.get("rentabilidade_24m", "—")

    # Faixa vertical esquerda
    add_rect(sl, 0, 0, Cm(1.2), SH, C_CARD)
    add_rect(sl, 0, 0, Cm(0.12), SH, C_GOLD)

    # Frase principal
    add_text(sl, "Antecipe seus sonhos.", Cm(1.7), Cm(2.5), Cm(28), Cm(3.5),
             size=54, bold=True, color=C_WHITE)
    add_text(sl, "Obrigado.", Cm(1.7), Cm(6.2), Cm(28), Cm(2.5),
             size=46, bold=True, color=C_GOLD)

    # Assessor
    add_text(sl, f"{assessor}  |  Assessor de Investimentos  |  XP Unique",
             Cm(1.7), Cm(9.2), Cm(28), Cm(0.7),
             size=14, color=C_DKGRAY)
    add_text(sl, "Braúna Ecossistema de Investimentos  |  braunainvestimentos.com.br",
             Cm(1.7), Cm(10.0), Cm(28), Cm(0.6),
             size=13, color=RGBColor(0x6A, 0x9E, 0x9A))
    add_text(sl, "/Braunainvestimentos",
             Cm(1.7), Cm(10.7), Cm(28), Cm(0.55),
             size=12, color=RGBColor(0x6A, 0x9E, 0x9A))

    # Card resumo da reunião
    card_x = Cm(1.7)
    card_y = Cm(11.7)
    card_w = Cm(30.0)
    card_h = Cm(4.5)
    add_rect(sl, card_x, card_y, card_w, card_h, C_CARD, C_TEAL_D, 0.5)
    add_rect(sl, card_x, card_y, Cm(0.1), card_h, C_TEAL)

    add_text(sl, "Resumo da Reunião",
             card_x + Cm(0.4), card_y + Cm(0.2), Cm(10), Cm(0.6),
             size=12, bold=True, color=C_TEAL)

    resumo_items = [
        ("💰", "Patrimônio:",   fmt_brl(pat)),
        ("📉", "Rent. Mês:",    str(r_mes)),
        ("📈", "Rent. 12M:",    str(r_12m)),
        ("📊", "Rent. 24M:",    str(r_24m)),
    ]

    item_w = card_w / len(resumo_items)
    for ii, (icone, lab, val) in enumerate(resumo_items):
        ix = card_x + Cm(0.4) + ii * item_w
        add_text(sl, icone, ix, card_y + Cm(0.9), Cm(0.7), Cm(0.75), size=14)
        add_text(sl, lab, ix + Cm(0.75), card_y + Cm(0.9), item_w - Cm(0.8), Cm(0.5),
                 size=10, color=RGBColor(0x8E, 0xCE, 0xCA))
        cor_v = pct_cor(val) if "%" in str(val) else C_WHITE
        add_text(sl, val, ix + Cm(0.75), card_y + Cm(1.5), item_w - Cm(0.8), Cm(0.6),
                 size=10, bold=True, color=cor_v)

    # Hashtag
    add_text(sl, "#AntecipeSeusSonhos",
             card_x, card_y + card_h + Cm(0.3), Cm(20), Cm(0.65),
             size=14, bold=True, color=C_TEAL_DP)

    footer(sl, b["data_ref"], b["conta"], b["tipo_conta"])


# ── SLIDE: ANÁLISE DE PRODUTOS & SUGESTÕES DE TROCA ─────────────────────────
def s_analise_produtos(prs, d, num=5):
    """NN | Análise de Produtos & Sugestões de Troca"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)
    header(sl, num, "Análise de Produtos & Sugestões de Troca",
           sub="Baseada no asset allocation do perfil e metas do assessor", **b)

    sugestoes = d.get("sugestoes_por_classe", []) or []

    if not sugestoes:
        add_rect(sl, MARGIN, Cm(2.65), SW - 2 * MARGIN, Cm(3.0), C_CARD, C_TEAL_D, 0.4)
        add_text(sl, "✅  Carteira alinhada ao modelo — nenhuma divergência relevante identificada.",
                 MARGIN + Cm(0.4), Cm(3.1), SW - 2 * MARGIN - Cm(0.8), Cm(1.5),
                 size=13, color=C_VLTEAL)
        return

    card_w   = Cm(14.8)
    card_gap = Cm(0.5)
    card_h   = Cm(7.3)
    top_y    = Cm(2.65)

    COR_PRIOR = {"alta": C_RED, "media": C_AMBER, "baixa": C_GREEN}
    LABEL_PRIOR = {"alta": "ALTA", "media": "MÉDIA", "baixa": "BAIXA"}

    for idx, sug in enumerate(sugestoes[:4]):
        col = idx % 2
        row = idx // 2
        x = MARGIN + col * (card_w + card_gap)
        y = top_y + row * (card_h + Cm(0.3))

        prior     = sug.get("prioridade", "baixa")
        cor_prior = COR_PRIOR.get(prior, C_TEAL)
        gap_pp    = float(sug.get("gap_pp", 0) or 0)
        gap_val   = float(sug.get("gap_valor", 0) or 0)
        label_cls = sug.get("label_classe", sug.get("classe", "—"))
        motivo    = sug.get("motivo", "")
        produtos  = sug.get("produtos", []) or []

        # Card — borda dourada se alta prioridade
        borda = C_GOLD if prior == "alta" else C_TEAL_D
        add_rect(sl, x, y, card_w, card_h, C_CARD, borda, 0.6 if prior == "alta" else 0.3)
        add_rect(sl, x, y, Cm(0.12), card_h, cor_prior)

        # Título da classe
        add_text(sl, label_cls, x + Cm(0.3), y + Cm(0.15), card_w - Cm(2.5), Cm(0.65),
                 size=14, bold=True, color=C_WHITE)

        # Badge de prioridade
        badge_x = x + card_w - Cm(2.1)
        add_rect(sl, badge_x, y + Cm(0.1), Cm(1.9), Cm(0.55), cor_prior)
        add_text(sl, LABEL_PRIOR.get(prior, prior.upper()),
                 badge_x, y + Cm(0.12), Cm(1.9), Cm(0.4),
                 size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Gap info
        gap_txt = f"Gap: {abs(gap_pp):.1f}pp  |  {fmt_brl(gap_val)} a alocar"
        add_text(sl, gap_txt, x + Cm(0.3), y + Cm(0.88), card_w - Cm(0.5), Cm(0.5),
                 size=10, color=C_GOLD)

        # Motivo
        add_text(sl, str(motivo)[:90], x + Cm(0.3), y + Cm(1.45), card_w - Cm(0.5), Cm(0.7),
                 size=9, color=C_LTEAL, italic=True)

        # Divisor
        add_rect(sl, x + Cm(0.3), y + Cm(2.2), card_w - Cm(0.6), Cm(0.04), C_TEAL_D)

        # Produtos sugeridos
        add_text(sl, "Produtos sugeridos:", x + Cm(0.3), y + Cm(2.3), card_w - Cm(0.5), Cm(0.45),
                 size=9, bold=True, color=C_TEAL)

        for pi, prod in enumerate(produtos[:3]):
            py = y + Cm(2.85) + pi * Cm(1.35)
            add_rect(sl, x + Cm(0.3), py, card_w - Cm(0.6), Cm(1.2), C_CARD2, C_TEAL_D, 0.25)
            nome_p = str(prod.get("nome", "—"))[:50]
            tipo_p = str(prod.get("tipo", ""))
            taxa_p = str(prod.get("taxa", ""))
            badge_txt = tipo_p or taxa_p
            add_text(sl, nome_p, x + Cm(0.55), py + Cm(0.1), card_w - Cm(2.5), Cm(0.55),
                     size=10, bold=True, color=C_WHITE)
            if badge_txt:
                add_rect(sl, x + card_w - Cm(2.1), py + Cm(0.1), Cm(1.7), Cm(0.45), C_TEAL_D)
                add_text(sl, badge_txt[:16], x + card_w - Cm(2.1), py + Cm(0.12),
                         Cm(1.7), Cm(0.35), size=8, color=C_TEAL, align=PP_ALIGN.CENTER)
            desc_p = str(prod.get("descricao", ""))[:65]
            if desc_p:
                add_text(sl, desc_p, x + Cm(0.55), py + Cm(0.7), card_w - Cm(1.0), Cm(0.45),
                         size=8, color=C_LGRAY)


# ── SLIDE: AÇÕES — ALERTAS DE PREÇO & ESTRATÉGIA ─────────────────────────────
def s_acoes_alerta(prs, d, num=8):
    """NN | Ações — Alertas de Preço & Estratégia"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)
    header(sl, num, "Ações — Alertas de Preço & Estratégia",
           sub="Posições acima do preço-teto do HP · Operações estruturadas recomendadas", **b)

    alertas = d.get("acoes_alerta", []) or []

    if not alertas:
        add_rect(sl, MARGIN, Cm(2.65), SW - 2 * MARGIN, Cm(3.0), C_CARD, C_GREEN, 0.5)
        add_rect(sl, MARGIN, Cm(2.65), Cm(0.12), Cm(3.0), C_GREEN)
        add_text(sl, "✅  Todas as posições dentro dos limites do HP",
                 MARGIN + Cm(0.4), Cm(3.15), SW - 2 * MARGIN - Cm(0.8), Cm(0.7),
                 size=14, bold=True, color=C_WHITE)
        add_text(sl, "Nenhuma ação da carteira está acima do preço-teto ou no stop configurado.",
                 MARGIN + Cm(0.4), Cm(3.95), SW - 2 * MARGIN - Cm(0.8), Cm(0.65),
                 size=11, color=C_VLTEAL)
        return

    card_w   = Cm(14.8)
    card_gap = Cm(0.5)
    card_h   = Cm(7.5)
    top_y    = Cm(2.65)

    for idx, alerta in enumerate(alertas[:4]):
        col = idx % 2
        row = idx // 2
        x = MARGIN + col * (card_w + card_gap)
        y = top_y + row * (card_h + Cm(0.3))

        situacao   = alerta.get("situacao", "monitorar")
        ticker     = str(alerta.get("ticker", "—"))
        saldo      = float(alerta.get("saldo", 0) or 0)
        pct_cart   = float(alerta.get("pct_carteira", 0) or 0)
        preco_alvo = alerta.get("preco_alvo")
        upside     = alerta.get("upside")
        acao_rec   = str(alerta.get("acao_recomendada", ""))
        estruturada = str(alerta.get("estruturada_sugerida", ""))

        cor_sit = C_RED if situacao == "acima_alvo" else C_AMBER
        label_sit = "ACIMA DO ALVO" if situacao == "acima_alvo" else "MONITORAR"
        icone_sit = "🔴" if situacao == "acima_alvo" else "⚠️"

        add_rect(sl, x, y, card_w, card_h, C_CARD, cor_sit, 0.5)
        add_rect(sl, x, y, Cm(0.12), card_h, cor_sit)

        preco_atual = alerta.get("preco_atual")
        variacao_d  = alerta.get("variacao_dia")
        preco_ent   = alerta.get("preco_entrada")
        tese        = str(alerta.get("tese", ""))

        # Ticker grande
        add_text(sl, ticker, x + Cm(0.3), y + Cm(0.1), Cm(5.5), Cm(1.0),
                 size=26, bold=True, color=C_WHITE)

        # Preço atual em tempo real (brapi.dev) — exibido próximo ao ticker
        if preco_atual is not None:
            var_cor  = C_GREEN if (variacao_d or 0) >= 0 else C_RED
            var_txt  = f"({variacao_d:+.2f}%)" if variacao_d is not None else ""
            add_text(sl, f"R$ {preco_atual:.2f}  {var_txt}",
                     x + Cm(0.3), y + Cm(1.0), Cm(7.0), Cm(0.55),
                     size=11, bold=True, color=var_cor)

        # Badge situação
        badge_x = x + card_w - Cm(3.5)
        add_rect(sl, badge_x, y + Cm(0.12), Cm(3.3), Cm(0.65), cor_sit)
        add_text(sl, f"{icone_sit} {label_sit}",
                 badge_x, y + Cm(0.14), Cm(3.3), Cm(0.5),
                 size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Saldo e % carteira
        add_text(sl, f"{fmt_brl(saldo)}   |   {pct_cart:.1f}% da carteira",
                 x + Cm(0.3), y + Cm(1.65), card_w - Cm(0.5), Cm(0.45),
                 size=10, color=C_LTEAL)

        # Preço entrada → alvo → upside
        info_parts = []
        if preco_ent is not None:
            info_parts.append(f"Entrada: R$ {preco_ent}")
        if preco_alvo is not None:
            info_parts.append(f"Alvo HP: R$ {preco_alvo}")
        if upside is not None:
            try:
                up_f = float(upside)
                info_parts.append(f"Upside: {up_f:+.1f}%")
            except Exception:
                info_parts.append(f"Upside: {upside}")
        if info_parts:
            add_text(sl, "  |  ".join(info_parts),
                     x + Cm(0.3), y + Cm(2.2), card_w - Cm(0.5), Cm(0.45),
                     size=9, color=C_GOLD)

        # Divisor
        add_rect(sl, x + Cm(0.3), y + Cm(2.75), card_w - Cm(0.6), Cm(0.04), C_TEAL_D)

        # Ação recomendada
        if acao_rec:
            add_text(sl, "Ação recomendada:", x + Cm(0.3), y + Cm(2.9), card_w - Cm(0.5), Cm(0.4),
                     size=9, bold=True, color=C_TEAL)
            add_text(sl, str(acao_rec)[:100], x + Cm(0.3), y + Cm(3.4),
                     card_w - Cm(0.5), Cm(0.85), size=10, color=C_WHITE)

        # Tese resumida (se disponível)
        if tese:
            add_text(sl, str(tese)[:80], x + Cm(0.3), y + Cm(4.3),
                     card_w - Cm(0.5), Cm(0.6), size=9, color=C_DKGRAY)

        # Estruturada sugerida
        if estruturada:
            add_rect(sl, x + Cm(0.3), y + Cm(5.0), card_w - Cm(0.6), Cm(1.8), C_CARD2, C_GOLD, 0.4)
            add_rect(sl, x + Cm(0.3), y + Cm(5.0), Cm(0.1), Cm(1.8), C_GOLD)
            add_text(sl, "📐 Estruturada sugerida:", x + Cm(0.55), y + Cm(5.1),
                     card_w - Cm(0.9), Cm(0.4), size=9, bold=True, color=C_GOLD)
            add_text(sl, str(estruturada)[:100], x + Cm(0.55), y + Cm(5.6),
                     card_w - Cm(0.9), Cm(1.0), size=10, color=C_AMBER)


# ── SLIDE: CARTEIRA FII — PORTFÓLIO AUTOMATIZADO ─────────────────────────────
def s_carteira_fiis(prs, d, num=9):
    """NN | Carteira FII — Portfólio Automatizado"""
    sl = new_slide(prs)
    bg(sl)
    b = _get_base(d)
    header(sl, num, "Carteira FII — Portfólio Automatizado",
           sub="Sugestão de alocação em Fundos Imobiliários conforme perfil", **b)

    fiis_data = d.get("carteira_fiis_sugerida", {}) or {}
    posicoes  = fiis_data.get("posicoes", []) or []
    gestora_nome   = fiis_data.get("gestora_nome", "Carteira Braúna FII")
    referencia     = fiis_data.get("referencia", "")
    total_pct      = float(fiis_data.get("total_sugerido_pct", 0) or 0)
    total_val      = float(fiis_data.get("total_sugerido_valor", 0) or 0)

    if not posicoes:
        add_rect(sl, MARGIN, Cm(2.65), SW - 2 * MARGIN, Cm(2.5), C_CARD, C_TEAL_D, 0.4)
        add_text(sl, "Dados de FIIs não disponíveis para este perfil.",
                 MARGIN + Cm(0.4), Cm(3.0), SW - 2 * MARGIN - Cm(0.8), Cm(1.2),
                 size=13, color=C_VLTEAL)
        return

    # KPIs topo (4 cards)
    kpi_w = Cm(7.5)
    kpi_h = Cm(2.1)
    kpi_y = Cm(2.65)
    kpis  = [
        ("Total Sugerido", fmt_brl(total_val), C_GOLD),
        ("% do Patrimônio", f"{total_pct:.1f}%",  C_TEAL),
        ("Nº de FIIs",      str(len(posicoes)),    C_GREEN),
        ("Referência",      gestora_nome[:22],     C_LGRAY),
    ]
    for ki, (lab, val, cor) in enumerate(kpis):
        kx = MARGIN + ki * (kpi_w + Cm(0.37))
        add_rect(sl, kx, kpi_y, kpi_w, kpi_h, C_CARD, C_TEAL_D, 0.3)
        add_rect(sl, kx, kpi_y, Cm(0.1), kpi_h, cor)
        add_text(sl, lab, kx + Cm(0.25), kpi_y + Cm(0.1), kpi_w - Cm(0.4), Cm(0.55),
                 size=9, color=C_LGRAY)
        add_text(sl, val, kx + Cm(0.25), kpi_y + Cm(0.7), kpi_w - Cm(0.4), Cm(1.1),
                 size=16, bold=True, color=cor)

    # Cabeçalho da tabela
    tbl_y = Cm(5.1)
    tbl_h_row = Cm(1.55)
    cols_x = [MARGIN, Cm(4.2), Cm(13.5), Cm(19.5), Cm(25.5), Cm(29.5)]
    cols_w = [Cm(3.8), Cm(9.0), Cm(5.7), Cm(5.7), Cm(3.7), Cm(3.0)]
    hdrs   = ["Ticker", "Fundo", "Tipo", "% Carteira FII", "Valor Sugerido", "Destaque"]

    add_rect(sl, MARGIN, tbl_y, SW - 2 * MARGIN, Cm(0.65), C_TEAL_D)
    for ci, (hdr, cx, cw) in enumerate(zip(hdrs, cols_x, cols_w)):
        add_text(sl, hdr, cx + Cm(0.1), tbl_y + Cm(0.05), cw, Cm(0.55),
                 size=9, bold=True, color=C_WHITE)

    tem_etf = any((p.get("tipo", "").upper() in ("ETF", "IFIX")) for p in posicoes)

    for ri, pos in enumerate(posicoes[:8]):
        ry   = tbl_y + Cm(0.65) + ri * tbl_h_row
        cor_bg = C_CARD2 if ri % 2 == 0 else C_CARD
        add_rect(sl, MARGIN, ry, SW - 2 * MARGIN, tbl_h_row - Cm(0.05), cor_bg)

        ticker = str(pos.get("ticker", "—"))
        nome   = str(pos.get("nome", "—"))[:40]
        tipo   = str(pos.get("tipo", "—"))
        pct_c  = float(pos.get("pct_carteira", 0) or 0)
        val_c  = float(pos.get("valor", 0) or 0)

        is_etf = tipo.upper() in ("ETF", "IFIX") or "IFIX" in ticker.upper()
        destaque = "ETF diversif." if is_etf else ""

        vals = [ticker, nome, tipo, f"{pct_c:.2f}%", fmt_brl(val_c), destaque]
        cors = [C_TEAL if is_etf else C_WHITE, C_WHITE, C_LTEAL, C_GOLD, C_GREEN,
                C_AMBER if destaque else C_LGRAY]

        for ci, (v, cx, cw, cv) in enumerate(zip(vals, cols_x, cols_w, cors)):
            add_text(sl, v, cx + Cm(0.1), ry + Cm(0.35), cw - Cm(0.15), Cm(0.75),
                     size=10, bold=(ci == 0), color=cv)

    # Linha de destaque ETF
    if tem_etf:
        nota_y = tbl_y + Cm(0.65) + len(posicoes[:8]) * tbl_h_row + Cm(0.2)
        add_rect(sl, MARGIN, nota_y, SW - 2 * MARGIN, Cm(0.65), C_TEAL_D, C_GOLD, 0.4)
        add_text(sl, "★  ETF IFIX11 — diversificação instantânea ao setor imobiliário listado na B3",
                 MARGIN + Cm(0.3), nota_y + Cm(0.08), SW - 2 * MARGIN - Cm(0.5), Cm(0.5),
                 size=9, color=C_GOLD)

    if referencia:
        add_text(sl, f"Fonte: {gestora_nome}  |  Ref.: {referencia}",
                 MARGIN, SH - Cm(0.8), Cm(20), Cm(0.4),
                 size=8, color=C_LGRAY, italic=True)


# ── FUNÇÃO PRINCIPAL ──────────────────────────────────────────────────────────
def gerar_apresentacao_pptx(d: dict) -> bytes:
    """
    Gera PPTX no modelo Reunião Trimestral Braúna. Retorna bytes.
    Slides: 1-Capa, 2-Agenda, 3-Macro, 4-Patrimônio, 5-Carteira vs Modelo,
            6-Composição, [7-Risco/Desvios], [8-Calls], [9-RV c/ Estruturada],
            [10-Internacional (perfis moderado+)], [11-Sugestões],
            12-Estratégias, 13-Ecossistema, 14-Próximos Passos, 15-Encerramento
    """
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    tem_desvios  = bool(d.get("desvios"))
    tem_calls    = bool(d.get("calls") or d.get("estruturadas"))
    tem_rv       = bool(d.get("acoes") or d.get("fiis") or d.get("sugerir_estruturada_rv"))
    tem_sugs     = bool(d.get("sugestoes_produtos"))

    alerta_i = d.get("alerta_internacional", {})
    tem_intl = (alerta_i.get("perfil_elegivel", False) and
                alerta_i.get("gap_pp", 0) > 1.0)

    # Novos slides inteligentes
    tem_analise_produtos = bool(d.get("sugestoes_por_classe"))
    tem_acoes_alerta     = bool(d.get("acoes_alerta"))
    tem_carteira_fiis    = bool(d.get("carteira_fiis_sugerida"))

    # Agenda dinâmica
    agenda_items = [
        (1, "Contexto Macro & Mercado",    "SELIC, IPCA, câmbio e perspectivas"),
        (2, "Visão Geral do Patrimônio",    "Evolução patrimonial e rentabilidade"),
        (3, "Carteira vs. Modelo",          "Alocação atual vs. perfil recomendado"),
        (4, "Composição da Carteira",       "Distribuição estratégica por classe de ativo"),
    ]
    n = 5
    if tem_desvios:
        agenda_items.append((n, "Análise de Risco & Diagnóstico", "Desvios e pontos de atenção")); n += 1
    if tem_calls:
        agenda_items.append((n, "Calls & Oportunidades", "Recomendações e operações estruturadas")); n += 1
    if tem_rv:
        agenda_items.append((n, "Renda Variável", "Posições, performance e operação estruturada sugerida")); n += 1
    if tem_intl:
        agenda_items.append((n, "Internacional — 15% em Dólar", "Proteção cambial e diversificação global")); n += 1
    if tem_sugs:
        agenda_items.append((n, "Sugestões de Produtos", "Produtos recomendados pelo Head de Produtos")); n += 1
    if tem_analise_produtos:
        agenda_items.append((n, "Análise de Produtos & Trocas", "Sugestões inteligentes por classe e assessor")); n += 1
    if tem_acoes_alerta:
        agenda_items.append((n, "Ações — Alertas de Preço", "Posições acima do preço-teto · Estruturadas")); n += 1
    if tem_carteira_fiis:
        agenda_items.append((n, "Carteira FII Automatizada", "Portfólio de FIIs sugerido por perfil")); n += 1
    agenda_items.append((n, "Estratégias & Recomendações", "Plano de ação para o próximo período")); n += 1
    agenda_items.append((n, "Ecossistema Braúna", "One Stop Shop — soluções integradas")); n += 1
    agenda_items.append((n, "Próximos Passos", "Imediato · 30 dias · 90 dias — ações acordadas")); n += 1

    # ── Numeração sequencial dos slides de conteúdo ──────────────────────────
    _n = [0]
    def _N():
        _n[0] += 1
        return _n[0]

    # Slides fixos
    s_capa(prs, d)
    s_agenda(prs, d, agenda_items)
    s_macro(prs, d,        num=_N())
    s_patrimonio(prs, d,   num=_N())
    s_vs_modelo(prs, d,    num=_N())
    s_composicao(prs, d,   num=_N())

    # Opcionais clássicos
    if tem_desvios:
        s_desvios(prs, d,  num=_N())
    if tem_calls:
        s_calls(prs, d,    num=_N())
    if tem_rv:
        s_rv(prs, d,       num=_N())
    if tem_intl:
        s_internacional(prs, d, num=_N())
    if tem_sugs:
        s_sugestoes(prs, d, num=_N())

    # Novos slides inteligentes
    if tem_analise_produtos:
        s_analise_produtos(prs, d, num=_N())
    if tem_acoes_alerta:
        s_acoes_alerta(prs, d, num=_N())
    if tem_carteira_fiis:
        s_carteira_fiis(prs, d, num=_N())

    s_recomendacoes(prs, d, num=_N())
    s_ecossistema(prs, d,   num=_N())
    s_proximos(prs, d,      num=_N())
    s_encerramento(prs, d)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
